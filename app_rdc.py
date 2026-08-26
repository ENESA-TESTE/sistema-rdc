import os
import re
import datetime
import zipfile
import io
import openpyxl
import pandas as pd
import streamlit as st
import altair as alt
import base64
import json
import time
import tempfile
import plotly.express as px

# ==========================================
# AUTO-RECUPERAR LOGO CASO O ARQUIVO SEJA DELETADO
# ==========================================
B64_LOGO_PADRAO = """iVBORw0KGgoAAAANSUhEUgAAANgAAAEsCAYAAAChVH1WAAAAAXNSR0IArs4c6QAAAARnQU1BAACxjwv8YQUAAAAJcEhZcwAADsMAAA7DAcdvqGQAAFV6SURBVHhe7b13dBTHtvZ9/rn3/da6733vPec45wTOOGeDbXJGIATYxgZMBiUEIoMjtsnKIgcRhXIO5JxBAgQiKqIcUM7s79k1MwqjklCYERO6WD9q1N3T1VW1n9pV1dU9//jH8D3/8V/Dvf/ZETw+2uvR/x256xlz5bEfd77x2He7Bv979I6Rj3y/08qcePR7b8t/jd793b+/9/7tX995rzJpRnuv5vj/frvjf//xyPc+vR4Z7R3zyA/el0SsN3bHPPK9d9Ijo/eQgjnjbfr84CPi/x2765l//Hv0npGPjvGnR8cGkoj1ih89+qOvglnjY/oIO/dRCezx732sHvthDz2GzItYQUGhfai19KRGYI/jj8exUcQKCgrtQ62lWoE9gT+ewEYRKygotA+1llQC+9HH6skf8ccYXxKxgoJC+9BoSRGYgoIe0BbYU/jjKWwUsYKCQvtQa0kI7CkI7Gn88TQ2ilhBQaF9qLVUK7BnfvShZ8b4kYgVFBTah1pLT44NUASmoKBzFIEpKOiR+gJ7FgJ7dowPPTvWj0Ss0GFwX/0JVMaTEjTb+RjZdxUMGLWWFIHpieeARjyPjgmgR8YE0r/GBIFg+p8fQ+n//RgOwiCgYOqM4zo1iR89/mMw/bf6+P/9MUSc498432M47+Nj/NFa+or0ZNeh8JCoL7DnIbDnx/rQ8+P8SMQKreY58MRYXszMQgqm/zcmHAUcTO+M9aau47ZQz5/W05DxLjR8/CoaO+FXmjxxHk2aOJ8WT3Kg9ZMn0joJYvukSbRgogONn7hAHP/9hL/Icvxq6v+TJ305biu9P3a3EOl/jwmD6ILosbEB9MxYX+k1KnQgai29pAistbDxwmOApyCoxyCof44JhbDCIKRtMH43Gj/hb5o3aQZ5Th5D/lOs6PC0fnR2ene6Yf0xJdq8RwW2z1GV3b+pwu5RKrN7jIrtHqeSZuBjKnEsH59l8yol2LxPl6Z3pUPT+lPwlOHkPHkczZg0l8aMX079IORX4dX+B9fzKATH18jXqrlueZ4UdE59gb0Mgb2Iwn9xnD+JWKERL8DlPw9jfQYe4gkY7tPwGt0gqEkTltMfk+1py9QRdGRab4qz/oxSbN6hUrunhIhKIY4iuyeo0O5Jumf7FOXbPk05EFi27fNtItf2GXGOAtsnxXlLcH5Op8D2BUqwfp/OT/+K/Kda0NLJ02naxL/o63Fe4lqfHBuEbos/8uCHvMjzqKBD1FqqFdhL/Ac2ilhBwKJ6TniAAEEnbBv+0xpaMtlRGPG56Z/Dm7wL4bwoPAwbewGE1F4RtYVc22eFgIvg8cpxLUVIP8mmi7hGn6mW9POkBdRz3GbRSGjyw/mT5VtBB9RqSS2wl8f50ss/+ZOIzZznUTjPjA1EKxRAr43bQ4MgqlUYKx2Y1pduWnehLNtXIKYnRfeNjZqNW2b0DxON4Lj7WWz3NGXYdqKr8HC+U0fShPEr6J1xu2AEAfTsuECRX1k5KLQDtZZUAhvnY/XKT370yvgAErGZ8vxPAfTUuBB6FYXTc/wGmjdxMUVOs4AneB0e6QV0w54BT1Ee4o72UO2Br5WvuRDd1nx8zrB5iY5P70l/T55DPcZvpM4/BSLfofQi8i0rF4U2oNaS2QusE3gBwnpsXBi9/ZMvDZvoSh5TptNVm8/gAZ6GQfKY5xm1lzIeUTUFi02TJ+ay9RfkOsWaBk7wgNCC6ImfQuglRWjtRxGYH9y4nxDW6xDY6IkraOfUH+Ct3sH4pW4cZUyeqrWoxPY0Fds+QXcwlnSC0PpM2IgGJ5iehVeTlZlCCzF3gT2D1voFdItGTXSBsL6juzZvUSnGK2xwMmM0dcQECTxarPWXYhLn3Z986KmfQsVYQlZ+Cg+gvsA6j/O36jzenzpPCCQRmyidwItombmFHoiWet2UiXTL+j0qF8J6Smp45gR7tGLbxzBee4mipg2iKZOWobuIhghlxmUnK1OFJlBryWwE9jJalBfHB9OH4/fQH1Pm0hXrj6kQ46pCdAdlxmbO8IRIKbqNSTZv05qpU6jr+J30/PgQUYayslWQoC2wV/HHq9goYhNCZHJ8EIwjmEZM9KSI6RaUZfsyukNPGOT0uuHwnOg23sPng9P703cT3UUZvoyylJWzghZqLQmBvQ6BvY4/XsdGEZsIrwqvFUIfT/Ch36fMpxs2H4gJDG6h5UaloA03Qlxmt23eg+efR++Nx7gCZfqapLwV6nhDraU3TFVgLK7nfwqjARM2055pP8BYXhYtsinPCuoLLjMuu2zbV2jr1PH0xYQ99ML4UEVkzdBAYNxF5MJ6DRtFbMRwpjqND6TnxofR2EnOdNK6J8ZZTyvdQR3A6yCLUJYh0yxp4MRNEFmdJ9OuB3OHnRXHKoFNgMAmYMdEbOTYSHkddJoQgM/BNGPyEoq3+ZiKbR9XvJYO4bIss3uUjlv3IotJG+mFCRCZVj0owBbVWhICewsCe2NiAL0xKYhEbKR0nhgEkYXQn1Mc6a7Nm/BcT0iNRKF9qKbzn6BT6B2MmLyWXpwYCoOS14m58qZGS1PVAnsLf7yFjSI2Ql6bGExvQVx/TZlNGbavqqffFc+lL3LE5McTdNnmC5o8ZSVEFg6DCpTWjTnytlpLJiGwzhND6EN0DZ2n2FM6xKUakMsNQ0F3sCcrtX2MLlp/RSPhyV5Wi+xtSR2ZG40ExoXCG0VsRAjPBYGtmDoDg/AXzXap08OCRVaEcS5PJg2fvBGNXajKwMycLmotGa3AROuAMVcXCGzpNEfKtO1MeYq4Hho8mcQi6zdpG70yMczsRdZAYG9P9h/11mQY7RS4NY6NgDcmB9Irk8Lpj6nzKAvdQsVzPXxKILKg6VbUY/Iuem1yiLTezIW31VpSCWxCwJdvw2C7TAkmjg2dtyYHUafJ4WQ3dSkl2XZRxlwGQi66i4W2z9C2aRPoA9TRG5ONw570gUZLQmD8Xxf88S42cmzovDo5jL6dsp7irD8XM1nKbKHhwA9xZqJH8de0efBi4TCyIHpHUoemzjtqLanGYDZ+T7+LP97DRo4Nmdcnh1KvKd50yKafchPZQOFbJLdt36Mp05wgsjBpPZo6Gi0JgX2A/96bEkjvTw3GDt5pmLw5JYQ+AFunT1Avf1IW7Romz2E89piY9Og/ZSfqLVRan6aMRkusLaMQ2DsYNL45JZwWT/9ZdEGUByQNG278+AWrXtMn0Udozd9CwyirV1PF6AT2OsQ1euoairf9BBWnLIEyBrgRTLd9jeZOXwKBRWCMHyStW1OkkcDeR+Y/mBpCHBsab00Jo8+n+FOw9SjxbJKsMhUME34FQYxNNxowdRu9NiVSWr+miEZLtQL7YGoQfTgNYxzEhsR7aAnenRpGv0z/RazUUB47MT6K0Sius55GH8Lo3gGyejY1NFr6dHLIU2qBBWNjKDYGGxRvTw2nIVO9KM7mU7qneC+jhJ8gz7DtTNOmO6E+I9B9YuOT17epoNHSe5MCPxAC+xB/fISNHBsKqooIpc02U9Gff1aZkjdieOr+kM1A+mqqn+iRyOrbpFBr6YPJIZZCYB/hj4+xkWNDoQtau7HT1lKi7ZvwXspSKGMmT3TtX0RX/3dRr9xdlNW5yaDW0oeTgkcYpMB4kPgZCLb+FuJSvJexw/XH7/U/Y9ODek/bI3omsno3GTRamhxkJQT28bRg+mR6KETGQnv4vDUtkmZZr6BUm9eVtYYmAo/Fcm1fIBebWfTmtChpvZsMai19Ok0tsE+m4w/rUIiMhfZw+XB6CHVFHGZjKablFe9lOvDa0fO23Wig9W56f3qYtP5NArWWDFJgb02PIkfrZZQqxl6K9zIlVL/Q+Tw52zjQG9P3SuvfFNBoqVZgn8JrfGYdRhw/TD5Eq9YV1xFg+614HF1WSQrGDS/SPmv7DQ2z3kHvTY+Q2oGxU6ulWoFZY6MNNiJ+mHSxjqapNm50x/Z9DIqV+16mCL8wp9DuWVph60hvW++X2oGxU6slFlhXCOwz/PE5NnL8sPgYbrWrdTh52UyE93pUWjkKpkEJ6veQ7WDqb+1H71tHSO3BmNFoyWAE9jl4xzqKRthspSu2XZVVGyYOj8WybDuTo81Setd6r9QmjJlGAvsCf3yJjRw/DD6D9/oYAvvNZjEV2BnXbyArtI0S20cw1h5NvWyCUPfhUrswVmq1pBJY5DNf2ITSl7bIJOKHwcc2EdTPxo8O2lkqj6OYCezFUu3eph9t1tK7NvukdmGs1GpJI7Av8UdXbOS4o+EL+cgmmibYrKEM207Kinkzgu+LOdnOhg1EiG6VzD6MEY2WVAJziHymm20ofWUXThx3NF/gYr6yDaONdnbivpfSPTQf+KHM03b9aJDtHvrUNlJqH8ZIrZZs1ALrCgPvZhcB5YV1OJ+hYAfY+lCs3VcocOU9G+YE91ZyEE+3c4PAotG1ktuIsaHREnpnD1dgXKCf20bRDDtXdA9fUn590szg3grX+WY7e7T4GLcAmZ0YG40E9pVdGH1tHwHXFtahdIMr7YqL2WQ/A4X9gjL+MjNYYNxrOW/Xnfrb+dMXdpFSOzE2NFrqphHY1/Zh9M2MCOzgnR1HV/tw6mUXQiftemP8pTzzZY7wbwpw78XB3o2+tI+iryR2YmxotNTNvlZg4dgYiY3hHcqX9pE0ZcYmSrJ9Q1nYa6bw0qkc2xfRi3FAgwsPhtZfZivGhEZL3ewjHp7AvgKf20eTm/18yrR9RdwXkVWAgulTAC921LYvbCISRtnxDb2uaSSwb2aEU3eHSOzgnR3D13CjX87YS+F2VuJNvcr0vPnCTzvfsHuPxjjsoK4zoqT2YkzUaokF1hsC644/emAjxx3FV1C5lYM/XbT7kvhH3GQFr2Ae8M9Ppdm9Rn86/AmBRUvtxZjQaKnnwxTYF/Becx1cKcHuXRSwsrjXnOFxGL97xdt+PH0644DUXoyJxgJziKAeM6Pg2iI6hB7gU4eD5DJjPmXbKeMvBdWDmIft+lEvGOfXDh1ni/pAo6WeM9UC64k/emEjxx0BC+xLh0Pka/8DClb5jS8FfnfiExRn9xlNmrmdujrsFV6A7USDzI4MlVot1QpsJjbOwkbEHcHXuIBhM4PpmH1f8diCrMAVzAu+TZNo9w4tmOlCn808BBuJFnzTCHgHAbphoAeQ2djDpJ6WVALrNSuSejtGY0dkh9B11n6ynbmB4tFiKQ9XKjA8Dsu2fYGO2w8kD4f55OawkH6ftYxmz/KksbN20IBZwSCE+s0Ko76zIqi3sKUo6jFrL30Ne/oKfD1rH9hL38yKpu6gB/bL7E/faLRUKzC+2D7YyHFH8OWsQ/TbzGWUZvemuJMvK3AF84Nv1fBvPWfZvkyZIM2uM6XYvU4JsJPbIMauKx2cMZxCHMbSdgdb8py5gP6AHU2fuZ6+nbmLRszyoWGzAmnwrFAIEQYOwXWddZi6zTooRPgNxNcDwpPZpC7RaKnPwxLYp7OO0Sq0UvdQmKqWS17gCuaHSmT8ktJnxEJgnsJnuAvJy+n4dwr4VdwsRF6Jz14vC3Gq7Wt02b4bHZ9hIQS4ZaYDrZj5K02duZYmzNxCo2d5kyW8YH94tS9mHRGNPHs9fQjuoQqM3Xo3x8O0fcYEKrf9V6MCVlBoDhZgHSw0FSxGbUHym8nycFyi3Xt0YUZvinL4nnY4TKNfHX6jGTM90PX0pr6Oe+lL2GN3eDeZvbaFxgJzxB+zsRGxvumOxC1nh1L0jBFUbKO8PUpBH9QXoeqRGH64k19HUWz7qPgN6RS7LnR2Rj/yd/iR/pz5B41wDKRejlFSm20tGi3VCqwvTtxv9l6oOUrvfOO4n8Y5etMZ+75UqKzgUOhANIJjj8eC4wa+yua/YYfPUrDDeBrqGEw9HfdJ7bY1aLSkEtgCCGw2Ns7BRsT65qvZh8jGcTNdt/8CmVRmEBV0j0pIdV1HTbdR1Y1Ujd94ci3F7g06NWOQENeCWU7UfzY8ELyPzG5bg0ZLfWarBdYPJ+0/Zx+UF613vpp9hOY4ulOqfReRSVkBKSi0BJWIWDQsINWYi1H9YCNPfrxEabadKdW2E12x/5yOOAwjv5lTyXXWLzTL0YPGOO6kUbMDyALC6jn7IASxV2qzrUWjJXQVO15g3WYfo19m/U35dqoCkhWcgkJ9NCLie6YFNhhH2TxGJeje5cMTZdu+THftXqdEu7fput0HdMG+K0XPHEV7Zk4nj1k/06LZHjRx9g6IKASE0SCIqd/sfdRzzkHqjt5Uz9kH4LV0a/uNBNZ/TjQNmLsPO3infuk29wQtnbWIymz+3aggFcwX1WQEeyKViIogolLYSJnNI5Rr9xKlQ0S37T+gq/afYfzenQ7a9yPvmZNp48w5tNRxGS2Ys5bGzvGlAbCxvjBwDX3m7KfecyAiwJ/7iO3wVFp2qUtqtdTRAhOZnnuQNs6yFy2QrKAVTBl4HOtnKGvaU5Q55XHKn/aomGgot/kXFaGLl44x0W37j+jyjK8wNupLEQ5W5O0whjwc59Nyx79oHkRkPXcHWc6NoO5zj4DD1GPuIeoFm+qlFpBGWDwO0tBfILdJfdBAYEMgsAFz99LAefuxY69e6TN3Pw2dG41+8ETRQskrQcH4gHBs6gERZU9/CkJ6krKmPEZZkx+hrIn/pKxJ/6KcGa9Q7sIPqHhFL4r92YrCHEbDC1mT86z59LvjCpo7m0W0ncbODaAB8w5St3kn6eu5x4SgeqrF1Bd21J+NGDYlQ2Z7HYlGS3UCm4eN87ERsT7pM+8AWc0Lp7CZY6jQRpmiNw7q7idl22DMbM1AQNZPCxHVComZ+gRAw2n/MuX+/BndWzGACjdNohLfBVS2z53KzwdR5dVDVHXnLFFWPB09eo3G/HKAumJczkL6Zh6ENA9CmneIekNc/ebtk9qRoaPRklpgh58ZiIwMmn8AytunV3qj4H6cF0jHZg4V/Wx5hSp0NPWntevPyvFqCH6tWoOlSXb8nRcox/ENyvuzFxW4jqTibTZUGvo3lZ/YTpW3TlFVUixVZ9ykmrwUul+ST1RVTrKQmV9NU5zOiNXxMnsxVjRaGlQrsPn4YyE2ItYnvecfpjHzA+jEzCGKwDoYFlED8dg8iV4ETyY8Ltb45YrjeDKhE6XavUaJdm/SDfv36ZTDAIqe+T3tmWVNSb5OVH3RDwKKoZrcVKopzBICul9RTFRTpZZNK8J9omU7LonlcwPYOCU2Y4wMWqDS0qB5aoENwh+DITCO9UkfCOyn+X50euZARWB6QCMijYAKIB4e6/KEEr+5KUfMxr1GSXZv0037DyjO/hM6Z/8V7Z9pSQGOk2nL7DnkNOdP+nmeB42f50uW6M4Pmx9JFvP5PRkRdOBChloZugsBBxLJYgGPqfZKbcYYGQyBic8dLbDeC47SNAxgr8/4XBiBzEgUmodFJLyQEJBqOrvY5hEqEzNxT6IL14lS7N+mW2I2riuddehJex2GkN/M8RCQI7nMWUK/zXWlmfO9aOyCIBqwAGOdBYdq6d+Ag9R//kEagC7P1w7RtOdAkloWugu3k4vohyVHxPo9bXsxVhoJbDBakCGLDmLHfr3Sa+Exspm3jW7P+EgYicyAFFQ/iiATEU9nF9o8Ten2b9IdlOFlh2/o5MwB8EBW6ML9RJsc7SCg3+jPuatp/rwNNH2BN/2wIJx6LzxOPReeQAN3jPqgkeu74AjEdFiIaCB3Z9Rw10bFfkH9uuvtuJdW7Y6j0rJqtTR0E6pr7pO92xmVwLTSNFaGLFRpyWKhxoPxDghMU7D6ggVmLQT2sZkJrG4mTkPteMhGMx6q687xi4ASuQvH94NmDqLoWd+Rn+Mk8nB0pOWzf6Nf5jrTnHnryXr+Lhq3IJhGLNhLPRaeom8WnkZ8gnpBSNxb6CtExPVaJ5q20mfOXlq87iLlF1SopaG74LbnGg2cqxm/yNM3JgZDYOKzRmCDF0J1i6E6xPqk9yJ4sPnb0PqaksA0otHMwtU9mySHF50+Tan2b9FVBxbQEIpy/IH2zLamtXPmo/u2lOZiDOQwfxNNn7+Txi8MoO8WRtAwdOF7LIKIFqlFhMaqz8Ij1HchRMQVKilvXdJ33j6yXnWK0rJL1bLQXQg5kkIWiw6gK8pDFXn6xoToDSKu9WBDkDmLnw9hxwG90mfxcbJdYHwCqxOQSiAqr/MUum+8sJQFo/qdq2zbFylTzMS9TKl2nemKQ1c6M3Mg7Z81mvwgoI3zFtKyectpxvyNNGXBTpqw0JfGLAym7xZF0vBFaPEWHaFei0+Bk2iMjlMfNEj9Fh2l/osO08BFh9DLkJdrRzAALfK4P49RUnqRWha6CzeTC2nk74chMDZOefrGhMVilZYsFu5XBKZBIyBNl43HPtxt44WlvDo71/YFyrDrRHd5Fg7e57b9u3QN3biTDgNor+No2jPHjtbPXUTL56+iWQs2QzjseULp24XhNGJRFFku2odWGp5mMbzO4hPgOMRzDOI5QgPUAhqE/bJyMwR4Eow9THxygVoWugslZVU0YdlxcX5Z2sZGI4Gxex4KgXGsT9io7B6WwNRLeHgdnFhxMPlRypr0b8qf8i/1yuznKUs9A8eLSuNnfEIXZnSjkzN6UJjjGPKZY01r5v0KAa2k+RDQ1IU+EM1eGgrPM1Tkj8XBXobFwqI5KuKB+HuQgPfxMQcblYsxMAQC40Wyl2/nq2Whu1BRWUO/bYzBGEzlwWTpGxNDITDx2eQExst4pj8NEUFAvAYOAsqa8E/K/On/inVwfEzunLco75dPKX9pbyp0taSbf39HIbPH0NbZDuQ59xf6e74zLVi4iewWetO38Dx94G36Lz5KAwQQTC31haNBJSBjFVFzsOH3xzhs79k0tSx0F2pq7tPW0FtiYoCFLEvfmGgksCE/wyB+PUwc65M+v5wg24Xb2yiwuoWktV5II6KJLKL/EvvFEp5fPqN7y/pSoedoKtpmS8VBf1FJlAuVHdtOFTFhVHXrBNVkXicqyyKfQ8nUfdFZ6olxT5+fT1C/n49T/58hpp+P0KCfuUzg7muR58scGAz6QQABh5PFCgxdhyCcdxCMcvBiiFkrbWPD4heVlix+ZoH9Dg8G4xn66xExDtMnfX85Cc/AAuP7YE+pJw+0Z+E0Ewl1U9j3MB7KmYIuHbp1ObNep7xfP6eC1UOo0OM7Kt7lSKVhy6js+A4qvxhKFXEHqSrhHNVkxNP9okw0j81PK2/Ze0e8pFJ2vQp1cAMzAN5la/htqqrWvcLiEvJVM9nc8mulbWwM/UWjJbXAhv5yiIb9dgQ7eKf+6PfrSbJftIOSZnwgJhFUi0hVsMhYXPx06lWHbnR+1gDaP3s0+cyeRt6L51NyhBdVxh+FeC5Qddo1qslNpvsFGUSV7Zs23hx1i7o7YiwluV6FOizAQBi/k/c1KivX7c1mDhl5ZTQCNjgYIpOlb0wMg7PiuM6D8Q5kjmN90v+X4zRpcQBFzxlHF2b1o0jEPnPtaOP8xbRy4UqasWgr9u+hnxYH0pifQ+n7X6JoyIJ9NNXpNCVnt2Ex6QPCfTTELgHx1HPOPun1KjRkAAS2Wk8Cy84vp2mrT6ObyF0sefrGgugN8meNwNh7Wf5+FMo7rFeGAstfD9GoX/fRd79E08hfD5AV/rbkC/r1GA387SQNAAN/PUGDfj1Og7Gt1/xDNGPdRcI4WOehpLSKlmy/TH3n75der0JD2Lv86XWZilFuug75RRU0b90FGrSYPYA8fWPB8jeVlix/PmT1D4ufDz4Nw68cBoGxAPQNT6awcFQcpSECFhjcqtaxTL8FB2jB5hh1Neg2FJeoBNYHApOlrdCQQRhXOLifE2LQdeB7Yct2XkE3FF1ESdrGxDAIjGOMw6z+YbnkyCNwabctfz/W6MCHDc/G8NS3s2+8uhp0G4qKq+iPbSoPJktfoSH6FFh5RTW5+8XTwIUmJrDhP5/4p+Vvhy4MX3K8kat72IiLRJfBDQWvj5CWVUqz116g/hCxLH2FhgyGwGZ6nKN7ehDY/RqiPfsThcBkaRsTtV3E39QCG/rb4RjLJSfERIchMQQezBLx7n2J6mrQbbiTWkQ2rmepP0QsS1+hIYNQHzM9z+tFYHxvzedQEho746+LYdwbRGzx21GVwIb9fjhm+J8nsIN3Gg5D0AqMWHKUwk/cVdeCbkPC3SKyQ5dnAFpmWfoKDRmM+vh+6XHKyi9Tl6BuQ/S5dNGbGCpJ25iw/OOYKq4T2NEYq79OYsdRg8ICrcDIP49RyLFUdRXoNrDA7CGwgewpJekrNITrwxINHt+z0keIuZ0vehPDJGkbEzzcEp+NQWCjIDC9ebA0CAxjCkVgLYPrwwr1oS+BXUowLYEN+/XYMIMW2BBU6OilJ+jMlRx1Feg2KAJrHfoWWOwdk/Ng44TAhv9xNGbE3yexg3caDlyhY1acpKsJ99RVoNuQcLeYZnicp0EYW8jSV2gIj41G/AWB6WkMxgLj8TB3Q2XpGwvcCHFs+efRv1UC+/NYzIilpwixQWHx+1Eas/IUxd3Rk8DSIDDPCzQIQpalr9CQoWiVR/x9XL8C496EJG1jwuqv4yK2/PPIakVgisBajCKwllErsCVHVv1j+MoT/8SGmJHLTosdhsTQP47R2FWnKU5PXcRECMxhzQUaDCHL0ldoyLAlx2gkxsR6E1hCPg38FY2dJG1jYsTfJ9SfjykCUwTWchSBtYxGAkOGYkZAYJwxQ8ICAhsDgV1L0v2LVjiwwGZAYIMgMFn6Cg0ZCoGN0LPABkBglpK0jQkrCIxjS0MXGFfoDxiDnbmmn2l6RWCtQxFYy2gkMKu/TsSMXHGGrFB4hsQwXOSoZScp5Lh+bjSLLuK6izQYhiNLX6EhXB8jl5/Uo8Du0UBu7CRpGxMjYLPis0Zg6DPGjILAuHUyJLgl+w4VGn5K928y4pCYXkwzIbAh6pZZoXm4PvQpsMuJBaI3IYzUiBkJgXFstVTjwZaeihm58ixZsfIMiGFwtd+uOKVXgTmsj6HBf8KtS9JXaAjXx3jns5R9T/6Deu0K94kizqbTYIy7ZWkbEyOWn1J9XnpcLbDlp2NGrjpHVrzDgLBUX+zug8nqWtBtUAkslgb/hdZGK22FxgxdepIcN8TSveJKdQnqMEBg/sdS0dgZf12MWHFa/VktsBEQ2CgIjI3ZkLBadgr98ZO0LvyOuhZ0G9Jzymj+1suiUmXpKzRkGOpizqZLehEYP3DpezSVhphAXYyEwDiuFZjVCniw1ecJsUExfPlpUakuQTfV1aDbUFJaTUv9rtMAbjUl6Ss0ZCgaPMdNl/UisPKKGvIIuUVD/kbXSpK2MTFi5RnV5+UnDVtgzCB03xbvilNXg26DIrDWwQKbs1k/AuMf9lvlf50s0KDK0jYmGglsxEp0EZ3OYwfvNCxYYLO9LlNlle7f21ZcUk1/+8bTQKQhS1uhIUPR7Vnue51KdPwrlxxYtIu3xdEw7mZJ0jYmRq46I2KrFUYgsCFo0Ww3xlJGru5nrvgV0O4Rt6n/kuPStBUaMgwezCP0lujO6TrkFJSTzboYsjRFgWFDzLfOF8QOQ4NbtClrLtLVRP0sl9p6OJn6/nFCmrZCQywxJnYPuS1esabrkJZXRt87nSUrtYEaM6NWn1V9XnHa8AXGlTrW9Twdu5Ktrgrdhu1Hkqk/xmCytBXqGAG4sdtzJFX83JCuQ1xSIcSFsYspCmzE6rMxo1wuEmKDYzgu1GrlWfI5mqquCt2G8LPpcOVnRDqy9BVUWIGhaOyCT6WrS063Ye+FLNQB17c8fWNipNM5EVutOmv4AuOKtVh2mjbp6d2IUahYvgdoqQisWUQ9QGAHL+m+J1GDsfCGqASyZA+mla4xIheYKwSGPrChMRIMWHqKlvjfED81qutw4moujXW7QMNYYFppK9QDBjMYAruepvsfQa+ovE9/+d5QNXKytI2Mkc4QGOJagWFDzLduMWKHIWKBLtzMrVcoJbN9vwUmC3EJBTR9fSwNReXK0lZQwQYzCi3zzbRidcnpLpSh4Zy+8RI8mMo4jZ1RLudVn52MRGCW6JdP9Iyh8zd0/wPciRkl5LDlClmgeyJLW0HFcHiwqWtjKTlb943crbQS+tH9gkhDlrax0VhgLhCY+yUayTsMECtxseco8kKmukp0F/IKKmnB7ngazAKTpK2gYtjqczRn21XKzNf9e+mjz2fRd64XVPUsSdvYGIW8iM9O59UCc4XAPCAwV2w0UAag/7/hQLL4VUpdBj7f38E3aQC6obJ0FVQMhXf5y++mXpZJeUYkCHGNYKOUpG1sjMKYXnx2MSKBDcYY6Ref65RboPsW1D0qgQYrAmuWIeime0Ym6vwmc3XNfZq946oQsCxdY8QoBTYMLdzkdZfoWrLuZ7H8TqSjz4yuqAu6KJK0Fc7TIDRAvifT4PF124W4jfHX5A2XRBdUlq4xoiWwK/8c4XohZpTHZUJs0FigEqL1cB8m+mIWjfa4SJboosjSNXtgMANXnqXD8bnqEtNd4PuQo91Nq+xHul0UsZVGYCMhsG89L2MH7zRcuJLXRifq/H7YlcQCmrzxEg3lQbYkXXNnOFrj7yGCmATdrwd1jrxDQ9BwsohlaRsjo1BW4nOtwNwuxny75gqN5B0GzDC0crO2X6W7ubp96UpmXjnZbYsjCye4dUm65s5QlPsMr6uUkFGiLjHdhOz8Cpq14xoN5u6hJF1jZZRHjOqz2wUIbD26iB4XY0atvUKIDR4LjJVO3tTtm3558eo8n+s0iJe4SNI0dwajzP8IuEX5hbqdQTx2NY/GrI2lYWjtZekaKyM9Y0Rs5V4rsBgILA4beYdhM9j5AnkduUvV1bodbLtEJtIwVy4UebrmTP/V58ljXzJVVumua85zJRsOpaBRO09WkjSNmZGesSKuFRgUF/PtujihPEPHEq6XuxW5hbqdrvc7mU7f4vxW6hZIoY6BEFiwjm/yp6NbPnf3dRrsgrGKJE1jZtSaWPVnIxQYtwxD4WnO3NLtsqmz6HaO3XBJCFiWrrkyHOUxZu0lOqXjZWrH4/NoNFr64aLFNy0aCWwEBDZq3VUawa7NCLBEN84tMgndRN11We6VVNHErXE01A0CQwHJ0jU7UA5D3GJoBnoMupzg4FngNQeSaSC6+5yGNG0jZuSaSyK28oxRC2wNBLYeAuPMGgGWaPUmbI6jzHu66ybymGCR300ahhbbSpKmuTIAXbgl4jUBumvMEjJLycbrKlnweFeSprEzEh6f4zqBrYXANkBga3GAETBcXHwsBZ7JUFeZbsL2I3eFuPj8snTNDS6LweiObz2m2x/fiLqcI84rDFKSrrEzch0EhthqrVpgI9fGxny74ZpQnjEwAgzziCXHPdeppFx3a+NO3yoQ52fD0k7THLFEI/bj+it0KE53Kzjyiippke9N4b1kaZoCo9ZdVn9WC8xqXWzMyI3XyArKMxaGQQTfb7xMR67lqauu/SG3uJLGbokT55alaW4MxkDdZmc8JerwQdcztwvQQ7hEw2GAsjRNgRHrL6s+awQ2YsPlmFGb4sUOY8EKDIEX+y3kDlXo6P4MP1m7LCyRhqHyZWmaG4NRvn+EJoj3R+oilKN8lwTfoaE8CSBJz1QYueGKiK3WXzJegTFD4WnGweOcQ6uoi8BvIwu5mE0W7MEk6ZkTwsOAnSd099NR1+4W0yic29LEG7BGArPacCVm5Kbr2HDFyEAr63mJXPalUJWOXq19C92h7zbE0XD0o+Vpmj4jgAVEMM7rGl1O0t3jQauiktE7MP1yHQH7UX3WCGwjBLYFAtuIjUbGEAhh3NZrFJuoG0MoKq2iP0ISaAi34JL0zIHhYAiEMDfgts66h1dTi+nHrVdpKDdcWumZGiM2QWD8eYNaYMM3xsWM2HIDBYuW2wgZsOYyrYpOplIdzCjy/bCQ2BzqD88oS8scGMYeHN2c7cd185JR7l38HoZGC+KSpWdqWG26KmLLDVdUAhu2KS7GausNsoTyjBELNghk6ISOxmJX+S1H8IoWaIVk6Zk6QyCu75H/G2m6mT08Ep9P3225Ks4rS8/UGL75qupzncCuQmA3sZF3GCcD1l2hRcEJlF9Spa7WtofC0mr6KzKZ+sIzytIydYagwZoTcEfM+rU3cJd7MbrcA1E/sDNpeqbG8M3XVJ/rC2w4BMYFYIxwZjgehIFl2KUc9PPUtduOEIzzDIJRDK2XjjkwdCPKEQILvpitk/dvBOI8ljinBXeZJOmZIpYQmIg1AoNLi7Hyggdj12bEDEFFTtp5gxKy2//EcxLOYbvnFg1YD8OQpGWqWMAwftp+nVLy2v97bFwPk3ffVJchWnVJeqbI8C2qvA7bpPFgm6/FDN92i4Zhh1GDShyA1tf1YPvXzvEDnWuOpVO/dXE0lFskWXomCJffhqPp7X6glSec3A7dpb7cNTSj8mMst8arPm+KUwts6/UYyx23aSh2GDPDwBBU5iiv63TsZvsnPK6kltB4b3gxeEY+tyxNU4HzZ8HGgfKLvdv+98/vj78nGqYhbGhaaZk6w2B/It6iFthQExGYhgGbrtF0v9uU1s7HWe7X3Kel+1Kp3wZ0nSTpmBIsgr5oSJYhv+1dQM3lPt3nNvXnMYlWOuZAI4FhQ8zwnbexg3caPxaARbbiwN12/3j66TuF9D3GdYO4JdZKx5QYDIOw2nqdzrbzhn1ZRQ2t2J9KA+G9hkrSMQu2XVfFGoFZQGDDdt6BYV43GQZuQQZhNAExOeLmcVtDTQ3R79Ep1H8zup+SdEyFXhuu0VIIo7isfd4r6HKeECsjS8cs2HZDFW+5pu4iQmCWEBhik6I/WtGxu29SbHL7xhSXUorhxW7SIBiNLB1jZwAao1E7btCZdnov/v537O1NtJxaDATG8TAWWG9veLBt8GC74MHg2kyNvhCZQ1BCu8dja49n0AAYzhBJGsbMkG3x1A/daZfD6e3qTidml9F0/zvUH11pWTpmxXZ4MI691AIbsu1GzNBdCSjsGybHYNAX3btfolLoXmnbV3mk5lXQBB64Y5wiS8dY6Y2yGe97h5Jy2n7fq6CkCt3oVOqzKZ4GwbBk6ZgV22+KeDALjLuIMMIYi90JNBjKM0UGgb7oBnkcTRcPVbY1RF7LR6FhfMeFp5WGMTIAeWFD8I9t+ysByspryPVYBsoXXU0+nyQds2PHTREPFALbk/lfg7bdiLbYndj4QBNiAIsCbD2dRffbqLHi8mpaeiCNeqHVl6VhTHCj0xtd3j/2pra5a8ivHN91PocGwqv3h7hk6Zgl9QXWx/nGfw7cftNryJ5kGoQdpkxfCMxi5y0KuXqvzSK7lV1Ok/wTqCeMSpaGsdDb6waN9rlDN7PatqyMV3rsupgLod6k/kCWhtkCG+N44PbrKoEN2Hlr6yCfFEJs8vTi/vHu2xQaxyJrW8t95E4hCvAW9dx2U5qGodMHlT901y3ae71tq134IcxdMXkwoFvUF8jSMGt23RYxGh61wHbc9DIXgQ0ELDKL3bfaLDLuGu2MyaX+EBkbK59TlpahwdfZf+dN6oOGYcuZrDatNxSeC3kfiLz3NaK8dyiNBLbrltdgv1QaiJbdHBgEVCK7LbqLbXm8hZ+VWnk0Q5ynP7yBLB1Dgyuer3cVrruwDTeUuVstxAXj6QtkaSgA7zsi7r9DLbB+ENggCGwANpoTPdEKW+y5QxvPZlNpG14NnVNcRUsOpFFvbs251ZKkYSj0A9233aL5e9Pa9FBqcXkNbTmfqxIpt9Ja51eoBwTGMXS1slZgA/zuohLumBX9QY+dtyGQ2+R+MrNNhsff+eNguhBrbxifLJ2HDV9Xd4yV5u69S+kFrf8Rvbzialp5PIt6II+9UF6yNBTq4Z2g/nx7r1kLTEOvXXeoDwxn8f40SmuTAVbRX4fTqTfOwYKVpfGw4OvhvK04ktkmcfGsqWNUqjhPLwNtQAwOjcB23TotBNZn9y2vfv53qQ9cm7nSEwXyDTzZ+NAUOpVc0upxWQm6UBvP5YrC7YbzyNLoaLpBFDwecD+dLVa5tybwAumTKIcJwcn0DbqEvVA+sjQUJOxJEHFv79snFIHVozdgoxzim0i7L+UL0bQqwCijbhbSj0FJ9CW6ZD0ehlFiTNkT8RdIf3RgEoXFF7S6seBur1dsPg2EoXwDr8XlIk1LQY62wHrvvuPVNwCDdexQgFFBGD3QbVxyNJNu5rR+kXB8Vhn9gS5Zz13wZhCsLA190At8yZMQuxPod6Qfl9H6m8j8nYUH0nHd8OhsJFppKLQAn0T15zsqgfXcc8erT2A69cIOBRQO6IkC+gJG9lNYKoXAC7S2i8XeL/pWEc2IThPn+Rpdxx44J5+31x55um1Cfc6vIKqv0ShMj7hLodcLqLSV3re4rIZ2X7lHo4JSqCuul6+Vy0GapkLz+CZpPmsEluDVJyhDtUOhlp7gSxhuX58kWnAok86mloofiGhNyCquokgIzT46HUJIpG/q0QPn5TRkaTcFH8/f617vPN33JNE0CCs4vpAyCls/E8pjrTkHUP84D+dXlq5CK/BLVn9WC+wbCKwXBNYDGxXq+AZ0gwF/gq7eZ2B4YAptvJjXppu0+aXVdCqlhFaeyKERIanU2z9ZCKWrdyJ9Ab4EXZEWp6cNb/8S8HHdQHd8rzcqcURwqugKHk8qoVycvzWBJzHiMstp6fFsGhSQIs79Fc4rKwdDgOuCy+EreIbuWvsMDtQNx93rC6wHBMaZUEiir8CnbPBgEET1A7zD3zDEgKsFYozSlpvSmsDr+HJLq+gkxLb5Qi45wnN8h/OPCE0lC3TPBgYkN4K3W2H/aBw392AGrTuXS4cTiykT3rGilcudOP2ke5WioRiK87JwGc43G652WTxs+Jq4Pr4Gw4JTaADK42s0BLJjDQYITPWZBRZ24z+7+SZ4fROcQV/5ITNmTDfwMYyNW8pvw+/SijM5dDyxlIraIaiWhGKcPym/ks6hC3rgdnED9gPefie3Ese1sn9aL/BKjNiMcnI/n0uWISn0KfL4BbxBV60yMDS4Tj7HddruT6dLaWWikRuChuFziEx2vEGA3gnH3RSB1cGV+CHGH2Mj0mjLpXxKyG3fKwYMJdyFt9oHb/frsWwajNb/oz0J9KUvupqSMjBEPkO99A9MpjPw+BwKSmpo2v4M+sg7QXq8QdBAYOgidvVL8Po6JEMUulnhr4o/QWs4EOOQ1edy4CmMX1g5xdV0DOMy95g8mhSdTl/AW33IHkuSd4MG1/gJBGYNQVWoHwotKKuh6Qcz6QO+oS/7jiEAgXHc1b9WYMleX4VmYUOKWfEFeH9PEo2PzqBT6Ia19fmwhx0qq4luZFdQ2PVCWnk2l6bsz6SBgakwwiT6wCdZ5FOWf0Pnc78UeIMU2nWl7rm1glKVwLjeZN8xCAJS1Z+TVQL7AgL7EgL7HBvNhY9ReZ/A+JacyaX0ovb/5FFHhnK05jHp5bQThrfoeDb9sDeDBoWm0dfI1/sQ1Tt70BX0TabPtPJsbHyIPPQNuUt369VP2r0qGof8voe6k33HIIDAVJ8hsOcwBvs0INnr87Bs+hQ7TJ1PwHt+qfQV4vWx94SxGltgR3sPXaVTKWW0+nQeDQxLo8+Qn4/QaLzjk4L8oQFR59dY4Xr6APn463RugxfHXrxbRhbI77vYJ/ueQRB4V/NZJbAPAlK9Pg7PoQ+xw9R52/8ufY54W1yhusqMO7DYyqvv08WMclobc48mHMyiXmHp1AWNyGu+EBwqWlYOhs77AainoDQ6ldZwuZf3jSL6yB+NJPbLvmcQ4Lo5/iDwrlpggRBYBAQWxDtNl3cxLvkSmfe62nJxVVbep/ySajFxwOSCsnZMl+s78H2u0/Bsf57Pp+8PZNHXwWn0KoT2BozyfUmZGCqvw0PNPZ5LZfV6GPy6goVn8qgT8iP7jsGAMue4gcA+gsA+4I0mynsMjMwD3cIHrS6/k1dJEUkl5IIxzm9n88n+aA5NP5xN1sAGLEK3zOVyAXndKqJjqWWUkl8JV6L+sgGFPDQM0Qkl9DvE9t3eTOrim0qd4dneQevKYpOVkyHwNhrCHhhTnoNXrh/icypoSEQGvQWvLPvew4LL8t2QjFq6hGZRl5AsejMo46QQGFp2r/cjc+ldKM9U6YxKmXMiV0zzSgM2s1h+hjGO3JcluiDPo6V8BQb5KoT5Kr6voRP+fgnG+oJPKnULTqfv92fRfAhxN7ov8VmGOc2fgEbD/3YxzTuVS9+EoDyQh9chtC4wjnfYQODZZeXWkbwHuuA63sS1bYorEC8Xqg34uCW+CN44Ve/X+g7qtAsLBbwdkklvqXkTvFFLFr0Rmk2vg7dDM6u+CUnNYvoEJWVbBVzPHuF/Net7v6uHhMC6QGDvQGBdcHJT5FUY0pDoTLoFI5OFG7mVNAteqSvGLs+hAl+BuN5EJcrOpc3rOI6PfwFCfBfpDIzKpLnoxkTCczQp5ocYeJV/bHYFbb9eTJMP5dBbgVw+yAd4A3l5S5LHjuJN0Blj5NknGzeEmUXV9AMasle4UQhJh+EDCEHE4C0Gf9cCcbSVT4PvFvUNTkzvHZSUYeV/PWuMT2zmWJ/YjKm7TxXb7zhA9jsO0gyvaJq5MZAcNvjTrA1+1X+s21X4x7rdBX+v21G4es2WQmfPzYUunpuPCIG9GXzX6+2oPJxcfaEmxOvgzcB08kPrrR2q0L/3hqENjM6iF/whFBiZ7Bwt5XUYyMs4Bwvuk9AMGnYwm7ZeLaLEgqpWP/TYEYEXCJ+H2NxiC2kQGiA21DdhnJ2Qh84Q2xvCUDuG15HW8wFpNA7d8YzixrdNtt0up5eCMum10MyaN0Iyqt+E13gT8dshaSC9qgvid4PhiTUEpcAzs7dLqekeeCd/WMCNuxYBN9Im+F7MnOp9Ns3a+0yG4879pQu2htN8sHCDLy1as0Pwu6dX5TLPzSVLPbaUrvTYWObitq6UcXdbU+np6k4CFzda4+QEVpMnYg9nF4E7cHNxJVfg7OSs6iK+BoG9AYGxMZoaz6LSbE/mUaHW81HV1USecUX0dlA6vQhjek3y3bbC53oZBvMiewbEPaOyyO1SISVBaIZ4M5t/Ay0HHuN4Wjn9ef4e9UUX+YOwDIgMRo88cF5e1cqjruCyejEYaSGdGSdzazKLGz8VkF1cWT0+Kone9r9N3fxvlg/3jcsY5Xvl7o++sRkzdx5Nm7XraOq87Qdyf98SVL1ksz/9scmP/vLYRH+5bwAbabnbuqpVbusqGBfXNZUurp4Vrq6elW4u7jUeEIo7EAJxchZoROIOVGJxE7gAZ2c1/NnFvVmcXNxUkxyvBqd7vR6dj9Yhw6R4Ce7+XfSX96c0nOplG/eAZ3kDFfsyjpF9V1d0wvnZQF9GWr3gKT2uFNKdewZ8Yxtlw/fYDqeU0y8Q26CDOfRZRCYai3R6KiiDXkB5dpLks7V0Bs+jTJ4LzKAvIjPLnWLvVZdUVDcawN6/f7/qyKlzJSuc3GiViwetBk4QhrOLRzWMuMbVxb1a4OxW4wajVwFhsEhqaSgSDSqRtEwsbaGBwF6DwF5Fpk2J52AUtqfyqUjLe4UklMKz1IlL9l1doTl/ZxYaWuvn4BX6QmhbbxbT3TY8HNnRIb8UYksuoz9jC2rGnMgrRiNR/Fpges0TEMbTwZn0PPL1Cqif16bg414AT6Kr9zy+33Nvdtnsc/dKL2ZW5ENI0gfakpKSy9zc19Q4rXYhiKiRSFxgyIzMwB82tQLrBIF1hsC4ZTIVXgBdgP9t1SpsTbidX0XdYeDcenZCN0j2XX3zHIzsaQjt28O5GBuWUJEBToZoBwigBv8KrmVX3vO+VVz92+WCIuvTeblD9meXdYFXeywgnR6DcJ6A6J5Uw581PIr9PBPXLzqr0vr0vfw18UWV13IqSnDehhVUL5SUlJTt2Lm7fBXEJTNgQ6dWYC+FZHq9vPcevRSWaTI8DgH1RvcmhScY6oVfLxTQ09ySSr6jL14My2rEC+AxGF7n0CyaDC+7/26FIc6DNBkgjOKqypp7CXmV5QcwdtuTXEqu14vzVl0rzlh2pSjr99jCnNXxRdmrrhVlrrxakLU9saRyP467ll1RVV15vxDfb/YljdhfeuDg4SKIq8ZJjHnkRmzI1AoM/WqvlyCwF4UxGD/PAx4vzDx3T11dqnAtt4o+hfd6CgKTfa85WBAvhKt4XpAt4ucE2VKeD8uq0YDva1CdC7DQng5F6w+hfR6RUeZ8ISc3Ob+8bb8nZACBvRFg8RTB2xUjZorUtPidBji2Mib2UpmLq0fNaidXqfEaA7UCeyYs0+vZfQX0DIzFFHiCjRfxxviGU/N/xRbRU/AYT2DfUxE5IFvET4In1DwuyK0jMk/EOK766YicyqcjsqsgqioIrfL58Myql8LSa14OSydB6F2MM1JF/F5IYnGfwJvpfYNupVkFxGVYe59Km+pzJtVh5+HCBV5hxCzcHEi/rPGixWu20Xz3LdW/uW6qWLPRq/r4iVNlpaWlFa0xSlMKcVevVbt5rCVjFhdTK7CnIjK9nt5fACNiwzM+nozMpsejIIaoHHoM8f9AFK/szaPwxLqlNvzUf8+TxfQf4fn0SGRezdNRuWXMs+CliMyyTvAgncLTy98MTaG3QpNBEr0VfIfeCrpNbwfequkbcC3XMuDy3RF+MenTvE+kz9516K7jrsNpv24Nrlyy2Z+WbPajvyCUv93W0d/u62mZ69qaVa5rqla7ra10cvWscnP1qHQFiKt5WliFa+3UME8Tu4JVq5xhWG4127bvqrhx42ZReXl5lTkJLTExqXLTpq2VK1EOMqM1JgxWYE9EwosIweTSI1F59Eg0RAH+FZVP/6zlHv0zuoD+d28hRJVf+WJkTtHzkXlFnSOyil8NvVvULTqj+GRaee0ALD6zuGpo2LXq93wv09c+sZVj95xKH7PndPrE3Scz5m/fl/GzV0T6r9siMpdv9K5ZtX4nrVq/g5wgFtU9D56pcqtxcVbBU8Kuzq4CngJuCs3UcNPTw01UDNLkgT3HQUFhxWlp6dzVMo13GDQRkL+apKTkik2bvSpMQVxMhwlMIxb2LiyYf0Ms/1YL5X81sVosEE3Nc1F5xZ2jcvI6RebkfxieWvBpWHL+Z+DroOtlfQKuUG/Qx/ci9fU+Tb33nCXLXcfLbHccyJi242DmHK+orN83+Ge67YlMz8kvqf1174uxV8s9XFyrPFbxnXfVzcTaO+/1RdGAhsKoQ16guoYFxsbm4bmu8vCRYyU5ubmlPD5RZ8lkAvJUdedOYvVmExIX06zAnmxEtojrexeNYIRowL9qgbeJVnmYR6LyqztHZOe/FZmd9U5ERk73kDu5vUJuZ/cNvpE/1O9ilZXPWWJG7TxMI7cfoFHb9tdYb40onLMlNGf2lrDcvzf65K1cvzt31YZdOW5rtpSu8VxPnh7raa2rB61bvUqwlperODsLRFeLu1xu7tXX42/kqyuwLDw8qmL5itW1HsmY4LEIX7vX9p1lZ06fLSspKW5yatvYAuqm/PLluPJ16zaLxkSWf2OlVmCPR2Z5PXGgCAN6jGEiIZwmeBTieSkyu/jDsPTMjyLSM78JScy2DLiSOSzwWsb3fheKJngfJWbizoM0eVMwTd4cQlM3B1fP3+Bb9MsGn3u/rvcuWLXWq9BpzdZ7Lmu2FHl6rKte676G1rp50nqIYj3Ewqx1cqoVjGbZiliyAq/iXEvzQuHWf/uO3TWnTp+lg4eOVKxZu6HS2AfNbIArVjrV7Pb2LbtwMba0srLKqIVWXV1dcvzEyWJXN4+alUZ6r6s5hMBWeJ/4/wYG3tz8dVhKnmXg1ayf/C6kjfW9kD7d+1ipw/Zoctixl2ZtCabZ63aR43pvWrh2Z+Wfa7xKmOWeW0pd3deXuLhvLHZ3W1exBl6FWYtulMq7rKa1QCUWVZdMdMXEUhZXcqkvGMkFthcWFAuN49VOxue5ZHB+Vqx0RuxetXOXT2V8/I3CmpoannE0/LvV6sDXmpubVx4cElbC97mM9UbygxAC27x48/9Z7b5+q9ParZUY1Fe4ua4pd3VbU+7u6lHNK4Y9caAnBKBaOazyLPXHLrWDeCEUGLEGrcQeFmyQjGyfMcONxgp4NFc3zyp//6DqxMSke6WlZdU8nlHbscEFFlZlZWX5lStXi7Zs2VYNj2zU97kehBCYs7Pzf7qudvZa67EGQoFX0QgG1InF9AzUVFB5Z2GkNegSV1yNv15w716BwQkN3cHKtLT00sCg0Ao0eMJrmWLDV59agUFEXh6e66QHKRgHbKxstCy2rV47Ks6eO1+Snp5Reb/mPs88PrTuI9Iuw3VUHTp0pAg9o2pzEJYGRWAmCBsvi2zZitW0bt2myn37DhTExV2rysvL53tpHbIMC+lUV1VVF925k1B28ODhIp5g4gkaU+4OylAEZuKwQa9Y6QTP5sorQ4r37T9YFht7ubKgsDAPImj4RhkdBJyzMCsrq/D0mbPVgYEhRe4eayv59oK5CUuDIjAzgb0aexBxH9DFo2rb9p3FgUEhlUeOHq+5cycxD2O2zJrqmiK1Tloc2CPm5+enX79+o/Dw4WM1/gHBpZu2eJVC0CItcxWWBkVgZohGbMuWrxIC8Fyzvnzjpq2lGLdV+vkHUVT0Pjp48AjFXrqSFXc1Pk2bo0dPVkTvPUAQE99nrNq4aUuJh8e6CiEqeMuVq8xnjPUgFIGZOZrxGk88cFdSM07i7a5unpWgQhsnZ3cxtc7HqrqfqokVRVSNUQSm0AgWSkuQfVehIYrAFBT0iCIwBQU9oghMQUGPKAJTUNAjisAUFPSIIjAFBT2iCExBQY8oAlNQ0COKwBQU9IgiMAUFPaIITEFBjygCU1DQI4rAFBT0SK3AXFzctykCU1DQLUJgixcv/g9nJ7ctisAUFHSLEBgR/cPZ1f1XzzXrpQcpKCi0DSEwDk4uLvaKwBQUdEutwJxdXWcqAlNQ0C2KwBQU9IgiMAUFPaIITEFBjygCU1DQI07ObkcVgSko6AknZ/cERWAKCnpEEZiCgh5RBKbQYviNvvyabOVV2S1HrwLTVIY+aGkFy76rQVdGUt/wNMiO00b7O7rEqZ2/e8150vwiy9JlK2nZ8tVim4urh8AJx/APSPA+fkc9pyk7jz6QlTejq/rUJXoTGFeCm/sacvdYK2Jdwud0dfOUplufB10DV4guKsXF1bPBeTk9Tlt2bH1c3fRZPmukaT4INtSly1aJctnqtYP8/AMpIjKaDh46QmfOnqcLF2IEZ/H5wMHDFB4eSbu9fckD9sPfY1HKzqtrXNX5bJhnTyF82fEPC70IjAvZa9sOunT5Cl2/cZOuX9cxOOeJk6dFobJByK6Bf/Vj02Yviom5JL2G+Os3aO++A/zjdG0WGX+PhRQeEd3w/EiP89+UsfE189MLp06d0WP5nKG16zaKcpBdgzarndyEJ1q3fhPt3bufrly5ShkZmVRaWqr6MbBmQsG9Arp56zYdO3aCtmzdIX5ds6l6aS9cpnv2+FHc1WsNyw6fT50+26o8dwR6ERh3K7jlq6rS3+9wc+XzNTdVkWwsO3Z6U1lZ07+aWlRUTDtxjOqH6Rqf40FoBHbx4iX1GesC57+p87IBrFm7gbKystVH6z5kZGTR5s3bWuRRNAa5b99Buns3jWpq2v6TzpynY8dOCq+iD0Pnc169Gq9OrWGoqblPAQFBQuCy7z4M9CYwH19/KikpUWdd9yElJfWBAtu+YzcVFBSovyEPly5dIWeIpC0trkZgZ89eUJ+tLnD+HyQwNmZ9hZSUu7Rpk9cDBcb72dPHxV2jykrdNIj37xPduHGLNm7a2iKBtxQ+F9fpPXjMpgLXJ3cV9eVBW4vRCiw9PUMnAmOjCg4OE8fLztMchiyw9PTMB3owvo6tW7dTWlq6+lvycB+KqayspPLyclAh4oqKCrG9ucDdcO4K68qT8aTKseMnm023pKQU3fOdqE/dCbs9dKjAuGAKCgopMzOLMrPax+XLcaLy2iswDklJyaLv3trWVh8CYwPhblZ7y4jLZ/2GzU0aN187t/Q8fmkqcBf6TkKiGCsGh4SLyQxvjH+89/iSr28gHYexc9mVNtMNP3L0uChXTk92HS2Fz7EBHvFBjQGHI0eOi/Tam6Yu6FCBVVRUUnh4FParpnXbiyxtDa0RGIdDh47CGFs34cHH6lpg589fFMesgEHJ8twatNOtD19DRORe6XiLvVV8/A0hKP7NZdm5GT4Hl1lAYDDl5OSqv90w8PYt8JJNlUVL4e/zbGb9cT032NXV1eq/6kJ6RoYo36Yal46kwwUWERFd26K1F1naGpoSGPff2UtoB26tN2/Z1qpK4WvQtcB4CpzPycdo57e1yNJmeB97L1n6NTXVdOLEKXLGMfXrSXYehvdx121HM41ZZNTeNnXBNayCmD081jXyttxVjY29LLqr9QP3IAODQkUjIDtfR9LhAotEq9kRLUtTArt9+w4dPYZ+vPrv+iEGlaUxbNk5teHj9CEwfQ/SOX0f3wAqLW3ctYuPvy7y1NruMueVp+mrqxt7xPbmaTnq0sfHX4z/6ofMzEwxQXPz5m31lrrAHpjTa2ld6guzE9iNGzdpw8YtlJycqt5SF/j6AgKCRbdIdk5tjFVgXDaHjxxr1L3iv0NCIsTqDNn3moOvl++BsVfRDux52L7aUu8agZw+c059trpw4sRpUcZ8e0E7FKNH4uW1U6/l2BLMTmDswTj9gMAQaf89KTmF3NybnjypjzEL7PTps+oU6wLfM/TzC2jyupuDy4JXj3ADxhMRKampAv588tQZckcXry154tlA7roXFhapr1IVyiDkbdt30d9LV4h6zr93T71HFXh8xgJsrSfWNWYnsDt3EsR+XrnA90y0QxVEt3//IfF97XNqY2oCY+8jBIYxlex7LYFndtfAltieGP7MK25kxz4ILl8uh4MHj6ivsC5wt5DPywLiOrh0OU69py5kZmThetbXesGHgVkKzA0GzN2g3d4+VFTUsGXkwDNf3Go+SGTGLLAzki4Xh4iIqDZ1ETXwdXP+eHZRhWoWUnbsg+DvsUBS795VX11d4BlFnvzgOuDrDcd1a/dIKjBmCwuLbJNH1hVmKzC+Bj6G1zTKpqrPnb8gDKS51s+YBcY3bHlpkXZISEyidRs2i5nBh9nyM1xO3JXXvrGcl5eHBrBulQofx+PqtLuN75Fdu3Zd5ONh5cVsBaZpafkG811JxfC1+/oFCkOTnZ8xVoFx+jt3+SCP8oW8vHB3124fMXvHeeDjO9pANenxtWiHk2gUuYtf/5r4Os+dPa8+oi7w7ReeMV2xou23CdpDhwtMV/fBZOnWpyUC4+P4Wvnmt/a9FA43b94iz2aW+vB16ENgfE6NUbcVWboaVMe4UyK8VVMhP/8eXYm7Snv3HhCzg3w9XKba1yU7vy7gdHixdonWan5uFPb4+Ddq+PjaeIF1cXFDm+PAK1Ee1gJgsxeYiGFssiVD3HVU3SSVL/XhbcYoMIbPz9f4oMdReOVETk4eJSQk0vETp8jHJ6DWw9anJWm2Bq6/c+cuNOoexqOeZOsbOX0XXBd3cbUDz2Ru3KzbhcctpUMFxgbLRhWHlpEfOWgLfCOUH/x70HM/LRUYwwW/bduuRtfLgRcVb8CYRFY5olJ1LLDc3Dy6hjxevXqtUd5bwvXrN8QEBq9kb658GN4fFh7Z7FrC+oGNvRQeJBdjIJ6OP3z4mJgo8lyzrrZRYNorNq4bHlNpL7/i9PfuO9hkt53rnO/vaY+peazJkyLtmbxpKx0qMF0FLvhND3gUojUCY4Pgv3ngLwv8IJ+TS2PPoA+B6SKwx+G8P6jF5jzzRE5IaLgQdmsDGzKvW8zJzYXnvSgmJLgx4vO2x1twue0/cKjR84T8DGBzj+BwurzusaCwUP2NusArO5pbHK4vFIGpYaPnp3nZY2kHXqLD4wE+pv53jF1gDOeBr4c93sWLsWLs1Z7Ay5dYHOyBOH2nVho014ub+1pKSkpRn1EV2HvxQmjNu0Fk39Xk5Xp84+4+d4XFxE0TdaIvjFJgbASbdSwwUTmrXCgkJFy6wkM13as6rv53DFFg+fkFokFoicA08DWxN+PvnT5zlpKTU1r0uoCmQmrqXQoOCRP3qlrjNbj7xwt1tZ9ELykuIT+/QFGvsu8xXB+cVnAw12HjWy9Hj52oPUb2fX3QoQLjVigvLx+FnwbjSm8TPGDlscZ6eBs2Cln6TGsFxvD5eADPT/dqB+4OhYZF4Lx1RqsPgfG0Mt/P4X2y/D+ItLQMlM910VVqrnzkqDwAj1X4+gLR5Tt24qR4cDInJ6fR2OZBoaysnKKi97XYqPk4bsRiJStseMaT64aPkX1XA6fj6bleXK92yEE3eP2GTa1qeNpLhwqM++uHDx8V3Qd+2rZNbNlGG/F9njFqrrDbIjCGr527EvxgqHbg7iNPrmi+qw+B8ctm+Clj7rJJ8/8AtmzBd1E+/E6MBxljc3Ae+fqXwqPwOsLt23eJ574OoPt34XyM6MLxhMeDAk+g7PEJEA3Tg65HlP2uPY3KvqqqWrzBSnPzuzn4urlbylPzssCTOuIYHCu7Bl3ToQLjafqwsChRUNyKtJWWtMxtFRgXPJ+fu0nagT0w3+Tka9BUqK4FxlPTGuPWzndL4fPr0oD4erg8ud44Zk/CT0tv37FLLBzgNzpVom6bCnw/kZc8NVdvmnLnbpx24CHB+g1bavP1ILgMdkKoFWjQtcOtW3f0fiO/Ph0uMENaySH7HsMVxM8Zpac3XuHBlc2Vx+fnytS1wDpiJUd74DzztXEeuAw4Zm8ZGhZJ2ZJuGQdumPb4+IlylZ2TEWUOr80zhdohOzuHjh49IWZ5+TUFzXMKxx6n06fPUbFkDoAXFHAPRRFYO2mPwBjOQ6R4RL3xhAev3HZ2Ud33MTeBacOC4/ww/LKZpmYhjxw5VuthZOdhgYWGRqqPbhhYoExrQnPf4QdruXybuhZdogisCXg/i+eG5GlZnrbnGTLuMhmjwPjcfG18/byESIDPfF0urvLvPAiN0PbtPyg1bJ5ibypf/D1+9OTmzTvqo/UbeKKNb8koAmsH7RUYw60qdye018Nx4IcJPdeuF8cYk8DYqHiixt8/SNySCAoOFYSFRYjJlZY+zS2Dy4JnL2W3OXiFSVNvYuby4Lf1ynoL+gh8A5sbArYR7WvRNYrAHgBfKy/NknU2+IYqn+fMmcaruA1VYCwCXl3OKzdU7zrk9xxWCFHwjWZOt60tO5+bZ4hbIzDN3xcuxqiPrAssuOLiYvE0Mz+312rwPV4cLLu9wGsWOZ/69mKKwB4AH7du3Sbp4DsvP1+02CdPNp4SNlSBcdq7vH1RLo1vQ9y+naCa7XNqW/00J7Cmuoh8c5u/o/1KAA78bkde68iTSrvRk2gtO3buFrcW+N6gdmDhsRfXtxdTBNYC2HD4V0RkgX9c4vz5xq2voQqMz7seBs0vNtUOvARt2/bd4jkw2Xebgz0BlxOPTWVjMF6Jrzmu/nf4eniBrnbgNy57e/vRX3+vEOXI79JsPavx/eVNrjG9GBMr7ESfXkwRWAvg6/XgQbjk4T/uz7OxahuVoQqMjYm9RmJSsjrFusB54J8p4gmP1hodi4uv+9Zt+USFv+RtXZwGL8CViZ3vV7FNttdW2A54+RdPbGgHTnfT5pb9QEZb6XCB6eJGswZZ2hp0KTCGz8evjNZ+N19Toa0C08WNZg1NiYTzwq9Bk3kaHpv54to1s4qy79dHc618L4zvQ8m6h5xPXtqmLTC+Dl7Jz3ZRP/CYiZ/Da6r8WgOXAacrW/7G+ecVIrpIpyk6VGA6WSrFbNkmbgTzFHlzRqRLgfHx/ONuTb0sRju0VWDtXSqlgc/Bkwqy8uG8iKeFm1jqlAeR8bpL9kh/L10p8sFLnTTC5bJlAfI+/vFBnjTh65a944ND9N79UsHz3zduNl75npaeoVPPwtfPQpa9s/E27IHror2esik6VGDcYrR3sS/DC36Tk1KafehS1wJjeGzCws7Klq9YqB/aKrD2LvbVwOsm+Z0iTV0D5/88uqNNBc376Q8cOEyBQSG0a/ceIUqGfwCC3+LEy5rYM/Ci3qYC/4yRrKvH4uF33/Nsn3Y4dfoM9utuuRefhxsbXuGvHfja+Tk2bjBk320vHSow3Yb7zf56iD4ExnBLzvdQtB8G1A5tFZguQ0Rk069g4/xzA8WN1YMC3wfkp5h5EoTh9/s/KP8cOH/cIHFdaKfP5cjjTe1uKtfXrt2+2K/bbhvXxSH0npqa4eRj2moTzWG0AuMW9mEIjNNzRbdINuFRPxiCwHjleHOPya/EdfAq+eSUxq8Rb09g0VyOi2vytQXsvTZv2U5ZWY17Avy6BO03RukCrm9evSG7HcDbVL8ppvspe70JzNcvAIPXlk0ItC082IOpxhkNRc6/Z9UegTEifxh3yN5gpAmc/wcJTDZ7psvwoPdQsBFzOfGYmFvx5n5utyWBvUNmZrZ4sRFPejQ1huJyYW+iHXhqnhuFpsqtPXBeuc75toos8PXwMboWtl4ExpXGfXZuofklkfqAf2SuuTEYVy7/Sn5ycnKD7126dLndAmP4/DxWqH/u+nD+m2oR+Zq5vPkFPrLv6gq+ydqSsQVfDxsWe/wLF2MpLzdfTIA86CdleT8fl4txNf8oOb9FlxsOLtumype389Q8GzrfqK9/vfzGqPY+x9YcnDbffGbbqZ8uw9ejj3d26EVg5oA+WruHiUYUDNsCv3sweu8BOnvuvJg5vXYtnpLQWPFsId9K4Hfb834+TrMESkNLyoVngLW3dUSZNnV+2fXoAr0KTFNg+kKWpjZt/V5LkJ1bg+x4bWTf0yWyNB8Ef49Fwl6NvbQG9sa8OoJj/vVNzXY+rqWiqk/966yP7FhdI0uXkR3bXvQqMAUFc0cRmIKCHlEEpqCgRxSBKSjoEUVgCgp6RBGYgoIeUQSmoKBHFIEpKOgRRWAKCnpEEZiCgh7RCGy2IjAFBd2jEpib2xRFYAoKukcIzMnJrfuatRv1tqJYQcFcUQvMtaciMAUF3aMIrBm4PBTMHU/xdq22oghMQUGPSAXGsZOzW42Ts2sC4uNOLm4nzI3Vzm6nnV3cIpyd3dycnDxWYJy6ypxwcXFftdrZfTniqSgCy1XOblbmxkon15FObm7dly9f94K7u/szrcf9mf8fuEFDCdxuDgoAAAAASUVORK5CYII="""
pasta_base = os.path.dirname(os.path.abspath(__file__))
caminho_logo_padrao = os.path.join(pasta_base, "logo.png")
if not os.path.exists(caminho_logo_padrao):
    try:
        with open(caminho_logo_padrao, "wb") as f:
            f.write(base64.b64decode(B64_LOGO_PADRAO))
    except: pass

try:
    from google import genai
    from pydantic import BaseModel
    
    class RDC_Schema(BaseModel):
        DATA: str
        DISCIPLINA: str
        ENCARREGADO: str
        TURNO: str
        DDS: str
        TRANSCRICAO: str
        ATIVIDADE: str
        PROBLEMAS: str
        LOCAL: str
        AREA: str
        CALDEIRA: str
        
    class RDC_CC_Schema(BaseModel):
        LOCAL: str
        AREA: str
        DISCIPLINA: str
        ENCARREGADO: str
        CALDEIRA: str

except ImportError:
    pass

# --- SISTEMA DE INTERNACIONALIZAÇÃO (i18n) ---
TRANSLATIONS = {
    "Dashboard": "Dashboard",
    "Resumo Diário": "Daily Summary",
    "Emissão de RDC": "RDC Issuance",
    "Controle de C.C": "C.C Control",
    "Competição F1": "F1 Competition",
    "Leitor de RDC (IA)": "RDC Reader (AI)",
    "IA - Atualizador de C.C": "AI - C.C Updater",
    "RDC Digital": "Digital RDC",
    "Painel de Configurações": "Settings Panel",
    "Ver Usuários e Senhas": "View Users & Passwords",
    "Nome da Empresa/Site:": "Company/Site Name:",
    "Salvar Nome": "Save Name",
    "Backup Seguro": "Secure Backup",
    "Baixar Backup (.zip)": "Download Backup (.zip)",
    "Gestão de Usuários": "User Management",
    "Adicionar / Editar Usuário": "Add / Edit User",
    "Salvar Usuário": "Save User",
    "Sair (Logout)": "Logout",
    "Trocar Logo (PNG/JPG):": "Change Logo (PNG/JPG):",
    "Idioma / Language": "Language",
    "Modelo de IA (Gemini)": "AI Model (Gemini)",
    "Gerenciar Lista F1": "Manage F1 List",
    "Selecione a Data do Resumo:": "Select Summary Date:",
    "GERAR EXCEL": "GENERATE EXCEL",
    "GERAR PDF": "GENERATE PDF",
    "GERAR TODOS (.ZIP)": "GENERATE ALL (.ZIP)",
    "Bem-vindo(a)": "Welcome",
    "Condição Climática": "Weather Condition",
    "Horas Perdidas": "Lost Hours",
    "Observação / Justificativa": "Note / Justification",
    
    # NOVAS TRADUÇÕES
    "Sistema RDC & PDE": "RDC & PDE System",
    "ACESSO RESTRITO": "RESTRICTED ACCESS",
    "Usuário (Login):": "Username (Login):",
    "Senha:": "Password:",
    "Manter conectado": "Keep me logged in",
    "Entrar no Sistema": "Login to System",
    "Usuário ou senha incorretos. Verifique espaços em branco ou letras erradas.": "Invalid username or password. Check for typos or blank spaces.",
    "Controle Operacional de Efetivo": "Operational Manpower Control",
    "Carregar Arquivo Excel / CSV:": "Upload Excel / CSV File:",
    "Efetivo": "Manpower",
    "Encarregados": "Foremen",
    "% MOD Global": "% Direct Labor",
    "Funções": "Roles",
    "Centros de Custo": "Cost Centers",
    "Total Alocados": "Total Allocated",
    "Funções Distintas": "Distinct Roles",
    "Span of Control": "Span of Control",
    "Total de RDCs Entregues": "Total RDCs Delivered",
    "Identificação": "Identification",
    "Localização": "Location",
    "Atividades e Envio": "Activities & Submission",
    "Sincronização de RDCs (Nuvem)": "RDC Sync (Cloud)",
    "Puxar Dados Automáticos (Google Sheets)": "Fetch Auto Data (Google Sheets)",
    "Aguardando Base de Dados": "Waiting for Database",
    "O sistema está pronto.<br>Para iniciar a gestão, <b>arraste o arquivo de Efetivo (.csv ou .xlsx)</b><br>para a área de upload na barra lateral.": "System is ready.<br>To start managing, <b>drag and drop the Manpower file (.csv or .xlsx)</b><br>to the upload area in the sidebar.",
    "Desenvolvido por": "Developed by",
    "Análise de Gargalos": "Bottleneck Analysis",
}

def t(texto):
    """Traduz texto baseado no idioma selecionado no session_state."""
    if st.session_state.get("idioma", "Português") == "English":
        return TRANSLATIONS.get(texto, texto)
    return texto

from streamlit_gsheets import GSheetsConnection

# --- CONFIGURAÇÃO DA PÁGINA ---
# Força a criação do arquivo de tema Escuro automaticamente
config_dir = ".streamlit"
os.makedirs(config_dir, exist_ok=True)
config_path = os.path.join(config_dir, "config.toml")
if not os.path.exists(config_path):
    with open(config_path, "w", encoding="utf-8") as f:
        f.write('[theme]\nbase="dark"\nprimaryColor="#f39c12"\nbackgroundColor="#1e1e1e"\nsecondaryBackgroundColor="#2b2b2b"\ntextColor="#e0e4ea"\n')

caminho_nome_site = "nome_empresa.txt"
if os.path.exists(caminho_nome_site):
    with open(caminho_nome_site, "r", encoding="utf-8") as f:
        nome_site = f.read().strip()
    if not nome_site:
        nome_site = "ENESA Engenharia"
else:
    nome_site = "ENESA Engenharia"

st.set_page_config(page_title=f"Sistema RDC & PDE - {nome_site}", layout="wide", initial_sidebar_state="expanded")

# Injeção de CSS para ajustes de interface
st.markdown("""
    <style>
        /* Tipografia Moderna */
        @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap');
        
        html, body, [class*="css"] {
            font-family: 'Inter', sans-serif;
        }

        /* Glassmorphism na Sidebar */
        [data-testid="stSidebar"] {
            background-color: rgba(15, 23, 42, 0.5) !important;
            backdrop-filter: blur(12px) !important;
            border-right: 1px solid rgba(255, 255, 255, 0.05) !important;
        }

        /* Melhorar visual dos Métricas (Cards) */
        [data-testid="stMetric"] {
            background: rgba(30, 41, 59, 0.5) !important;
            border-radius: 12px !important;
            padding: 15px 20px !important;
            border: 1px solid rgba(255, 255, 255, 0.05) !important;
            box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1) !important;
            transition: all 0.3s ease !important;
        }
        [data-testid="stMetric"]:hover {
            transform: translateY(-2px) !important;
            box-shadow: 0 8px 15px rgba(0, 0, 0, 0.2) !important;
            border-color: rgba(14, 165, 233, 0.3) !important;
        }

        /* Botões Arredondados com Hover Vivo */
        .stButton button {
            border-radius: 8px !important;
            transition: all 0.2s ease !important;
            font-weight: 600 !important;
            border: 1px solid rgba(255, 255, 255, 0.1) !important;
        }
        .stButton button:hover {
            transform: scale(1.02) !important;
            box-shadow: 0 4px 12px rgba(14, 165, 233, 0.3) !important;
            border-color: rgba(14, 165, 233, 0.5) !important;
        }
        
        /* Headers com gradiente sutil */
        h1, h2, h3 {
            font-weight: 700 !important;
        }
        
        /* Ocultar barra de topo do Streamlit */
        header[data-testid="stHeader"] {
            background: transparent !important;
        }

        /* ================================================ */
        /* 3. TABELAS ZEBRADAS (linhas alternadas)          */
        /* ================================================ */
        [data-testid="stDataFrame"] table tbody tr:nth-child(even) {
            background-color: rgba(30, 41, 59, 0.4) !important;
        }
        [data-testid="stDataFrame"] table tbody tr:nth-child(odd) {
            background-color: rgba(15, 23, 42, 0.6) !important;
        }
        [data-testid="stDataFrame"] table tbody tr:hover {
            background-color: rgba(14, 165, 233, 0.12) !important;
            transition: background-color 0.2s ease !important;
        }
        [data-testid="stDataFrame"] table thead tr {
            background: linear-gradient(135deg, rgba(14, 165, 233, 0.2), rgba(139, 92, 246, 0.15)) !important;
            border-bottom: 2px solid rgba(14, 165, 233, 0.3) !important;
        }
        [data-testid="stDataFrame"] table thead th {
            font-weight: 700 !important;
            text-transform: uppercase !important;
            font-size: 0.75rem !important;
            letter-spacing: 0.5px !important;
            color: #e2e8f0 !important;
        }

        /* ================================================ */
        /* 4. TABS/MENUS COM ESTILO PREMIUM                 */
        /* ================================================ */
        .stTabs [data-baseweb="tab-list"] {
            gap: 4px !important;
            background: rgba(15, 23, 42, 0.6) !important;
            border-radius: 12px !important;
            padding: 4px !important;
            border: 1px solid rgba(255, 255, 255, 0.06) !important;
        }
        .stTabs [data-baseweb="tab"] {
            border-radius: 8px !important;
            padding: 8px 16px !important;
            font-weight: 600 !important;
            font-size: 0.85rem !important;
            color: #94a3b8 !important;
            transition: all 0.25s ease !important;
            border: none !important;
        }
        .stTabs [data-baseweb="tab"]:hover {
            background: rgba(14, 165, 233, 0.1) !important;
            color: #e2e8f0 !important;
        }
        .stTabs [aria-selected="true"] {
            background: linear-gradient(135deg, #0ea5e9, #3b82f6) !important;
            color: #ffffff !important;
            box-shadow: 0 4px 12px rgba(14, 165, 233, 0.35) !important;
        }
        .stTabs [data-baseweb="tab-highlight"] {
            display: none !important;
        }
        .stTabs [data-baseweb="tab-border"] {
            display: none !important;
        }

        /* ================================================ */
        /* 6. SPINNER DE CARREGAMENTO PREMIUM               */
        /* ================================================ */
        .stSpinner > div {
            border-radius: 12px !important;
            background: rgba(15, 23, 42, 0.8) !important;
            backdrop-filter: blur(8px) !important;
            border: 1px solid rgba(14, 165, 233, 0.2) !important;
            padding: 20px !important;
            box-shadow: 0 8px 32px rgba(0, 0, 0, 0.3) !important;
        }
        .stSpinner > div > div {
            border-top-color: #0ea5e9 !important;
        }

        /* Tentar forçar o carregamento da fonte oficial de ícones do Google */
        @import url('https://fonts.googleapis.com/css2?family=Material+Symbols+Rounded:opsz,wght,FILL,GRAD@24,400,0,0');
        
        /* Esconder o botão de collapse da sidebar para evitar texto quebrado */
        [data-testid="stSidebarCollapseButton"] {
            display: none !important;
        }
        
        /* Ocultar texto quebrado "arrow_down" do st.expander caso a fonte não carregue */
        summary .material-symbols-rounded,
        .st-emotion-cache-1t8fpt5 .material-symbols-rounded,
        [data-testid="stExpander"] .material-symbols-rounded {
            display: none !important;
            color: transparent !important;
        }
        
        /* Esconder o menu superior chato do Streamlit (Deploy, Rerun, etc) */
        header { visibility: hidden !important; display: none !important; }
        [data-testid="stHeader"] { display: none !important; }
        [data-testid="stToolbar"] { display: none !important; }
        #MainMenu { display: none !important; }
        footer { display: none !important; }
        .stApp > header { display: none !important; }
        
        /* Remover a linha colorida no topo e subir o layout */
        [data-testid="stDecoration"] { display: none !important; }
        .stApp { margin-top: -60px; }
        
        /* Travar a largura da barra lateral (esconder o arrastador) */
        [data-testid="stSidebarResizer"] {
            display: none !important;
        }
    </style>
""", unsafe_allow_html=True)

# Paleta Premium Dark Mode (Glassmorphism & Neon Subtle)
cor_fundo = "#0f172a" # Slate 900 (azul-marinho elegante)
cor_fundo_grad = "linear-gradient(135deg, #0f172a 0%, #1e293b 100%)"
cor_card = "rgba(30, 41, 59, 0.55)" # Slate 800 glass
cor_borda = "rgba(255, 255, 255, 0.08)"
cor_texto = "#f8fafc"
cor_texto_sub = "#94a3b8"
cor_azul = "#3b82f6"
cor_azul_hover = "#2563eb"
cor_destaque = "#0ea5e9" # Light blue neon
cor_verde = "#10b981"

st.markdown(f"""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&family=Outfit:wght@500;600;700&display=swap');
    
    @keyframes fadeIn {{
        from {{ opacity: 0; transform: translateY(15px); }}
        to {{ opacity: 1; transform: translateY(0); }}
    }}
    
    /* === BASE === */
    .stApp {{
        font-family: 'Inter', sans-serif !important;
        background: {cor_fundo_grad};
        background-attachment: fixed;
        color: {cor_texto};
    }}
    
    .block-container {{
        animation: fadeIn 0.8s cubic-bezier(0.16, 1, 0.3, 1);
        max-width: 1200px;
        padding-top: 1rem;
    }}

    /* === CABEÇALHO === */
    .enesa-header {{
        text-align: center;
        margin-top: -30px;
        margin-bottom: 28px;
        padding: 32px 20px;
        background: {cor_card};
        backdrop-filter: blur(12px);
        -webkit-backdrop-filter: blur(12px);
        border-radius: 16px;
        border: 1px solid {cor_borda};
        border-left: 5px solid {cor_destaque};
        box-shadow: 0 8px 32px rgba(0,0,0,0.3);
        transition: transform 0.3s ease;
    }}
    .enesa-header:hover {{
        transform: translateY(-2px);
    }}
    
    /* === FUNDO === */
    @keyframes panBackground {{
        0% {{ background-position: 0 0; }}
        100% {{ background-position: 100vw 100vh; }}
    }}
    .stApp {{
        background: radial-gradient(ellipse at top, #1e293b 0%, {cor_fundo} 60%, #0f172a 100%) !important;
        background-image: 
            linear-gradient(rgba(148, 163, 184, 0.03) 1px, transparent 1px),
            linear-gradient(90deg, rgba(148, 163, 184, 0.03) 1px, transparent 1px) !important;
        background-size: 50px 50px !important;
        animation: panBackground 200s linear infinite !important;
    }}
    
    /* === DATAFRAME GLOW === */
    .stDataFrame {{
        border-radius: 12px;
        overflow: hidden;
        border: 1px solid rgba(14, 165, 233, 0.1);
        box-shadow: 0 4px 15px rgba(0,0,0,0.2);
        transition: all 0.3s ease;
    }}
    .stDataFrame:hover {{
        border-color: rgba(14, 165, 233, 0.4);
        box-shadow: 0 8px 30px rgba(14, 165, 233, 0.15);
    }}
    
    /* === GLOWING ACCENT PARA TABS === */
    .stTabs [data-baseweb="tab-list"]::before {{
        content: '';
        position: absolute;
        top: 0; left: 0; right: 0; bottom: 0;
        background: linear-gradient(90deg, transparent, rgba(14, 165, 233, 0.1), transparent);
        pointer-events: none;
        z-index: 0;
        animation: panBackground 5s linear infinite;
    }}
    
    /* === TIPOGRAFIA === */
    h1, h2, h3 {{
        color: {cor_texto} !important;
        font-family: 'Outfit', sans-serif !important;
        font-weight: 600 !important;
        letter-spacing: -0.5px;
    }}
    p, span, label, div {{
        font-family: 'Inter', sans-serif !important;
    }}
    
    /* === BOTÃO PRINCIPAL === */
    div.stButton > button[data-baseweb="button"] {{
        background: linear-gradient(135deg, {cor_azul} 0%, {cor_destaque} 100%);
        color: white;
        border: none;
        font-family: 'Outfit', sans-serif !important;
        font-weight: 600;
        font-size: 0.95rem;
        letter-spacing: 0.5px;
        border-radius: 10px;
        padding: 12px 24px;
        transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1);
        box-shadow: 0 4px 15px rgba(14, 165, 233, 0.2);
    }}
    div.stButton > button[data-baseweb="button"]:hover {{
        transform: translateY(-2px) scale(1.02);
        box-shadow: 0 8px 25px rgba(14, 165, 233, 0.4);
    }}
    div.stButton > button[data-baseweb="button"]:active {{
        transform: translateY(1px) scale(0.98);
    }}
    
    /* === SIDEBAR === */
    [data-testid="stSidebar"] {{
        background-color: rgba(15, 23, 42, 0.8) !important;
        backdrop-filter: blur(20px);
        -webkit-backdrop-filter: blur(20px);
        border-right: 1px solid rgba(14, 165, 233, 0.2);
        box-shadow: 4px 0 25px rgba(14, 165, 233, 0.15);
    }}
    [data-testid="stSidebar"] h1,
    [data-testid="stSidebar"] h2,
    [data-testid="stSidebar"] h3 {{
        color: {cor_texto} !important;
        font-family: 'Outfit', sans-serif !important;
        text-shadow: 0 0 15px rgba(14, 165, 233, 0.5);
    }}
    
    /* === ADMIN AVATAR GLOW === */
    .admin-avatar {{
        width: 60px; height: 60px;
        border-radius: 50%;
        background: {cor_azul};
        display: flex; align-items: center; justify-content: center;
        margin: 20px auto;
        box-shadow: 0 0 20px {cor_azul};
        color: white; font-weight: bold; font-size: 20px;
        border: 2px solid white;
    }}
    
    /* === INPUTS & CONTAINERS === */
    .stTextInput input, .stSelectbox > div > div, .stTextArea textarea {{
        border-radius: 10px !important;
        border: 1px solid {cor_borda} !important;
        background-color: rgba(15, 23, 42, 0.5) !important;
        color: {cor_texto} !important;
        transition: all 0.3s ease !important;
        backdrop-filter: blur(4px);
    }}
    .stTextInput input:focus, .stSelectbox > div > div:focus-within, .stTextArea textarea:focus {{
        border-color: {cor_destaque} !important;
        box-shadow: 0 0 0 3px rgba(14, 165, 233, 0.2) !important;
        background-color: rgba(30, 41, 59, 0.8) !important;
    }}
    
    /* === ABAS === */
    .stTabs [data-baseweb="tab-list"] {{
        gap: 8px;
        background-color: {cor_card};
        backdrop-filter: blur(10px);
        border-radius: 12px;
        padding: 6px;
        border: 1px solid {cor_borda};
        box-shadow: 0 4px 6px rgba(0,0,0,0.1);
    }}
    .stTabs [data-baseweb="tab"] {{
        border-radius: 8px;
        font-family: 'Outfit', sans-serif !important;
        font-weight: 500;
        font-size: 0.95rem;
        padding: 10px 20px;
        color: {cor_texto_sub};
        transition: all 0.3s ease;
        border: 1px solid transparent;
    }}
    .stTabs [data-baseweb="tab"]:hover {{
        color: {cor_texto};
        background-color: rgba(255, 255, 255, 0.05);
    }}
    .stTabs [aria-selected="true"] {{
        background: linear-gradient(135deg, {cor_azul} 0%, {cor_destaque} 100%) !important;
        color: white !important;
        box-shadow: 0 4px 15px rgba(14, 165, 233, 0.3);
        border: 1px solid rgba(255,255,255,0.1);
    }}
    
    /* === MÉTRICAS (GLASS CARDS) === */
    [data-testid="stMetric"] {{
        background: {cor_card};
        backdrop-filter: blur(12px);
        -webkit-backdrop-filter: blur(12px);
        border: 1px solid {cor_borda};
        border-radius: 16px;
        padding: 20px;
        box-shadow: 0 4px 15px rgba(0,0,0,0.15);
        transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1);
        position: relative;
        overflow: hidden;
    }}
    [data-testid="stMetric"]::before {{
        content: '';
        position: absolute;
        top: 0; left: 0; width: 100%; height: 100%;
        background: linear-gradient(135deg, rgba(255,255,255,0.05) 0%, rgba(255,255,255,0) 100%);
        pointer-events: none;
    }}
    [data-testid="stMetric"]:hover {{
        transform: translateY(-5px);
        box-shadow: 0 12px 25px rgba(0,0,0,0.3);
        border-color: rgba(14, 165, 233, 0.3);
    }}
    [data-testid="stMetricLabel"] {{
        color: {cor_texto_sub} !important;
        font-family: 'Outfit', sans-serif !important;
        font-weight: 500 !important;
        letter-spacing: 0.5px;
    }}
    [data-testid="stMetricValue"] {{
        color: {cor_texto} !important;
        font-family: 'Outfit', sans-serif !important;
        font-weight: 700 !important;
        font-size: 2.2rem !important;
        background: linear-gradient(to right, #fff, #94a3b8);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
    }}
    
    /* === DATAFRAMES === */
    .stDataFrame {{
        border-radius: 12px;
        overflow: hidden;
        border: 1px solid {cor_borda};
        box-shadow: 0 4px 15px rgba(0,0,0,0.1);
    }}
    
    /* === FILE UPLOADER === */
    [data-testid="stFileUploader"] {{
        border-radius: 12px;
        background: {cor_card};
        border: 1px dashed rgba(14, 165, 233, 0.4);
        padding: 10px;
        transition: all 0.3s ease;
    }}
    [data-testid="stFileUploader"]:hover {{
        border-color: {cor_destaque};
        background: rgba(30, 41, 59, 0.6);
    }}
    [data-testid="stFileUploaderDropzone"] button {{
        color: transparent !important;
        position: relative;
        min-width: 140px !important;
    }}
    [data-testid="stFileUploaderDropzone"] button::after {{
        content: "Procurar arquivos";
        color: white;
        position: absolute;
        left: 50%;
        top: 50%;
        transform: translate(-50%, -50%);
        font-family: 'Outfit', sans-serif !important;
        font-weight: 600;
        font-size: 14px;
        white-space: nowrap;
    }}
    
    /* === EXPANDERS === */
    [data-testid="stExpander"] {{
        background: {cor_card};
        backdrop-filter: blur(10px);
        border-radius: 12px !important;
        border: 1px solid {cor_borda} !important;
        overflow: hidden;
        transition: all 0.3s ease;
    }}
    [data-testid="stExpander"]:hover {{
        border-color: rgba(14, 165, 233, 0.3) !important;
        box-shadow: 0 4px 15px rgba(14, 165, 233, 0.1);
    }}
    [data-testid="stExpander"] summary {{
        font-family: 'Outfit', sans-serif !important;
        font-weight: 500;
    }}
    
    /* === ALERTS (SUCCESS, ERROR, WARNING, INFO) === */
    [data-testid="stAlert"] {{
        background: rgba(15, 23, 42, 0.6) !important;
        backdrop-filter: blur(12px) !important;
        border-radius: 12px !important;
        border: 1px solid rgba(255,255,255,0.05) !important;
        color: {cor_texto} !important;
        box-shadow: 0 4px 20px rgba(0,0,0,0.2) !important;
        overflow: hidden;
        position: relative;
    }}
    [data-testid="stAlert"]::before {{
        content: '';
        position: absolute;
        top: 0; left: 0; width: 4px; height: 100%;
        background: rgba(255,255,255,0.2);
    }}
    /* Success Alert Green Glow */
    [data-testid="stAlert"]:has(.st-emotion-cache-1cvow4s) {{
        border-left: 1px solid rgba(16, 185, 129, 0.5) !important;
    }}
    [data-testid="stAlert"]:has(.st-emotion-cache-1cvow4s)::before {{
        background: {cor_verde};
        box-shadow: 0 0 15px {cor_verde};
    }}
    /* Error Alert Red Glow */
    [data-testid="stAlert"]:has(.st-emotion-cache-1eq3m7k) {{
        border-left: 1px solid rgba(239, 68, 68, 0.5) !important;
    }}
    [data-testid="stAlert"]:has(.st-emotion-cache-1eq3m7k)::before {{
        background: #ef4444;
        box-shadow: 0 0 15px #ef4444;
    }}
    /* Warning Alert Yellow/Orange Glow */
    [data-testid="stAlert"]:has(.st-emotion-cache-1gtyq0b) {{
        border-left: 1px solid rgba(245, 158, 11, 0.5) !important;
    }}
    [data-testid="stAlert"]:has(.st-emotion-cache-1gtyq0b)::before {{
        background: #f59e0b;
        box-shadow: 0 0 15px #f59e0b;
    }}
    
    /* === DOWNLOAD BUTTONS (PREMIUM) === */
    .stDownloadButton > button {{
        background: linear-gradient(135deg, {cor_verde} 0%, #059669 100%) !important;
        color: white !important;
        border: none !important;
        border-radius: 10px !important;
        font-family: 'Outfit', sans-serif !important;
        font-weight: 600 !important;
        font-size: 0.9rem !important;
        letter-spacing: 0.5px;
        transition: all 0.3s ease !important;
        box-shadow: 0 4px 15px rgba(16, 185, 129, 0.25) !important;
    }}
    .stDownloadButton > button:hover {{
        transform: translateY(-2px) !important;
        box-shadow: 0 8px 25px rgba(16, 185, 129, 0.4) !important;
    }}
    
    /* === TOGGLE PREMIUM === */
    .stToggle label {{
        font-family: 'Inter', sans-serif !important;
        font-weight: 500 !important;
    }}
    
    /* === FORMS (GLASS) === */
    [data-testid="stForm"] {{
        background: {cor_card} !important;
        backdrop-filter: blur(12px) !important;
        border-radius: 16px !important;
        border: 1px solid {cor_borda} !important;
        padding: 20px !important;
        box-shadow: 0 4px 20px rgba(0,0,0,0.2) !important;
    }}
    
    /* === SEPARADORES PREMIUM === */
    hr {{
        border: none !important;
        height: 1px !important;
        background: linear-gradient(90deg, transparent 0%, rgba(14, 165, 233, 0.3) 50%, transparent 100%) !important;
        margin: 16px 0 !important;
    }}
    
    /* === SELECTBOX PREMIUM === */
    [data-testid="stSelectbox"] label {{
        font-family: 'Outfit', sans-serif !important;
        font-weight: 500 !important;
        color: {cor_texto_sub} !important;
        letter-spacing: 0.3px;
    }}
    
    /* === NUMBER INPUT === */
    .stNumberInput {{
        border-radius: 10px !important;
    }}
    
    /* === CAPTIONS (subtle) === */
    .stCaption {{
        font-style: italic !important;
        opacity: 0.7;
    }}
    
    /* === ANIMATED GLOW for primary buttons === */
    @keyframes pulseGlow {{
        0% {{ box-shadow: 0 0 5px rgba(14, 165, 233, 0.4); }}
        50% {{ box-shadow: 0 0 20px rgba(14, 165, 233, 0.6); }}
        100% {{ box-shadow: 0 0 5px rgba(14, 165, 233, 0.4); }}
    }}
    
    div.stButton > button[kind="primary"] {{
        animation: pulseGlow 2s ease-in-out infinite;
    }}
    
    /* === MULTISELECT PREMIUM === */
    [data-testid="stMultiSelect"] {{
        font-family: 'Inter', sans-serif !important;
    }}
    [data-testid="stMultiSelect"] span[data-baseweb="tag"] {{
        background: linear-gradient(135deg, {cor_azul}, {cor_destaque}) !important;
        border-radius: 8px !important;
        color: white !important;
    }}
    
    /* === STATUS BARS === */
    .stProgress > div > div {{
        background: linear-gradient(90deg, {cor_azul}, {cor_destaque}) !important;
        border-radius: 10px !important;
    }}
    
    /* === MARCA D'ÁGUA EDSON GARCIA === */
    .watermark-edson {{
        position: fixed;
        top: 20px;
        right: 30px;
        z-index: 999999;
        font-family: 'Outfit', sans-serif;
        font-size: 13px;
        font-weight: 700;
        background: linear-gradient(135deg, rgba(255,255,255,0.4), rgba(14, 165, 233, 0.4));
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        letter-spacing: 3px;
        text-transform: uppercase;
        pointer-events: none;
        user-select: none;
        text-shadow: 0 0 20px rgba(14, 165, 233, 0.2);
        display: flex;
        align-items: center;
        gap: 8px;
        opacity: 0.8;
        animation: floatGlow 4s ease-in-out infinite;
    }}
    .watermark-edson::before {{
        content: '✦';
        font-size: 14px;
        -webkit-text-fill-color: rgba(14, 165, 233, 0.5);
    }}
    
    @keyframes floatGlow {{
        0% {{ transform: translateY(0px); opacity: 0.7; }}
        50% {{ transform: translateY(-3px); opacity: 1; }}
        100% {{ transform: translateY(0px); opacity: 0.7; }}
    }}
    
    /* ================================================ */
    /* 📱 DESIGN MOBILE - EXPERIÊNCIA DE APP            */
    /* ================================================ */
    
    /* === TABLETS (até 1024px) === */
    @media screen and (max-width: 1024px) {{
        .block-container {{
            max-width: 100% !important;
            padding-left: 1rem !important;
            padding-right: 1rem !important;
        }}
        
        /* Sidebar vira overlay em tablet */
        [data-testid="stSidebar"] {{
            width: 280px !important;
            min-width: 280px !important;
        }}
    }}
    
    /* === CELULARES (até 768px) === */
    @media screen and (max-width: 768px) {{
        
        /* --- Layout Geral --- */
        .block-container {{
            max-width: 100% !important;
            padding: 0.5rem 0.75rem !important;
            padding-top: 0.5rem !important;
        }}
        
        .stApp {{
            margin-top: -30px !important;
        }}
        
        /* --- Esconder Sidebar por padrão no celular --- */
        [data-testid="stSidebar"] {{
            width: 260px !important;
            min-width: 260px !important;
            transform: translateX(-100%);
            transition: transform 0.3s ease !important;
            z-index: 999 !important;
        }}
        [data-testid="stSidebar"][aria-expanded="true"] {{
            transform: translateX(0) !important;
        }}
        
        /* --- Cabeçalho Compacto --- */
        .enesa-header {{
            padding: 18px 14px !important;
            margin-top: -15px !important;
            margin-bottom: 16px !important;
            border-radius: 12px !important;
        }}
        .enesa-header h1 {{
            font-size: 1.3rem !important;
        }}
        .enesa-header p {{
            font-size: 0.8rem !important;
        }}
        
        /* --- Tipografia Mobile --- */
        h1 {{
            font-size: 1.4rem !important;
        }}
        h2 {{
            font-size: 1.2rem !important;
        }}
        h3 {{
            font-size: 1.05rem !important;
        }}
        p, span, label, div {{
            font-size: 0.9rem !important;
        }}
        
        /* --- Tabs: Scroll Horizontal (Swipe) --- */
        .stTabs [data-baseweb="tab-list"] {{
            overflow-x: auto !important;
            overflow-y: hidden !important;
            flex-wrap: nowrap !important;
            -webkit-overflow-scrolling: touch !important;
            scroll-snap-type: x mandatory !important;
            gap: 4px !important;
            padding: 4px !important;
            border-radius: 10px !important;
            scrollbar-width: none !important;
            -ms-overflow-style: none !important;
        }}
        .stTabs [data-baseweb="tab-list"]::-webkit-scrollbar {{
            display: none !important;
        }}
        .stTabs [data-baseweb="tab"] {{
            flex-shrink: 0 !important;
            scroll-snap-align: start !important;
            padding: 8px 14px !important;
            font-size: 0.78rem !important;
            min-width: max-content !important;
            white-space: nowrap !important;
        }}
        
        /* --- Botões Grandes e Touch-Friendly --- */
        .stButton button {{
            min-height: 52px !important;
            font-size: 0.95rem !important;
            border-radius: 12px !important;
            padding: 14px 20px !important;
            width: 100% !important;
            touch-action: manipulation !important;
            -webkit-tap-highlight-color: transparent !important;
        }}
        .stButton button:active {{
            transform: scale(0.97) !important;
            transition: transform 0.1s ease !important;
        }}
        
        /* --- Download Buttons --- */
        .stDownloadButton button {{
            min-height: 52px !important;
            font-size: 0.95rem !important;
            border-radius: 12px !important;
            padding: 14px 20px !important;
            width: 100% !important;
        }}
        
        /* --- Inputs Maiores (Dedo-Friendly) --- */
        .stTextInput input {{
            min-height: 48px !important;
            font-size: 16px !important;
            border-radius: 10px !important;
            padding: 12px 14px !important;
        }}
        .stSelectbox > div > div {{
            min-height: 48px !important;
            font-size: 16px !important;
        }}
        .stTextArea textarea {{
            font-size: 16px !important;
            min-height: 100px !important;
        }}
        .stDateInput input {{
            min-height: 48px !important;
            font-size: 16px !important;
        }}
        .stNumberInput input {{
            min-height: 48px !important;
            font-size: 16px !important;
        }}
        
        /* --- Métricas: Cards Compactos --- */
        [data-testid="stMetric"] {{
            padding: 12px 14px !important;
            border-radius: 12px !important;
        }}
        [data-testid="stMetricValue"] {{
            font-size: 1.6rem !important;
        }}
        [data-testid="stMetricLabel"] {{
            font-size: 0.75rem !important;
        }}
        
        /* --- Colunas: Empilhar Verticalmente --- */
        [data-testid="stHorizontalBlock"] {{
            flex-wrap: wrap !important;
            gap: 8px !important;
        }}
        [data-testid="stHorizontalBlock"] > [data-testid="stVerticalBlock"] {{
            flex: 1 1 100% !important;
            min-width: 100% !important;
        }}
        
        /* --- Tabelas: Scroll Horizontal --- */
        .stDataFrame {{
            border-radius: 10px !important;
            overflow-x: auto !important;
            -webkit-overflow-scrolling: touch !important;
        }}
        [data-testid="stDataFrame"] table thead th {{
            font-size: 0.65rem !important;
            padding: 6px 8px !important;
        }}
        [data-testid="stDataFrame"] table tbody td {{
            font-size: 0.75rem !important;
            padding: 6px 8px !important;
        }}
        
        /* --- Expanders Touch-Friendly --- */
        [data-testid="stExpander"] {{
            border-radius: 10px !important;
        }}
        [data-testid="stExpander"] summary {{
            min-height: 48px !important;
            padding: 12px 16px !important;
            font-size: 0.9rem !important;
            display: flex !important;
            align-items: center !important;
        }}
        
        /* --- File Uploader Maior --- */
        [data-testid="stFileUploader"] {{
            border-radius: 12px !important;
            padding: 16px !important;
        }}
        [data-testid="stFileUploaderDropzone"] {{
            min-height: 80px !important;
        }}
        
        /* --- Alerts Compactos --- */
        [data-testid="stAlert"] {{
            border-radius: 10px !important;
            padding: 12px !important;
            font-size: 0.85rem !important;
        }}
        
        /* --- Forms dentro de Expanders --- */
        .stForm {{
            padding: 8px !important;
        }}
        
        /* --- Popover Mobile --- */
        [data-testid="stPopover"] {{
            width: 90vw !important;
            max-width: 90vw !important;
        }}
        
        /* --- Watermark menor no mobile --- */
        .watermark-edson {{
            font-size: 8px !important;
            bottom: 4px !important;
            right: 4px !important;
        }}
        
        /* --- Toast/Notificações no mobile --- */
        [data-testid="stToast"] {{
            bottom: 70px !important;
            right: 10px !important;
            left: 10px !important;
            max-width: calc(100vw - 20px) !important;
        }}
        
        /* --- Scrollbar elegante no mobile --- */
        * {{
            scrollbar-width: thin !important;
            scrollbar-color: rgba(14, 165, 233, 0.3) transparent !important;
        }}
        *::-webkit-scrollbar {{
            width: 4px !important;
            height: 4px !important;
        }}
        *::-webkit-scrollbar-thumb {{
            background: rgba(14, 165, 233, 0.4) !important;
            border-radius: 4px !important;
        }}
        *::-webkit-scrollbar-track {{
            background: transparent !important;
        }}
    }}
    
    /* === CELULARES PEQUENOS (até 480px) === */
    @media screen and (max-width: 480px) {{
        .block-container {{
            padding: 0.3rem 0.5rem !important;
        }}
        
        .enesa-header {{
            padding: 14px 10px !important;
            margin-bottom: 12px !important;
        }}
        .enesa-header h1 {{
            font-size: 1.1rem !important;
        }}
        
        h1 {{
            font-size: 1.2rem !important;
        }}
        h2 {{
            font-size: 1.05rem !important;
        }}
        h3 {{
            font-size: 0.95rem !important;
        }}
        
        [data-testid="stMetricValue"] {{
            font-size: 1.3rem !important;
        }}
        
        .stTabs [data-baseweb="tab"] {{
            padding: 6px 10px !important;
            font-size: 0.72rem !important;
        }}
        
        .stButton button {{
            min-height: 48px !important;
            font-size: 0.88rem !important;
            padding: 12px 16px !important;
        }}
    }}
    
    /* === MODO PAISAGEM NO CELULAR === */
    @media screen and (max-height: 500px) and (orientation: landscape) {{
        .block-container {{
            padding-top: 0.3rem !important;
        }}
        .enesa-header {{
            padding: 10px !important;
            margin-bottom: 8px !important;
        }}
    }}
    
    /* ================================================ */
    /* 🌊 GLASSMORPHISM EXTREMO - VIDRO FOSCO PREMIUM   */
    /* ================================================ */
    
    /* Cards de Métrica com Vidro Profundo */
    [data-testid="stMetric"] {{
        background: rgba(15, 23, 42, 0.35) !important;
        backdrop-filter: blur(20px) saturate(180%) !important;
        -webkit-backdrop-filter: blur(20px) saturate(180%) !important;
        border: 1px solid rgba(255, 255, 255, 0.12) !important;
        box-shadow: 
            0 8px 32px rgba(0, 0, 0, 0.25),
            inset 0 1px 0 rgba(255, 255, 255, 0.1),
            inset 0 -1px 0 rgba(0, 0, 0, 0.1) !important;
    }}
    
    /* Tabs com Vidro */
    .stTabs [data-baseweb="tab-list"] {{
        background: rgba(15, 23, 42, 0.3) !important;
        backdrop-filter: blur(16px) saturate(160%) !important;
        -webkit-backdrop-filter: blur(16px) saturate(160%) !important;
        border: 1px solid rgba(255, 255, 255, 0.08) !important;
        box-shadow: 0 4px 24px rgba(0, 0, 0, 0.2), inset 0 1px 0 rgba(255, 255, 255, 0.06) !important;
    }}
    
    /* Expanders com Vidro */
    [data-testid="stExpander"] {{
        background: rgba(15, 23, 42, 0.3) !important;
        backdrop-filter: blur(16px) saturate(150%) !important;
        -webkit-backdrop-filter: blur(16px) saturate(150%) !important;
        border: 1px solid rgba(255, 255, 255, 0.08) !important;
        box-shadow: 0 4px 20px rgba(0, 0, 0, 0.15), inset 0 1px 0 rgba(255, 255, 255, 0.05) !important;
    }}
    
    /* Alerts com Vidro */
    [data-testid="stAlert"] {{
        background: rgba(15, 23, 42, 0.35) !important;
        backdrop-filter: blur(16px) saturate(150%) !important;
        -webkit-backdrop-filter: blur(16px) saturate(150%) !important;
        border: 1px solid rgba(255, 255, 255, 0.08) !important;
        box-shadow: 0 4px 20px rgba(0, 0, 0, 0.2), inset 0 1px 0 rgba(255, 255, 255, 0.06) !important;
    }}
    
    /* Inputs com Vidro */
    .stTextInput input, .stSelectbox > div > div, .stTextArea textarea,
    .stDateInput input, .stNumberInput input {{
        background: rgba(15, 23, 42, 0.35) !important;
        backdrop-filter: blur(12px) !important;
        -webkit-backdrop-filter: blur(12px) !important;
        border: 1px solid rgba(255, 255, 255, 0.1) !important;
        box-shadow: inset 0 1px 0 rgba(255, 255, 255, 0.05) !important;
    }}
    .stTextInput input:focus, .stSelectbox > div > div:focus-within, .stTextArea textarea:focus {{
        border-color: rgba(14, 165, 233, 0.5) !important;
        box-shadow: 0 0 0 3px rgba(14, 165, 233, 0.15), 0 0 20px rgba(14, 165, 233, 0.1) !important;
    }}
    
    /* DataFrames com Vidro */
    .stDataFrame {{
        background: rgba(15, 23, 42, 0.25) !important;
        backdrop-filter: blur(12px) !important;
        -webkit-backdrop-filter: blur(12px) !important;
        border: 1px solid rgba(255, 255, 255, 0.08) !important;
        box-shadow: 0 8px 32px rgba(0, 0, 0, 0.2), inset 0 1px 0 rgba(255, 255, 255, 0.04) !important;
    }}
    
    /* Sidebar com Vidro Profundo */
    [data-testid="stSidebar"] {{
        background: rgba(10, 15, 30, 0.6) !important;
        backdrop-filter: blur(24px) saturate(200%) !important;
        -webkit-backdrop-filter: blur(24px) saturate(200%) !important;
        border-right: 1px solid rgba(14, 165, 233, 0.15) !important;
        box-shadow: 4px 0 30px rgba(0, 0, 0, 0.3), inset -1px 0 0 rgba(255, 255, 255, 0.04) !important;
    }}
    
    /* File Uploader com Vidro */
    [data-testid="stFileUploader"] {{
        background: rgba(15, 23, 42, 0.3) !important;
        backdrop-filter: blur(12px) !important;
        -webkit-backdrop-filter: blur(12px) !important;
        border: 1px dashed rgba(14, 165, 233, 0.35) !important;
        box-shadow: inset 0 1px 0 rgba(255, 255, 255, 0.04) !important;
    }}
    
    /* Botões com Vidro e Brilho */
    div.stButton > button[data-baseweb="button"] {{
        background: linear-gradient(135deg, rgba(59, 130, 246, 0.85) 0%, rgba(14, 165, 233, 0.85) 100%) !important;
        backdrop-filter: blur(8px) !important;
        -webkit-backdrop-filter: blur(8px) !important;
        border: 1px solid rgba(255, 255, 255, 0.15) !important;
        box-shadow: 0 4px 15px rgba(14, 165, 233, 0.25), inset 0 1px 0 rgba(255, 255, 255, 0.2) !important;
    }}
    div.stButton > button[data-baseweb="button"]:hover {{
        box-shadow: 0 8px 30px rgba(14, 165, 233, 0.4), inset 0 1px 0 rgba(255, 255, 255, 0.25), 0 0 40px rgba(14, 165, 233, 0.15) !important;
    }}
    
    /* ================================================ */
    /* ✨ MICRO-ANIMAÇÕES E TRANSIÇÕES SUAVES           */
    /* ================================================ */
    
    /* Fade-Up: Elementos aparecem flutuando */
    @keyframes fadeUpIn {{
        from {{
            opacity: 0;
            transform: translateY(20px);
        }}
        to {{
            opacity: 1;
            transform: translateY(0);
        }}
    }}
    
    @keyframes fadeUpInDelayed {{
        0% {{
            opacity: 0;
            transform: translateY(25px);
        }}
        30% {{
            opacity: 0;
            transform: translateY(25px);
        }}
        100% {{
            opacity: 1;
            transform: translateY(0);
        }}
    }}
    
    /* Pulse Glow para botões */
    @keyframes pulseGlow {{
        0% {{
            box-shadow: 0 4px 15px rgba(14, 165, 233, 0.2);
        }}
        50% {{
            box-shadow: 0 4px 25px rgba(14, 165, 233, 0.45), 0 0 40px rgba(14, 165, 233, 0.1);
        }}
        100% {{
            box-shadow: 0 4px 15px rgba(14, 165, 233, 0.2);
        }}
    }}
    
    @keyframes shimmer {{
        0% {{
            background-position: -200% 0;
        }}
        100% {{
            background-position: 200% 0;
        }}
    }}
    
    /* Aplicar fade-up nos blocos */
    .block-container {{
        animation: fadeUpIn 0.6s cubic-bezier(0.16, 1, 0.3, 1) !important;
    }}
    
    /* Métricas aparecem com delay escalonado */
    [data-testid="stMetric"] {{
        animation: fadeUpIn 0.5s cubic-bezier(0.16, 1, 0.3, 1) !important;
        transition: all 0.35s cubic-bezier(0.4, 0, 0.2, 1) !important;
    }}
    [data-testid="stHorizontalBlock"] > [data-testid="stVerticalBlock"]:nth-child(1) [data-testid="stMetric"] {{
        animation-delay: 0.05s !important;
    }}
    [data-testid="stHorizontalBlock"] > [data-testid="stVerticalBlock"]:nth-child(2) [data-testid="stMetric"] {{
        animation-delay: 0.12s !important;
    }}
    [data-testid="stHorizontalBlock"] > [data-testid="stVerticalBlock"]:nth-child(3) [data-testid="stMetric"] {{
        animation-delay: 0.19s !important;
    }}
    [data-testid="stHorizontalBlock"] > [data-testid="stVerticalBlock"]:nth-child(4) [data-testid="stMetric"] {{
        animation-delay: 0.26s !important;
    }}
    
    /* Botão Principal com pulso neon suave */
    div.stButton > button[kind="primary"],
    div.stButton > button[data-baseweb="button"] {{
        transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1) !important;
    }}
    div.stButton > button[kind="primary"]:hover,
    div.stButton > button[data-baseweb="button"]:hover {{
        animation: pulseGlow 2s ease-in-out infinite !important;
        transform: translateY(-2px) scale(1.02) !important;
    }}
    div.stButton > button[data-baseweb="button"]:active {{
        transform: translateY(1px) scale(0.97) !important;
        transition: transform 0.1s ease !important;
    }}
    
    /* Tabs com transição suave */
    .stTabs [data-baseweb="tab"] {{
        transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1) !important;
    }}
    .stTabs [aria-selected="true"] {{
        transition: all 0.3s cubic-bezier(0.16, 1, 0.3, 1) !important;
        box-shadow: 0 4px 20px rgba(14, 165, 233, 0.35), 0 0 30px rgba(14, 165, 233, 0.1) !important;
    }}
    
    /* Expanders com transição suave */
    [data-testid="stExpander"] {{
        transition: all 0.3s ease !important;
    }}
    [data-testid="stExpander"]:hover {{
        transform: translateY(-1px) !important;
        border-color: rgba(14, 165, 233, 0.25) !important;
        box-shadow: 0 8px 25px rgba(14, 165, 233, 0.1), inset 0 1px 0 rgba(255, 255, 255, 0.06) !important;
    }}
    
    /* DataFrames com hover suave */
    .stDataFrame {{
        transition: all 0.3s ease !important;
    }}
    .stDataFrame:hover {{
        transform: translateY(-2px) !important;
        box-shadow: 0 12px 40px rgba(0, 0, 0, 0.25), 0 0 30px rgba(14, 165, 233, 0.08) !important;
        border-color: rgba(14, 165, 233, 0.3) !important;
    }}
    
    /* Inputs com transição de foco suave */
    .stTextInput input, .stSelectbox > div > div, .stTextArea textarea {{
        transition: all 0.3s ease !important;
    }}
    .stTextInput input:focus, .stTextArea textarea:focus {{
        transform: scale(1.005) !important;
    }}
    
    /* File Uploader com animação */
    [data-testid="stFileUploader"] {{
        transition: all 0.3s ease !important;
    }}
    [data-testid="stFileUploader"]:hover {{
        border-color: rgba(14, 165, 233, 0.6) !important;
        background: rgba(14, 165, 233, 0.05) !important;
        box-shadow: 0 0 25px rgba(14, 165, 233, 0.1) !important;
        transform: translateY(-1px) !important;
    }}
    
    /* Alerts com entrada suave */
    [data-testid="stAlert"] {{
        animation: fadeUpIn 0.4s cubic-bezier(0.16, 1, 0.3, 1) !important;
    }}
    
    /* Shimmer decorativo no cabeçalho */
    .enesa-header {{
        position: relative;
        overflow: hidden;
        transition: all 0.4s cubic-bezier(0.4, 0, 0.2, 1) !important;
    }}
    .enesa-header::after {{
        content: '';
        position: absolute;
        top: 0;
        left: -200%;
        width: 200%;
        height: 100%;
        background: linear-gradient(
            90deg,
            transparent,
            rgba(14, 165, 233, 0.04),
            rgba(255, 255, 255, 0.06),
            rgba(14, 165, 233, 0.04),
            transparent
        );
        background-size: 200% 100%;
        animation: shimmer 8s linear infinite;
        pointer-events: none;
    }}
    .enesa-header:hover {{
        transform: translateY(-3px) !important;
        box-shadow: 0 12px 40px rgba(0, 0, 0, 0.35), 0 0 30px rgba(14, 165, 233, 0.08) !important;
    }}
    
    /* ================================================ */
    /* 📱 BARRA DE NAVEGAÇÃO INFERIOR (MOBILE ONLY)     */
    /* ================================================ */
    @media screen and (max-width: 768px) {{
        
        /* Reposicionar as tabs principais para o rodapé */
        .stTabs [data-baseweb="tab-list"] {{
            position: fixed !important;
            bottom: 0 !important;
            left: 0 !important;
            right: 0 !important;
            z-index: 998 !important;
            border-radius: 20px 20px 0 0 !important;
            background: rgba(10, 15, 30, 0.85) !important;
            backdrop-filter: blur(24px) saturate(200%) !important;
            -webkit-backdrop-filter: blur(24px) saturate(200%) !important;
            border: none !important;
            border-top: 1px solid rgba(14, 165, 233, 0.2) !important;
            box-shadow: 0 -4px 30px rgba(0, 0, 0, 0.4), 0 -1px 0 rgba(255, 255, 255, 0.05) !important;
            padding: 6px 8px 10px 8px !important;
            gap: 2px !important;
            overflow-x: auto !important;
            overflow-y: hidden !important;
            flex-wrap: nowrap !important;
            -webkit-overflow-scrolling: touch !important;
            scrollbar-width: none !important;
        }}
        .stTabs [data-baseweb="tab-list"]::-webkit-scrollbar {{
            display: none !important;
        }}
        
        /* Abas do rodapé com estilo de ícone */
        .stTabs [data-baseweb="tab"] {{
            flex-direction: column !important;
            align-items: center !important;
            justify-content: center !important;
            padding: 6px 10px !important;
            font-size: 0.65rem !important;
            min-width: 60px !important;
            min-height: 50px !important;
            border-radius: 12px !important;
            gap: 2px !important;
            flex-shrink: 0 !important;
            white-space: nowrap !important;
            line-height: 1.1 !important;
        }}
        
        /* Tab ativa com brilho neon */
        .stTabs [aria-selected="true"] {{
            background: linear-gradient(135deg, rgba(14, 165, 233, 0.25), rgba(59, 130, 246, 0.2)) !important;
            color: #0ea5e9 !important;
            border: 1px solid rgba(14, 165, 233, 0.3) !important;
            box-shadow: 0 0 20px rgba(14, 165, 233, 0.2), inset 0 1px 0 rgba(255, 255, 255, 0.1) !important;
        }}
        
        /* Espaço extra no final da página para não ficar escondido pela barra */
        .block-container {{
            padding-bottom: 85px !important;
        }}
        
        /* Indicador de brilho na barra inferior */
        .stTabs [data-baseweb="tab-list"]::before {{
            content: '' !important;
            position: absolute !important;
            top: -1px !important;
            left: 10% !important;
            right: 10% !important;
            height: 2px !important;
            background: linear-gradient(90deg, transparent, rgba(14, 165, 233, 0.5), transparent) !important;
            border-radius: 2px !important;
            z-index: 1 !important;
            animation: none !important;
        }}
    }}
    </style>
    
    <div class="watermark-edson">EDSON GARCIA DE ARAUJO</div>
""", unsafe_allow_html=True)

# --- CHECAR LOGIN POR LINK RÁPIDO (QR CODE) ---
try:
    if "pwd" in st.query_params and st.query_params["pwd"] == "Campo@2026":
        st.session_state.logged_in = True
        st.session_state.role = "encarregado"
except Exception:
    pass

# --- LOGIN / AUTENTICAÇÃO (MIGRAÇÃO) ---
if "logged_in" not in st.session_state:
    st.session_state.logged_in = True  # Bypass do login antigo para usar o novo (usuario_logado)
    st.session_state.role = "admin"

# Removido o header global daqui para aparecer apenas após o login.

# =================================================================
# CAMINHOS E CONFIGURAÇÕES
# =================================================================
pasta_base = os.path.dirname(os.path.abspath(__file__))
caminho_logo = os.path.join(pasta_base, "logo.png")
caminho_pde_padrao = os.path.join(pasta_base, "PDE.csv")
caminho_modelo_padrao = os.path.join(pasta_base, "MODELO.xlsx")
caminho_modelo_salvo = os.path.join(pasta_base, "MODELO_SALVO.xlsx")

# Caminhos para PERSISTÊNCIA (salvar a base do usuário)
caminho_base_salva_csv = os.path.join(pasta_base, "BASE_ATUAL.csv")
caminho_hist_cc = os.path.join(pasta_base, "historico_cc.csv")
caminho_historico_f1_csv = os.path.join(pasta_base, "historico_f1_local.csv")
caminho_base_salva_xlsx = os.path.join(pasta_base, "BASE_ATUAL.xlsx")
caminho_escala_csv = os.path.join(pasta_base, "escala_diaria.csv")
caminho_rdc_registros_csv = os.path.join(pasta_base, "rdc_registros.csv")

celula_encarregado = "I4"
celula_matricula = "B9"
celula_nome = "C9"
celula_funcao = "H9"



# =================================================================
# FUNÇÕES UTILITÁRIAS
# =================================================================
def extrair_coordenadas(celula_str):
    match = re.match(r"([A-Z]+)([0-9]+)", celula_str.strip().upper())
    if match:
        return match.group(1), int(match.group(2))
    return None, None

def salvar_historico(wb, nome_arquivo):
    try:
        hoje = datetime.datetime.now()
        pasta_hist = os.path.join(pasta_base, "Historico_RDC", str(hoje.year), f"{hoje.month:02d}_{hoje.strftime('%B')}")
        os.makedirs(pasta_hist, exist_ok=True)
        caminho_hist = os.path.join(pasta_hist, f"{hoje.strftime('%d_%H%M')}_{nome_arquivo}")
        wb.save(caminho_hist)
    except Exception:
        pass

def salvar_modelo_no_disco(arquivo_modelo_up):
    try:
        arquivo_modelo_up.seek(0)
        conteudo = arquivo_modelo_up.read()
        arquivo_modelo_up.seek(0)
        with open(caminho_modelo_salvo, "wb") as f:
            f.write(conteudo)
        return True
    except Exception:
        return False

def obter_caminho_modelo():
    if os.path.exists(caminho_modelo_salvo):
        return caminho_modelo_salvo
    elif os.path.exists(caminho_modelo_padrao):
        return caminho_modelo_padrao
    return None

def preencher_excel(equipe, encarregado_selecionado, data_rdc=""):
    try:
        caminho = obter_caminho_modelo()
        if caminho is None:
            st.error("⚠️ Ficheiro MODELO.xlsx não encontrado! Faça upload na barra lateral.")
            return None
        
        wb = openpyxl.load_workbook(caminho)
            
        from copy import copy
        ws = wb.active
        celula_enc = ws[celula_encarregado]
        celula_enc.value = encarregado_selecionado
        
        # Inserir data se fornecida
        if data_rdc:
            for c_dt in ["O2", "P2", "Q2", "N2", "M2", "O3", "P3"]:
                if ws[c_dt].value is None or "DATA" in str(ws[c_dt].value).upper():
                    pass
        
        # Ajuste dinâmico de tamanho para nomes longos
        tamanho = 16
        if len(encarregado_selecionado) > 25:
            tamanho = 14
        if len(encarregado_selecionado) > 35:
            tamanho = 12
            
        if celula_enc.font:
            nova_fonte = copy(celula_enc.font)
            nova_fonte.size = tamanho
            nova_fonte.bold = True
            celula_enc.font = nova_fonte
        letra_mat, _ = extrair_coordenadas(celula_matricula)
        letra_nom, linha_nom = extrair_coordenadas(celula_nome)
        letra_fun, _ = extrair_coordenadas(celula_funcao)
        linha_atual = linha_nom if linha_nom else 9
        
        for _, row in equipe.iterrows():
            if letra_mat:
                ws[f"{letra_mat}{linha_atual}"] = str(row.get("MATRICULA", ""))
            if letra_nom:
                ws[f"{letra_nom}{linha_atual}"] = str(row.get("NOME", ""))
            if letra_fun:
                ws[f"{letra_fun}{linha_atual}"] = str(row.get("FUNÇÃO", ""))
            linha_atual += 1
            
        ws.print_area = "A1:R61"
        ws.page_setup.orientation = ws.ORIENTATION_PORTRAIT
        ws.page_setup.paperSize = ws.PAPERSIZE_A4
        ws.page_setup.horizontalCentered = True
        ws.page_setup.verticalCentered = True
        ws.page_margins.left = 0.2
        ws.page_margins.right = 0.2
        ws.page_margins.top = 0.5
        ws.page_margins.bottom = 0.5
        
        ws.sheet_properties.pageSetUpPr.fitToPage = True
        ws.page_setup.fitToWidth = 1
        ws.page_setup.fitToHeight = 1
        return wb
    except Exception as e:
        st.error(f"⚠️ Erro ao preencher o modelo Excel: {e}")
        return None

def gerar_pdf_rdc(equipe, encarregado_selecionado, nome_empresa="", logo_path="", data_rdc=""):
    """Gera um PDF que replica fielmente o MODELO.xlsx da Apropriacao de Campo."""
    try:
        from fpdf import FPDF
    except ImportError:
        return None

    def safe(text):
        return str(text).encode("latin-1", "replace").decode("latin-1")

    pdf = FPDF(orientation="P", unit="mm", format="A4")
    pdf.add_page()
    pdf.set_auto_page_break(auto=False, margin=5)
    
    # Margens
    ml = 5  # margem esquerda
    page_w = 200  # largura util
    
    # ================================================================
    # CABECALHO COMPLETO (Logo esquerda + info direita)
    # ================================================================
    y_start = 5
    logo_w = 28  # largura da coluna da logo
    info_w = page_w - logo_w  # largura da area de info
    h_row = 5.5  # altura de cada linha do cabecalho
    h_titulo = 7  # altura do titulo
    header_h = h_titulo + h_row * 4  # altura total do cabecalho (titulo + 4 linhas)
    
    pdf.set_font("Helvetica", "B", 10)
    
    # --- LOGO (celula esquerda, altura total do cabecalho) ---
    has_logo = logo_path and os.path.exists(logo_path)
    pdf.set_xy(ml, y_start)
    pdf.cell(logo_w, header_h, "", 1, 0)  # Borda da celula da logo
    
    if has_logo:
        try:
            # Centralizar logo dentro da celula
            logo_img_w = logo_w - 4
            logo_img_h = header_h - 4
            pdf.image(logo_path, x=ml + 2, y=y_start + 2, w=logo_img_w, h=logo_img_h)
        except Exception:
            pass
    
    # --- TITULO: APROPRIACAO DE CAMPO ---
    pdf.set_xy(ml + logo_w, y_start)
    pdf.set_font("Helvetica", "B", 10)
    pdf.cell(info_w, h_titulo, "APROPRIACAO DE CAMPO", 1, 1, "C")
    
    y = y_start + h_titulo
    
    # --- LINHA 2: OBRA | DISCIPLINA/SETOR | DATA ---
    x_info = ml + logo_w
    col_left = 75
    col_mid = 60
    col_right = info_w - col_left - col_mid  # 172 - 75 - 60 = 37
    
    pdf.set_xy(x_info, y)
    pdf.set_font("Helvetica", "", 5)
    pdf.cell(12, h_row, "OBRA:", 1, 0, "R")
    pdf.set_font("Helvetica", "", 7)
    obra_txt = safe(f" 125 - {nome_empresa}") if nome_empresa else " 125 - ARAUCO"
    pdf.cell(col_left - 12, h_row, obra_txt, 1, 0, "L")
    pdf.set_font("Helvetica", "", 5)
    pdf.cell(28, h_row, "DISCIPLINA/SETOR:", 1, 0, "R")
    pdf.cell(col_mid - 28, h_row, "", 1, 0, "L")
    pdf.cell(10, h_row, "DATA:", 1, 0, "R")
    pdf.set_font("Helvetica", "B", 7)
    pdf.cell(col_right - 10, h_row, str(data_rdc) if data_rdc else "", 1, 1, "C")
    pdf.set_font("Helvetica", "", 5)
    
    y += h_row
    
    # --- LINHA 3: TURNO | COORDENADOR/SUPERVISOR ---
    pdf.set_xy(x_info, y)
    pdf.set_font("Helvetica", "", 5)
    pdf.cell(12, h_row, "TURNO:", 1, 0, "R")
    pdf.cell(col_left - 12, h_row, "", 1, 0, "L")
    pdf.cell(38, h_row, "COORDENADOR / SUPERVISOR:", 1, 0, "R")
    pdf.cell(info_w - col_left - 38, h_row, "", 1, 1, "L")
    
    y += h_row
    
    # --- LINHA 4: HORARIO | ENCARREGADO OU MESTRE ---
    pdf.set_xy(x_info, y)
    pdf.set_font("Helvetica", "", 5)
    pdf.cell(14, h_row, "HORARIO:", 1, 0, "R")
    pdf.set_font("Helvetica", "", 5)
    pdf.cell(col_left - 14, h_row, " Seg a Sex= 07:00 as 17:00 h", 1, 0, "L")
    pdf.set_font("Helvetica", "", 5)
    pdf.cell(38, h_row, "ENCARREGADO OU MESTRE:", 1, 0, "R")
    pdf.set_font("Helvetica", "B", 8)
    pdf.cell(info_w - col_left - 38, h_row, safe(f" {encarregado_selecionado}"), 1, 1, "L")
    
    y += h_row
    
    # --- LINHA 5: CONDICOES CLIMATICAS | EQUIPAMENTO | COMPONENTE ---
    pdf.set_xy(x_info, y)
    pdf.set_font("Helvetica", "", 5)
    pdf.cell(24, h_row, "COND. CLIMATICAS:", 1, 0, "L")
    pdf.cell(15, h_row, "Bom (  )", 1, 0, "C")
    pdf.cell(18, h_row, "Chuva Leve(  )", 1, 0, "C")
    pdf.cell(18, h_row, "Chuva Forte(  )", 1, 0, "C")
    pdf.cell(20, h_row, "EQUIPAMENTO:", 1, 0, "R")
    pdf.cell(col_mid - 20, h_row, "", 1, 0, "L")
    pdf.cell(17, h_row, "COMPONENTE:", 1, 0, "R")
    pdf.cell(col_right - 17, h_row, "", 1, 1, "L")
    
    y += h_row
    
    # ================================================================
    # LINHA 6: espaco/separador
    # ================================================================
    pdf.set_xy(ml, y)
    pdf.cell(page_w, 2, "", 0, 1)
    y = pdf.get_y()
    
    # ================================================================
    # CABECALHO DA TABELA DE EFETIVO
    # ================================================================
    # Colunas: ITEM | MATRICULA | NOME | FUNCAO | HORARIO(INICIO|TERMINO) | 1|2|3|4|5|6 | TOTAL
    col_item = 8
    col_mat = 20
    col_nome = 60
    col_func = 45
    col_h_ini = 11
    col_h_ter = 11
    col_num = 6
    col_total = 9
    # Total: 8+20+60+45+11+11+36+9 = 200
    
    pdf.set_xy(ml, y)
    pdf.set_font("Helvetica", "B", 6)
    pdf.set_fill_color(220, 220, 220)
    
    # Linha superior do cabecalho (bordas Top, Left, Right para as primeiras colunas)
    pdf.cell(col_item, 4, "ITEM", "LTR", 0, "C", True)
    pdf.cell(col_mat, 4, "MATRICULA", "LTR", 0, "C", True)
    pdf.cell(col_nome, 4, "NOME", "LTR", 0, "C", True)
    pdf.cell(col_func, 4, "FUNCAO", "LTR", 0, "C", True)
    pdf.cell(col_h_ini + col_h_ter, 4, "HORARIO", 1, 0, "C", True)
    for n in range(1, 7):
        pdf.cell(col_num, 4, str(n), 1, 0, "C", True)
    pdf.cell(col_total, 4, "TOTAL", 1, 1, "C", True)
    
    # Sub-cabecalho (bordas Left, Right, Bottom para as primeiras colunas)
    y = pdf.get_y()
    pdf.set_xy(ml, y)
    pdf.set_font("Helvetica", "", 5)
    pdf.cell(col_item, 3, "", "LBR", 0, "C", True)
    pdf.cell(col_mat, 3, "", "LBR", 0, "C", True)
    pdf.cell(col_nome, 3, "", "LBR", 0, "C", True)
    pdf.cell(col_func, 3, "", "LBR", 0, "C", True)
    pdf.cell(col_h_ini, 3, "INICIO", 1, 0, "C", True)
    pdf.cell(col_h_ter, 3, "TERMINO", 1, 0, "C", True)
    for _ in range(6):
        pdf.cell(col_num, 3, "HN", 1, 0, "C", True)
    pdf.cell(col_total, 3, "HN", 1, 1, "C", True)
    
    y = pdf.get_y()
    
    # ================================================================
    # TABELA DE EFETIVO (25 linhas)
    # ================================================================
    pdf.set_font("Helvetica", "", 6)
    num_linhas_efetivo = 25
    h_row = 4.2
    
    for idx in range(num_linhas_efetivo):
        pdf.set_xy(ml, y)
        
        if idx < len(equipe):
            row = equipe.iloc[idx]
            mat = safe(str(row.get("MATRICULA", "")))
            nome_col = safe(str(row.get("NOME", "")))
            funcao = safe(str(row.get("FUNCAO", row.get("FUNÇÃO", ""))))
        else:
            mat, nome_col, funcao = "", "", ""
        
        pdf.cell(col_item, h_row, str(idx + 1), 1, 0, "C")
        pdf.cell(col_mat, h_row, mat, 1, 0, "C")
        pdf.cell(col_nome, h_row, nome_col, 1, 0, "L")
        pdf.cell(col_func, h_row, funcao, 1, 0, "L")
        pdf.cell(col_h_ini, h_row, "", 1, 0, "C")
        pdf.cell(col_h_ter, h_row, "", 1, 0, "C")
        for _ in range(6):
            pdf.cell(col_num, h_row, "", 1, 0, "C")
        pdf.cell(col_total, h_row, "", 1, 1, "C")
        
        y = pdf.get_y()
    
    # Linha TOTAL
    pdf.set_xy(ml, y)
    pdf.set_font("Helvetica", "B", 6)
    pdf.cell(col_item + col_mat + col_nome, h_row, "", 1, 0)
    pdf.cell(col_func, h_row, "TOTAL", 1, 0, "C")
    pdf.cell(col_h_ini, h_row, "", 1, 0)
    pdf.cell(col_h_ter, h_row, "", 1, 0)
    for _ in range(6):
        pdf.cell(col_num, h_row, "", 1, 0)
    pdf.cell(col_total, h_row, "", 1, 1)
    y = pdf.get_y()
    
    # ================================================================
    # MAQUINAS E EQUIPAMENTOS (2 linhas)
    # ================================================================
    pdf.set_xy(ml, y)
    pdf.set_font("Helvetica", "B", 5)
    pdf.cell(page_w, 3.5, "MAQUINAS E EQUIPAMENTOS (quando aplicavel)", 1, 1, "L")
    y = pdf.get_y()
    
    pdf.set_xy(ml, y)
    col_qt = 8
    col_placa = 18
    col_equip = 62
    col_tag = 30
    col_hm_ini = 12
    col_hm_ter = 12
    rest = page_w - col_qt - col_placa - col_equip - col_tag - col_hm_ini - col_hm_ter
    
    pdf.cell(col_qt, 3.5, "Qt.", 1, 0, "C", True)
    pdf.cell(col_placa, 3.5, "PLACA", 1, 0, "C", True)
    pdf.cell(col_equip, 3.5, "EQUIPAMENTO", 1, 0, "C", True)
    pdf.cell(col_tag, 3.5, "TAG", 1, 0, "C", True)
    pdf.cell(col_hm_ini, 3.5, "Inicio", 1, 0, "C", True)
    pdf.cell(col_hm_ter, 3.5, "Termino", 1, 0, "C", True)
    pdf.cell(rest, 3.5, "Total", 1, 1, "C", True)
    y = pdf.get_y()
    
    for mq in range(2):
        pdf.set_xy(ml, y)
        pdf.set_font("Helvetica", "", 5)
        pdf.cell(col_qt, 3.5, str(mq + 1), 1, 0, "C")
        pdf.cell(col_placa, 3.5, "", 1, 0)
        pdf.cell(col_equip, 3.5, "", 1, 0)
        pdf.cell(col_tag, 3.5, "", 1, 0)
        pdf.cell(col_hm_ini, 3.5, "", 1, 0)
        pdf.cell(col_hm_ter, 3.5, "", 1, 0)
        pdf.cell(rest, 3.5, "", 1, 1)
        y = pdf.get_y()
    
    # ================================================================
    # ATIVIDADES (10 linhas)
    # ================================================================
    pdf.set_xy(ml, y)
    pdf.set_font("Helvetica", "B", 6)
    col_item_a = 8
    col_ativ = page_w - col_item_a - 20
    col_ca = 20
    
    pdf.cell(col_item_a, 5, "ITEM", 1, 0, "C", True)
    pdf.cell(col_ativ, 5, "ATIVIDADES", 1, 0, "C", True)
    pdf.cell(col_ca, 5, "C. ATIVID.", 1, 1, "C", True)
    y = pdf.get_y()
    
    pdf.set_font("Helvetica", "", 6)
    for ai in range(10):
        pdf.set_xy(ml, y)
        pdf.cell(col_item_a, 5, str(ai + 1), 1, 0, "C")
        pdf.cell(col_ativ, 5, "", 1, 0)
        pdf.cell(col_ca, 5, "", 1, 1)
        y = pdf.get_y()
    
    # ================================================================
    # AREA / LOCAL DE TRABALHO
    # ================================================================
    pdf.set_xy(ml, y)
    pdf.set_font("Helvetica", "B", 6)
    pdf.cell(page_w, 5, "          PB (  )            RB (  )                                                                                    OBS", 1, 1, "L")
    y = pdf.get_y()
    
    areas = [
        "DUTO (  )      EQUIPAMENTO (  )     TUBULACAO (  )",
        "ESTRUTURA MET (  )      PRECIPITADOR (  )",
        "PRESSAO - MEC (  )      PRESSAO - TUBULACAO (  )",
        "PRESSAO - FORNALHA (  )      PINTURA (  )",
        "SOPRAGEM (  )      ANDAIME (  )",
    ]
    pdf.set_font("Helvetica", "", 6)
    for area_txt in areas:
        pdf.set_xy(ml, y)
        pdf.cell(page_w, 4.5, safe(area_txt), 1, 1, "L")
        y = pdf.get_y()
    
    pdf.set_xy(ml, y)
    pdf.set_font("Helvetica", "B", 5)
    pdf.cell(page_w, 4, " MARCAR CONFORME O LOCAL DE TRABALHO", 1, 1, "L")
    y = pdf.get_y()
    
    # Espaco
    pdf.ln(3)
    y = pdf.get_y()
    
    # ================================================================
    # ASSINATURAS (3 colunas)
    # ================================================================
    col_ass = page_w / 3
    pdf.set_xy(ml, y)
    pdf.set_font("Helvetica", "", 5)
    
    # Linhas pontilhadas
    pdf.ln(8)
    y = pdf.get_y()
    
    x1 = ml + 5
    x2 = ml + col_ass + 5
    x3 = ml + 2 * col_ass + 5
    line_len = col_ass - 10
    
    pdf.line(x1, y, x1 + line_len, y)
    pdf.line(x2, y, x2 + line_len, y)
    pdf.line(x3, y, x3 + line_len, y)
    
    pdf.set_xy(ml, y + 1)
    pdf.cell(col_ass, 4, "ENCARREGADO / MESTRE", 0, 0, "C")
    pdf.cell(col_ass, 4, "ENGENHEIRO / COORDENADOR", 0, 0, "C")
    pdf.cell(col_ass, 4, "PLANEJAMENTO", 0, 1, "C")

    return bytes(pdf.output())

def ler_arquivo_seguro(arquivo, nome_arquivo=""):
    try:
        if nome_arquivo.endswith(".xlsx") or nome_arquivo.endswith(".xls"):
            return pd.read_excel(arquivo)
        else:
            try:
                return pd.read_csv(arquivo, sep=";", encoding="latin-1")
            except Exception:
                pass
            try:
                if hasattr(arquivo, 'seek'):
                    arquivo.seek(0)
                return pd.read_csv(arquivo, sep=",", encoding="latin-1")
            except Exception:
                pass
            try:
                if hasattr(arquivo, 'seek'):
                    arquivo.seek(0)
                return pd.read_csv(arquivo, sep=";", encoding="utf-8")
            except Exception:
                pass
            if hasattr(arquivo, 'seek'):
                arquivo.seek(0)
            return pd.read_csv(arquivo)
    except Exception as e:
        st.sidebar.error(f"❌ Não foi possível ler o arquivo: {e}")
        return None

def salvar_base_localmente(arquivo_upload):
    try:
        nome = arquivo_upload.name
        arquivo_upload.seek(0)
        conteudo = arquivo_upload.read()
        arquivo_upload.seek(0)
        
        if nome.endswith(".xlsx") or nome.endswith(".xls"):
            with open(caminho_base_salva_xlsx, "wb") as f:
                f.write(conteudo)
            if os.path.exists(caminho_base_salva_csv):
                os.remove(caminho_base_salva_csv)
        else:
            with open(caminho_base_salva_csv, "wb") as f:
                f.write(conteudo)
            if os.path.exists(caminho_base_salva_xlsx):
                os.remove(caminho_base_salva_xlsx)
        return True
    except Exception:
        return False

def salvar_f1_seguro(conn, df_f1, caminho_csv):
    """Salva o Histórico F1 no Google Sheets COM PROTEÇÃO ANTI-PERDA.
    NUNCA sobrescreve a nuvem com dados vazios ou muito menores.
    """
    if df_f1 is None or df_f1.empty:
        return False, "Bloqueado: tentativa de salvar dados vazios na nuvem."
    
    # Filtrar linhas inválidas antes de salvar
    if 'ENCARREGADO' in df_f1.columns:
        df_f1 = df_f1[df_f1['ENCARREGADO'].notna() & (df_f1['ENCARREGADO'] != '')]
    
    if df_f1.empty:
        return False, "Bloqueado: todos os registros são inválidos."
    
    # Verificar quantos registros existem na nuvem antes de sobrescrever
    try:
        df_nuvem = conn.read(worksheet="Historico_F1", ttl=0)
        if df_nuvem is not None:
            df_nuvem = df_nuvem.dropna(how='all')
            qtd_nuvem = len(df_nuvem)
        else:
            qtd_nuvem = 0
    except Exception:
        qtd_nuvem = 0
    
    qtd_novo = len(df_f1)
    
    # PROTEÇÃO: Se os dados novos têm MUITO MENOS registros que a nuvem (mais de 50% menor),
    # algo está errado. Bloquear a escrita para não perder dados.
    if qtd_nuvem > 10 and qtd_novo < qtd_nuvem * 0.5:
        return False, f"Bloqueado: tentativa de reduzir de {qtd_nuvem} para {qtd_novo} registros (perda >50%)."
    
    # Tudo OK, salvar
    try:
        conn.update(worksheet="Historico_F1", data=df_f1)
        # Backup local
        df_f1.to_csv(caminho_csv, index=False)
        return True, f"OK: {qtd_novo} registros salvos."
    except Exception as e:
        return False, f"Erro ao salvar: {e}"

def preparar_dataframe(df):
    # Auto-detect header row if the file has title rows above the headers
    unnamed_cols = [c for c in df.columns if str(c).startswith('Unnamed')]
    if len(unnamed_cols) > len(df.columns) / 2:
        for i, row in df.head(15).iterrows():
            row_str = " ".join([str(x).upper() for x in row.values])
            if "NOME" in row_str or "MATRICULA" in row_str or "CHAPA" in row_str or "CRACHA" in row_str or "COLABORADOR" in row_str:
                df.columns = row
                df = df.iloc[i+1:].reset_index(drop=True)
                break

    mapeamento = {}
    for col in df.columns:
        col_clean = str(col).strip().upper()
        
        # Ignorar matriculas de líderes/encarregados para não sobrescrever a matricula principal
        if "MATRÍCULA" in col_clean or "MATRICULA" in col_clean or "MAT." in col_clean or "CHAPA" in col_clean or "CRACHÁ" in col_clean or "CRACHA" in col_clean or "RE " in col_clean or "RE" == col_clean:
            if "LÍDER" in col_clean or "LIDER" in col_clean or "ENCARREGADO" in col_clean or "COORDENADOR" in col_clean or "SUPERVISOR" in col_clean:
                continue
            else:
                mapeamento[col] = "MATRICULA"
                
        elif "ENCARREGADO" in col_clean or "LÍDER" in col_clean or "LIDER" in col_clean or "SUPERVISOR" in col_clean or "COORDENADOR" in col_clean:
            mapeamento[col] = "ENCARREGADO"
            
        elif "NOME" in col_clean or "COLABORADOR" in col_clean or "FUNCIONÁRIO" in col_clean or "FUNCIONARIO" in col_clean or "EMPREGADO" in col_clean:
            mapeamento[col] = "NOME"
            
        elif "FUNÇÃO" in col_clean or "FUNCAO" in col_clean or "CARGO" in col_clean or "CBO" in col_clean:
            mapeamento[col] = "FUNÇÃO"
            
        elif "CENTRO DE CUSTO" in col_clean or col_clean == "C.C" or col_clean == "CC" or col_clean == "CECO":
            mapeamento[col] = "C.C"
            
        elif "DISCIPLINA" in col_clean or "ÁREA" in col_clean or "AREA" in col_clean or "SETOR" in col_clean:
            mapeamento[col] = "DISCIPLINA"
            
        elif "MÃO DE OBRA" in col_clean or "MAO DE OBRA" in col_clean or "TIPO" in col_clean:
            mapeamento[col] = "MÃO DE OBRA"
            
        elif "TURNO" in col_clean or "HORÁRIO" in col_clean or "HORARIO" in col_clean:
            mapeamento[col] = "TURNO"
            
        elif "STATUS" in col_clean or "SITUAÇÃO" in col_clean or "SITUACAO" in col_clean:
            mapeamento[col] = "STATUS"

    df = df.rename(columns=mapeamento)
    
    # Remover colunas duplicadas mantendo a primeira encontrada (que geralmente é a principal da esquerda pra direita)
    df = df.loc[:, ~df.columns.duplicated(keep='first')]
    
    for c in ["MATRICULA", "NOME", "FUNÇÃO", "ENCARREGADO", "TURNO", "STATUS"]:
        if c not in df.columns:
            df[c] = ""
        df[c] = df[c].fillna("").astype(str).str.strip()
        df[c] = df[c].replace(["nan", "NaN", "None", "0.0", "0", "#N/D", "#N/A", "#REF!", "-"], "")

    for c in ["C.C", "DISCIPLINA", "MÃO DE OBRA"]:
        if c not in df.columns:
            df[c] = ""
        df[c] = df[c].fillna("").astype(str).str.strip()
        df[c] = df[c].replace(["nan", "NaN", "None", "0.0", "0", "#N/D", "#N/A", "#REF!", "-"], "")

    df["MATRICULA"] = df["MATRICULA"].str.replace(".0", "", regex=False)
    
    # -------------------------------------------------------------------------
    # NOVO: NORMALIZAR NOMES DE C.C VINDOS DO ERP EM FORMATO DE TEXTO
    # Ex: "PB - Dutos" -> "125.02.002"
    # -------------------------------------------------------------------------
    def normalizar_cc(valor):
        val = str(valor).upper()
        if not val or val.startswith("125."):
            return valor
            
        prefixo = ""
        if "PB" in val: prefixo = "125.02"
        elif "RB" in val: prefixo = "125.01"
            
        sufixo = ""
        if "DUTO" in val: sufixo = "002"
        elif "EQUIPA" in val: sufixo = "001"
        elif "TUBU" in val and "PRESS" not in val: sufixo = "003"
        elif "ESTRUTURA" in val: sufixo = "004"
        elif "PRECIP" in val: sufixo = "005"
        elif "PRESS" in val and "MEC" in val: sufixo = "006"
        elif "PRESS" in val and "TUBU" in val: sufixo = "007"
        elif "PRESS" in val and "FORN" in val: sufixo = "008"
        elif "PINTURA" in val: sufixo = "009"
        elif "COMIS" in val: sufixo = "010"
        elif "ASSISTIDA" in val: sufixo = "011"
        elif "LAVAGEM" in val: sufixo = "012"
        elif "SOPRAGEM" in val: sufixo = "013"
        elif "ANDAIME" in val: sufixo = "014"
        elif "OPERADOR" in val: sufixo = "015"
        elif "ESCOPO" in val: sufixo = "016"
        elif "GERENC" in val: sufixo = "101"
        elif "PRODUC" in val or "PRODUÇ" in val: sufixo = "102"
        elif "QUALIDADE" in val: sufixo = "103"
        elif "PLANEJAMENTO" in val: sufixo = "104"
        elif "ADMINISTRA" in val: sufixo = "105"
        elif "MEDICINA" in val or "SEGURAN" in val: sufixo = "106"
        elif "INFRAESTRUTURA" in val: sufixo = "107"
        elif "ALMOXARIFADO" in val and "ENESA" in val: sufixo = "108"
        elif "ALMOXARIFADO" in val: sufixo = "109"
        elif "ELETRICA PROV" in val: sufixo = "110"
        elif "TOPOGRAFIA" in val: sufixo = "111"
        elif "CARGA" in val or "MOVIMENTA" in val: sufixo = "112"
        elif "MEDICAO" in val or "CUSTO" in val: sufixo = "113"
        
        if prefixo and sufixo:
            return f"{prefixo}.{sufixo}"
        return valor

    df["C.C"] = df["C.C"].apply(normalizar_cc)
    
    # -------------------------------------------------------------------------
    # NOVO: FORÇAR A DISCIPLINA A SER SEMPRE CORRETA DE ACORDO COM O C.C ATUAL
    # -------------------------------------------------------------------------
    mapa_sufixo_disciplina = {
        '001': 'EQUIPAMENTOS', '002': 'DUTOS', '003': 'TUBULACAO', 
        '004': 'ESTRUTURA METALICA', '005': 'PRECIPITADOR', '006': 'PRESSAO - MECANICA', 
        '007': 'PRESSAO - TUBULACAO', '008': 'PRESSAO - FORNALHA', '009': 'PINTURA', 
        '010': 'COMISSIONAMENTO', '011': 'OP. ASSISTIDA', '012': 'LAVAGEM QUIMICA', 
        '013': 'SOPRAGEM', '014': 'ANDAIME', '015': 'OPERADORES', '016': 'FORA DE ESCOPO',
        '101': 'GERENCIA', '102': 'PRODUCAO', '103': 'GARANTIA DA QUALIDADE',
        '104': 'PLANEJAMENTO', '105': 'ADMINISTRACAO', '106': 'SEGURANCA E MEDICINA DO TRABALHO',
        '107': 'INFRAESTRUTURA', '108': 'ALMOXARIFADO ENESA', '109': 'ALMOXARIFADO MATERIAIS',
        '110': 'MANUT. ELETRICA PROVISORIA', '111': 'TOPOGRAFIA', '112': 'MOVIMENTACAO DE CARGAS',
        '113': 'MEDICAO/CUSTO/CONTRATOS'
    }

    def corrigir_disciplina(row):
        cc_val = str(row["C.C"]).strip()
        sufixo = cc_val.split('.')[-1] if '.' in cc_val else cc_val
        if sufixo in mapa_sufixo_disciplina:
            return mapa_sufixo_disciplina[sufixo]
        return str(row.get("DISCIPLINA", "")).upper()

    df["DISCIPLINA"] = df.apply(corrigir_disciplina, axis=1)
    # -------------------------------------------------------------------------
    colunas_ordenadas = ["MATRICULA", "NOME", "FUNÇÃO", "C.C", "ENCARREGADO"]
    outras_cols = [c for c in df.columns if c not in colunas_ordenadas]
    df = df[colunas_ordenadas + outras_cols]
    df = df[df["NOME"].str.strip() != ""]
    
    return df

# =================================================================
# INTEGRAÇÃO GOOGLE DRIVE (BACKUP NUVEM)
# =================================================================
def backup_google_drive(file_path, mime_type, file_name):
    try:
        from google.oauth2 import service_account
        from googleapiclient.discovery import build
        from googleapiclient.http import MediaFileUpload
        
        # Tentar pegar as credenciais que já usamos pro Sheets
        if "connections" in st.secrets and "gsheets" in st.secrets["connections"]:
            creds_info = st.secrets["connections"]["gsheets"]
            # Precisamos do escopo do Drive
            scopes = ['https://www.googleapis.com/auth/drive.file']
            creds = service_account.Credentials.from_service_account_info(creds_info, scopes=scopes)
            
            drive_service = build('drive', 'v3', credentials=creds)
            
            # Buscar pasta "RDO_Backups"
            pasta_nome = "RDO_Backups"
            query = f"name='{pasta_nome}' and mimeType='application/vnd.google-apps.folder' and trashed=false"
            response = drive_service.files().list(q=query, spaces='drive', fields='files(id, name)').execute()
            pastas = response.get('files', [])
            
            if not pastas:
                # Criar pasta se não existe
                folder_metadata = {
                    'name': pasta_nome,
                    'mimeType': 'application/vnd.google-apps.folder'
                }
                folder = drive_service.files().create(body=folder_metadata, fields='id').execute()
                folder_id = folder.get('id')
            else:
                folder_id = pastas[0].get('id')
                
            file_metadata = {
                'name': file_name,
                'parents': [folder_id]
            }
            media = MediaFileUpload(file_path, mimetype=mime_type, resumable=True)
            arquivo_salvo = drive_service.files().create(body=file_metadata, media_body=media, fields='id').execute()
            return True, f"Backup salvo no Drive com ID: {arquivo_salvo.get('id')}"
        else:
            return False, "Credenciais do Google não encontradas no secrets.toml"
    except ImportError:
        return False, "Biblioteca google-api-python-client não instalada."
    except Exception as e:
        return False, f"Erro no Drive: {e}"

# =================================================================
# SESSION STATE
# =================================================================
if 'df' not in st.session_state:
    st.session_state.df = None
if 'df_ia' not in st.session_state:
    st.session_state.df_ia = pd.DataFrame(columns=['ITEM', 'SUB', 'DATA', 'DISCIPLINA', 'ENCARREGADO', 'TURNO', 'DDS', 'TRANSCRICAO', 'ATIVIDADE', 'SUB_ATIVIDADE', 'LOCAL_ESPECIFICO', 'EFETIVO_ATIVIDADE', 'PROBLEMAS', 'LOCAL', 'AREA', 'CALDEIRA'])
if 'df_historico_f1' not in st.session_state or st.session_state.df_historico_f1.empty:
    _f1_carregado = False
    # PRIORIDADE 1: Tentar carregar da NUVEM (Google Sheets) para nunca perder dados
    try:
        _conn_f1 = st.connection("gsheets", type=GSheetsConnection)
        _df_f1_nuvem = _conn_f1.read(worksheet="Historico_F1", ttl=0)
        if _df_f1_nuvem is not None:
            _df_f1_nuvem = _df_f1_nuvem.dropna(how='all')
            # SÓ aceitar se tiver dados reais (pelo menos 1 linha com DATA e ENCARREGADO)
            if not _df_f1_nuvem.empty and 'DATA' in _df_f1_nuvem.columns and 'ENCARREGADO' in _df_f1_nuvem.columns:
                _df_f1_nuvem = _df_f1_nuvem[_df_f1_nuvem['ENCARREGADO'].notna() & (_df_f1_nuvem['ENCARREGADO'] != '')]
                if len(_df_f1_nuvem) > 0:
                    st.session_state.df_historico_f1 = _df_f1_nuvem
                    # Salvar cópia local como backup
                    _df_f1_nuvem.to_csv(caminho_historico_f1_csv, index=False)
                    _f1_carregado = True
    except Exception as e:
        st.sidebar.caption(f"⚠️ F1 nuvem: {e}")
    # PRIORIDADE 2: Se a nuvem falhou, tentar o CSV local
    if not _f1_carregado:
        if os.path.exists(caminho_historico_f1_csv):
            try:
                _df_local = pd.read_csv(caminho_historico_f1_csv)
                if not _df_local.empty and 'ENCARREGADO' in _df_local.columns:
                    st.session_state.df_historico_f1 = _df_local
                    _f1_carregado = True
            except Exception:
                pass
    # PRIORIDADE 3: Se nada funcionou, criar DataFrame vazio
    if not _f1_carregado:
        st.session_state.df_historico_f1 = pd.DataFrame(columns=["DATA", "ENCARREGADO"])
if 'mostrar_upload' not in st.session_state:
    st.session_state.mostrar_upload = False

# =================================================================
# SISTEMA DE LOGIN (BLOQUEIO GLOBAL) E COOKIES
# =================================================================
import extra_streamlit_components as stx

cookie_manager = stx.CookieManager()

caminho_usuarios = "usuarios.json"
import json
if not os.path.exists(caminho_usuarios):
    with open(caminho_usuarios, "w", encoding="utf-8") as f:
        json.dump({"admin": {"senha": "123", "nome": "Administrador", "role": "admin"}}, f)


def carregar_usuarios():
    if not os.path.exists(caminho_usuarios): return {}
    with open(caminho_usuarios, "r", encoding="utf-8") as f:
        return json.load(f)

def salvar_usuarios(users):
    with open(caminho_usuarios, "w", encoding="utf-8") as f:
        json.dump(users, f)

if "usuario_logado" not in st.session_state:
    st.session_state.usuario_logado = None
if "role_usuario" not in st.session_state:
    st.session_state.role_usuario = None
if "nome_completo" not in st.session_state:
    st.session_state.nome_completo = None

usuarios_db = carregar_usuarios()

# Tentativa de auto-login via Cookie
cookie_user = cookie_manager.get("rdc_user_session")
if st.session_state.usuario_logado is None and cookie_user and cookie_user in usuarios_db:
    st.session_state.usuario_logado = cookie_user
    st.session_state.role_usuario = usuarios_db[cookie_user].get("role", "user")
    st.session_state.nome_completo = usuarios_db[cookie_user].get("nome", cookie_user)

if st.session_state.usuario_logado is None:
    st.markdown("<br><br>", unsafe_allow_html=True)
    
    col1, col2, col3 = st.columns([1, 1.2, 1])
    with col2:
        # Se houver logo, mostra logo acima da caixa
        if os.path.exists(caminho_logo):
            col_l1, col_l2, col_l3 = st.columns([1, 2, 1])
            with col_l2:
                st.image(caminho_logo, use_container_width=True)
            st.markdown("<br>", unsafe_allow_html=True)
            
        with st.container():
            st.markdown("""
                <style>
                    .login-premium-card {
                        background: rgba(15, 23, 42, 0.7);
                        backdrop-filter: blur(24px);
                        -webkit-backdrop-filter: blur(24px);
                        border-radius: 24px;
                        border: 1px solid rgba(255,255,255,0.08);
                        padding: 48px 40px;
                        box-shadow: 0 25px 60px rgba(0,0,0,0.5), 0 0 40px rgba(14, 165, 233, 0.08);
                        text-align: center;
                        margin-bottom: 20px;
                    }
                    .login-premium-icon {
                        width: 80px; height: 80px;
                        border-radius: 50%;
                        background: linear-gradient(135deg, #0ea5e9, #8b5cf6);
                        display: flex; align-items: center; justify-content: center;
                        margin: 0 auto 24px;
                        font-size: 36px;
                        box-shadow: 0 8px 30px rgba(14, 165, 233, 0.3);
                        animation: pulseGlow 3s ease-in-out infinite;
                    }
                </style>
                <div class="login-premium-card">
                    <div class="login-premium-icon">🔐</div>
                    <h3 style='color: #f8fafc; margin-bottom: 5px; font-weight: 700; font-size: 26px; font-family: "Outfit", sans-serif;'>{t("Sistema RDC & PDE")}</h3>
                    <p style='color: #0ea5e9; font-size: 12px; font-weight: 600; letter-spacing: 2px;'>{t("ACESSO RESTRITO")}</p>
                </div>
            """, unsafe_allow_html=True)
            
            user_input = st.text_input(t("Usuário (Login):"), placeholder="Digite sua credencial")
            pass_input = st.text_input(t("Senha:"), type="password", placeholder="••••••••")
            lembrar_me = st.checkbox(t("Manter conectado"), value=True)
            
            st.markdown("<br>", unsafe_allow_html=True)
            if st.button("Entrar no Sistema", type="primary", use_container_width=True):
                user_clean = user_input.strip().upper()
                pass_clean = pass_input.strip()
                
                user_encontrado = None
                for key_db in usuarios_db.keys():
                    if key_db.strip().upper() == user_clean:
                        user_encontrado = key_db
                        break
                
                if user_encontrado and usuarios_db[user_encontrado]["senha"] == pass_clean:
                    st.session_state.usuario_logado = user_encontrado
                    st.session_state.role_usuario = usuarios_db[user_encontrado].get("role", "user")
                    st.session_state.nome_completo = usuarios_db[user_encontrado].get("nome", user_encontrado)
                    
                    if lembrar_me:
                        cookie_manager.set("rdc_user_session", user_encontrado, expires_at=datetime.datetime.now() + datetime.timedelta(days=30))
                        
                    time.sleep(1) # Tempo para o cookie assentar
                    st.rerun()
                else:
                    st.error("Usuário ou senha incorretos. Verifique espaços em branco ou letras erradas.")
    st.stop() # Bloqueia todo o resto do sistema!


# =================================================================


# =================================================================
# CABEÇALHO GLOBAL (Mostrado apenas se logado)
# =================================================================
# Status bar with live info
import datetime as dt_mod
hora_agora = dt_mod.datetime.now().strftime("%H:%M")
data_agora = dt_mod.datetime.now().strftime("%d/%m/%Y")
nome_user_logado = st.session_state.get('nome_completo', 'Admin')

st.markdown(f"""
    <div class="enesa-header">
        <div style="display: flex; justify-content: space-between; align-items: center; flex-wrap: wrap; gap: 12px;">
            <div>
                <div style="display: flex; align-items: center; gap: 10px; margin-bottom: 4px;">
                    <h1 style="margin: 0; font-size: 1.7rem; font-weight: 700;">
                        <span style="background: linear-gradient(135deg, #0ea5e9, #8b5cf6); -webkit-background-clip: text; -webkit-text-fill-color: transparent;">Sistema de Gestao RDC & PDE</span>
                    </h1>
                    <span style="background: rgba(14, 165, 233, 0.15); border: 1px solid rgba(14, 165, 233, 0.25); border-radius: 6px; padding: 2px 8px; font-size: 10px; color: #0ea5e9; font-weight: 700; letter-spacing: 1px;">v7.0</span>
                </div>
                <p style="color: {cor_texto_sub}; font-size: 0.82rem; margin: 0; letter-spacing: 0.5px;">{nome_site} — Controle Operacional de Efetivo</p>
            </div>
            <div style="display: flex; gap: 10px; align-items: center; flex-wrap: wrap;">
                <div style="background: rgba(16, 185, 129, 0.12); border: 1px solid rgba(16, 185, 129, 0.25); border-radius: 20px; padding: 5px 14px; font-size: 11px; color: #10b981; font-weight: 600; letter-spacing: 0.5px; display: flex; align-items: center; gap: 6px;">
                    <span style="width: 7px; height: 7px; border-radius: 50%; background: #10b981; display: inline-block; box-shadow: 0 0 6px #10b981; animation: pulse 2s infinite;"></span>
                    ONLINE
                </div>
                <div style="background: rgba(14, 165, 233, 0.08); border: 1px solid rgba(14, 165, 233, 0.15); border-radius: 10px; padding: 6px 14px; font-size: 12px; color: #94a3b8;">
                    <span style="color: #0ea5e9; font-weight: 600;">{nome_user_logado}</span> · {data_agora}
                </div>
            </div>
        </div>
    </div>
    <style>
        @keyframes pulse {{
            0%, 100% {{ opacity: 1; }}
            50% {{ opacity: 0.4; }}
        }}
    </style>
""", unsafe_allow_html=True)

# =================================================================
# BARRA LATERAL
# =================================================================
with st.sidebar:
    if os.path.exists(caminho_logo):
        col1, col2, col3 = st.columns([1.5, 2, 1.5]) 
        with col2:
            st.image(caminho_logo, use_container_width=True)
        st.markdown("<br>", unsafe_allow_html=True)
        
    html_avatar = """
    <div style="display: flex; align-items: center; gap: 15px; padding: 15px; background: rgba(14, 165, 233, 0.1); border-radius: 12px; border: 1px solid rgba(14, 165, 233, 0.3); margin-bottom: 25px; box-shadow: 0 0 20px rgba(14, 165, 233, 0.2); transition: transform 0.3s;" onmouseover="this.style.transform='scale(1.05)'" onmouseout="this.style.transform='scale(1)'">
        <div style="width: 50px; height: 50px; border-radius: 50%; background: linear-gradient(135deg, #0ea5e9, #8b5cf6); display: flex; justify-content: center; align-items: center; font-size: 24px; color: white; box-shadow: 0 0 15px rgba(14, 165, 233, 0.5);">
            👨‍💻
        </div>
        <div>
            <div style="font-size: 14px; font-weight: 800; color: #f8fafc; text-shadow: 0 0 10px rgba(255,255,255,0.3); letter-spacing: 0.5px;">ADMINISTRADOR</div>
            <div style="font-size: 12px; color: #0ea5e9; font-weight: bold; margin-top: 2px; text-shadow: 0 0 5px rgba(14,165,233,0.5);">Acesso Supremo</div>
        </div>
    </div>
    """
    st.markdown(html_avatar, unsafe_allow_html=True)
    
    st.header("📂 Arquivos Base")
    
    if st.button("➕ Enviar Nova Base (PDE)", use_container_width=True):
        st.session_state.mostrar_upload = not st.session_state.mostrar_upload
        
    arquivo_pde = None
    arquivo_modelo = None
    
    if st.session_state.mostrar_upload:
        st.markdown("<div style='background-color: #22262e; padding: 10px; border-radius: 8px;'>", unsafe_allow_html=True)
        arquivo_pde = st.file_uploader("Base de Efetivo (.csv/.xlsx):", type=["csv", "xlsx"])
        arquivo_modelo = st.file_uploader("📄 Layout MODELO.xlsx:", type=["xlsx"])
        st.markdown("</div>", unsafe_allow_html=True)
        
        if arquivo_pde is not None:
            if salvar_base_localmente(arquivo_pde):
                st.success("💾 Base de Efetivo salva localmente!")
                
        if arquivo_modelo is not None:
            if salvar_modelo_no_disco(arquivo_modelo):
                st.success("💾 Modelo salvo!")

    
    st.markdown("---")
    
    base_existe = os.path.exists(caminho_base_salva_csv) or os.path.exists(caminho_base_salva_xlsx)
    if base_existe:
        st.success("✅ Base salva no sistema.")
    else:
        st.info("ℹ️ Nenhuma base salva ainda.")
    
    st.markdown("---")
    
    st.markdown("---")
    
    # === MODO TV ===
    if st.button("📺 Modo TV (Apresentação)", use_container_width=True, type="secondary"):
        st.session_state.modo_tv = True
        st.session_state.tv_slide = 0
        st.rerun()
    
    st.markdown(f"👤 Bem-vindo(a), **{st.session_state.nome_completo}**")
    
    if st.button("Sair (Logout)", use_container_width=True):
        cookie_manager.delete("rdc_user_session")
        st.session_state.usuario_logado = None
        st.session_state.role_usuario = None
        st.session_state.nome_completo = None
        time.sleep(1)
        st.rerun()
        
    if st.session_state.role_usuario == "admin":
        st.markdown("---")
        st.markdown("#### ⚙️ Painel de Configurações")
        
        # --- Seletor de Idioma ---
        idioma_opcoes = ["Português", "English"]
        idioma_atual = st.session_state.get("idioma", "Português")
        idx_idioma = idioma_opcoes.index(idioma_atual) if idioma_atual in idioma_opcoes else 0
        idioma_sel = st.selectbox("🌐 Idioma / Language", idioma_opcoes, index=idx_idioma, key="sel_idioma")
        if idioma_sel != st.session_state.get("idioma", "Português"):
            st.session_state.idioma = idioma_sel
            st.rerun()
        
        # --- Seletor de Modelo Gemini ---
        modelos_gemini = ["gemini-2.5-flash", "gemini-2.5-pro", "gemini-1.5-pro", "gemini-1.5-flash"]
        modelo_atual = st.session_state.get("modelo_gemini", "gemini-2.5-flash")
        idx_modelo = modelos_gemini.index(modelo_atual) if modelo_atual in modelos_gemini else 0
        modelo_sel = st.selectbox("🤖 Modelo de IA (Gemini)", modelos_gemini, index=idx_modelo, key="sel_modelo_gemini")
        if modelo_sel != st.session_state.get("modelo_gemini", "gemini-2.5-flash"):
            st.session_state.modelo_gemini = modelo_sel
            st.rerun()
        
        st.markdown("---")
        
        if st.toggle("🔑 Ver Usuários e Senhas"):
            usuarios_carregados = carregar_usuarios()
            dados_usuarios = []
            for u_nome, u_dados in sorted(usuarios_carregados.items()):
                dados_usuarios.append({
                    "Login": u_nome,
                    "Senha": u_dados.get("senha", ""),
                    "Acesso": u_dados.get("role", "user")
                })
            df_usuarios = pd.DataFrame(dados_usuarios)
            st.dataframe(df_usuarios, hide_index=True, use_container_width=True)
            
        novo_logo = st.file_uploader("Trocar Logo (PNG/JPG):", type=["png", "jpg", "jpeg"])
        if novo_logo:
            with open(caminho_logo, "wb") as f:
                f.write(novo_logo.getbuffer())
            st.success("Logo atualizado! Recarregue a página.")
            
        novo_nome_site = st.text_input("Nome da Empresa/Site:", value=nome_site)
        if st.button("Salvar Nome"):
            with open(caminho_nome_site, "w", encoding="utf-8") as f:
                f.write(novo_nome_site)
            st.success("Nome atualizado!")
            time.sleep(1)
            st.rerun()
                
        st.markdown("---")
        st.markdown("**💾 Backup Seguro**")
        
        # Função para gerar backup ZIP
        buffer_zip = io.BytesIO()
        with zipfile.ZipFile(buffer_zip, "w") as z:
            # Backup da Base de Efetivo
            if st.session_state.df is not None:
                buffer_pde = io.BytesIO()
                st.session_state.df.to_excel(buffer_pde, index=False, engine='openpyxl')
                z.writestr("BASE_EFETIVO.xlsx", buffer_pde.getvalue())
            
            # Backup do Histórico F1
            if "df_historico_f1" in st.session_state and not st.session_state.df_historico_f1.empty:
                buffer_f1 = io.BytesIO()
                st.session_state.df_historico_f1.to_excel(buffer_f1, index=False, engine='openpyxl')
                z.writestr("HISTORICO_F1.xlsx", buffer_f1.getvalue())
        
        st.download_button(
            label="📥 Baixar Backup (.zip)",
            data=buffer_zip.getvalue(),
            file_name=f"Backup_RDC_{datetime.datetime.now().strftime('%Y%m%d')}.zip",
            mime="application/zip",
            use_container_width=True
        )

        st.markdown("---")
        st.markdown("#### 👥 Gestão de Usuários")
        with st.form("form_novo_usuario"):
            st.markdown("**Adicionar / Editar Usuário**")
            novo_user = st.text_input("Usuário (Login):")
            nova_senha = st.text_input("Senha:")
            novo_nome = st.text_input("Nome Completo:")
            nova_role = st.selectbox("Nível de Acesso:", ["user", "admin", "apontador"])
            submit_user = st.form_submit_button("Salvar Usuário")
            if submit_user and novo_user and nova_senha:
                usuarios_db[novo_user] = {"senha": nova_senha, "nome": novo_nome, "role": nova_role}
                salvar_usuarios(usuarios_db)
                st.success(f"Usuário '{novo_user}' salvo!")
                time.sleep(1)
                st.rerun()
        
        st.markdown("**Usuários Cadastrados:**")
        for u, dados in sorted(usuarios_db.items()):
            col_u, col_del = st.columns([4, 1])
            if u == "admin":
                col_u.markdown(f"👤 **{u}** (admin)")
            else:
                current_role = dados.get('role', 'user')
                roles_options = ["user", "admin", "apontador"]
                idx = roles_options.index(current_role) if current_role in roles_options else 0
                
                new_role = col_u.selectbox(f"👤 {u}", roles_options, index=idx, key=f"role_{u}")
                
                if new_role != current_role:
                    usuarios_db[u]['role'] = new_role
                    salvar_usuarios(usuarios_db)
                    st.rerun()
                    
                with col_del:
                    st.markdown("<div style='height:28px'></div>", unsafe_allow_html=True)
                    if st.button("❌", key=f"del_{u}"):
                        del usuarios_db[u]
                        salvar_usuarios(usuarios_db)
                        st.rerun()

        st.markdown("---")
        if st.toggle("🏎️ Gerenciar Lista F1", key="toggle_f1_config"):
            import json as json_mod
            caminho_f1 = "encarregados_f1.json"
            try:
                with open(caminho_f1, "r", encoding="utf-8") as f_f1:
                    lista_f1 = json_mod.load(f_f1)
            except:
                lista_f1 = []
            
            col_add_f1, col_rem_f1 = st.columns(2)
            with col_add_f1:
                novo_enc = st.text_input("➕ Adicionar Encarregado ao F1:", key="add_enc_f1_config")
                if st.button("Adicionar", key="btn_add_f1_config", use_container_width=True):
                    if novo_enc.strip():
                        nome_up = novo_enc.strip().upper()
                        if nome_up not in [e.upper() for e in lista_f1]:
                            lista_f1.append(nome_up)
                            with open(caminho_f1, "w", encoding="utf-8") as f_f1:
                                json_mod.dump(lista_f1, f_f1, ensure_ascii=False, indent=2)
                            st.success(f"'{nome_up}' adicionado!")
                            time.sleep(1)
                            st.rerun()
                        else:
                            st.warning("Já existe na lista.")
            with col_rem_f1:
                enc_rem = st.multiselect("🗑️ Remover do F1:", sorted([e.upper() for e in lista_f1]), key="rem_enc_f1_config")
                if st.button("Remover", key="btn_rem_f1_config", use_container_width=True):
                    if enc_rem:
                        lista_f1 = [e for e in lista_f1 if e.upper() not in enc_rem]
                        with open(caminho_f1, "w", encoding="utf-8") as f_f1:
                            json_mod.dump(lista_f1, f_f1, ensure_ascii=False, indent=2)
                        st.success(f"{len(enc_rem)} removido(s)!")
                        time.sleep(1)
                        st.rerun()

    st.markdown("---")
    st.markdown(
        f"""
        <div style='text-align: center; margin-top: 30px; padding: 24px 16px; background: linear-gradient(135deg, rgba(15, 23, 42, 0.7), rgba(30, 41, 59, 0.4)); border-radius: 16px; border: 1px solid rgba(255,255,255,0.06); backdrop-filter: blur(8px);'>
            <div style='display: flex; justify-content: center; align-items: center; gap: 24px; flex-wrap: wrap; margin-bottom: 12px;'>
                <div style='display: flex; align-items: center; gap: 6px;'>
                    <span style='width: 8px; height: 8px; border-radius: 50%; background: #10b981; display: inline-block; box-shadow: 0 0 8px #10b981;'></span>
                    <span style='font-size: 11px; color: #94a3b8; font-weight: 500;'>Sistema Operacional</span>
                </div>
                <span style='font-size: 10px; color: #334155;'>|</span>
                <span style='font-size: 11px; color: #64748b;'>📅 Última att: 15/07/2026</span>
                <span style='font-size: 10px; color: #334155;'>|</span>
                <div style='display: inline-block; background: rgba(14, 165, 233, 0.1); border: 1px solid rgba(14, 165, 233, 0.2); border-radius: 20px; padding: 2px 12px;'>
                    <span style='font-size: 10px; color: #0ea5e9; font-weight: 700; letter-spacing: 1px;'>v7.0</span>
                </div>
            </div>
            <div style='border-top: 1px solid rgba(255,255,255,0.04); padding-top: 12px;'>
                <p style='font-size: 10px; color: #475569; letter-spacing: 2px; text-transform: uppercase; margin: 0 0 4px 0;'>Desenvolvido por</p>
                <p style='font-size: 14px; font-weight: 700; margin: 0; background: linear-gradient(135deg, #0ea5e9, #8b5cf6); -webkit-background-clip: text; -webkit-text-fill-color: transparent;'>Edson Garcia · {nome_site}</p>
            </div>
        </div>
        """, 
        unsafe_allow_html=True
    )

# =================================================================
# LÓGICA DE CARREGAMENTO DA NUVEM (GOOGLE SHEETS)
# =================================================================
try:
    conn = st.connection("gsheets", type=GSheetsConnection)
except Exception:
    conn = None

if arquivo_pde is not None:
    df_carregado = ler_arquivo_seguro(arquivo_pde, arquivo_pde.name)
    if df_carregado is not None:
        # PREPARAR PRIMEIRO para padronizar nomes de colunas (MATRÍCULA -> MATRICULA, etc.)
        df_carregado = preparar_dataframe(df_carregado)
        
        # PRESERVAR O C.C DA BASE ANTIGA DO APP
        if st.session_state.df is not None:
            # Prepara a base antiga para garantir que o nome da coluna é exatamente "C.C" e "MATRICULA"
            df_antigo = preparar_dataframe(st.session_state.df)
            
            if "C.C" in df_antigo.columns and "MATRICULA" in df_carregado.columns:
                try:
                    df_com_cc = df_antigo[df_antigo["C.C"].astype(str).str.strip() != ""]
                    
                    # O C.C do PDE agora é a verdade absoluta.
                    # Apenas herdamos o C.C do Encarregado caso o colaborador venha sem C.C no PDE.
                    # Passo 2: Herdar o C.C do Encarregado para novos colaboradores (ou transferidos sem C.C)
                    if "ENCARREGADO" in df_carregado.columns and "ENCARREGADO" in df_com_cc.columns:
                        mapa_enc_cc = df_com_cc.groupby("ENCARREGADO")["C.C"].agg(lambda x: x.value_counts().index[0] if len(x) > 0 else "").to_dict()
                        
                        def herdar_cc(row):
                            cc_atual = str(row.get("C.C", "")).strip()
                            if not cc_atual:
                                return mapa_enc_cc.get(row.get("ENCARREGADO", ""), "")
                            return cc_atual
                            
                        df_carregado["C.C"] = df_carregado.apply(herdar_cc, axis=1)
                        
                except Exception:
                    pass
        
        st.session_state.df = df_carregado
        if conn and not st.session_state.get('force_use_local', False):
            try:
                conn.update(worksheet="PDE", data=st.session_state.df)
                st.toast("☁️ Base salva! C.Cs preservados com sucesso!", icon="✅")
            except Exception as e:
                st.sidebar.error(f"Erro Nuvem: {e}")

elif st.session_state.df is None:
    carregado_nuvem = False
    
    # 1. LER DA NOVA PLANILHA MESTRE DO GOOGLE SHEETS
    if not st.session_state.get('force_use_local', False):
        if conn:
            try:
                df_gsheets = conn.read(worksheet="PDE", ttl=5)
                df_gsheets = df_gsheets.dropna(how='all')
                if not df_gsheets.empty:
                    st.session_state.df = preparar_dataframe(df_gsheets)
                    carregado_nuvem = True
                    st.toast(f"PDE Mestre carregado! {len(df_gsheets)} funcionários.", icon="☁️")
            except Exception as e:
                st.sidebar.warning(f"⚠️ Erro ao ler PDE Mestre autenticado: {e}")
            
    # Resetar a flag (dentro do elif st.session_state.df is None)
    if st.session_state.get('force_use_local', False):
        carregado_nuvem = True
        st.session_state.force_use_local = False
            
    if not carregado_nuvem:
        if os.path.exists(caminho_base_salva_xlsx):
            df_carregado = ler_arquivo_seguro(caminho_base_salva_xlsx, "BASE_ATUAL.xlsx")
            if df_carregado is not None:
                st.session_state.df = df_carregado
        elif os.path.exists(caminho_base_salva_csv):
            df_carregado = ler_arquivo_seguro(caminho_base_salva_csv, "BASE_ATUAL.csv")
            if df_carregado is not None:
                st.session_state.df = df_carregado
        elif os.path.exists(caminho_pde_padrao):
            df_carregado = ler_arquivo_seguro(caminho_pde_padrao, "PDE.csv")
            if df_carregado is not None:
                st.session_state.df = df_carregado

# =================================================================
# SEMPRE VERIFICAR O HISTÓRICO F1 NA NUVEM (COM PROTEÇÃO ANTI-PERDA)
# =================================================================
if conn and not st.session_state.get('force_use_local', False):
    try:
        df_f1 = conn.read(worksheet="Historico_F1", ttl=0)
        if df_f1 is not None:
            df_f1 = df_f1.dropna(how='all')
            # Filtrar apenas registros válidos (com ENCARREGADO preenchido)
            if not df_f1.empty and 'ENCARREGADO' in df_f1.columns:
                df_f1 = df_f1[df_f1['ENCARREGADO'].notna() & (df_f1['ENCARREGADO'] != '')]
            
            qtd_nuvem = len(df_f1) if not df_f1.empty else 0
            qtd_local = len(st.session_state.df_historico_f1) if not st.session_state.df_historico_f1.empty else 0
            
            # SÓ ATUALIZAR SE: nuvem tem dados E (nuvem tem MAIS dados OU local está vazio)
            # Isso IMPEDE que uma leitura vazia da nuvem apague dados bons
            if qtd_nuvem > 0 and (qtd_nuvem >= qtd_local or qtd_local == 0):
                st.session_state.df_historico_f1 = df_f1
                df_f1.to_csv(caminho_historico_f1_csv, index=False)
            elif qtd_nuvem == 0 and qtd_local > 0:
                # Nuvem está vazia mas local tem dados - REENVIAR para a nuvem!
                ok, msg = salvar_f1_seguro(conn, st.session_state.df_historico_f1, caminho_historico_f1_csv)
                if ok:
                    st.sidebar.caption(f"☁️ F1: {qtd_local} registros reenviados à nuvem (estava vazia).")
    except Exception:
        pass

# =================================================================
# CONTEUDO PRINCIPAL
# =================================================================
if st.session_state.df is not None:
    df_atual = preparar_dataframe(st.session_state.df.copy())
    lista_encarregados_base = sorted([str(e) for e in df_atual["ENCARREGADO"].unique() if str(e).strip() != ""])

    # ====== HISTÓRICO DE C.C (FOTOGRAFIA DIÁRIA) ======
    try:
        hoje_hist = datetime.date.today().strftime("%Y-%m-%d")
        df_cc_valido = df_atual[df_atual["C.C"].astype(str).str.strip() != ""]
        if not df_cc_valido.empty:
            cc_counts = df_cc_valido["C.C"].value_counts().reset_index()
            cc_counts.columns = ["C.C", "Efetivo"]
            cc_counts["DATA"] = hoje_hist
            
            hist_cc_existente = pd.DataFrame()
            if os.path.exists(caminho_hist_cc):
                try:
                    hist_cc_existente = pd.read_csv(caminho_hist_cc)
                except:
                    pass
                
            if not hist_cc_existente.empty and "DATA" in hist_cc_existente.columns:
                hist_cc_existente = hist_cc_existente[hist_cc_existente["DATA"] != hoje_hist]
                hist_cc_novo = pd.concat([hist_cc_existente, cc_counts], ignore_index=True)
            else:
                hist_cc_novo = cc_counts
                
            hist_cc_novo.to_csv(caminho_hist_cc, index=False)
    except Exception as e:
        pass # Falha silenciosa para não quebrar o app
    # ==================================================

    # === LISTA DINÂMICA DE ENCARREGADOS (Carrega do JSON ou cria com a lista padrão) ===
    import json
    caminho_f1_json = os.path.join(os.path.dirname(__file__), "encarregados_f1.json")
    caminho_f1_excecoes = os.path.join(os.path.dirname(__file__), "f1_excecoes.csv")
    
    encarregados_f1_padrao = [
        "ABMAEL PEREIRA PAIVA", "JEAN PEDRO", "ANANIAS DE SOUSA NETO", "GILDO GONCALVES DA SILVA",
        "SIDNEI FERNANDES DA SILVA", "BARTOLOMEU FERNANDES", "FRANCINALDO DE SOUSA", "IZAIAS BAIA BELO",
        "SANDRO LIMA DE SOUZA", "ALOISIO FERREIRA SOUZA", "ARLINDO PEREIRA DA SILVA", "FAUZE CELIS RODRIGUES COSTA",
        "FRANCISCO PEREIRA LIMA", "JOAO PAULO DA COSTA QUARESMA", "JOSE ORLANDO DAS NEVES MADEIRA",
        "JOSE TARCISIO ARAUJO DA SILVA", "LEANDRO DA CRUZ DE SOUZA", "CLAUDIO LUCIANO ARGELINO",
        "EDVALDO CARVALHO ANGELIM", "ELDER MENDES JUNIOR", "MANOEL MARIA SARGES SOARES", "CLAUDIO CRUZ SOUSA",
        "CLIDENILDO GOMES DE ALMEIDA", "GRACINEI PEREIRA DOS SANTOS", "JAILSON MENDES DE OLIVEIRA",
        "JARBAS DA ROCHA GOMES", "JOSE MAURICIO RODRIGUES DA SILVA", "JOSE SARAIVA LOPES NETO",
        "JOSMAEL RODRIGUES PEREIRA", "ALEX PANTOJA DE OLIVEIRA", "ARILSON DIAS DO PRADO", "ELTON GOMES DOS SANTOS",
        "RICARDO SARMENTO FERREIRA", "WENISON DA SILVA CUNHA CORREIA", "FRANCISCO ALVES DA PENHA",
        "IVAN DO NASCIMENTO RAMOS", "ELDER MENDES", "GEAN LENO JOSE DE FREITAS", "JOSE EDUARDO FARIAS FERREIRA",
        "EDIMILSON NUNES VASCONCELOS", "LOURISVALDO AMARAL ARAUJO", "VALDEMIR BARBOSA REIS",
        "LUZINALDO AMARAL DE ARAUJO", "MAURO DE QUEIROZ ANDRADE", "ELIAS SOUSA DA COSTA", "ISAIAS SOUSA LISBOA",
        "ISMAEL CARLOS GOMES DA SILVA", "RAIMUNDO DA SILVA DOS SANTOS", "RAIMUNDO EUDE DA SILVA FREITAS",
        "RODOLFO DOS SANTOS COSTA", "ELISEU DA SILVA BISPO", "IRON MARQUES MOREIRA", "LUIZ CARLOS DE SOUZA",
        "ANTONIO TEIXEIRA BORBA", "JOSE FRANCIVAN MONTEIRO SANTOS", "JOSE WALKER CARNEIRO OLIVEIRA",
        "LEANDRO DA SILVA QUEIROZ", "SILVIO MANOEL DE ANDRADE", "EVERALDO DOS SANTOS SOARES",
        "FRANCISCO GRACIEL DE SOUSA MARTINS", "JAILSON SILVA DE GOIS", "JORGINALDO NUNES DA SILVA",
        "CLAUDIVAN OLIVEIRA DOS SANTOS", "GUILHERME HENRIQUE DE ARAUJO SOUSA", "LEANDRO MARTINS DA SILVA BORGES",
        "WEVERTON FERNANDES MARIANO", "JORGE DA COSTA SILVA", "RAIMUNDO FRAZAO DOS SANTOS",
        "JOSE RIBEIRO DO NASCIMENTO JUNIOR", "JOSE ROBERTO SALVADOR FILHO", "MARCUS ANTONIO DE SOUZA",
        "RAIMUNDO ROGERIO LEITE", "ROUBERVAL SANTOS DOS SANTOS", "CARLOS ALBERTO DA COSTA MOREIRA",
        "JOSE FELIPE DOS SANTOS", "JOSE GERIARDI FONSECA DE SENA", "JOSE HENRIQUE SILVA VIEIRA",
        "ODAIR MENEZES DA SILVA", "SIDNALDO SANTOS DE JESUS", "ANDERSON VICTALINO",
        "FRANCISCO AUGUSTO DE SOUSA BARROS", "GENILSON PEREIRA DE SOUSA", "HELENO MARQUES DE SOUZA NETO",
        "HEMERSON MONTEIRO DE OLIVEIRA", "JACKSON DEIBSON FELICIANO DA SILVA", "JARDELINO PEREIRA DA COSTA",
        "JOAO TIAGO OLIVEIRA DE AMORIM", "JOSE MARIA DA SILVA PESSOA", "LUCIO FABIO DA SILVA LEANDRO",
        "RAIMUNDO GONCALVES DOS SANTOS", "FABRICIO FIGUEIREDO", "RHOKSONY FERREIRA SILVEIRA",
        "FERNANDO DA CONCEIÇÃO", "ROGERIO BARROS DOS SANTOS", "SIQUEU SANTOS SOLEDADE",
        "SEBASTIAO CARLOS DE OLIVEIRA", "MANOEL NEPOMUCENO DOS SANTOS", "LUIZ RAMOS DE LIMA",
        "JORGE LUIS LOPES", "VALDINEI GOMES OLIVEIRA", "CARLOS DA SILVA OLIVEIRA"
    ]
    
    # Carregar ou criar o JSON
    if os.path.exists(caminho_f1_json):
        try:
            with open(caminho_f1_json, "r", encoding="utf-8") as f:
                encarregados_f1_oficial = json.load(f)
        except Exception:
            encarregados_f1_oficial = encarregados_f1_padrao
    else:
        encarregados_f1_oficial = encarregados_f1_padrao
        with open(caminho_f1_json, "w", encoding="utf-8") as f:
            json.dump(encarregados_f1_padrao, f, ensure_ascii=False, indent=2)
    
    lista_completa_encarregados = sorted([str(e).upper().strip() for e in df_atual["ENCARREGADO"].unique() if str(e).strip() != ""])
    
    # Carregar exceções (Abonos)
    if "df_f1_excecoes" not in st.session_state:
        if os.path.exists(caminho_f1_excecoes):
            try:
                st.session_state.df_f1_excecoes = pd.read_csv(caminho_f1_excecoes)
            except Exception:
                st.session_state.df_f1_excecoes = pd.DataFrame(columns=["DATA", "ENCARREGADO", "MOTIVO"])
        else:
            st.session_state.df_f1_excecoes = pd.DataFrame(columns=["DATA", "ENCARREGADO", "MOTIVO"])

    # =================================================================
    # MODO ENCARREGADO (Lançamento Nativo com Formatação Original)
    # =================================================================
    if st.session_state.get("role") == "encarregado":
        st.markdown("### <span class='material-symbols-rounded' style='vertical-align: middle; color: #0ea5e9; font-size: 32px;'>edit_document</span> Lançamento de RDC Digital", unsafe_allow_html=True)
        st.caption("Preencha as informações do seu dia de trabalho seguindo as 3 etapas abaixo. Os dados serão salvos na nuvem.")
        
        with st.form("form_rdc_digital_encarregado"):
            tab_id, tab_local, tab_ativ = st.tabs(["1️⃣ Identificação", "2️⃣ Localização", "3️⃣ Atividades e Envio"])
            
            with tab_id:
                st.markdown("<p style='color: #94a3b8; font-size: 14px;'>Quem é você e qual seu turno?</p>", unsafe_allow_html=True)
                rdc_encarregado = st.selectbox("Selecione seu Nome (Encarregado):", [""] + lista_completa_encarregados)
                rdc_turno = st.selectbox("Turno de Trabalho:", ["DIURNO", "NOTURNO", "MISTO"])
                
            with tab_local:
                import datetime
                st.markdown("<p style='color: #94a3b8; font-size: 14px;'>Onde você trabalhou hoje?</p>", unsafe_allow_html=True)
                
                rdc_data = st.date_input("Data do Relatório:", datetime.date.today())
                
                area_options = ["PB", "RB", "ESP", "LAYDOWN 1", "LAYDOWN 2", "OUTRO (DIGITAR)"]
                area_sel = st.selectbox("Área / Local de Trabalho:", area_options)
                rdc_area = area_sel
                if area_sel == "OUTRO (DIGITAR)":
                    rdc_area = st.text_input("Qual Área/Local?", placeholder="Ex: Escritório, Almoxarifado...")
                
                disc_options = [
                    "EQUIPAMENTOS", "DUTOS", "TUBULACAO", "ESTRUTURA METALICA", "PRECIPITADOR", 
                    "PRESSAO - MECANICA", "PRESSAO - TUBULACAO", "PRESSAO - FORNALHA", "PINTURA", 
                    "COMISSIONAMENTO", "OP. ASSISTIDA", "LAVAGEM QUIMICA", "SOPRAGEM", "ANDAIME", 
                    "OPERADORES", "FORA DE ESCOPO", "GERENCIA", "PRODUCAO", "GARANTIA DA QUALIDADE", 
                    "PLANEJAMENTO", "ADMINISTRACAO", "SEGURANCA E MEDICINA DO TRABALHO", "INFRAESTRUTURA", 
                    "ALMOXARIFADO ENESA", "ALMOXARIFADO MATERIAIS", "MANUT. ELETRICA PROVISORIA", 
                    "TOPOGRAFIA", "MOVIMENTACAO DE CARGAS", "MEDICAO/CUSTO/CONTRATOS", "CIVIL", "MECÂNICA", "ELÉTRICA", "INSTRUMENTAÇÃO", "ISOLAMENTO", "OUTRA (DIGITAR)"
                ]
                disc_sel = st.selectbox("Disciplina Principal:", disc_options)
                
                rdc_disciplina = disc_sel
                if disc_sel == "OUTRA (DIGITAR)":
                    rdc_disciplina = st.text_input("Qual Disciplina?", placeholder="Ex: Tubulação, Solda...")
                    
            with tab_ativ:
                st.markdown("<p style='color: #94a3b8; font-size: 14px;'>O que foi executado?</p>", unsafe_allow_html=True)
                rdc_dds = st.text_input("Tópico do DDS do dia:")
                rdc_atividades = st.text_area("Atividades Executadas (Detalhe os serviços feitos pela equipe):", height=150)
                rdc_problemas = st.text_area("Problemas / Interrupções / Ocorrências (Opcional):", height=68)
                
                st.markdown("<br>", unsafe_allow_html=True)
                submit_rdc = st.form_submit_button("🚀 Salvar e Enviar RDC na Nuvem", use_container_width=True, type="primary")
            
            if submit_rdc:
                if not rdc_encarregado:
                    st.error("⚠️ Por favor, selecione o nome do Encarregado.")
                elif not rdc_atividades.strip():
                    st.error("⚠️ Por favor, preencha as Atividades Executadas.")
                elif disc_sel == "OUTRA (DIGITAR)" and not rdc_disciplina.strip():
                    st.error("⚠️ Digite a disciplina na caixa 'Qual Disciplina?'.")
                else:
                    rdc_json = [{
                        "ENCARREGADO": rdc_encarregado,
                        "DATA": rdc_data.strftime("%Y/%m/%d"),
                        "TURNO": rdc_turno,
                        "AREA": rdc_area.strip().upper(),
                        "DISCIPLINA": rdc_disciplina.strip().upper(),
                        "DDS": rdc_dds.strip(),
                        "ATIVIDADE": rdc_atividades.strip(),
                        "CALDEIRA": rdc_problemas.strip(),
                        "PROBLEMAS": rdc_problemas.strip()
                    }]
                    
                    import json
                    import requests
                    
                    WEBHOOK_URL = "https://script.google.com/macros/s/AKfycbxfE96gE7ckdmapBLBHJuoX2bvAt-2d76OUJNiSRsLgFCOiySeQhFOopp3DoC5Fn95D/exec"
                    
                    try:
                        with st.spinner("Enviando dados para a nuvem..."):
                            res = requests.post(WEBHOOK_URL, json=rdc_json, allow_redirects=True)
                        if res.status_code == 200:
                            st.toast(f"RDC Digital de {rdc_encarregado} salvo com sucesso na Nuvem!", icon="✅")
                            st.info("O relatório já foi enviado. Você pode fechar esta página.")
                        else:
                            st.error(f"❌ Erro ao enviar. Servidor retornou: {res.text}")
                    except Exception as e:
                        st.error(f"❌ Falha de conexão: {e}")
                        
        st.stop() # Finaliza o script para não mostrar as abas do admin

    def gerar_relatorio_pdf(df):
        import io
        from fpdf import FPDF
        import datetime as dt_mod
        
        class PDF(FPDF):
            def header(self):
                self.set_font('helvetica', 'B', 15)
                self.cell(0, 10, 'Relatorio Executivo - Sistema RDO & PDE', new_x="LMARGIN", new_y="NEXT", align='C')
                self.set_font('helvetica', '', 10)
                agora = dt_mod.datetime.now().strftime("%d/%m/%Y %H:%M")
                self.cell(0, 5, f'Gerado em: {agora}', new_x="LMARGIN", new_y="NEXT", align='C')
                self.ln(10)
                
            def footer(self):
                self.set_y(-15)
                self.set_font('helvetica', 'I', 8)
                self.cell(0, 10, f'Pagina {self.page_no()}', align='C')

        pdf = PDF()
        pdf.add_page()
        
        # 1. Resumo de Efetivo
        pdf.set_font('helvetica', 'B', 12)
        pdf.set_fill_color(240, 240, 240)
        pdf.cell(0, 10, ' 1. Resumo de Efetivo Global', new_x="LMARGIN", new_y="NEXT", fill=True)
        pdf.ln(2)
        
        total = len(df)
        mod = len(df[df["MÃO DE OBRA"].str.strip().str.upper() == "MOD"])
        moi = len(df[df["MÃO DE OBRA"].str.strip().str.upper() == "MOI"])
        
        pdf.set_font('helvetica', '', 11)
        pdf.cell(50, 8, f'Efetivo Total:', border=0)
        pdf.set_font('helvetica', 'B', 11)
        pdf.cell(50, 8, f'{total}', border=0, new_x="LMARGIN", new_y="NEXT")
        
        pdf.set_font('helvetica', '', 11)
        pdf.cell(50, 8, f'Mao de Obra Direta (MOD):', border=0)
        pdf.set_font('helvetica', 'B', 11)
        pdf.cell(50, 8, f'{mod}', border=0, new_x="LMARGIN", new_y="NEXT")
        
        pdf.set_font('helvetica', '', 11)
        pdf.cell(50, 8, f'Mao de Obra Indireta (MOI):', border=0)
        pdf.set_font('helvetica', 'B', 11)
        pdf.cell(50, 8, f'{moi}', border=0, new_x="LMARGIN", new_y="NEXT")
        pdf.ln(5)
        
        # 2. Distribuição por Área
        pdf.set_font('helvetica', 'B', 12)
        pdf.cell(0, 10, ' 2. Distribuicao por Area', new_x="LMARGIN", new_y="NEXT", fill=True)
        pdf.ln(2)
        
        df_area = df.copy()
        df_area['AREA'] = df_area['C.C'].apply(lambda x: 'PB' if '125.02' in str(x) and '.005' not in str(x) else ('RB' if '125.01' in str(x) and '.005' not in str(x) else ('ESP' if '.005' in str(x) else 'OUTROS')))
        contagem = df_area[df_area['AREA'] != 'OUTROS'].groupby('AREA').size()
        
        pdf.set_font('helvetica', '', 10)
        for area, qtd in contagem.items():
            pdf.cell(100, 8, f'Area {area}:', border=1)
            pdf.cell(40, 8, f'{qtd} funcionarios', border=1, new_x="LMARGIN", new_y="NEXT", align="R")
        pdf.ln(5)
        
        return bytes(pdf.output())

    def gerar_relatorio_pptx_dashboard(df, nome_site, caminho_logo, lista_completa_encarregados):
        import io
        import datetime as dt_mod
        from pptx import Presentation as PptxPresentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        prs = PptxPresentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Paleta de Cores
        COR_FUNDO = RGBColor(15, 23, 42)
        COR_AZUL = RGBColor(14, 165, 233)
        COR_BRANCO = RGBColor(224, 228, 234)
        COR_CINZA = RGBColor(148, 163, 184)
        COR_VERDE = RGBColor(34, 197, 94)
        COR_AMARELO = RGBColor(245, 158, 11)
        COR_CARD_BG = RGBColor(30, 41, 59)
        
        def set_slide_bg(slide, cor):
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = cor
        
        def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=COR_BRANCO, alignment=PP_ALIGN.LEFT):
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = color
            p.alignment = alignment
            return txBox
        
        def add_card(slide, left, top, width, height, titulo, valor, cor_valor=COR_AZUL):
            shape = slide.shapes.add_shape(1, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = COR_CARD_BG
            shape.line.color.rgb = RGBColor(51, 65, 85)
            shape.line.width = Pt(1)
            
            txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), Inches(0.35))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = titulo
            p.font.size = Pt(11)
            p.font.color.rgb = COR_CINZA
            
            txBox2 = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.5), width - Inches(0.3), Inches(0.55))
            tf2 = txBox2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = str(valor)
            p2.font.size = Pt(26)
            p2.font.bold = True
            p2.font.color.rgb = cor_valor

        # Cálculos de Métricas
        total_efetivo = len(df)
        qtd_mod = len(df[df["MÃO DE OBRA"].astype(str).str.strip().str.upper() == "MOD"])
        qtd_moi = len(df[df["MÃO DE OBRA"].astype(str).str.strip().str.upper() == "MOI"])
        total_mo = qtd_mod + qtd_moi
        pct_mod = round((qtd_mod / total_mo * 100), 1) if total_mo > 0 else 0
        
        encs_validos = [e for e in df["ENCARREGADO"].unique() if str(e).strip() != "" and str(e) in lista_completa_encarregados]
        qtd_enc = len(encs_validos)
        span = round(total_efetivo / qtd_enc, 1) if qtd_enc > 0 else 0
        total_funcoes = df["FUNÇÃO"].nunique() if "FUNÇÃO" in df.columns else 0

        # ============================================
        # SLIDE 1: CAPA
        # ============================================
        slide1 = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide1, COR_FUNDO)
        
        if os.path.exists(caminho_logo):
            try:
                slide1.shapes.add_picture(caminho_logo, Inches(5.4), Inches(1.0), height=Inches(1.5))
            except:
                pass
        
        line = slide1.shapes.add_shape(1, Inches(3), Inches(3.0), Inches(7.333), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = COR_AZUL
        line.line.fill.background()
        
        add_text_box(slide1, Inches(1.5), Inches(3.2), Inches(10.333), Inches(1.0),
            "CENTRO DE COMANDO & CONTROLE OPERACIONAL", font_size=32, bold=True, color=COR_BRANCO, alignment=PP_ALIGN.CENTER)
        add_text_box(slide1, Inches(1.5), Inches(4.2), Inches(10.333), Inches(0.8),
            nome_site, font_size=24, bold=False, color=COR_AZUL, alignment=PP_ALIGN.CENTER)
        add_text_box(slide1, Inches(1.5), Inches(5.0), Inches(10.333), Inches(0.6),
            f"Relatório de Efetivo e Produtividade — {dt_mod.datetime.now().strftime('%d/%m/%Y')}", font_size=16, color=COR_CINZA, alignment=PP_ALIGN.CENTER)
        add_text_box(slide1, Inches(1.5), Inches(5.7), Inches(10.333), Inches(0.5),
            f"Gerado em: {dt_mod.datetime.now().strftime('%d/%m/%Y às %H:%M')}", font_size=12, color=COR_CINZA, alignment=PP_ALIGN.CENTER)

        # ============================================
        # SLIDE 2: PAINEL DE EFETIVO E ESTRUTURA
        # ============================================
        slide2 = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide2, COR_FUNDO)
        
        add_text_box(slide2, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
            "INDICADORES DE EFETIVO GLOBAL", font_size=26, bold=True, color=COR_BRANCO)
        
        line2 = slide2.shapes.add_shape(1, Inches(0.5), Inches(0.95), Inches(12.333), Inches(0.03))
        line2.fill.solid()
        line2.fill.fore_color.rgb = COR_AZUL
        line2.line.fill.background()
        
        card_w = Inches(2.3)
        card_h = Inches(1.15)
        card_y = Inches(1.2)
        gap = Inches(0.2)
        start_x = Inches(0.5)
        
        add_card(slide2, start_x, card_y, card_w, card_h, "Efetivo Total", total_efetivo, COR_AZUL)
        add_card(slide2, start_x + card_w + gap, card_y, card_w, card_h, "Encarregados", qtd_enc, COR_VERDE)
        add_card(slide2, start_x + 2*(card_w + gap), card_y, card_w, card_h, "% MOD Global", f"{pct_mod}%", COR_AZUL)
        add_card(slide2, start_x + 3*(card_w + gap), card_y, card_w, card_h, "Funções Distintas", total_funcoes, COR_AMARELO)
        add_card(slide2, start_x + 4*(card_w + gap), card_y, card_w, card_h, "Span of Control", span, RGBColor(139, 92, 246))
        
        # Tabela 1: Mão de Obra (MOD vs MOI)
        add_text_box(slide2, Inches(0.5), Inches(2.7), Inches(5.8), Inches(0.5),
            "COMPOSIÇÃO DA MÃO DE OBRA", font_size=16, bold=True, color=COR_AZUL)
        
        tbl_mo = slide2.shapes.add_table(3, 3, Inches(0.5), Inches(3.3), Inches(5.8), Inches(1.2)).table
        tbl_mo.columns[0].width = Inches(2.8)
        tbl_mo.columns[1].width = Inches(1.5)
        tbl_mo.columns[2].width = Inches(1.5)
        
        for j, h in enumerate(["Classificação", "Efetivo", "Proporção"]):
            c = tbl_mo.cell(0, j)
            c.text = h
            c.fill.solid()
            c.fill.fore_color.rgb = COR_AZUL
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = COR_FUNDO
        
        mo_data = [
            ("Mão de Obra Direta (MOD)", str(qtd_mod), f"{pct_mod}%"),
            ("Mão de Obra Indireta (MOI)", str(qtd_moi), f"{round(100 - pct_mod, 1)}%")
        ]
        for i, (tipo, qtd_s, prop_s) in enumerate(mo_data):
            row_idx = i + 1
            for j, val in enumerate([tipo, qtd_s, prop_s]):
                c = tbl_mo.cell(row_idx, j)
                c.text = val
                c.fill.solid()
                c.fill.fore_color.rgb = COR_CARD_BG if row_idx % 2 == 0 else COR_FUNDO
                for p in c.text_frame.paragraphs:
                    p.font.size = Pt(10)
                    p.font.color.rgb = COR_BRANCO

        # Tabela 2: Efetivo por Área
        add_text_box(slide2, Inches(6.8), Inches(2.7), Inches(6.0), Inches(0.5),
            "DISTRIBUIÇÃO DE EFETIVO POR ÁREA", font_size=16, bold=True, color=COR_AZUL)
        
        df_area_dash = df.copy()
        df_area_dash['ÁREA_RESUMO'] = df_area_dash['C.C'].apply(lambda x: 'PB' if '125.02' in str(x) and '.005' not in str(x) else ('RB' if '125.01' in str(x) and '.005' not in str(x) else ('ESP' if '.005' in str(x) else 'OUTROS')))
        area_counts = df_area_dash[df_area_dash['ÁREA_RESUMO'] != 'OUTROS'].groupby('ÁREA_RESUMO').size().reset_index(name='Quantidade')
        
        if not area_counts.empty:
            rows_area = len(area_counts) + 1
            tbl_area = slide2.shapes.add_table(rows_area, 3, Inches(6.8), Inches(3.3), Inches(6.0), Inches(0.4 * rows_area)).table
            tbl_area.columns[0].width = Inches(2.5)
            tbl_area.columns[1].width = Inches(1.8)
            tbl_area.columns[2].width = Inches(1.7)
            
            for j, h in enumerate(["Área", "Colaboradores", "% do Efetivo"]):
                c = tbl_area.cell(0, j)
                c.text = h
                c.fill.solid()
                c.fill.fore_color.rgb = COR_AMARELO
                for p in c.text_frame.paragraphs:
                    p.font.size = Pt(11)
                    p.font.bold = True
                    p.font.color.rgb = COR_FUNDO
            
            total_area_s = area_counts['Quantidade'].sum()
            for i, r in area_counts.iterrows():
                row_idx = i + 1
                prop_area = f"{round(r['Quantidade'] / total_area_s * 100, 1)}%" if total_area_s > 0 else "0%"
                for j, val in enumerate([f"Área {r['ÁREA_RESUMO']}", str(r['Quantidade']), prop_area]):
                    c = tbl_area.cell(row_idx, j)
                    c.text = val
                    c.fill.solid()
                    c.fill.fore_color.rgb = COR_CARD_BG if row_idx % 2 == 0 else COR_FUNDO
                    for p in c.text_frame.paragraphs:
                        p.font.size = Pt(10)
                        p.font.color.rgb = COR_BRANCO

        # ============================================
        # SLIDE 3: TOP 10 EQUIPES E FUNÇÕES
        # ============================================
        slide3 = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide3, COR_FUNDO)
        
        add_text_box(slide3, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
            "DISTRIBUIÇÃO DE EQUIPES & FUNÇÕES", font_size=26, bold=True, color=COR_BRANCO)
        
        line3 = slide3.shapes.add_shape(1, Inches(0.5), Inches(0.95), Inches(12.333), Inches(0.03))
        line3.fill.solid()
        line3.fill.fore_color.rgb = COR_AZUL
        line3.line.fill.background()
        
        # Top 10 Encarregados
        df_enc_pptx = df[(df["ENCARREGADO"].str.strip() != "") & (df["ENCARREGADO"].isin(lista_completa_encarregados))]
        if not df_enc_pptx.empty:
            top_enc_pptx = df_enc_pptx["ENCARREGADO"].value_counts().head(10).reset_index()
            top_enc_pptx.columns = ["Encarregado", "Efetivo"]
            
            add_text_box(slide3, Inches(0.5), Inches(1.2), Inches(6.0), Inches(0.45),
                "TOP 10 MAIORES EQUIPES", font_size=16, bold=True, color=COR_VERDE)
            
            rows_top = len(top_enc_pptx) + 1
            tbl_top = slide3.shapes.add_table(rows_top, 2, Inches(0.5), Inches(1.75), Inches(6.0), Inches(0.38 * rows_top)).table
            tbl_top.columns[0].width = Inches(4.2)
            tbl_top.columns[1].width = Inches(1.8)
            
            for j, h in enumerate(["Encarregado", "Colaboradores"]):
                c = tbl_top.cell(0, j)
                c.text = h
                c.fill.solid()
                c.fill.fore_color.rgb = COR_VERDE
                for p in c.text_frame.paragraphs:
                    p.font.size = Pt(11)
                    p.font.bold = True
                    p.font.color.rgb = COR_FUNDO
            
            for i, r in top_enc_pptx.iterrows():
                row_idx = i + 1
                for j, val in enumerate([str(r["Encarregado"]), str(r["Efetivo"])]):
                    c = tbl_top.cell(row_idx, j)
                    c.text = val
                    c.fill.solid()
                    c.fill.fore_color.rgb = COR_CARD_BG if row_idx % 2 == 0 else COR_FUNDO
                    for p in c.text_frame.paragraphs:
                        p.font.size = Pt(9.5)
                        p.font.color.rgb = COR_BRANCO

        # Top Funções Mais Demandadas
        if "FUNÇÃO" in df.columns:
            top_func = df["FUNÇÃO"].value_counts().head(10).reset_index()
            top_func.columns = ["Função", "Quantidade"]
            
            add_text_box(slide3, Inches(6.8), Inches(1.2), Inches(6.0), Inches(0.45),
                "FUNÇÕES COM MAIOR CONTINGENTE", font_size=16, bold=True, color=COR_AMARELO)
            
            rows_func = len(top_func) + 1
            tbl_func = slide3.shapes.add_table(rows_func, 2, Inches(6.8), Inches(1.75), Inches(6.0), Inches(0.38 * rows_func)).table
            tbl_func.columns[0].width = Inches(4.2)
            tbl_func.columns[1].width = Inches(1.8)
            
            for j, h in enumerate(["Cargo / Função", "Efetivo"]):
                c = tbl_func.cell(0, j)
                c.text = h
                c.fill.solid()
                c.fill.fore_color.rgb = COR_AMARELO
                for p in c.text_frame.paragraphs:
                    p.font.size = Pt(11)
                    p.font.bold = True
                    p.font.color.rgb = COR_FUNDO
            
            for i, r in top_func.iterrows():
                row_idx = i + 1
                for j, val in enumerate([str(r["Função"]), str(r["Quantidade"])]):
                    c = tbl_func.cell(row_idx, j)
                    c.text = val
                    c.fill.solid()
                    c.fill.fore_color.rgb = COR_CARD_BG if row_idx % 2 == 0 else COR_FUNDO
                    for p in c.text_frame.paragraphs:
                        p.font.size = Pt(9.5)
                        p.font.color.rgb = COR_BRANCO

        # Rodapé
        add_text_box(slide3, Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4),
            f"{nome_site} — Relatório Executivo Gerado Automaticamente — {dt_mod.datetime.now().strftime('%d/%m/%Y %H:%M')}",
            font_size=9, color=COR_CINZA, alignment=PP_ALIGN.CENTER)

        buffer_pptx = io.BytesIO()
        prs.save(buffer_pptx)
        buffer_pptx.seek(0)
        return buffer_pptx.getvalue()

    def gerar_briefing_matinal_ia(df_rdcs_dia, data_str, nome_site):
        """Gera um briefing executivo de 3 tópicos usando Gemini com fallback automático."""
        if df_rdcs_dia is None or df_rdcs_dia.empty:
            return {
                "avancos": ["Ainda não há RDCs registrados para esta data."],
                "atencao": ["Verifique se os encarregados já enviaram os relatórios no sistema."],
                "bloqueios": ["Nenhum bloqueio registrado até o momento."],
                "origem": "Base de Dados"
            }
        
        # Preparar dados para o prompt
        linhas_resumo = []
        for _, row in df_rdcs_dia.iterrows():
            enc = row.get("ENCARREGADO", "N/I")
            disc = row.get("DISCIPLINA", "Geral")
            ativ = row.get("ATIVIDADE", "")
            prob = row.get("PROBLEMAS", "")
            area = row.get("AREA", "")
            cald = row.get("CALDEIRA", "")
            
            info = f"- Encarregado: {enc} ({disc}) | Área: {area} {cald} | Atividade: {ativ}"
            if prob and str(prob).strip().lower() not in ["", "nan", "nenhum", "não informado", "nao informado", "sem problemas", "-", "n/a"]:
                info += f" | PROBLEMA REPORTADO: {prob}"
            linhas_resumo.append(info)
        
        texto_dados = "\n".join(linhas_resumo[:40])
        
        chave_api = ""
        try:
            chave_api = st.secrets.get("GEMINI_API_KEY", "")
        except Exception:
            pass
            
        if chave_api:
            try:
                from google import genai
                client = genai.Client(api_key=chave_api.split(",")[0].strip())
                prompt = f"""
Você é o Engenheiro Coordenador Geral da obra da empresa {nome_site}.
Analise os apontamentos reais de RDC (Relatório Diário de Campo) das equipes de campo na data {data_str} e monte um Briefing Executivo direto para a Reunião Matinal com encarregados e diretoria.

Dados reais das frentes de serviço:
{texto_dados}

Retorne ESTRITAMENTE um JSON puro válido com as 3 chaves abaixo (cada uma com uma lista de 2 a 4 tópicos objetivos em português):
{{
  "avancos": [
    "Frase curta sobre avanço ou frente concluída 1",
    "Frase curta sobre avanço 2"
  ],
  "atencao": [
    "Frase sobre ponto de atenção ou frente com rendimento menor 1",
    "Frase sobre ponto de atenção 2"
  ],
  "bloqueios": [
    "Frase sobre impedimento, falta de liberação ou decisão urgente necessária hoje 1",
    "Frase sobre impedimento 2"
  ]
}}
Retorne apenas o JSON sem crases ou formatação markdown."""
                
                resp = client.models.generate_content(
                    model='gemini-2.5-flash',
                    contents=prompt
                )
                
                resp_text = resp.text.strip()
                if resp_text.startswith("```"):
                    partes = resp_text.split("```")
                    if len(partes) > 1:
                        resp_text = partes[1]
                        if resp_text.startswith("json"):
                            resp_text = resp_text[4:]
                resp_text = resp_text.strip()
                
                import json as json_m
                dados_json = json_m.loads(resp_text)
                
                return {
                    "avancos": dados_json.get("avancos", []),
                    "atencao": dados_json.get("atencao", []),
                    "bloqueios": dados_json.get("bloqueios", []),
                    "origem": "Inteligência Artificial (Gemini 2.5)"
                }
            except Exception:
                pass
                
        # FALLBACK INTELIGENTE
        avancos = []
        atencao = []
        bloqueios = []
        
        ativs_validas = df_rdcs_dia[df_rdcs_dia["ATIVIDADE"].astype(str).str.strip().str.len() > 5]
        for _, r in ativs_validas.head(3).iterrows():
            avancos.append(f"**{r.get('DISCIPLINA', 'Frente')} ({r.get('ENCARREGADO', '')}):** {str(r.get('ATIVIDADE', ''))[:95]}")
        if not avancos:
            avancos.append(f"{len(df_rdcs_dia)} equipes operacionais em andamento nas frentes de serviço.")
            
        probs_validos = df_rdcs_dia[df_rdcs_dia["PROBLEMAS"].notna() & (df_rdcs_dia["PROBLEMAS"].astype(str).str.strip().str.lower() != "nenhum") & (df_rdcs_dia["PROBLEMAS"].astype(str).str.strip().str.lower() != "não informado") & (df_rdcs_dia["PROBLEMAS"].astype(str).str.strip().str.len() > 3)]
        
        if not probs_validos.empty:
            for _, r in probs_validos.head(3).iterrows():
                bloqueios.append(f"⚠️ **{r.get('DISCIPLINA', '')} ({r.get('ENCARREGADO', '')}):** {str(r.get('PROBLEMAS', ''))[:100]}")
        else:
            bloqueios.append("Nenhum impedimento crítico reportado pelas equipes no período.")
            
        discs_count = df_rdcs_dia["DISCIPLINA"].value_counts()
        disc_maior = discs_count.index[0] if not discs_count.empty else "Geral"
        atencao.append(f"Maior concentração de mão de obra na disciplina **{disc_maior}** ({discs_count.iloc[0]} frentes).")
        if len(probs_validos) > 0:
            atencao.append(f"{len(probs_validos)} equipe(s) apontaram restrições que demandam acompanhamento.")
        else:
            atencao.append("Ritmo de produção estável em todas as frentes ativas.")
            
        return {
            "avancos": avancos,
            "atencao": atencao,
            "bloqueios": bloqueios,
            "origem": "Base de Dados Operacional"
        }


    mapa_area_sufixo = {
        'EQUIPAMENTO': '001', 'EQUIPAMENTOS': '001',
        'DUTO': '002', 'DUTOS': '002',
        'TUBULACAO': '003', 'TUBULAÇÃO': '003',
        'ESTRUTURA MET': '004', 'ESTRUTURA METALICA': '004', 'ESTRUTURA METÁLICA': '004',
        'PRECIPITADOR': '005', 'ESP': '005',
        'PRESSAO-MEC': '006', 'PRESSAO - MEC': '006', 'PARTE DE PRESSAO - MECANICA': '006',
        'PRESSAO-TUBULACAO': '007', 'PRESSAO - TUBULACAO': '007', 'PARTE DE PRESSAO - TUBULACAO': '007',
        'PRESSAO-FORNALHA': '008', 'PRESSAO - FORNALHA': '008', 'PARTE DE PRESSAO - FORNALHA': '008', 'PARTE DE PRESSAO - FORNALIA': '008',
        'PINTURA': '009',
        'COMISSIONAMENTO': '010', 'APOIO AO COMISSIONAMENTO': '010',
        'OPERACAO ASSISTIDA': '011', 'OPERAÇÃO ASSISTIDA': '011',
        'LAVAGEM QUIMICA': '012', 'LAVAGEM QUÍMICA': '012',
        'SOPRAGEM': '013',
        'ANDAIME': '014', 'ANDAIMES': '014',
        'OPERADOR': '015', 'OPERADORES E MOTORISTAS': '015', 'MOTORISTA': '015',
        'FORA DE ESCOPO': '016', 'SERVICOS FORA DE ESCOPO': '016', 'SERVIÇOS FORA DE ESCOPO': '016'
    }
    # =================================================================
    # MODO TV (APRESENTAÇÃO AUTOMÁTICA)
    # =================================================================
    if st.session_state.get("modo_tv", False):
        import streamlit.components.v1 as components
        
        slide_atual = st.session_state.get("tv_slide", 0)
        proximo_slide = 1 - slide_atual  # alterna entre 0 e 1
        
        # CSS para esconder sidebar e maximizar conteúdo
        st.markdown("""
        <style>
            [data-testid="stSidebar"] { display: none !important; }
            [data-testid="stSidebarCollapseButton"] { display: none !important; }
            .block-container { max-width: 100% !important; padding: 1rem 2rem !important; }
            .stApp { margin-top: -80px; }
            @keyframes fadeSlideIn {
                from { opacity: 0; transform: translateY(20px); }
                to { opacity: 1; transform: translateY(0); }
            }
            .tv-container { animation: fadeSlideIn 0.8s ease; }
        </style>
        """, unsafe_allow_html=True)
        
        # Barra superior do Modo TV
        col_tv_tit, col_tv_btn = st.columns([5, 1])
        with col_tv_tit:
            slide_nome = "📊 Dashboard" if slide_atual == 0 else "🏎️ Competição F1"
            st.markdown(f"""
            <div style="display: flex; align-items: center; gap: 15px; margin-top: 10px;">
                <div style="background: linear-gradient(135deg, #ef4444, #f97316); border-radius: 8px; padding: 4px 12px; font-size: 11px; color: white; font-weight: 700; letter-spacing: 1px; animation: pulse 1.5s infinite;">
                    📺 AO VIVO
                </div>
                <span style="color: #94a3b8; font-size: 14px;">{slide_nome} · Atualiza em 20s</span>
            </div>
            <style>@keyframes pulse {{ 0%,100%{{ opacity:1; }} 50%{{ opacity:0.5; }} }}</style>
            """, unsafe_allow_html=True)
        with col_tv_btn:
            if st.button("❌ Sair do Modo TV", type="primary"):
                st.session_state.modo_tv = False
                st.rerun()
        
        st.markdown("<div class='tv-container'>", unsafe_allow_html=True)
        
        if slide_atual == 0:
            # ====== SLIDE 0: DASHBOARD ======
            st.markdown(f"""
            <div class="enesa-header" style="margin-top: 10px;">
                <div style="display: flex; justify-content: space-between; align-items: center;">
                    <div>
                        <h1 style="margin: 0; font-size: 2rem; font-weight: 700;">
                            <span style="background: linear-gradient(135deg, #0ea5e9, #8b5cf6); -webkit-background-clip: text; -webkit-text-fill-color: transparent;">Painel de Controle — {nome_site}</span>
                        </h1>
                        <p style="color: {cor_texto_sub}; font-size: 0.9rem; margin: 4px 0 0 0;">Efetivo Operacional · {data_agora}</p>
                    </div>
                </div>
            </div>
            """, unsafe_allow_html=True)
            
            # Filtros do Modo TV
            col_f1_tv, col_f2_tv = st.columns(2)
            with col_f1_tv:
                filtro_tv_local = st.segmented_control(
                    "📍 Filtrar por Local:", 
                    ["Todas", "PB (Caldeira)", "RB (Retorta)", "ESP (Precipitador)"], 
                    default="Todas",
                    key="tv_filtro_local"
                )
                if not filtro_tv_local:
                    filtro_tv_local = "Todas"
            with col_f2_tv:
                enc_lista_tv = ["Todos"] + sorted(lista_completa_encarregados)
                filtro_tv_enc = st.selectbox("👷 Filtrar por Encarregado:", enc_lista_tv, key="tv_filtro_enc")
            
            # Aplicar filtros
            df_tv = df_atual.copy()
            if "PB" in filtro_tv_local:
                df_tv = df_tv[df_tv["C.C"].apply(lambda x: "125.02" in str(x) and ".005" not in str(x))]
            elif "RB" in filtro_tv_local:
                df_tv = df_tv[df_tv["C.C"].apply(lambda x: "125.01" in str(x) and ".005" not in str(x))]
            elif "ESP" in filtro_tv_local:
                df_tv = df_tv[df_tv["C.C"].apply(lambda x: ".005" in str(x))]
            
            if filtro_tv_enc != "Todos":
                df_tv = df_tv[df_tv["ENCARREGADO"] == filtro_tv_enc]
            
            label_filtro = filtro_tv_local if filtro_tv_enc == "Todos" else f"{filtro_tv_enc} ({filtro_tv_local})"
            
            # KPIs
            total_tv = len(df_tv)
            mod_tv = len(df_tv[df_tv["MÃO DE OBRA"].str.strip().str.upper() == "MOD"])
            moi_tv = len(df_tv[df_tv["MÃO DE OBRA"].str.strip().str.upper() == "MOI"])
            pct_mod_tv = round(mod_tv / (mod_tv + moi_tv) * 100, 1) if (mod_tv + moi_tv) > 0 else 0
            enc_tv = len([e for e in df_tv["ENCARREGADO"].unique() if str(e).strip() != "" and str(e) in lista_completa_encarregados])
            funcoes_tv = df_tv["FUNÇÃO"].nunique()
            
            def card_tv(titulo, valor, cor):
                return f"""
                <div style="background: rgba(30, 41, 59, 0.5); backdrop-filter: blur(10px); border-radius: 16px; border: 1px solid rgba(255,255,255,0.06); padding: 24px; text-align: center; position: relative; overflow: hidden;">
                    <p style="margin: 0; font-size: 14px; color: #94a3b8; font-weight: 600; text-transform: uppercase; letter-spacing: 1px;">{titulo}</p>
                    <h2 style="margin: 8px 0 0 0; font-size: 42px; font-weight: 700; color: #f8fafc; text-shadow: 0 0 20px {cor}60;">{valor}</h2>
                    <div style="position: absolute; bottom: 0; left: 0; width: 100%; height: 4px; background: linear-gradient(90deg, {cor}, transparent);"></div>
                </div>
                """
            
            c1, c2, c3, c4, c5 = st.columns(5)
            with c1: st.markdown(card_tv("Efetivo Total", total_tv, "#3b82f6"), unsafe_allow_html=True)
            with c2: st.markdown(card_tv("MOD", mod_tv, "#10b981"), unsafe_allow_html=True)
            with c3: st.markdown(card_tv("MOI", moi_tv, "#ef4444"), unsafe_allow_html=True)
            with c4: st.markdown(card_tv("% MOD", f"{pct_mod_tv}%", "#0ea5e9"), unsafe_allow_html=True)
            with c5: st.markdown(card_tv("Encarregados", enc_tv, "#8b5cf6"), unsafe_allow_html=True)
            
            st.markdown("<br>", unsafe_allow_html=True)
            
            # Gráficos lado a lado
            col_g1, col_g2 = st.columns(2)
            with col_g1:
                st.markdown("#### Efetivo por Área")
                df_area_tv = df_tv.copy()
                df_area_tv['AREA'] = df_area_tv['C.C'].apply(lambda x: 'PB' if '125.02' in str(x) and '.005' not in str(x) else ('RB' if '125.01' in str(x) and '.005' not in str(x) else ('ESP' if '.005' in str(x) else 'OUTROS')))
                df_area_count = df_area_tv[df_area_tv['AREA'] != 'OUTROS'].groupby('AREA').size().reset_index(name='Quantidade')
                if not df_area_count.empty:
                    import plotly.express as px
                    fig_tv1 = px.pie(df_area_count, values='Quantidade', names='AREA', hole=0.55, color_discrete_sequence=["#3b82f6", "#10b981", "#f59e0b", "#ef4444"])
                    fig_tv1.update_layout(margin=dict(l=10, r=10, t=10, b=10), paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", font=dict(color="#e0e4ea", size=14), height=350, showlegend=True, legend=dict(orientation="h", yanchor="bottom", y=-0.15, xanchor="center", x=0.5))
                    st.plotly_chart(fig_tv1, use_container_width=True)
            
            with col_g2:
                st.markdown("#### Top 10 Maiores Equipes")
                df_top_enc = df_tv[df_tv["ENCARREGADO"].isin(lista_completa_encarregados)]
                top10 = df_top_enc.groupby("ENCARREGADO").size().nlargest(10).reset_index(name="Qtd")
                if not top10.empty:
                    import plotly.express as px
                    fig_tv2 = px.bar(top10, y="ENCARREGADO", x="Qtd", orientation="h", color_discrete_sequence=["#3b82f6"], text="Qtd")
                    fig_tv2.update_layout(margin=dict(l=10, r=10, t=10, b=10), paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", font=dict(color="#e0e4ea", size=12), height=350, yaxis=dict(autorange="reversed"), xaxis_title="", yaxis_title="")
                    fig_tv2.update_traces(textposition="outside")
                    st.plotly_chart(fig_tv2, use_container_width=True)
        
        else:
            # ====== SLIDE 1: F1 RANKING ======
            st.markdown(f"""
            <div class="enesa-header" style="margin-top: 10px;">
                <div style="text-align: center;">
                    <h1 style="margin: 0; font-size: 2rem; font-weight: 700;">
                        <span style="background: linear-gradient(135deg, #f59e0b, #ef4444); -webkit-background-clip: text; -webkit-text-fill-color: transparent;">🏎️ Competição F1 — Ranking de Entregas</span>
                    </h1>
                    <p style="color: {cor_texto_sub}; font-size: 0.9rem; margin: 4px 0 0 0;">{nome_site} · {data_agora}</p>
                </div>
            </div>
            """, unsafe_allow_html=True)
            
            # Montar ranking do F1 (mesmo mês atual)
            df_hist_tv = st.session_state.get("df_historico_f1", pd.DataFrame())
            if not df_hist_tv.empty:
                mes_atual_str = datetime.date.today().strftime("%Y-%m")
                df_hist_tv["DATA"] = df_hist_tv["DATA"].astype(str)
                df_mes = df_hist_tv[df_hist_tv["DATA"].str.startswith(mes_atual_str)]
                
                if not df_mes.empty:
                    ranking_tv = df_mes.groupby("ENCARREGADO").size().reset_index(name="ENTREGAS").sort_values("ENTREGAS", ascending=False)
                    
                    # Pódio Top 3
                    top3_tv = ranking_tv.head(3)
                    if len(top3_tv) >= 3:
                        def nome_curto(nome):
                            p = str(nome).split()
                            return p[0] + " " + (p[-1] if len(p) > 1 else "")
                        
                        n1, t1 = nome_curto(top3_tv.iloc[0]["ENCARREGADO"]), top3_tv.iloc[0]["ENTREGAS"]
                        n2, t2 = nome_curto(top3_tv.iloc[1]["ENCARREGADO"]), top3_tv.iloc[1]["ENTREGAS"]
                        n3, t3 = nome_curto(top3_tv.iloc[2]["ENCARREGADO"]), top3_tv.iloc[2]["ENTREGAS"]
                        
                        st.markdown(f"""
                        <div style="display: flex; justify-content: center; align-items: flex-end; height: 250px; gap: 25px; margin: 30px 0;">
                            <div style="display: flex; flex-direction: column; align-items: center; width: 180px;">
                                <span style="font-size: 40px;">🥈</span>
                                <div style="background: linear-gradient(180deg, #94a3b8, #64748b); width: 100%; height: 120px; border-radius: 12px 12px 0 0; display: flex; flex-direction: column; justify-content: center; align-items: center;">
                                    <span style="font-size: 28px; font-weight: 700; color: white;">{t2}</span>
                                    <span style="font-size: 11px; color: #e2e8f0; margin-top: 4px;">{n2}</span>
                                </div>
                            </div>
                            <div style="display: flex; flex-direction: column; align-items: center; width: 200px;">
                                <span style="font-size: 50px;">🥇</span>
                                <div style="background: linear-gradient(180deg, #f59e0b, #d97706); width: 100%; height: 170px; border-radius: 12px 12px 0 0; display: flex; flex-direction: column; justify-content: center; align-items: center; box-shadow: 0 0 30px rgba(245, 158, 11, 0.4);">
                                    <span style="font-size: 36px; font-weight: 700; color: white;">{t1}</span>
                                    <span style="font-size: 13px; color: #fef3c7; margin-top: 4px; font-weight: 600;">{n1}</span>
                                </div>
                            </div>
                            <div style="display: flex; flex-direction: column; align-items: center; width: 180px;">
                                <span style="font-size: 40px;">🥉</span>
                                <div style="background: linear-gradient(180deg, #b45309, #92400e); width: 100%; height: 100px; border-radius: 12px 12px 0 0; display: flex; flex-direction: column; justify-content: center; align-items: center;">
                                    <span style="font-size: 28px; font-weight: 700; color: white;">{t3}</span>
                                    <span style="font-size: 11px; color: #e2e8f0; margin-top: 4px;">{n3}</span>
                                </div>
                            </div>
                        </div>
                        """, unsafe_allow_html=True)
                    
                    # Tabela completa do ranking
                    st.markdown("#### 📊 Ranking Completo do Mês")
                    ranking_tv["POS"] = range(1, len(ranking_tv) + 1)
                    ranking_tv = ranking_tv[["POS", "ENCARREGADO", "ENTREGAS"]]
                    st.dataframe(ranking_tv, use_container_width=True, height=350, hide_index=True)
                else:
                    st.info("Nenhuma entrega registrada neste mês ainda.")
            else:
                st.info("Histórico F1 não carregado.")
        
        st.markdown("</div>", unsafe_allow_html=True)
        
        # JavaScript: Fullscreen + Auto-rotação a cada 20s
        st.session_state.tv_slide = proximo_slide
        components.html(f"""
        <script>
            // Tentar fullscreen
            if (!document.fullscreenElement) {{
                document.documentElement.requestFullscreen().catch(e => {{}});
            }}
            // Recarregar em 20 segundos
            setTimeout(function() {{
                window.parent.location.reload();
            }}, 20000);
        </script>
        """, height=0)
        
        st.stop()  # Impede o resto da página de renderizar

    tab_dashboard, tab_resumo, tab_emissao, tab_escala, tab_cc, tab_f1, tab_ia, tab_ia_cc, tab_rdc_digital, tab_pde, tab_banco_rdc, tab_gargalos, tab_banco_dados, tab_admin = st.tabs([f"📊 {t('Dashboard')}", f"📅 {t('Resumo Diário')}", f"📝 {t('Emissão de RDC')}", f"📋 {t('Escala')}", f"💰 {t('Controle de C.C')}", f"🏎️ {t('Competição F1')}", f"🤖 {t('Leitor de RDC (IA)')}", f"🤖 {t('IA - Atualizador de C.C')}", f"📱 {t('RDC Digital')}", f"👷 {t('Gerenciar PDE')}", f"📑 {t('Banco de RDCs')}", f"🔍 {t('Análise de Gargalos')}", f"📊 {t('Banco de Dados')}", f"⚙️ {t('Admin')}"])

    if st.session_state.get("role_usuario") == "apontador":
        st.markdown("""
        <style>
            div[data-baseweb="tab-list"] button:nth-child(1),
            div[data-baseweb="tab-list"] button:nth-child(2),
            div[data-baseweb="tab-list"] button:nth-child(5),
            div[data-baseweb="tab-list"] button:nth-child(7),
            div[data-baseweb="tab-list"] button:nth-child(8),
            div[data-baseweb="tab-list"] button:nth-child(9),
            div[data-baseweb="tab-list"] button:nth-child(10),
            div[data-baseweb="tab-list"] button:nth-child(11),
            div[data-baseweb="tab-list"] button:nth-child(12),
            div[data-baseweb="tab-list"] button:nth-child(13),
            div[data-baseweb="tab-list"] button:nth-child(14) {
                display: none !important;
            }
        </style>
        """, unsafe_allow_html=True)

    with tab_dashboard:
        # === RELÓGIO DIGITAL ===
        import streamlit.components.v1 as components
        html_relogio = """
        <div id="clock_container" style="font-family: 'Courier New', Courier, monospace; font-size: 28px; color: #0ea5e9; font-weight: bold; text-shadow: 0 0 10px rgba(14, 165, 233, 0.8); text-align: center; background: linear-gradient(145deg, rgba(15, 23, 42, 0.9), rgba(30, 41, 59, 0.9)); padding: 10px 20px; border-radius: 12px; border: 1px solid rgba(14, 165, 233, 0.4); width: fit-content; margin: 0 auto 20px auto; box-shadow: 0 4px 15px rgba(0,0,0,0.5), inset 0 0 10px rgba(14,165,233,0.1);">
            <div id="time" style="letter-spacing: 2px;">--:--:--</div>
            <div id="date" style="font-size: 14px; color: #94a3b8; font-weight: normal; text-shadow: none; font-family: 'Inter', sans-serif; margin-top: 5px; text-transform: uppercase; letter-spacing: 1px;">Carregando...</div>
        </div>
        <script>
            function updateClock() {
                const now = new Date();
                const timeStr = now.toLocaleTimeString('pt-BR');
                const dateOptions = { weekday: 'long', year: 'numeric', month: 'long', day: 'numeric' };
                let dateStr = now.toLocaleDateString('pt-BR', dateOptions);
                document.getElementById('time').innerText = timeStr;
                document.getElementById('date').innerText = dateStr;
            }
            setInterval(updateClock, 1000);
            updateClock();
        </script>
        """
        components.html(html_relogio, height=110)
        
        col_dash_tit, col_dash_btn_pdf, col_dash_btn_pptx = st.columns([2, 1, 1])
        with col_dash_tit:
            st.markdown("### 🎛️ Centro de Comando (Overview)")
        with col_dash_btn_pdf:
            st.markdown("<br>", unsafe_allow_html=True)
            pdf_bytes = gerar_relatorio_pdf(df_atual)
            st.download_button(
                label="📥 Baixar PDF",
                data=pdf_bytes,
                file_name=f"Relatorio_Executivo_{datetime.date.today().strftime('%d_%m_%Y')}.pdf",
                mime="application/pdf",
                use_container_width=True,
                type="secondary"
            )
        with col_dash_btn_pptx:
            st.markdown("<br>", unsafe_allow_html=True)
            try:
                pptx_dash_bytes = gerar_relatorio_pptx_dashboard(df_atual, nome_site, caminho_logo, lista_completa_encarregados)
                st.download_button(
                    label="📑 Baixar PowerPoint (.pptx)",
                    data=pptx_dash_bytes,
                    file_name=f"Dashboard_Executivo_{datetime.date.today().strftime('%d_%m_%Y')}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                    type="primary"
                )
            except Exception as e_pptx_dash:
                st.error(f"Erro ao gerar PPTX: {e_pptx_dash}")
        
        # Filtro de MOI / MOD, Local, Turno e Status
        col_filtros1, col_filtros2, col_filtros3, col_filtros4 = st.columns(4)
        with col_filtros1:
            filtro_dash_mo = st.segmented_control(
                "Filtrar Visão por Tipo de Mão de Obra:", 
                ["Ambas", "MOD", "MOI"], 
                default="Ambas"
            )
            if not filtro_dash_mo:
                filtro_dash_mo = "Ambas"
                
        with col_filtros2:
            filtro_dash_local = st.segmented_control(
                "Filtrar Dados por Local:", 
                ["Ambas", "PB", "RB", "ESP"], 
                default="Ambas",
                key="filtro_dash_local_key"
            )
            if not filtro_dash_local:
                filtro_dash_local = "Ambas"
                
        with col_filtros3:
            # Pegar todos os turnos únicos do PDE, ou padronizar
            turnos_disponiveis = ["Todos"]
            if "TURNO" in df_atual.columns:
                turnos_reais = [t for t in df_atual["TURNO"].unique() if str(t).strip() and str(t) != "nan"]
                turnos_disponiveis.extend(sorted(turnos_reais))
            
            filtro_dash_turno = st.selectbox(
                "Filtrar por Turno:", 
                turnos_disponiveis,
                index=0
            )
            
        with col_filtros4:
            # Pegar todos os status únicos do PDE, ou padronizar
            status_disponiveis = ["Todos"]
            if "STATUS" in df_atual.columns:
                status_reais = [s for s in df_atual["STATUS"].unique() if str(s).strip() and str(s) != "nan"]
                status_disponiveis.extend(sorted(status_reais))
                
            filtro_dash_status = st.selectbox(
                "Filtrar por Status:", 
                status_disponiveis,
                index=status_disponiveis.index("ATIVO") if "ATIVO" in status_disponiveis else 0
            )
            
        df_dash = df_atual.copy()
        
        # Aplicar filtro MOI/MOD
        if filtro_dash_mo == "MOD":
            df_dash = df_dash[df_dash["MÃO DE OBRA"].astype(str).str.strip().str.upper() == "MOD"]
        elif filtro_dash_mo == "MOI":
            df_dash = df_dash[df_dash["MÃO DE OBRA"].astype(str).str.strip().str.upper() == "MOI"]
            
        # Aplicar filtro Local
        df_dash = df_dash[df_dash["C.C"].str.strip() != ""]
        if filtro_dash_local == "PB":
            df_dash = df_dash[df_dash["C.C"].apply(lambda x: "125.02" in str(x) and ".005" not in str(x))]
        elif filtro_dash_local == "RB":
            df_dash = df_dash[df_dash["C.C"].apply(lambda x: "125.01" in str(x) and ".005" not in str(x))]
        elif filtro_dash_local == "ESP":
            df_dash = df_dash[df_dash["C.C"].apply(lambda x: ".005" in str(x))]
            
        # Aplicar filtro Turno
        if filtro_dash_turno != "Todos" and "TURNO" in df_dash.columns:
            df_dash = df_dash[df_dash["TURNO"] == filtro_dash_turno]
            
        # Aplicar filtro Status
        if filtro_dash_status != "Todos" and "STATUS" in df_dash.columns:
            df_dash = df_dash[df_dash["STATUS"] == filtro_dash_status]
        
        # Linha 1: Cartões de KPI Customizados (Premium)
        qtd_encarregados_dash = len([e for e in df_dash["ENCARREGADO"].unique() if str(e).strip() != "" and str(e) in lista_completa_encarregados])
        qtd_mod_g = len(df_atual[df_atual["MÃO DE OBRA"].str.strip().str.upper() == "MOD"])
        qtd_moi_g = len(df_atual[df_atual["MÃO DE OBRA"].str.strip().str.upper() == "MOI"])
        total_mo_g = qtd_mod_g + qtd_moi_g
        pct_mod_g = round((qtd_mod_g / total_mo_g * 100), 1) if total_mo_g > 0 else 0
        span_control = round(len(df_dash) / qtd_encarregados_dash, 1) if qtd_encarregados_dash > 0 else 0
        
        def card_kpi(titulo, valor, icone, cor):
            return f"""
            <div style="background: rgba(30, 41, 59, 0.45); backdrop-filter: blur(10px); border-radius: 16px; border: 1px solid rgba(255,255,255,0.05); padding: 18px; box-shadow: 0 4px 15px rgba(0,0,0,0.2); position: relative; overflow: hidden; height: 110px; transition: transform 0.3s ease;" onmouseover="this.style.transform='translateY(-5px)'" onmouseout="this.style.transform='translateY(0px)'">
                <p style="margin: 0; font-size: 13px; color: #94a3b8; font-weight: 600; text-transform: uppercase; letter-spacing: 0.5px;">{titulo}</p>
                <h2 style="margin: 5px 0 0 0; font-size: 34px; font-weight: 700; color: #f8fafc; text-shadow: 0 0 15px {cor}60;">{valor}</h2>
                <div style="position: absolute; bottom: 0; left: 0; width: 100%; height: 4px; background: linear-gradient(90deg, {cor}, transparent); box-shadow: 0 -2px 10px {cor}80;"></div>
            </div>
            """
            
        st.markdown("<br>", unsafe_allow_html=True)
        m1, m2, m3, m4, m5 = st.columns(5)
        with m1: st.markdown(card_kpi(f"{t('Efetivo')} ({filtro_dash_mo})", len(df_dash), "engineering", "#3b82f6"), unsafe_allow_html=True)
        with m2: st.markdown(card_kpi(t("Encarregados"), qtd_encarregados_dash, "shield_person", "#10b981"), unsafe_allow_html=True)
        with m3: st.markdown(card_kpi(t("% MOD Global"), f"{pct_mod_g}%", "pie_chart", "#0ea5e9"), unsafe_allow_html=True)
        with m4: st.markdown(card_kpi(t("Funções"), df_dash["FUNÇÃO"].nunique(), "build", "#f59e0b"), unsafe_allow_html=True)
        with m5: st.markdown(card_kpi(t("Span of Control"), span_control, "groups", "#8b5cf6"), unsafe_allow_html=True)
        
        # ==============================================================
        # BRIEFING MATINAL COM IA (SOB DEMANDA - RÁPIDO)
        # ==============================================================
        st.markdown("<br>", unsafe_allow_html=True)
        with st.container(border=True):
            st.markdown("""
            <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 12px;">
                <h4 style="margin: 0; color: #f8fafc; font-size: 16px; font-weight: 700; display: flex; align-items: center; gap: 8px;">
                    🤖 Briefing Matinal com IA — Resumo Executivo para Reunião Diária
                </h4>
            </div>
            """, unsafe_allow_html=True)

            # Carregar banco de RDCs para o briefing
            df_rdc_briefing = None
            if os.path.exists(caminho_rdc_registros_csv):
                try:
                    df_rdc_briefing = pd.read_csv(caminho_rdc_registros_csv)
                except:
                    pass
            
            datas_disponiveis_br = []
            if df_rdc_briefing is not None and not df_rdc_briefing.empty and "DATA" in df_rdc_briefing.columns:
                df_rdc_briefing["_DATA_DT"] = pd.to_datetime(df_rdc_briefing["DATA"], errors='coerce')
                datas_disponiveis_br = sorted(df_rdc_briefing["_DATA_DT"].dropna().dt.strftime("%d/%m/%Y").unique().tolist(), key=lambda x: pd.to_datetime(x, format="%d/%m/%Y"), reverse=True)
            
            col_b1, col_b2, col_b3 = st.columns([2, 2, 2])
            with col_b1:
                data_brief_sel = st.selectbox(
                    "📅 Selecionar Data do Briefing:",
                    options=datas_disponiveis_br if datas_disponiveis_br else [datetime.datetime.now().strftime("%d/%m/%Y")],
                    index=0,
                    key="select_data_briefing"
                )
            with col_b2:
                st.markdown("<div style='height: 28px;'></div>", unsafe_allow_html=True)
                btn_gerar_br = st.button("⚡ Gerar / Carregar Briefing com IA", type="primary", use_container_width=True, key="btn_gerar_briefing_ia")
            with col_b3:
                st.markdown("<div style='height: 28px;'></div>", unsafe_allow_html=True)
                st.caption("💡 Clique no botão para carregar sob demanda sem lentidão no site")

            # Filtrar dados para a data selecionada
            df_dia_br = pd.DataFrame()
            if df_rdc_briefing is not None and not df_rdc_briefing.empty and "_DATA_DT" in df_rdc_briefing.columns:
                try:
                    dt_alvo = pd.to_datetime(data_brief_sel, format="%d/%m/%Y").date()
                    df_dia_br = df_rdc_briefing[df_rdc_briefing["_DATA_DT"].dt.date == dt_alvo]
                except:
                    df_dia_br = df_rdc_briefing
            
            # Gerar briefing APENAS quando o usuário clicar no botão
            cache_key = f"briefing_cache_{data_brief_sel}"
            if btn_gerar_br:
                with st.spinner("🤖 IA analisando RDCs e montando síntese da reunião..."):
                    res_brief = gerar_briefing_matinal_ia(df_dia_br, data_brief_sel, nome_site)
                    st.session_state[cache_key] = res_brief

            if cache_key in st.session_state:
                briefing_atual = st.session_state[cache_key]
                avancos_l = briefing_atual.get("avancos", [])
                atencao_l = briefing_atual.get("atencao", [])
                bloqueios_l = briefing_atual.get("bloqueios", [])
                origem_b = briefing_atual.get("origem", "IA")

                st.markdown(f"<p style='color: #64748b; font-size: 12px; margin: 6px 0 14px 0;'>Fonte: <b style='color: #0ea5e9;'>{origem_b}</b> · Total de {len(df_dia_br)} RDCs analisados nesta data</p>", unsafe_allow_html=True)

                # 3 Cards Coloridos (Verde, Amarelo, Vermelho)
                col_c1, col_c2, col_c3 = st.columns(3)
                
                with col_c1:
                    st.markdown(f"""
                    <div style="background: rgba(34, 197, 94, 0.08); border: 1px solid rgba(34, 197, 94, 0.3); border-radius: 12px; padding: 16px; height: 100%;">
                        <h4 style="color: #22c55e; margin: 0 0 10px 0; font-size: 15px; display: flex; align-items: center; gap: 6px;">
                            🟢 PRINCIPAIS AVANÇOS
                        </h4>
                        <ul style="color: #e2e8f0; font-size: 13px; margin: 0; padding-left: 18px; line-height: 1.6;">
                            {''.join([f'<li style="margin-bottom: 6px;">{item}</li>' for item in avancos_l])}
                        </ul>
                    </div>
                    """, unsafe_allow_html=True)
                    
                with col_c2:
                    st.markdown(f"""
                    <div style="background: rgba(245, 158, 11, 0.08); border: 1px solid rgba(245, 158, 11, 0.3); border-radius: 12px; padding: 16px; height: 100%;">
                        <h4 style="color: #f59e0b; margin: 0 0 10px 0; font-size: 15px; display: flex; align-items: center; gap: 6px;">
                            🟡 PONTOS DE ATENÇÃO
                        </h4>
                        <ul style="color: #e2e8f0; font-size: 13px; margin: 0; padding-left: 18px; line-height: 1.6;">
                            {''.join([f'<li style="margin-bottom: 6px;">{item}</li>' for item in atencao_l])}
                        </ul>
                    </div>
                    """, unsafe_allow_html=True)
                    
                with col_c3:
                    st.markdown(f"""
                    <div style="background: rgba(239, 68, 68, 0.08); border: 1px solid rgba(239, 68, 68, 0.3); border-radius: 12px; padding: 16px; height: 100%;">
                        <h4 style="color: #ef4444; margin: 0 0 10px 0; font-size: 15px; display: flex; align-items: center; gap: 6px;">
                            🔴 BLOQUEIOS & AÇÕES URGENTES
                        </h4>
                        <ul style="color: #e2e8f0; font-size: 13px; margin: 0; padding-left: 18px; line-height: 1.6;">
                            {''.join([f'<li style="margin-bottom: 6px;">{item}</li>' for item in bloqueios_l])}
                        </ul>
                    </div>
                    """, unsafe_allow_html=True)

                # Texto formatado para WhatsApp
                import urllib.parse
                texto_zap_lista = [
                    f"🏗️ *BRIEFING MATINAL DE OBRA — {nome_site}*",
                    f"📅 *Data de Referência:* {data_brief_sel}",
                    "",
                    "🟢 *PRINCIPAIS AVANÇOS:*",
                    *[f"• {item.replace('**', '')}" for item in avancos_l],
                    "",
                    "🟡 *PONTOS DE ATENÇÃO:*",
                    *[f"• {item.replace('**', '')}" for item in atencao_l],
                    "",
                    "🔴 *BLOQUEIOS / AÇÕES URGENTES:*",
                    *[f"• {item.replace('**', '')}" for item in bloqueios_l],
                    "",
                    f"📊 _Gerado pelo Sistema RDC & PDE em {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}_"
                ]
                texto_zap_completo = "\n".join(texto_zap_lista)
                link_zap = f"https://api.whatsapp.com/send?text={urllib.parse.quote(texto_zap_completo)}"

                st.markdown("<br>", unsafe_allow_html=True)
                col_z1, col_z2 = st.columns([1, 3])
                with col_z1:
                    st.link_button("📲 Enviar Briefing no WhatsApp", link_zap, use_container_width=True, type="secondary")
                with col_z2:
                    if st.toggle("📋 Ver Texto Formatado para Copiar", key="tgl_ver_texto_zap"):
                        st.text_area("Texto do Briefing:", value=texto_zap_completo, height=140, key="txt_area_briefing_zap")
            else:
                st.info("👆 Selecione a data desejada e clique em **'⚡ Gerar / Carregar Briefing com IA'** para carregar os indicadores sob demanda.")
        st.markdown("---")
        
        col_dash1, col_dash2, col_dash3 = st.columns([3, 3, 4])
        
        with col_dash1:
            st.markdown("**Status Operacional (Global)**")
            if total_mo_g > 0:
                df_mo_global = pd.DataFrame({"Tipo": ["MOD", "MOI"], "Quantidade": [qtd_mod_g, qtd_moi_g]})
                fig_mo_g = px.pie(df_mo_global, values="Quantidade", names="Tipo", hole=0.6, color_discrete_sequence=["#10b981", "#ef4444"])
                fig_mo_g.update_layout(margin=dict(l=20, r=20, t=10, b=10), paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", font=dict(color="#e0e4ea"), height=280, showlegend=True, legend=dict(orientation="h", yanchor="bottom", y=-0.2, xanchor="center", x=0.5))
                st.plotly_chart(fig_mo_g, use_container_width=True)
            else:
                st.info("Classificação de Mão de Obra não encontrada.")
                
        with col_dash2:
            st.markdown("**Efetivo por Área**")
            df_area = df_dash.copy()
            df_area['ÁREA_RESUMO'] = df_area['C.C'].apply(lambda x: 'PB' if '125.02' in str(x) and '.005' not in str(x) else ('RB' if '125.01' in str(x) and '.005' not in str(x) else ('ESP' if '.005' in str(x) else 'OUTROS')))
            df_area_count = df_area[df_area['ÁREA_RESUMO'] != 'OUTROS'].groupby('ÁREA_RESUMO').size().reset_index(name='Quantidade')
            
            if not df_area_count.empty and df_area_count['Quantidade'].sum() > 0:
                cores_areas = {'PB': '#3498db', 'RB': '#e67e22', 'ESP': '#9b59b6'}
                fig_area = px.pie(df_area_count, values="Quantidade", names="ÁREA_RESUMO", hole=0.6, color="ÁREA_RESUMO", color_discrete_map=cores_areas)
                fig_area.update_layout(margin=dict(l=20, r=20, t=10, b=10), paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", font=dict(color="#e0e4ea"), height=280, showlegend=True, legend=dict(orientation="h", yanchor="bottom", y=-0.2, xanchor="center", x=0.5))
                st.plotly_chart(fig_area, use_container_width=True)
            else:
                st.info("Áreas não identificadas.")
                
        with col_dash3:
            st.markdown(f"**Top 10 Maiores Equipes ({filtro_dash_mo})**")
            df_enc_dash = df_dash[(df_dash["ENCARREGADO"].str.strip() != "") & (df_dash["ENCARREGADO"].isin(lista_completa_encarregados))]
            if not df_enc_dash.empty:
                top_enc = df_enc_dash["ENCARREGADO"].value_counts().head(10).reset_index()
                top_enc.columns = ["Encarregado", "Efetivo"]
                fig_top_enc = px.bar(top_enc, x="Efetivo", y="Encarregado", orientation="h", color="Efetivo", color_continuous_scale=[(0, "#0f172a"), (1, "#0ea5e9")], text="Efetivo")
                fig_top_enc.update_layout(showlegend=False, xaxis_title="", yaxis_title="", margin=dict(l=0, r=40, t=10, b=0), paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", font=dict(color="#e0e4ea"), height=280)
                fig_top_enc.update_yaxes(categoryorder="total ascending")
                fig_top_enc.update_xaxes(visible=False)
                fig_top_enc.update_coloraxes(showscale=False)
                fig_top_enc.update_traces(textposition='outside', cliponaxis=False)
                st.plotly_chart(fig_top_enc, use_container_width=True)
                
        st.markdown("---")
        
        col_evolucao, col_gauge = st.columns([6, 4])
        
        with col_evolucao:
            st.markdown("**📈 Evolução Diária de Entregas de RDC (Mês Atual)**")
            if "df_historico_f1" in st.session_state and not st.session_state.df_historico_f1.empty:
                df_hist_dash = st.session_state.df_historico_f1.copy()
                df_hist_dash["DATA"] = pd.to_datetime(df_hist_dash["DATA"], errors='coerce')
                mes_atual = datetime.date.today().strftime("%Y-%m")
                df_hist_dash = df_hist_dash[df_hist_dash["DATA"].dt.strftime("%Y-%m") == mes_atual]
                
                if not df_hist_dash.empty:
                    entregas_por_dia = df_hist_dash.groupby(df_hist_dash["DATA"].dt.strftime("%Y-%m-%d")).size().reset_index(name="Entregas")
                    entregas_por_dia.columns = ["Data", "Qtd Entregue"]
                    
                    fig_evolucao = px.line(entregas_por_dia, x="Data", y="Qtd Entregue", markers=True, 
                                           title="", line_shape="spline", color_discrete_sequence=["#0ea5e9"])
                    fig_evolucao.update_layout(
                        xaxis_title="Dia", yaxis_title="RDCs Entregues",
                        yaxis=dict(dtick=1),
                        margin=dict(l=0, r=20, t=10, b=0),
                        paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
                        font=dict(color="#e0e4ea"), height=250
                    )
                    fig_evolucao.update_traces(line=dict(width=3), marker=dict(size=8))
                    st.plotly_chart(fig_evolucao, use_container_width=True)
                else:
                    st.info("Ainda não há entregas neste mês.")
            else:
                st.info("Sem histórico de F1.")
                
        with col_gauge:
            st.markdown("**🌡️ Termômetro de Engajamento**")
            if "df_historico_f1" in st.session_state and not st.session_state.df_historico_f1.empty:
                dias_unicos = df_hist_dash["DATA"].nunique()
                dias_unicos = dias_unicos if dias_unicos > 0 else 1
                rdcs_esperados = dias_unicos * len(lista_completa_encarregados)
                rdcs_entregues = len(df_hist_dash)
                
                pct_engajamento = round((rdcs_entregues / rdcs_esperados) * 100, 1) if rdcs_esperados > 0 else 0
                
                import plotly.graph_objects as go
                fig_gauge = go.Figure(go.Indicator(
                    mode = "gauge+number",
                    value = pct_engajamento,
                    number = {'suffix': "%", 'font': {'size': 30, 'color': '#e0e4ea'}},
                    domain = {'x': [0, 1], 'y': [0, 1]},
                    gauge = {
                        'axis': {'range': [None, 100], 'tickwidth': 1, 'tickcolor': "white"},
                        'bar': {'color': "#2c3e50"},
                        'bgcolor': "rgba(0,0,0,0)",
                        'borderwidth': 2,
                        'bordercolor': "gray",
                        'steps': [
                            {'range': [0, 60], 'color': '#ef4444'},
                            {'range': [60, 85], 'color': '#f59e0b'},
                            {'range': [85, 100], 'color': '#10b981'}],
                        'threshold': {
                            'line': {'color': "white", 'width': 4},
                            'thickness': 0.75,
                            'value': pct_engajamento}
                    }
                ))
                fig_gauge.update_layout(height=250, margin=dict(l=20, r=20, t=30, b=20), paper_bgcolor="rgba(0,0,0,0)", font=dict(color="#e0e4ea"))
                st.plotly_chart(fig_gauge, use_container_width=True)
            else:
                st.info("Sem dados suficientes.")
                
        st.markdown("---")
        st.markdown("**📦 Raio-X da Mão de Obra Indireta (MOI)**")
        col_moi1, col_moi2 = st.columns([5, 5])
        with col_moi1:
            df_moi = df_atual[df_atual["MÃO DE OBRA"].astype(str).str.strip().str.upper() == "MOI"].copy()
            if not df_moi.empty:
                moi_count = df_moi.groupby("DISCIPLINA").size().reset_index(name="Quantidade")
                moi_count = moi_count.sort_values(by="Quantidade", ascending=False).head(8)
                fig_moi = px.pie(moi_count, values="Quantidade", names="DISCIPLINA", hole=0.5, color_discrete_sequence=px.colors.sequential.YlOrRd[::-1])
                fig_moi.update_layout(margin=dict(l=20, r=20, t=10, b=10), paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", font=dict(color="#e0e4ea"), height=280)
                st.plotly_chart(fig_moi, use_container_width=True)
            else:
                st.info("Nenhuma MOI na base atual.")

        st.markdown("---")
        
        st.markdown(f"**👥 Liderança: Resumo Geral de Encarregados ({filtro_dash_mo})**")
        df_enc_full = df_dash[(df_dash["ENCARREGADO"].str.strip() != "") & (df_dash["ENCARREGADO"].isin(lista_completa_encarregados))]
        if not df_enc_full.empty:
            resumo_enc = df_enc_full["ENCARREGADO"].value_counts().reset_index()
            resumo_enc.columns = ["Encarregado", "Tamanho da Equipe"]
            
            termo_enc = st.text_input("🔍 Procurar Encarregado específico:")
            if termo_enc:
                resumo_enc = resumo_enc[resumo_enc["Encarregado"].astype(str).str.contains(termo_enc, case=False, na=False)]
                
            st.dataframe(resumo_enc, hide_index=True, use_container_width=True)
            
        st.markdown("---")
        st.markdown(f"**🔍 Base Completa ({filtro_dash_mo})**")
        termo_busca = st.text_input("Buscar funcionário (Nome, Matrícula ou Função):")
        df_exibicao = df_dash[["MATRICULA", "NOME", "FUNÇÃO", "ENCARREGADO", "C.C"]].copy()
        if termo_busca:
            mask = (
                df_exibicao["NOME"].astype(str).str.contains(termo_busca, case=False, na=False) |
                df_exibicao["MATRICULA"].astype(str).str.contains(termo_busca, case=False, na=False) |
                df_exibicao["FUNÇÃO"].astype(str).str.contains(termo_busca, case=False, na=False)
            )
            df_exibicao = df_exibicao[mask]
        st.dataframe(df_exibicao, hide_index=True, use_container_width=True)
        

    with tab_resumo:
        st.markdown("### 📅 Resumo Diário")
        
        # --- FIX: Filtro de data para ver o resumo de qualquer dia ---
        data_resumo = st.date_input("Selecione a Data do Resumo:", datetime.date.today())
        data_filtro_str = data_resumo.strftime("%Y-%m-%d")
        
        st.markdown("<br>", unsafe_allow_html=True)
        
        st.markdown("---")

        if st.toggle("➕ Lançar RDC Manualmente (Para papéis ilegíveis ou atrasados)", key="toggle_manual_resumo"):
            with st.form("form_resumo_manual"):
                st.info("💡 Você pode colar a lista inteira de encarregados aqui (um por linha ou separados por vírgula). O robô vai verificar: se a IA já tiver lido, ele ignora. Se faltou, ele adiciona!")
                col_m1, col_m2 = st.columns([1, 2])
                with col_m1:
                    data_manual_resumo = st.date_input("Data do RDC a lançar:", value=data_resumo, key="data_manual_resumo")
                with col_m2:
                    nomes_colados_resumo = st.text_area("Cole os nomes dos Encarregados", height=120, key="nomes_colados_resumo")
                
                btn_manual_resumo = st.form_submit_button("Processar Lista e Lançar no Sistema")
                if btn_manual_resumo and nomes_colados_resumo.strip():
                    import re
                    import difflib
                    data_str_resumo = data_manual_resumo.strftime("%Y-%m-%d")
                    
                    lista_suja = [n.strip().upper() for n in re.split(r'[\n,;]', nomes_colados_resumo) if n.strip()]
                    
                    novos_registros = []
                    nomes_ja_existentes = 0
                    nomes_nao_encontrados = []
                    
                    for nome_sujo in lista_suja:
                        match = difflib.get_close_matches(nome_sujo, lista_completa_encarregados, n=1, cutoff=0.55)
                        if match:
                            nome_oficial = match[0]
                            ja_existe = ((st.session_state.df_historico_f1["DATA"] == data_str_resumo) & (st.session_state.df_historico_f1["ENCARREGADO"] == nome_oficial)).any()
                            if ja_existe:
                                nomes_ja_existentes += 1
                            else:
                                novos_registros.append({"DATA": data_str_resumo, "ENCARREGADO": nome_oficial})
                        else:
                            nomes_nao_encontrados.append(nome_sujo)
                            
                    if novos_registros:
                        df_novos = pd.DataFrame(novos_registros)
                        
                        if conn and not st.session_state.get('force_use_local', False):
                            try:
                                df_fresco = conn.read(worksheet="Historico_F1", ttl=0)
                                if not df_fresco.empty:
                                    df_fresco = df_fresco.dropna(how='all')
                                    df_final = pd.concat([df_fresco, df_novos], ignore_index=True).drop_duplicates(subset=["DATA", "ENCARREGADO"])
                                else:
                                    df_final = pd.concat([st.session_state.df_historico_f1, df_novos], ignore_index=True).drop_duplicates(subset=["DATA", "ENCARREGADO"])
                                
                                ok, msg = salvar_f1_seguro(conn, df_final, caminho_historico_f1_csv)
                                if ok:
                                    st.session_state.df_historico_f1 = df_final
                                st.cache_data.clear()
                                st.toast(f"{len(novos_registros)} novos RDCs sincronizados com a nuvem! ({nomes_ja_existentes} já constavam).", icon="✅")
                            except Exception as e:
                                st.error(f"Erro ao salvar na nuvem: {e}")
                                st.session_state.df_historico_f1 = pd.concat([st.session_state.df_historico_f1, df_novos], ignore_index=True).drop_duplicates(subset=["DATA", "ENCARREGADO"])
                                st.session_state.df_historico_f1.to_csv(caminho_historico_f1_csv, index=False)
                        else:
                            st.session_state.df_historico_f1 = pd.concat([st.session_state.df_historico_f1, df_novos], ignore_index=True).drop_duplicates(subset=["DATA", "ENCARREGADO"])
                            st.session_state.df_historico_f1.to_csv(caminho_historico_f1_csv, index=False)
                            st.toast(f"{len(novos_registros)} novos RDCs adicionados localmente! ({nomes_ja_existentes} já constavam).", icon="✅")
                    elif nomes_ja_existentes > 0:
                        st.warning(f"⚠️ Todos os nomes reconhecidos ({nomes_ja_existentes}) já estavam devidamente lançados neste dia!")
                        
                    if nomes_nao_encontrados:
                        st.error(f"❌ Não encontrei na lista oficial (verifique a escrita): {', '.join(nomes_nao_encontrados)}")
                        
                    time.sleep(3)
                    st.rerun()
                    
        st.markdown("<hr style='margin-top:0px; margin-bottom:20px'>", unsafe_allow_html=True)

        
        df_hoje = pd.DataFrame()
        if "df_historico_f1" in st.session_state and not st.session_state.df_historico_f1.empty:
            df_hist = st.session_state.df_historico_f1.copy()
            df_hist["DATA_STR"] = pd.to_datetime(df_hist["DATA"], errors="coerce").dt.strftime("%Y-%m-%d")
            df_hoje = df_hist[df_hist["DATA_STR"] == data_filtro_str]
            
        encarregados_esperados = len(lista_completa_encarregados)
        entregues_hoje_lista = [e for e in df_hoje["ENCARREGADO"].unique() if e in lista_completa_encarregados] if not df_hoje.empty else []
        encarregados_entregues = len(entregues_hoje_lista)
        encarregados_pendentes = encarregados_esperados - encarregados_entregues
        
        c1, c2, c3 = st.columns(3)
        c1.metric("🎯 Esperados", encarregados_esperados)
        c2.metric("✅ Entregues", encarregados_entregues)
        c3.metric("⏳ Pendentes", encarregados_pendentes)
        
        st.markdown("---")
        if encarregados_pendentes > 0:
            st.error(f"**Atenção:** {encarregados_pendentes} encarregados ainda não entregaram o RDC nesta data ({data_filtro_str}).")
            entregues_list = df_hoje["ENCARREGADO"].unique() if not df_hoje.empty else []
            pendentes_list = [e for e in lista_completa_encarregados if e not in entregues_list]
            
            # Botão de Gerar PDF
            col_pdf, col_gap = st.columns([2, 3])
            with col_pdf:
                if st.button("📄 Gerar PDF de Cobrança", use_container_width=True):
                    from fpdf import FPDF
                    import tempfile
                    
                    class PDF(FPDF):
                        def header(self):
                            self.set_font('Helvetica', 'B', 15)
                            self.set_text_color(0, 0, 0)
                            self.cell(0, 10, 'Relatorio de Pendencias - RDC', 0, 1, 'C')
                            self.set_font('Helvetica', 'I', 10)
                            self.cell(0, 10, f'Data Referencia: {data_filtro_str} (Gerado em: {datetime.datetime.now().strftime("%d/%m/%Y %H:%M")})', 0, 1, 'C')
                            self.ln(5)
                    
                    pdf = PDF()
                    pdf.add_page()
                    pdf.set_font('Helvetica', 'B', 12)
                    pdf.set_text_color(200, 0, 0)
                    pdf.cell(0, 10, f'{len(pendentes_list)} Encarregados nao entregaram o RDC nesta data ({data_filtro_str}):', 0, 1, 'L')
                    pdf.ln(2)
                    
                    pdf.set_font('Helvetica', '', 10)
                    pdf.set_text_color(0, 0, 0)
                    for pendente in sorted(pendentes_list):
                        # Evitar problemas de encoding no PDF básico
                        nome_p = str(pendente).encode('latin-1', 'replace').decode('latin-1')
                        pdf.cell(0, 8, f'- {nome_p}', 0, 1, 'L')
                        
                    nome_pdf = f"Cobranca_RDC_{data_filtro_str}.pdf"
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                        pdf.output(tmp.name)
                        with open(tmp.name, "rb") as f:
                            pdf_bytes = f.read()
                            
                        # Backup Drive
                        success, msg = backup_google_drive(tmp.name, "application/pdf", nome_pdf)
                        if success:
                            st.toast("☁️ Relatório PDF salvo no Drive!")
                    
                    st.download_button(
                        label="⬇️ Baixar PDF",
                        data=pdf_bytes,
                        file_name=nome_pdf,
                        mime="application/pdf",
                        type="primary",
                        use_container_width=True
                    )
                    
            st.markdown("<br>", unsafe_allow_html=True)
            df_pend = pd.DataFrame({f"Encarregados Pendentes ({data_filtro_str})": sorted(pendentes_list)})
            st.dataframe(df_pend, hide_index=True, use_container_width=True)
        else:
            st.success(f"🎉 Todos os RDCs desta data ({data_filtro_str}) já foram entregues!")

    with tab_emissao:
        st.markdown("### Emissão de RDC")
        if not lista_encarregados_base:
            st.warning("Nenhum encarregado encontrado na base.")
        else:
            encarregados_sel = st.multiselect("Escolha o(s) Encarregado(s) (deixe vazio para gerar todos):", lista_encarregados_base, default=[lista_encarregados_base[0]] if lista_encarregados_base else [])
            lista_alvo = encarregados_sel if encarregados_sel else lista_encarregados_base
            
            equipe = df_atual[df_atual["ENCARREGADO"].isin(lista_alvo)]
            st.markdown("")
            texto_enc = ", ".join(lista_alvo) if len(lista_alvo) <= 3 else f"{len(lista_alvo)} Encarregados Selecionados"
            st.markdown(f"""<div style="background: {cor_card}; border-radius: 10px; padding: 20px; border: 1px solid {cor_borda}; margin-bottom: 16px;"><div style="text-align: center; border-bottom: 2px solid {cor_azul}; padding-bottom: 12px; margin-bottom: 12px;"><h3 style="margin: 0; font-size: 1.2rem; color: {cor_texto} !important;">RDC - Relatório Diário de Campo</h3><p style="color: {cor_texto_sub}; margin: 4px 0 0 0; font-size: 0.85rem;">{nome_site}</p></div><table style="width: 100%; color: {cor_texto}; font-size: 0.9rem;"><tr><td style="padding: 4px 0;"><strong>Encarregado(s):</strong></td><td>{texto_enc}</td></tr><tr><td style="padding: 4px 0;"><strong>Data:</strong></td><td>{datetime.datetime.now().strftime("%d/%m/%Y")}</td></tr><tr><td style="padding: 4px 0;"><strong>Efetivo:</strong></td><td>{len(equipe)} colaborador(es)</td></tr></table></div>""", unsafe_allow_html=True)
            
            if "ENCARREGADO" not in equipe.columns:
                st.dataframe(equipe[["MATRICULA", "NOME", "FUNÇÃO"]].reset_index(drop=True), hide_index=True, use_container_width=True)
            else:
                st.dataframe(equipe[["MATRICULA", "NOME", "FUNÇÃO", "ENCARREGADO"]].reset_index(drop=True), hide_index=True, use_container_width=True)
            
            st.markdown("")
            
            col_btn1, col_btn2, col_btn3, col_btn4 = st.columns(4)
            with col_btn1:
                if st.button("🟢 GERAR EXCEL", type="primary", use_container_width=True):
                    if len(lista_alvo) == 1:
                        enc = lista_alvo[0]
                        eq = df_atual[df_atual["ENCARREGADO"] == enc]
                        wb = preencher_excel(eq, enc)
                        if wb:
                            nome_limpo = enc.replace(" ", "_")
                            nome_arquivo = f"RDC_{nome_limpo}.xlsx"
                            buffer = io.BytesIO()
                            wb.save(buffer)
                            wb.close()
                            try:
                                hoje = datetime.datetime.now()
                                pasta_hist = os.path.join(pasta_base, "Historico_RDC", str(hoje.year), f"{hoje.month:02d}_{hoje.strftime('%B')}")
                                os.makedirs(pasta_hist, exist_ok=True)
                                caminho_local = os.path.join(pasta_hist, f"{hoje.strftime('%d_%H%M')}_{nome_arquivo}")
                                with open(caminho_local, "wb") as f:
                                    f.write(buffer.getvalue())
                                success, msg = backup_google_drive(caminho_local, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", f"{hoje.strftime('%d_%H%M')}_{nome_arquivo}")
                                if success:
                                    st.toast("☁️ Backup salvo no Google Drive!")
                            except Exception:
                                pass
                            buffer.seek(0)
                            st.download_button("⬇️ Baixar Planilha", data=buffer, file_name=nome_arquivo, mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", use_container_width=True)
                            st.success("✅ Gerado!")
                        else:
                            st.error("Modelo não encontrado.")
                    else:
                        with st.spinner("Gerando planilhas Excel em Lote..."):
                            zip_buffer = io.BytesIO()
                            qtd = 0
                            with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
                                for enc in lista_alvo:
                                    eq = df_atual[df_atual["ENCARREGADO"] == enc]
                                    if len(eq) > 0:
                                        wb_e = preencher_excel(eq, enc)
                                        if wb_e:
                                            buf = io.BytesIO()
                                            wb_e.save(buf)
                                            wb_e.close()
                                            n = enc.replace(" ", "_")
                                            zf.writestr(f"RDC_{n}.xlsx", buf.getvalue())
                                            qtd += 1
                                            try:
                                                hoje = datetime.datetime.now()
                                                pasta_hist = os.path.join(pasta_base, "Historico_RDC", str(hoje.year), f"{hoje.month:02d}_{hoje.strftime('%B')}")
                                                os.makedirs(pasta_hist, exist_ok=True)
                                                with open(os.path.join(pasta_hist, f"{hoje.strftime('%d_%H%M')}_RDC_{n}.xlsx"), "wb") as f:
                                                    f.write(buf.getvalue())
                                            except Exception:
                                                pass
                            if qtd > 0:
                                zip_buffer.seek(0)
                                nome_zip = f"LOTE_EXCEL_{datetime.datetime.now().strftime('%d_%m_%Y')}.zip"
                                st.download_button(f"⬇️ Baixar Planilhas ({qtd})", data=zip_buffer, file_name=nome_zip, mime="application/zip", use_container_width=True)
                                st.success(f"✅ {qtd} Planilhas Excel geradas no ZIP!")
                            else:
                                st.warning("Nenhuma planilha gerada.")

            with col_btn2:
                if st.button("📄 GERAR PDF", use_container_width=True):
                    if len(lista_alvo) == 1:
                        enc = lista_alvo[0]
                        eq = df_atual[df_atual["ENCARREGADO"] == enc]
                        pdf_bytes = gerar_pdf_rdc(eq, enc, nome_empresa=nome_site, logo_path=caminho_logo)
                        if pdf_bytes:
                            nome_limpo = enc.replace(" ", "_")
                            nome_pdf = f"RDC_{nome_limpo}.pdf"
                            st.download_button("⬇️ Baixar PDF", data=pdf_bytes, file_name=nome_pdf, mime="application/pdf", use_container_width=True)
                            st.success("✅ PDF gerado!")
                        else:
                            st.error("Erro ao gerar PDF.")
                    else:
                        with st.spinner("Gerando PDFs em Lote..."):
                            zip_buffer = io.BytesIO()
                            qtd = 0
                            with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
                                for enc in lista_alvo:
                                    eq = df_atual[df_atual["ENCARREGADO"] == enc]
                                    if len(eq) > 0:
                                        pdf_b = gerar_pdf_rdc(eq, enc, nome_empresa=nome_site, logo_path=caminho_logo)
                                        if pdf_b:
                                            zf.writestr(f"RDC_{enc.replace(' ', '_')}.pdf", pdf_b)
                                            qtd += 1
                            if qtd > 0:
                                zip_buffer.seek(0)
                                nome_zip = f"LOTE_PDF_{datetime.datetime.now().strftime('%d_%m_%Y')}.zip"
                                st.download_button(f"⬇️ Baixar PDFs ({qtd})", data=zip_buffer, file_name=nome_zip, mime="application/zip", use_container_width=True)
                                st.success(f"✅ {qtd} PDFs gerados no ZIP!")
                            else:
                                st.warning("Nenhum PDF gerado.")

            with col_btn3:
                if st.button("🚀 GERAR LOTE (.ZIP)", use_container_width=True):
                    with st.spinner("Gerando PDFs e planilhas..."):
                        try:
                            zip_buffer = io.BytesIO()
                            qtd = 0
                            with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
                                for enc in lista_alvo:
                                    eq = df_atual[df_atual["ENCARREGADO"] == enc]
                                    if len(eq) > 0:
                                        wb_e = preencher_excel(eq, enc)
                                        if wb_e:
                                            buf = io.BytesIO()
                                            wb_e.save(buf)
                                            wb_e.close()
                                            n = enc.replace(" ", "_")
                                            zf.writestr(f"RDC_{n}.xlsx", buf.getvalue())
                                            try:
                                                hoje = datetime.datetime.now()
                                                pasta_hist = os.path.join(pasta_base, "Historico_RDC", str(hoje.year), f"{hoje.month:02d}_{hoje.strftime('%B')}")
                                                os.makedirs(pasta_hist, exist_ok=True)
                                                with open(os.path.join(pasta_hist, f"{hoje.strftime('%d_%H%M')}_RDC_{n}.xlsx"), "wb") as f:
                                                    f.write(buf.getvalue())
                                            except Exception:
                                                pass
                                        
                                        pdf_b = gerar_pdf_rdc(eq, enc, nome_empresa=nome_site, logo_path=caminho_logo)
                                        if pdf_b:
                                            zf.writestr(f"RDC_{enc.replace(' ', '_')}.pdf", pdf_b)
                                        
                                        qtd += 1
                            if qtd > 0:
                                zip_buffer.seek(0)
                                nome_zip = f"LOTE_RDC_COMPLETO_{datetime.datetime.now().strftime('%d_%m_%Y')}.zip"
                                st.download_button(f"⬇️ Baixar Lote ({qtd} pastas/encarregados)", data=zip_buffer, file_name=nome_zip, mime="application/zip", use_container_width=True)
                                
                                try:
                                    with tempfile.NamedTemporaryFile(delete=False, suffix=".zip") as tmp_zip:
                                        tmp_zip.write(zip_buffer.getvalue())
                                        tmp_zip_path = tmp_zip.name
                                    success, msg = backup_google_drive(tmp_zip_path, "application/zip", nome_zip)
                                    if success:
                                        st.toast("☁️ Lote salvo no Google Drive!")
                                    os.remove(tmp_zip_path)
                                except:
                                    pass
                                    
                                st.success(f"✅ Arquivos de {qtd} encarregados gerados no ZIP!")
                            else:
                                st.warning("Nenhuma planilha gerada.")
                        except Exception as e:
                            st.error(f"Erro: {e}")

            with col_btn4:
                if st.button("🖨️ PDF ÚNICO (Impressão)", use_container_width=True):
                    with st.spinner("Gerando PDF único..."):
                        try:
                            import fitz  # PyMuPDF
                            pdf_final = fitz.open()
                            qtd_enc = 0
                            for enc in lista_alvo:
                                eq = df_atual[df_atual["ENCARREGADO"] == enc]
                                if len(eq) > 0:
                                    pdf_b = gerar_pdf_rdc(eq, enc, nome_empresa=nome_site, logo_path=caminho_logo)
                                    if pdf_b:
                                        pdf_individual = fitz.open(stream=pdf_b, filetype="pdf")
                                        pdf_final.insert_pdf(pdf_individual)
                                        pdf_individual.close()
                                        qtd_enc += 1
                            if qtd_enc > 0:
                                buffer_pdf_unico = io.BytesIO(pdf_final.tobytes(deflate=True))
                                pdf_final.close()
                                nome_pdf_unico = f"TODOS_RDC_{datetime.datetime.now().strftime('%d_%m_%Y')}.pdf"
                                st.download_button(
                                    f"⬇️ Baixar PDF Único ({qtd_enc} encarregados)",
                                    data=buffer_pdf_unico,
                                    file_name=nome_pdf_unico,
                                    mime="application/pdf",
                                    use_container_width=True
                                )
                                st.success(f"✅ PDF único gerado com {qtd_enc} RDCs!")
                            else:
                                pdf_final.close()
                                st.warning("⚠️ Nenhum RDC encontrado.")
                        except Exception as e:
                            st.error(f"Erro ao gerar PDF único: {e}")

    with tab_escala:
        st.markdown("### 📋 Escala Diária de Efetivo")
        st.markdown("Marque quem da equipe está escalado para trabalhar no dia selecionado. Os dados são salvos para controle do apontamento.")
        
        col_esc1, col_esc2 = st.columns(2)
        with col_esc1:
            data_escala = st.date_input("Data da Escala:", datetime.date.today(), key="data_escala_input")
            data_esc_str = data_escala.strftime("%Y-%m-%d")
        with col_esc2:
            if not lista_encarregados_base:
                st.warning("Nenhum encarregado encontrado na base.")
                enc_escala_sel = None
            else:
                enc_escala_sel = st.selectbox("Escolha o Encarregado:", lista_encarregados_base, key="enc_escala_sel")
                
        if enc_escala_sel:
            # Pegar a equipe atual do encarregado selecionado da base PDE (df_atual)
            equipe_base = df_atual[df_atual["ENCARREGADO"] == enc_escala_sel][["MATRICULA", "NOME", "FUNÇÃO", "ENCARREGADO"]].copy()
            
            # Carregar o arquivo de escala, se existir
            if os.path.exists(caminho_escala_csv):
                try:
                    df_escala_hist = pd.read_csv(caminho_escala_csv)
                except Exception:
                    df_escala_hist = pd.DataFrame(columns=["DATA", "MATRICULA", "NOME", "ENCARREGADO", "ESCALADO"])
            else:
                df_escala_hist = pd.DataFrame(columns=["DATA", "MATRICULA", "NOME", "ENCARREGADO", "ESCALADO"])
                
            # Filtrar histórico de escala para a data e encarregado específicos
            df_escala_hoje = df_escala_hist[(df_escala_hist["DATA"] == data_esc_str) & (df_escala_hist["ENCARREGADO"] == enc_escala_sel)]
            
            # Se já houver salvamento para este dia, usar o valor salvo. Senão, todos "Sim" (True)
            if not df_escala_hoje.empty:
                # Forçar tipo string na matrícula para evitar erros de merge no pandas
                equipe_base['MATRICULA'] = equipe_base['MATRICULA'].astype(str).str.replace(r'\.0$', '', regex=True).str.strip()
                df_escala_hoje_tmp = df_escala_hoje[["MATRICULA", "ESCALADO"]].copy()
                df_escala_hoje_tmp['MATRICULA'] = df_escala_hoje_tmp['MATRICULA'].astype(str).str.replace(r'\.0$', '', regex=True).str.strip()
                
                # Merge base com histórico
                equipe_edit = pd.merge(equipe_base, df_escala_hoje_tmp, on="MATRICULA", how="left")
                # Quem não estava no histórico, recebe True como padrão
                equipe_edit["ESCALADO"] = equipe_edit["ESCALADO"].fillna(True)
            else:
                equipe_edit = equipe_base.copy()
                equipe_edit["ESCALADO"] = True
                
            # Organizar colunas
            equipe_edit = equipe_edit[["ESCALADO", "MATRICULA", "NOME", "FUNÇÃO", "ENCARREGADO"]]
            
            # Mostrar editor interativo (data_editor)
            st.markdown(f"**Equipe de {enc_escala_sel} ({len(equipe_edit)} pessoas)**")
            
            # Desabilitar edição das colunas de identificação
            config_colunas = {
                "ESCALADO": st.column_config.CheckboxColumn("Escalado (Sim/Não)?", help="Marque se o funcionário vai trabalhar hoje.", default=True),
                "MATRICULA": st.column_config.TextColumn("Matrícula", disabled=True),
                "NOME": st.column_config.TextColumn("Nome", disabled=True),
                "FUNÇÃO": st.column_config.TextColumn("Função", disabled=True),
                "ENCARREGADO": None # Ocultar coluna
            }
            
            edited_df = st.data_editor(
                equipe_edit,
                column_config=config_colunas,
                hide_index=True,
                use_container_width=True,
                key=f"editor_escala_{enc_escala_sel}_{data_esc_str}"
            )
            
            # Botões de Ação
            col_action1, col_action2 = st.columns(2)
            with col_action1:
                if st.button("💾 Salvar Escala da Equipe", type="primary", use_container_width=True):
                    try:
                        # Preparar os dados editados
                        edited_df["DATA"] = data_esc_str
                        edited_df["ENCARREGADO"] = enc_escala_sel
                        
                        novos_dados = edited_df[["DATA", "MATRICULA", "NOME", "ENCARREGADO", "ESCALADO"]].copy()
                        
                        if not df_escala_hist.empty:
                            # Remover os dados antigos deste encarregado neste dia
                            mask_remover = (df_escala_hist["DATA"] == data_esc_str) & (df_escala_hist["ENCARREGADO"] == enc_escala_sel)
                            df_escala_hist = df_escala_hist[~mask_remover]
                            
                            # Concatenar
                            df_final = pd.concat([df_escala_hist, novos_dados], ignore_index=True)
                        else:
                            df_final = novos_dados
                            
                        # Salvar CSV Local
                        df_final.to_csv(caminho_escala_csv, index=False)
                        
                        # Backup Google Drive
                        success, msg = backup_google_drive(caminho_escala_csv, "text/csv", f"escala_diaria_{datetime.datetime.now().strftime('%d%m%Y')}.csv")
                        if success:
                            st.toast("☁️ Escala salva na nuvem com sucesso!")
                            
                        st.success(f"✅ Escala da equipe de **{enc_escala_sel}** para o dia **{data_escala.strftime('%d/%m/%Y')}** salva com sucesso!")
                        time.sleep(2)
                        st.rerun()
                    except Exception as e:
                        st.error(f"Erro ao salvar a escala: {e}")
                        
            with col_action2:
                # Gerar Excel em memória para download
                try:
                    buffer_escala = io.BytesIO()
                    with pd.ExcelWriter(buffer_escala, engine='openpyxl') as writer:
                        # Formatar DataFrame para Excel (remover colunas indesejadas, traduzir booleanos)
                        df_excel = edited_df.copy()
                        df_excel.insert(0, "DATA", data_esc_str) # Adiciona a data como primeira coluna
                        df_excel["ESCALADO"] = df_excel["ESCALADO"].apply(lambda x: "SIM" if x else "NÃO")
                        df_excel.to_excel(writer, index=False, sheet_name="Escala")
                    buffer_escala.seek(0)
                    
                    st.download_button(
                        label="⬇️ Baixar Escala em Excel",
                        data=buffer_escala,
                        file_name=f"Escala_{enc_escala_sel.replace(' ', '_')}_{data_esc_str}.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        use_container_width=True
                    )
                except Exception as e:
                    st.error(f"Erro ao gerar Excel: {e}")

        # === VISUALIZADOR DE ESCALAS SALVAS ===
        st.markdown("---")
        with st.expander("📂 Ver Banco de Dados da Escala (Histórico)", expanded=False):
            if os.path.exists(caminho_escala_csv):
                try:
                    df_historico_view = pd.read_csv(caminho_escala_csv)
                    if not df_historico_view.empty:
                        col_filtro1, col_filtro2 = st.columns(2)
                        with col_filtro1:
                            filtro_data = st.date_input("Filtrar por Data:", value=None, key="filtro_data_hist")
                        with col_filtro2:
                            encarregados_hist = ["Todos"] + sorted(df_historico_view['ENCARREGADO'].dropna().unique().tolist())
                            filtro_enc = st.selectbox("Filtrar por Encarregado:", encarregados_hist, key="filtro_enc_hist")
                            
                        df_view_filtrado = df_historico_view.copy()
                        if filtro_data:
                            df_view_filtrado = df_view_filtrado[df_view_filtrado["DATA"] == filtro_data.strftime("%Y-%m-%d")]
                        if filtro_enc != "Todos":
                            df_view_filtrado = df_view_filtrado[df_view_filtrado["ENCARREGADO"] == filtro_enc]
                            
                        st.dataframe(df_view_filtrado, use_container_width=True, hide_index=True)
                        
                        # Resumo
                        if 'ESCALADO' in df_view_filtrado.columns:
                            qtd_escalados = len(df_view_filtrado[df_view_filtrado['ESCALADO'] == True])
                            st.caption(f"Mostrando {len(df_view_filtrado)} registros no total. (Pessoas escaladas: {qtd_escalados})")
                    else:
                        st.info("Nenhuma escala salva ainda.")
                except Exception as e:
                    st.error(f"Erro ao ler histórico: {e}")
            else:
                st.info("O banco de dados de escalas ainda não existe. Salve uma escala primeiro.")

    with tab_f1:
        st.markdown("### 🏎️ Competição F1 — Entrega de RDC")
        st.markdown("Acompanhamento mensal, ranking de pontualidade e assiduidade dos Encarregados na entrega dos Relatórios Diários de Campo.")

        # === SINCRONIZAÇÃO AUTOMÁTICA DO HISTÓRICO COM BANCO DE RDC ===
        if os.path.exists(caminho_rdc_registros_csv):
            try:
                df_banco_rdc = pd.read_csv(caminho_rdc_registros_csv)
                if not df_banco_rdc.empty and "DATA" in df_banco_rdc.columns and "ENCARREGADO" in df_banco_rdc.columns:
                    df_banco_val = df_banco_rdc[["DATA", "ENCARREGADO"]].dropna().copy()
                    df_banco_val["ENCARREGADO"] = df_banco_val["ENCARREGADO"].astype(str).str.strip().str.upper()
                    df_banco_val = df_banco_val[df_banco_val["ENCARREGADO"] != ""]
                    
                    if not df_banco_val.empty:
                        if st.session_state.df_historico_f1.empty:
                            st.session_state.df_historico_f1 = df_banco_val.drop_duplicates(subset=["DATA", "ENCARREGADO"])
                        else:
                            st.session_state.df_historico_f1 = pd.concat([st.session_state.df_historico_f1, df_banco_val], ignore_index=True).drop_duplicates(subset=["DATA", "ENCARREGADO"])
            except Exception:
                pass

        # === MAPEAMENTO DE DISCIPLINAS POR ENCARREGADO ===
        dict_enc_disciplina = {}
        if "DISCIPLINA" in df_atual.columns and "ENCARREGADO" in df_atual.columns:
            for enc_n in df_atual["ENCARREGADO"].dropna().unique():
                enc_n_str = str(enc_n).strip().upper()
                if enc_n_str:
                    sub_d = df_atual[df_atual["ENCARREGADO"].astype(str).str.upper().str.strip() == enc_n_str]
                    if not sub_d.empty:
                        d_val = sub_d["DISCIPLINA"].dropna()
                        if not d_val.empty:
                            dict_enc_disciplina[enc_n_str] = str(d_val.iloc[0]).strip().upper()
                        else:
                            dict_enc_disciplina[enc_n_str] = "GERAL"

        # === PAINEL: GERENCIAR LISTA DE ENCARREGADOS ===
        if st.session_state.get("role_usuario") != "apontador" and st.toggle("👥 Gerenciar Lista de Encarregados do F1", key="toggle_gerenciar_lista_f1"):
            st.markdown("""
            <div style="background: rgba(14, 165, 233, 0.1); border: 1px solid rgba(14, 165, 233, 0.3); border-radius: 12px; padding: 15px; margin-bottom: 15px;">
                <p style="margin: 0; color: #94a3b8; font-size: 14px;">⚙️ Gerencie a lista oficial de encarregados. Você pode <b>adicionar</b>, <b>remover</b> ou <b>sincronizar automaticamente</b> com base na função 'ENCARREGADO' do PDE.</p>
            </div>
            """, unsafe_allow_html=True)
            
            # Botão de Sincronização Automática por Função
            col_auto_sync, col_auto_info = st.columns([1, 2])
            with col_auto_sync:
                if st.button("🔄 Auto-Sincronizar do PDE (Apenas Função 'ENCARREGADO')", key="btn_sync_enc_funcao", type="primary", use_container_width=True):
                    if "FUNÇÃO" in df_atual.columns:
                        mask_enc_func = df_atual["FUNÇÃO"].astype(str).str.upper().str.contains(r"ENCARREGAD|ENC\b|LIDER|LÍDER")
                        df_encs_auto = df_atual[mask_enc_func]
                        
                        nomes_pde_func = []
                        if "NOME" in df_encs_auto.columns:
                            nomes_pde_func.extend(df_encs_auto["NOME"].dropna().astype(str).str.strip().str.upper().unique().tolist())
                        if "ENCARREGADO" in df_encs_auto.columns:
                            nomes_pde_func.extend(df_encs_auto["ENCARREGADO"].dropna().astype(str).str.strip().str.upper().unique().tolist())
                            
                        nomes_pde_func = [n for n in set(nomes_pde_func) if n and len(n) > 3 and n != "AJUSTAR NOME"]
                        
                        if nomes_pde_func:
                            lista_unificada = sorted(list(set(encarregados_f1_oficial + nomes_pde_func)))
                            with open(caminho_f1_json, "w", encoding="utf-8") as f:
                                json.dump(lista_unificada, f, ensure_ascii=False, indent=2)
                            st.success(f"✅ Sincronização concluída! {len(lista_unificada)} encarregados mapeados a partir da função no PDE.")
                            time.sleep(2)
                            st.rerun()
                        else:
                            st.warning("⚠️ Nenhum colaborador com 'ENCARREGADO' na coluna Função foi encontrado no PDE.")
                    else:
                        st.error("Coluna 'FUNÇÃO' não encontrada no PDE.")
            with col_auto_info:
                st.caption("Puxa automaticamente todos os colaboradores que possuem 'ENCARREGADO' ou 'LÍDER' cadastrados na coluna Função do PDE.")

            st.markdown("<br>", unsafe_allow_html=True)
            col_add, col_rem = st.columns(2)
            
            with col_add:
                with st.form("form_add_enc_f1"):
                    st.markdown("**➕ Adicionar Encarregado Manualmente**")
                    novo_nome = st.text_input("Nome completo do Encarregado:", placeholder="Ex: JOÃO DA SILVA SOUZA")
                    btn_add = st.form_submit_button("Adicionar à Lista", type="secondary", use_container_width=True)
                    if btn_add and novo_nome.strip():
                        nome_upper = novo_nome.strip().upper()
                        if nome_upper in lista_completa_encarregados:
                            st.warning(f"⚠️ '{nome_upper}' já está na lista!")
                        else:
                            encarregados_f1_oficial.append(nome_upper)
                            with open(caminho_f1_json, "w", encoding="utf-8") as f:
                                json.dump(encarregados_f1_oficial, f, ensure_ascii=False, indent=2)
                            st.success(f"✅ '{nome_upper}' adicionado com sucesso!")
                            time.sleep(2)
                            st.rerun()
            
            with col_rem:
                with st.form("form_rem_enc_f1"):
                    st.markdown("**🗑️ Remover Encarregado**")
                    enc_remover = st.multiselect("Selecione quem remover:", lista_completa_encarregados)
                    btn_rem = st.form_submit_button("Remover da Lista", type="secondary", use_container_width=True)
                    if btn_rem and enc_remover:
                        lista_atualizada = [e for e in encarregados_f1_oficial if e.upper() not in enc_remover]
                        with open(caminho_f1_json, "w", encoding="utf-8") as f:
                            json.dump(lista_atualizada, f, ensure_ascii=False, indent=2)
                        st.success(f"✅ {len(enc_remover)} encarregado(s) removido(s)!")
                        time.sleep(2)
                        st.rerun()
            
            st.caption(f"📋 Total atual na lista: **{len(lista_completa_encarregados)}** encarregados oficiais")

        # === PAINEL: ABONAR FALTAS ===
        if st.session_state.get("role_usuario") != "apontador" and st.toggle("⏸️ Abonar Faltas (Folga / Atestado / Feriado)", key="toggle_abono_f1"):
            st.markdown("""
            <div style="background: rgba(245, 158, 11, 0.1); border: 1px solid rgba(245, 158, 11, 0.3); border-radius: 12px; padding: 15px; margin-bottom: 15px;">
                <p style="margin: 0; color: #94a3b8; font-size: 14px;">📝 Marque os dias em que o encarregado <b style="color: #f59e0b;">não precisava</b> entregar o RDC. Esses dias aparecerão como <b style="color: #f59e0b;">⏸️</b> e não penalizam a taxa de aproveitamento.</p>
            </div>
            """, unsafe_allow_html=True)
            
            with st.form("form_abono_f1"):
                col_ab1, col_ab2 = st.columns(2)
                with col_ab1:
                    datas_abono = st.date_input("📅 Data(s) do Abono:", value=datetime.date.today(), key="datas_abono_input")
                with col_ab2:
                    motivo_abono = st.selectbox("Motivo:", ["FOLGA", "ATESTADO MÉDICO", "FERIADO", "CHUVA / INTEMPÉRIE", "FALTA JUSTIFICADA", "OUTRO"])
                
                encs_abono = st.multiselect("Selecione os Encarregados para abonar:", lista_completa_encarregados, key="encs_abono_multi")
                
                btn_abono = st.form_submit_button("✅ Registrar Abono", type="primary", use_container_width=True)
                if btn_abono and encs_abono:
                    novos_abonos = []
                    if isinstance(datas_abono, (list, tuple)):
                        lista_datas = [d.strftime("%Y-%m-%d") for d in datas_abono]
                    else:
                        lista_datas = [datas_abono.strftime("%Y-%m-%d")]
                    
                    for data_ab in lista_datas:
                        for enc_ab in encs_abono:
                            ja_existe = False
                            if not st.session_state.df_f1_excecoes.empty:
                                ja_existe = ((st.session_state.df_f1_excecoes["DATA"] == data_ab) & (st.session_state.df_f1_excecoes["ENCARREGADO"] == enc_ab)).any()
                            if not ja_existe:
                                novos_abonos.append({"DATA": data_ab, "ENCARREGADO": enc_ab, "MOTIVO": motivo_abono})
                    
                    if novos_abonos:
                        df_novos_ab = pd.DataFrame(novos_abonos)
                        st.session_state.df_f1_excecoes = pd.concat([st.session_state.df_f1_excecoes, df_novos_ab], ignore_index=True)
                        st.session_state.df_f1_excecoes.to_csv(caminho_f1_excecoes, index=False)
                        st.success(f"✅ {len(novos_abonos)} abono(s) registrado(s) com sucesso!")
                        time.sleep(2)
                        st.rerun()
                    else:
                        st.info("ℹ️ Todos os abonos selecionados já estavam cadastrados.")
            
            if not st.session_state.df_f1_excecoes.empty:
                st.markdown("**Abonos registrados:**")
                df_exc_show = st.session_state.df_f1_excecoes.copy()
                df_exc_show = df_exc_show.sort_values("DATA", ascending=False).head(20)
                st.dataframe(df_exc_show, hide_index=True, use_container_width=True)
                
                if st.button("🗑️ Limpar Todos os Abonos", key="btn_limpar_abonos"):
                    st.session_state.df_f1_excecoes = pd.DataFrame(columns=["DATA", "ENCARREGADO", "MOTIVO"])
                    if os.path.exists(caminho_f1_excecoes):
                        os.remove(caminho_f1_excecoes)
                    st.success("✅ Todos os abonos foram removidos!")
                    time.sleep(2)
                    st.rerun()

        # --- LANÇAMENTO MANUAL ---
        if st.session_state.get("role_usuario") != "apontador" and st.toggle("➕ Lançar RDC Manualmente (Para papéis ilegíveis ou atrasados)"):
            with st.form("form_f1_manual"):
                st.info("💡 Você pode colar a lista inteira de encarregados aqui (um por linha ou separados por vírgula). O robô vai verificar: se a IA já tiver lido, ele ignora. Se faltou, ele adiciona!")
                col_m1, col_m2 = st.columns([1, 2])
                with col_m1:
                    data_manual = st.date_input("Data do RDC")
                with col_m2:
                    nomes_colados = st.text_area("Cole os nomes dos Encarregados", height=120)
                
                btn_manual = st.form_submit_button("Processar Lista e Lançar no F1")
                if btn_manual and nomes_colados.strip():
                    import re
                    import difflib
                    data_str = data_manual.strftime("%Y-%m-%d")
                    
                    lista_suja = [n.strip().upper() for n in re.split(r'[\n,;]', nomes_colados) if n.strip()]
                    novos_registros = []
                    nomes_ja_existentes = 0
                    nomes_nao_encontrados = []
                    
                    for nome_sujo in lista_suja:
                        match = difflib.get_close_matches(nome_sujo, lista_completa_encarregados, n=1, cutoff=0.55)
                        if match:
                            nome_oficial = match[0]
                            ja_existe = ((st.session_state.df_historico_f1["DATA"] == data_str) & (st.session_state.df_historico_f1["ENCARREGADO"] == nome_oficial)).any()
                            if ja_existe:
                                nomes_ja_existentes += 1
                            else:
                                novos_registros.append({"DATA": data_str, "ENCARREGADO": nome_oficial})
                        else:
                            nomes_nao_encontrados.append(nome_sujo)
                            
                    if novos_registros:
                        df_novos = pd.DataFrame(novos_registros)
                        if conn and not st.session_state.get('force_use_local', False):
                            try:
                                df_fresco = conn.read(worksheet="Historico_F1", ttl=0)
                                if not df_fresco.empty:
                                    df_fresco = df_fresco.dropna(how='all')
                                    df_final = pd.concat([df_fresco, df_novos], ignore_index=True).drop_duplicates(subset=["DATA", "ENCARREGADO"])
                                else:
                                    df_final = pd.concat([st.session_state.df_historico_f1, df_novos], ignore_index=True).drop_duplicates(subset=["DATA", "ENCARREGADO"])
                                
                                ok, msg = salvar_f1_seguro(conn, df_final, caminho_historico_f1_csv)
                                if ok:
                                    st.session_state.df_historico_f1 = df_final
                                st.cache_data.clear()
                                st.toast(f"{len(novos_registros)} novos RDCs sincronizados com a nuvem! ({nomes_ja_existentes} já constavam).", icon="✅")
                            except Exception as e:
                                st.error(f"Erro ao salvar na nuvem: {e}")
                                st.session_state.df_historico_f1 = pd.concat([st.session_state.df_historico_f1, df_novos], ignore_index=True).drop_duplicates(subset=["DATA", "ENCARREGADO"])
                                st.session_state.df_historico_f1.to_csv(caminho_historico_f1_csv, index=False)
                        else:
                            st.session_state.df_historico_f1 = pd.concat([st.session_state.df_historico_f1, df_novos], ignore_index=True).drop_duplicates(subset=["DATA", "ENCARREGADO"])
                            st.session_state.df_historico_f1.to_csv(caminho_historico_f1_csv, index=False)
                            st.toast(f"{len(novos_registros)} novos RDCs adicionados localmente! ({nomes_ja_existentes} já constavam).", icon="✅")
                    elif nomes_ja_existentes > 0:
                        st.warning(f"⚠️ Todos os nomes reconhecidos ({nomes_ja_existentes}) já estavam devidamente lançados neste dia!")
                        
                    if nomes_nao_encontrados:
                        st.error(f"❌ Não encontrei na lista oficial: {', '.join(nomes_nao_encontrados)}")
                        
                    time.sleep(3)
                    st.rerun()

        # === FILTROS PRINCIPAIS: MÊS E DISCIPLINA ===
        df_hist = st.session_state.df_historico_f1.copy()
        if not df_hist.empty:
            df_hist["DATA"] = pd.to_datetime(df_hist["DATA"], format="%Y-%m-%d", errors="coerce")
            df_hist = df_hist.dropna(subset=["DATA"])
            df_hist["MES_ANO"] = df_hist["DATA"].dt.strftime("%Y-%m")
            meses_disponiveis = sorted(df_hist["MES_ANO"].unique(), reverse=True)
        else:
            meses_disponiveis = [datetime.date.today().strftime("%Y-%m")]
            
        col_f1_mes, col_f1_disc = st.columns([1, 1])
        with col_f1_mes:
            mes_selecionado = st.selectbox("📅 Selecione o Mês para Análise:", meses_disponiveis)
        with col_f1_disc:
            disciplinas_disponiveis = sorted(list(set([d for d in dict_enc_disciplina.values() if d and d != "GERAL"])))
            filtro_f1_disc = st.selectbox("🎯 Filtrar por Disciplina:", ["Todas as Disciplinas"] + disciplinas_disponiveis)

        if not df_hist.empty:
            df_mes = df_hist[df_hist["MES_ANO"] == mes_selecionado]
        else:
            df_mes = pd.DataFrame(columns=["DATA", "ENCARREGADO"])
            
        import calendar
        ano, mes = map(int, mes_selecionado.split('-'))
        num_dias = calendar.monthrange(ano, mes)[1]
        
        nomes_no_mes = df_mes["ENCARREGADO"].dropna().unique().tolist() if not df_mes.empty else []
        todos_encarregados_matriz = sorted(list(set(lista_completa_encarregados + encarregados_f1_oficial + nomes_no_mes)))
        todos_encarregados_matriz = [e for e in todos_encarregados_matriz if str(e).strip() != "" and str(e).upper() != "AJUSTAR NOME"]
        
        # Aplicar filtro por disciplina se selecionado
        if filtro_f1_disc != "Todas as Disciplinas":
            todos_encarregados_matriz = [e for e in todos_encarregados_matriz if dict_enc_disciplina.get(e, "GERAL") == filtro_f1_disc]
        
        dias_str = [str(d) for d in range(1, num_dias + 1)]
        
        # Identificar fins de semana
        dias_fim_de_semana = set()
        for d in range(1, num_dias + 1):
            data_check = datetime.date(ano, mes, d)
            if data_check.weekday() >= 5:
                dias_fim_de_semana.add(str(d))
        
        dias_uteis = [d for d in dias_str if d not in dias_fim_de_semana]
        
        if not todos_encarregados_matriz:
            st.info("ℹ️ Nenhum encarregado encontrado para os filtros selecionados.")
        else:
            # Montar Matriz
            matriz = pd.DataFrame(index=todos_encarregados_matriz, columns=dias_str)
            for col in dias_str:
                if col in dias_fim_de_semana:
                    matriz[col] = "➖"
                else:
                    matriz[col] = "❌"
            
            for _, row in df_mes.iterrows():
                dia = str(row["DATA"].day)
                enc = str(row["ENCARREGADO"]).strip().upper()
                if enc in matriz.index and dia not in dias_fim_de_semana:
                    matriz.loc[enc, dia] = "✅"
            
            # Aplicar Abonos
            if not st.session_state.df_f1_excecoes.empty:
                df_exc_mes = st.session_state.df_f1_excecoes.copy()
                df_exc_mes["DATA"] = pd.to_datetime(df_exc_mes["DATA"], errors='coerce')
                df_exc_mes = df_exc_mes.dropna(subset=["DATA"])
                df_exc_mes = df_exc_mes[df_exc_mes["DATA"].dt.strftime("%Y-%m") == mes_selecionado]
                
                for _, row_exc in df_exc_mes.iterrows():
                    dia_exc = str(row_exc["DATA"].day)
                    enc_exc = str(row_exc["ENCARREGADO"]).strip().upper()
                    if enc_exc in matriz.index and dia_exc not in dias_fim_de_semana:
                        if matriz.loc[enc_exc, dia_exc] == "❌":
                            matriz.loc[enc_exc, dia_exc] = "⏸️"
            
            # Identificar dias decorridos no mês
            hoje_ref = datetime.date.today()
            if ano == hoje_ref.year and mes == hoje_ref.month:
                dias_decorridos_uteis = [d for d in dias_uteis if int(d) <= hoje_ref.day]
            elif (ano < hoje_ref.year) or (ano == hoje_ref.year and mes < hoje_ref.month):
                dias_decorridos_uteis = dias_uteis.copy()
            else:
                dias_decorridos_uteis = []

            # Totais e Aproveitamento (%)
            matriz["Total"] = (matriz[dias_uteis] == "✅").sum(axis=1)
            
            lista_aprov = []
            lista_badges = []
            lista_disciplinas_col = []
            
            for enc_i in matriz.index:
                lista_disciplinas_col.append(dict_enc_disciplina.get(enc_i, "GERAL"))
                
                # Cálculo de aproveitamento considerando apenas dias úteis decorridos e sem abonos
                if dias_decorridos_uteis:
                    dias_esp = sum(1 for d in dias_decorridos_uteis if matriz.loc[enc_i, d] != "⏸️")
                    entregues_ate_agora = sum(1 for d in dias_decorridos_uteis if matriz.loc[enc_i, d] == "✅")
                    aprov_val = round((entregues_ate_agora / dias_esp * 100), 1) if dias_esp > 0 else 0.0
                else:
                    aprov_val = 0.0
                lista_aprov.append(f"{aprov_val:.1f}%")
                
                # Badges / Conquistas Gamificadas
                badges_enc = []
                if aprov_val >= 100.0 and len(dias_decorridos_uteis) >= 2:
                    badges_enc.append("🚀 100% Imbatível")
                elif aprov_val >= 90.0:
                    badges_enc.append("⭐ Padrão Ouro")
                    
                if len(dias_decorridos_uteis) >= 3:
                    ultimos_3_dias = dias_decorridos_uteis[-3:]
                    if all(matriz.loc[enc_i, d] == "✅" for d in ultimos_3_dias):
                        badges_enc.append("🔥 Em Chamas")
                    elif all(matriz.loc[enc_i, d] == "❌" for d in ultimos_3_dias):
                        badges_enc.append("⚠️ Box / Pit Stop")
                        
                lista_badges.append(" ".join(badges_enc) if badges_enc else "—")

            matriz["Disciplina"] = lista_disciplinas_col
            matriz["Aproveitamento"] = lista_aprov
            matriz["Conquistas"] = lista_badges

            # Cabeçalhos com contagem por dia
            total_por_dia = (matriz[dias_str] == "✅").sum(axis=0)
            novas_colunas = {}
            for dia in dias_str:
                if dia in dias_fim_de_semana:
                    data_check = datetime.date(ano, mes, int(dia))
                    nome_dia = "SAB" if data_check.weekday() == 5 else "DOM"
                    novas_colunas[dia] = f"{dia}\n({nome_dia})"
                else:
                    novas_colunas[dia] = f"{dia}\n({total_por_dia[dia]})"
            
            # Cópia para ranking antes de renomear as colunas
            df_rank = matriz[["Total", "Aproveitamento", "Disciplina", "Conquistas"]].copy().reset_index()
            df_rank.columns = ["ENCARREGADO", "ENTREGAS", "APROVEITAMENTO_STR", "DISCIPLINA", "CONQUISTAS"]
            df_rank["APROV_NUM"] = df_rank["APROVEITAMENTO_STR"].str.replace('%', '').astype(float)
            df_rank = df_rank.sort_values(by=["ENTREGAS", "APROV_NUM"], ascending=[False, False]).reset_index(drop=True)

            matriz.rename(columns=novas_colunas, inplace=True)

            total_entregue = matriz["Total"].sum()
            
            # Média de aproveitamento geral
            media_aprov = 0.0
            if lista_aprov:
                vals_num = [float(a.replace('%', '')) for a in lista_aprov]
                media_aprov = round(sum(vals_num) / len(vals_num), 1)

            st.markdown("---")
            col_tit, col_met1, col_met2, col_met3 = st.columns([3, 1, 1, 1])
            with col_tit:
                st.markdown(f"#### 📊 Matriz de Entregas & Assiduidade — {mes_selecionado}")
            with col_met1:
                st.metric("📄 RDCs Entregues", total_entregue)
            with col_met2:
                st.metric("🎯 Aproveitamento Médio", f"{media_aprov}%")
            with col_met3:
                st.metric("👷 Encarregados", len(matriz))

            # Alerta de Devedores Críticos
            if mes_selecionado == datetime.date.today().strftime("%Y-%m") and len(dias_decorridos_uteis) >= 3:
                ultimos_3 = dias_decorridos_uteis[-3:]
                devedores = []
                for enc in matriz.index:
                    if all(matriz.loc[enc, novas_colunas[dia]] == "❌" for dia in ultimos_3):
                        devedores.append(enc)
                if devedores:
                    st.error(f"🚨 **ALERTA CRÍTICO:** {len(devedores)} encarregados não entregaram RDC nos últimos 3 dias úteis consecutivos.")
                    if st.toggle("👀 Ver encarregados com pendência crítica (Box)", key="tgl_dev_criticos"):
                        dados_dev = [{"Encarregado": enc, "Disciplina": dict_enc_disciplina.get(enc, 'GERAL'), "Faltas no Mês": len(dias_decorridos_uteis) - sum(1 for d in dias_decorridos_uteis if matriz.loc[enc, novas_colunas[d]] == "✅")} for enc in devedores]
                        df_dev = pd.DataFrame(dados_dev).sort_values(by="Faltas no Mês", ascending=False)
                        st.dataframe(df_dev, hide_index=True, use_container_width=True)

            def cor_fundo(valor):
                if valor == "✅":
                    return "background-color: rgba(74, 222, 128, 0.2); color: #4ade80; font-weight: bold;"
                elif valor == "❌":
                    return "background-color: rgba(255, 75, 75, 0.2); color: #ff4b4b; font-weight: bold;"
                elif valor == "⏸️":
                    return "background-color: rgba(245, 158, 11, 0.2); color: #f59e0b;"
                elif valor == "➖":
                    return "background-color: rgba(128, 128, 128, 0.2); color: #888;"
                return ""
                
            try:
                matriz_estilizada = matriz.style.map(cor_fundo)
            except AttributeError:
                matriz_estilizada = matriz.style.applymap(cor_fundo)
                
            st.dataframe(matriz_estilizada, use_container_width=True)

            # === MARCAR / DESMARCAR ENTREGA MANUALMENTE ===
            if st.toggle("✏️ Marcar ou Desmarcar Entrega de um Dia", key="toggle_marcar_dia_f1"):
                st.markdown("""
                <div style="background: rgba(16, 185, 129, 0.1); border: 1px solid rgba(16, 185, 129, 0.3); border-radius: 12px; padding: 15px; margin-bottom: 15px;">
                    <p style="margin: 0; color: #94a3b8; font-size: 14px;">Selecione os encarregados e o dia para colocar <b style="color: #10b981;">✅</b> ou tirar (voltar para <b style="color: #ef4444;">❌</b>).</p>
                </div>
                """, unsafe_allow_html=True)
                
                col_mk1, col_mk2, col_mk3 = st.columns([3, 1, 1])
                with col_mk1:
                    encs_marcar = st.multiselect("Encarregado(s):", todos_encarregados_matriz, key="encs_marcar_dia")
                with col_mk2:
                    dia_marcar = st.selectbox("Dia:", [int(d) for d in dias_uteis], key="dia_marcar_sel")
                with col_mk3:
                    acao_marcar = st.selectbox("Ação:", ["✅ Marcar Entregue", "❌ Desmarcar"], key="acao_marcar_sel")
                
                if st.button("Aplicar", type="primary", use_container_width=True, key="btn_aplicar_marcar"):
                    if encs_marcar:
                        data_str = f"{ano}-{str(mes).zfill(2)}-{str(dia_marcar).zfill(2)}"
                        if "✅" in acao_marcar:
                            novos = []
                            for enc_mk in encs_marcar:
                                ja_existe = ((st.session_state.df_historico_f1["DATA"] == data_str) & (st.session_state.df_historico_f1["ENCARREGADO"] == enc_mk)).any()
                                if not ja_existe:
                                    novos.append({"DATA": data_str, "ENCARREGADO": enc_mk})
                            if novos:
                                df_novos_mk = pd.DataFrame(novos)
                                st.session_state.df_historico_f1 = pd.concat([st.session_state.df_historico_f1, df_novos_mk], ignore_index=True)
                                st.session_state.df_historico_f1.to_csv(caminho_historico_f1_csv, index=False)
                                if conn and not st.session_state.get('force_use_local', False):
                                    try:
                                        salvar_f1_seguro(conn, st.session_state.df_historico_f1, caminho_historico_f1_csv)
                                        st.cache_data.clear()
                                    except Exception:
                                        pass
                                st.success(f"✅ {len(novos)} entrega(s) marcada(s) no dia {dia_marcar}!")
                            else:
                                st.info("ℹ️ Todos já estavam marcados nesse dia.")
                        else:
                            removidos = 0
                            for enc_mk in encs_marcar:
                                mask = (st.session_state.df_historico_f1["DATA"] == data_str) & (st.session_state.df_historico_f1["ENCARREGADO"] == enc_mk)
                                if mask.any():
                                    st.session_state.df_historico_f1 = st.session_state.df_historico_f1[~mask]
                                    st.session_state.df_historico_f1.to_csv(caminho_historico_f1_csv, index=False)
                                    removidos += 1
                            if removidos > 0:
                                if conn and not st.session_state.get('force_use_local', False):
                                    try:
                                        salvar_f1_seguro(conn, st.session_state.df_historico_f1, caminho_historico_f1_csv)
                                        st.cache_data.clear()
                                    except Exception:
                                        pass
                                st.success(f"❌ {removidos} entrega(s) desmarcada(s) no dia {dia_marcar}!")
                            else:
                                st.info("ℹ️ Nenhum deles estava marcado nesse dia.")
                        time.sleep(2)
                        st.rerun()

            # --- EXPORTAÇÃO E NUVEM ---
            col_exp1, col_exp2 = st.columns(2)
            with col_exp1:
                buffer_rh = io.BytesIO()
                matriz_export = matriz.reset_index().rename(columns={"index": "ENCARREGADO"})
                matriz_export.to_excel(buffer_rh, index=False, engine='openpyxl')
                buffer_rh.seek(0)
                st.download_button(
                    label="📥 Baixar Planilha do Mês para o RH (.xlsx)",
                    data=buffer_rh,
                    file_name=f"Relatorio_RH_F1_{mes_selecionado}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    use_container_width=True
                )
            with col_exp2:
                if st.button("☁️ Puxar Histórico F1 da Nuvem", key="btn_puxar_f1_nuvem", use_container_width=True):
                    try:
                        if conn:
                            df_nuvem = conn.read(worksheet="Historico_F1", ttl=0)
                            df_nuvem = df_nuvem.dropna(how='all')
                            if not df_nuvem.empty:
                                st.session_state.df_historico_f1 = df_nuvem
                                df_nuvem.to_csv(caminho_historico_f1_csv, index=False)
                                st.success(f"✅ Histórico F1 atualizado da nuvem! ({len(df_nuvem)} registros)")
                                st.rerun()
                            else:
                                st.warning("⚠️ A aba Historico_F1 na nuvem está vazia.")
                        else:
                            st.error("❌ Sem conexão com o Google Sheets.")
                    except Exception as e:
                        st.error(f"❌ Erro ao puxar da nuvem: {e}")

            # ==============================================================
            # PÓDIO & RANKING GAMIFICADO
            # ==============================================================
            st.markdown("---")
            st.markdown(f"#### 🏆 Pódio & Ranking de Campeões F1 — {mes_selecionado}")

            top3 = df_rank.head(3)
            if len(top3) >= 3:
                n1 = top3.iloc[0]["ENCARREGADO"].split()[0] + " " + (top3.iloc[0]["ENCARREGADO"].split()[-1] if len(top3.iloc[0]["ENCARREGADO"].split())>1 else "")
                t1 = f"{top3.iloc[0]['ENTREGAS']} RDCs ({top3.iloc[0]['APROVEITAMENTO_STR']})"
                n2 = top3.iloc[1]["ENCARREGADO"].split()[0] + " " + (top3.iloc[1]["ENCARREGADO"].split()[-1] if len(top3.iloc[1]["ENCARREGADO"].split())>1 else "")
                t2 = f"{top3.iloc[1]['ENTREGAS']} RDCs ({top3.iloc[1]['APROVEITAMENTO_STR']})"
                n3 = top3.iloc[2]["ENCARREGADO"].split()[0] + " " + (top3.iloc[2]["ENCARREGADO"].split()[-1] if len(top3.iloc[2]["ENCARREGADO"].split())>1 else "")
                t3 = f"{top3.iloc[2]['ENTREGAS']} RDCs ({top3.iloc[2]['APROVEITAMENTO_STR']})"
                
                html_podio = f"""
                <div style="display: flex; justify-content: center; align-items: flex-end; height: 210px; gap: 15px; margin-top: 30px; margin-bottom: 25px;">
                    <!-- 2 Lugar -->
                    <div style="display: flex; flex-direction: column; align-items: center; width: 140px; transition: transform 0.3s;" onmouseover="this.style.transform='translateY(-5px)'" onmouseout="this.style.transform='translateY(0)'">
                        <div style="font-size: 13px; color: #cbd5e1; font-weight: bold; text-align: center; margin-bottom: 5px; white-space: nowrap; overflow: hidden; text-overflow: ellipsis; max-width: 100%;">{n2}</div>
                        <div style="font-size: 28px; margin-bottom: -5px;">🥈</div>
                        <div style="background: linear-gradient(180deg, rgba(148,163,184,0.8), rgba(71,85,105,0.8)); backdrop-filter: blur(5px); width: 100%; height: 95px; border-radius: 12px 12px 0 0; display: flex; flex-direction: column; justify-content: center; align-items: center; color: white; font-weight: 800; font-size: 14px; box-shadow: 0 -5px 20px rgba(148,163,184,0.3); border: 1px solid rgba(255,255,255,0.3); border-bottom: none; padding: 5px;">
                            <span style="font-size: 20px;">{top3.iloc[1]['ENTREGAS']}</span>
                            <span style="font-size: 11px; color: #cbd5e1;">{top3.iloc[1]['APROVEITAMENTO_STR']}</span>
                        </div>
                    </div>
                    <!-- 1 Lugar -->
                    <div style="display: flex; flex-direction: column; align-items: center; width: 150px; transform: translateY(-15px); transition: transform 0.3s;" onmouseover="this.style.transform='translateY(-20px)'" onmouseout="this.style.transform='translateY(-15px)'">
                        <div style="font-size: 15px; color: #fbbf24; font-weight: bold; text-align: center; margin-bottom: 5px; text-shadow: 0 0 10px rgba(251, 191, 36, 0.6); white-space: nowrap; overflow: hidden; text-overflow: ellipsis; max-width: 100%;">{n1}</div>
                        <div style="font-size: 38px; margin-bottom: -5px;">👑</div>
                        <div style="background: linear-gradient(180deg, rgba(251,191,36,0.9), rgba(180,83,9,0.9)); backdrop-filter: blur(5px); width: 100%; height: 135px; border-radius: 12px 12px 0 0; display: flex; flex-direction: column; justify-content: center; align-items: center; color: white; font-weight: 800; font-size: 16px; box-shadow: 0 -5px 25px rgba(251,191,36,0.5); border: 1px solid rgba(255,255,255,0.5); border-bottom: none; padding: 5px;">
                            <span style="font-size: 24px;">{top3.iloc[0]['ENTREGAS']}</span>
                            <span style="font-size: 12px; color: #fef08a;">{top3.iloc[0]['APROVEITAMENTO_STR']}</span>
                        </div>
                    </div>
                    <!-- 3 Lugar -->
                    <div style="display: flex; flex-direction: column; align-items: center; width: 140px; transition: transform 0.3s;" onmouseover="this.style.transform='translateY(-5px)'" onmouseout="this.style.transform='translateY(0)'">
                        <div style="font-size: 13px; color: #d97706; font-weight: bold; text-align: center; margin-bottom: 5px; white-space: nowrap; overflow: hidden; text-overflow: ellipsis; max-width: 100%;">{n3}</div>
                        <div style="font-size: 28px; margin-bottom: -5px;">🥉</div>
                        <div style="background: linear-gradient(180deg, rgba(217,119,6,0.8), rgba(120,53,15,0.8)); backdrop-filter: blur(5px); width: 100%; height: 75px; border-radius: 12px 12px 0 0; display: flex; flex-direction: column; justify-content: center; align-items: center; color: white; font-weight: 800; font-size: 14px; box-shadow: 0 -5px 20px rgba(217,119,6,0.3); border: 1px solid rgba(255,255,255,0.2); border-bottom: none; padding: 5px;">
                            <span style="font-size: 18px;">{top3.iloc[2]['ENTREGAS']}</span>
                            <span style="font-size: 11px; color: #fed7aa;">{top3.iloc[2]['APROVEITAMENTO_STR']}</span>
                        </div>
                    </div>
                </div>
                """
                st.markdown(html_podio, unsafe_allow_html=True)

            # === BOTÃO: DIVULGAR RANKING NO WHATSAPP ===
            import urllib.parse
            texto_f1_zap = [
                f"🏎️ *RANKING F1 DE ENTREGA DE RDC — {nome_site}*",
                f"📅 *Mês de Referência:* {mes_selecionado}",
                f"🎯 *Visão:* {filtro_f1_disc}",
                "",
                "🏆 *TOP 5 LÍDERES EM PONTUALIDADE & ASSIDUIDADE:*",
            ]
            
            for idx_r, row_r in df_rank.head(5).iterrows():
                pos_icon = "🥇" if idx_r == 0 else ("🥈" if idx_r == 1 else ("🥉" if idx_r == 2 else f"{idx_r+1}º"))
                conq_txt = f" {row_r['CONQUISTAS']}" if row_r['CONQUISTAS'] != "—" else ""
                texto_f1_zap.append(f"{pos_icon} *{row_r['ENCARREGADO']}* ({row_r['DISCIPLINA']}): {row_r['ENTREGAS']} RDCs · *{row_r['APROVEITAMENTO_STR']}*{conq_txt}")
                
            texto_f1_zap.extend([
                "",
                f"📊 *Taxa Média de Entrega da Obra:* {media_aprov}%",
                f"📄 *Total Geral de RDCs Entregues:* {total_entregue}",
                "",
                "💪 _Parabéns a toda a liderança pelo compromisso diário com a medição e os registros da obra!_",
                f"_Atualizado em {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}_"
            ])
            
            zap_f1_msg = "\n".join(texto_f1_zap)
            zap_f1_url = f"https://api.whatsapp.com/send?text={urllib.parse.quote(zap_f1_msg)}"

            col_w1, col_w2 = st.columns([1, 2])
            with col_w1:
                st.link_button("📲 Divulgar Ranking F1 no WhatsApp", zap_f1_url, type="primary", use_container_width=True)
            with col_w2:
                with st.expander("📋 Ver Texto do Ranking Formatado para Copiar"):
                    st.text_area("Texto do WhatsApp:", value=zap_f1_msg, height=130, key="txt_f1_zap_copy")

            # Tabela Completa do Ranking
            st.markdown("<br>", unsafe_allow_html=True)
            st.markdown("**📋 Classificação Geral dos Encarregados:**")
            df_rank_display = df_rank[["ENCARREGADO", "DISCIPLINA", "ENTREGAS", "APROVEITAMENTO_STR", "CONQUISTAS"]].copy()
            df_rank_display.columns = ["Encarregado", "Disciplina", "RDCs Entregues", "Aproveitamento", "Badges / Conquistas"]
            st.dataframe(df_rank_display, hide_index=False, use_container_width=True)

            st.markdown("---")
            st.markdown("#### 📈 Evolução Mensal de Entregas")
            if not df_hist.empty and "MES_ANO" in df_hist.columns:
                df_evolucao = df_hist.groupby("MES_ANO").size().reset_index(name="RDCs Entregues")
            else:
                df_evolucao = pd.DataFrame()
            if not df_evolucao.empty:
                fig_ev = px.line(df_evolucao, x="MES_ANO", y="RDCs Entregues", text="RDCs Entregues", markers=True)
                fig_ev.update_traces(textposition="top center", line_color="#0ea5e9", marker=dict(size=8))
                fig_ev.update_layout(paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", font=dict(color="#e0e4ea"), xaxis_title="Mês", yaxis_title="Total de RDCs")
                st.plotly_chart(fig_ev, use_container_width=True)
                
            st.markdown("---")
            if st.button("📄 Gerar Relatório Mensal em PDF", type="secondary", use_container_width=True):
                try:
                    from fpdf import FPDF
                    pdf = FPDF()
                    pdf.add_page()
                    pdf.set_font("Arial", 'B', 16)
                    pdf.cell(200, 10, txt=f"Relatorio Mensal F1 - {mes_selecionado}", ln=True, align='C')
                    pdf.ln(10)
                    
                    pdf.set_font("Arial", 'B', 12)
                    pdf.cell(200, 10, txt=f"Total de RDCs Entregues no Mes: {total_entregue}", ln=True)
                    pdf.ln(10)
                    
                    pdf.set_font("Arial", 'B', 14)
                    pdf.cell(200, 10, txt="Os 3 Melhores do Mes:", ln=True)
                    pdf.set_font("Arial", '', 12)
                    for i, row in top3.iterrows():
                        pdf.cell(200, 10, txt=f"{i+1} Lugar: {row['ENCARREGADO']} - {row['ENTREGAS']} RDCs ({row['APROVEITAMENTO_STR']})", ln=True)
                    pdf.ln(5)
                    
                    pdf_output = bytes(pdf.output())
                    st.download_button("📥 Clique aqui para baixar o PDF", data=pdf_output, file_name=f"Relatorio_{mes_selecionado}.pdf", mime="application/pdf", type="primary")
                except ImportError:
                    st.error("Biblioteca FPDF não encontrada.")
            
            st.markdown("<br><br>", unsafe_allow_html=True)

    with tab_ia:
        st.markdown("### 🤖 Robô de Extração Inteligente (Google Gemini)")
        st.markdown("<p style='margin-top: -15px; font-size: 14px; color: #888;'>Uma ideia original por <b>Caio Farisco</b></p>", unsafe_allow_html=True)
        st.markdown("Arraste os formulários RDC físicos escaneados abaixo. A inteligência artificial irá extrair as informações e padronizar com a sua base de Encarregados.")
        
        try:
            from google import genai
            HAS_GENAI = True
        except ImportError:
            HAS_GENAI = False
            st.error("A biblioteca `google-genai` não está instalada no servidor. Instale usando `pip install google-genai`.")

        if HAS_GENAI:
            # Tentar ler a chave do cofre secreto (.streamlit/secrets.toml)
            chave_padrao = ""
            try:
                chave_padrao = st.secrets.get("GEMINI_API_KEY", "")
            except Exception:
                pass
            
            st.markdown("#### Configuração e Upload")
            if not chave_padrao:
                chave_padrao = st.text_input("🔑 Cole suas Chaves da API Gemini (separadas por vírgula):", type="password", help="Se usar múltiplas chaves, o robô alterna em caso de limite.")
            
            if not chave_padrao:
                st.info("☝️ Cole a(s) sua(s) chave(s) de API do Gemini acima para ativar o robô de leitura.")
            
            arquivos_scan = st.file_uploader("Upload de RDCs Escaneados (PDF, JPG, PNG)", type=["png", "jpg", "jpeg", "pdf"], accept_multiple_files=True)
                
            btn_processar = st.button("🚀 Processar Arquivos com IA", type="primary", use_container_width=True)
            
            if btn_processar and arquivos_scan and chave_padrao:
                # --- FIX: Evitar que o Gemini tente usar o Service Account do Google Sheets ---
                old_cred = os.environ.pop("GOOGLE_APPLICATION_CREDENTIALS", None)
                
                lista_chaves = [c.strip() for c in chave_padrao.split(",") if c.strip()]
                idx_chave_atual = 0
                client = genai.Client(api_key=lista_chaves[idx_chave_atual])
                nomes_para_prompt = ", ".join(lista_encarregados_base)
                
                prompt_ia = f"""
                Analise este documento (que pode ter várias páginas). Para CADA formulário de obra (RDC) encontrado no arquivo, extraia os dados.
                REGRA IMPORTANTÍSSIMA: Retorne APENAS UM objeto JSON por formulário/página. NÃO separe os formulários.
                DICA DE OURO: Todos os RDCs dentro deste arquivo PDF pertencem EXATAMENTE ao mesmo dia. Portanto, a DATA extraída deve ser idêntica para todos os formulários.
                
                === GLOSSÁRIO DE TERMOS DA CONSTRUÇÃO (USE PARA CORRIGIR ERROS) ===
                Os encarregados escrevem à mão e cometem erros. Quando encontrar palavras estranhas, use este glossário para deduzir a palavra correta:
                - TUBULASÃO, TUBULASAO, TUBULAÇAO, TUBLAÇÃO → TUBULAÇÃO
                - SOLDAJEM, SOLDAGEN, SOUDAGEM → SOLDAGEM
                - CALDEIRARIA, CALDEIRRARIA, CALDERARIA → CALDEIRARIA
                - SUPERAQUECEDOR, SUPER AQUECEDOR, SUPERAQUESSEDOR → SUPERAQUECEDOR
                - ECONOMISADOR, ECONOMIÇADOR, ECONOMIZADRO → ECONOMIZADOR
                - PRECIPITADRO, PRESIPITADOR, PRECIPTADOR → PRECIPITADOR
                - ESTRUTURA METALICA, ESTRURA MET, METALICA → ESTRUTURA METÁLICA
                - EQUIPAMENTO, EKIPAMENTO, EQUIPAMNETO → EQUIPAMENTO
                - ANDAINE, ANDAME, HANDAIME → ANDAIME
                - MONTAJEM, MONTAGEN, MONTAGEN → MONTAGEM
                - DESMONTAJEM, DISMONTAGEM → DESMONTAGEM
                - ESMERILHAMENTO, ESMERILHAMNETO → ESMERILHAMENTO
                - TRAÇAJEM, TRASAGEM → TRAÇAGEM
                - HIDROJAETEAMENTO, HIDRO JATO → HIDROJATEAMENTO
                - MAÇARICO, MASARICO, MAÇARIKO → MAÇARICO
                - PRE AQUECIMENTO, PRÉ AQUECIMENTO → PRÉ-AQUECIMENTO
                - ELEVAÇAO, ELEVASSÃO → ELEVAÇÃO
                - REVESTIMNTO, REVESTIMENTO → REVESTIMENTO
                - ISOLAMNTO, IZOLAMENTO → ISOLAMENTO
                - JUNTA DE ESPANÇÃO, JUNTA ESPANSÃO → JUNTA DE EXPANSÃO
                Se encontrar qualquer palavra estranha ou ilegível não listada acima, use o contexto para tentar deduzir o que o encarregado quis escrever. SEMPRE corrija a ortografia nos campos ATIVIDADE e SUB_ATIVIDADE.
                ===================================================================

                Retorne APENAS um array (lista) em formato JSON válido. Exemplo do formato exato esperado:
                [
                  {{
                    "DATA": "YYYY-MM-DD",
                    "DISCIPLINA": "...",
                    "ENCARREGADO": "...",
                    "TURNO": "...",
                    "DDS": "...",
                    "TRANSCRICAO": "TEXTO EXATAMENTE COMO ESTÁ ESCRITO NO RDC, SEM CORRIGIR NADA",
                    "ATIVIDADE": "RESUMO GERAL CURTO CORRIGIDO DE TODAS AS ATIVIDADES",
                    "PROBLEMAS": "...",
                    "LOCAL": "...",
                    "AREA": "...",
                    "SUBNIVEIS": [
                      {{"ATIVIDADE": "SOLDAGEM DE TUBULAÇÃO DN 8", "LOCAL_ESPECIFICO": "ELEVAÇÃO 35M - MÓDULO 3", "EFETIVO": "3 SOLDADORES, 1 AJUDANTE"}},
                      {{"ATIVIDADE": "MONTAGEM DE ANDAIME", "LOCAL_ESPECIFICO": "ÁREA DO SUPERAQUECEDOR", "EFETIVO": "2 MONTADORES"}}
                    ]
                  }}
                ]

                Regras de negócio:
                - DATA: Extraia a data em que o RDC foi preenchido. Retorne RIGOROSAMENTE no formato YYYY-MM-DD (Ano-Mês-Dia).
                - DISCIPLINA: Extraia a disciplina ou função do topo, mas RETORNE APENAS A PRIMEIRA PALAVRA OU A PALAVRA PRINCIPAL (ex: MECÂNICA, SOLDA, TOPOGRAFIA, CALDEIRARIA). Se for montador de andaime escreva ANDAIME. Sempre apenas 1 palavra.
                - ENCARREGADO: É QUASE PROIBIDO retornar 'AJUSTAR NOME'. Você DEVE escolher o nome da lista oficial [{nomes_para_prompt}] que for mais parecido fonética ou visualmente com o que está escrito à mão, mesmo que a letra seja péssima, haja apenas o primeiro nome, iniciais (ex: "J. Silva") ou erros grosseiros. Faça o cruzamento lógico e retorne EXATAMENTE o nome completo da lista. Só use 'AJUSTAR NOME' se o campo estiver 100% em branco ou completamente rasurado sem nenhuma letra legível.
                - TURNO: Analise os horários. De dia (ex: 07:00 as 17:00) = 'DIURNO'. De noite = 'NOTURNO'.
                - DDS: Extraia o tema principal de Segurança mencionado no relatório (DDS, Diálogo de Segurança). (ex: Trabalho a quente, Bloqueio, etc). Se não tiver, retorne 'Não Informado'.
                - TRANSCRICAO: Leia TUDO o que está escrito na seção de ATIVIDADES do RDC e transcreva o CONTEÚDO COMPLETO de forma LEGÍVEL e COMPREENSÍVEL. Corrija a ortografia usando o glossário acima, mas NÃO resuma e NÃO elimine detalhes. Inclua TODAS as informações que o encarregado anotou.
                - ATIVIDADE: Crie um RESUMO GERAL de no máximo 35 palavras contendo as principais atividades executadas em todo o RDC. TUDO EM MAIÚSCULAS. CORRIJA a ortografia usando o glossário acima. NÃO CRIE SUBNÍVEIS, APENAS UM ÚNICO RESUMO TEXTUAL.
                - CALDEIRA: Se mencionar 'caldeira de recuperação' = 'RB'. Se 'caldeira de potência' = 'PB'. Se a descrição da atividade mencionar 'PRECIPITADOR' ou 'ESP' = 'ESP'. Se nenhum = ''.
                - LOCAL: Analise a imagem CUIDADOSAMENTE. Procure as opções 'PB ( )' e 'RB ( )'. Verifique se há um 'X', um rabisco, um visto ou qualquer marcação (mesmo que mal desenhada) dentro, em cima ou do lado dos parênteses. Retorne APENAS 'PB' ou 'RB' correspondente ao que estiver marcado. Se nenhum, retorne ''.
                - AREA: Analise as caixinhas de área na imagem com LUPA. Procure por qualquer marcação (X, visto, círculo, rabisco) dentro ou sobre os parênteses. As opções são exatamente: DUTO, EQUIPAMENTO, TUBULAÇÃO, ESTRUTURA MET, PRECIPITADOR, PRESSAO - MEC, PRESSAO - TUBULACAO, PRESSAO - FORNALHA, PINTURA, SOPRAGEM, ANDAIME. Retorne EXATAMENTE o nome da área que estiver marcada. Se nenhuma estiver marcada, retorne ''.

                Não inclua crases, formatação markdown ou texto adicional, apenas o JSON puro começando com [ e terminando com ].
                """

                # === ANIMAÇÃO PREMIUM DE LOADING ===
                animacao_html = """
                <div style="display: flex; flex-direction: column; align-items: center; justify-content: center; padding: 25px; background: rgba(15,23,42,0.9); border: 1px solid #0ea5e9; border-radius: 15px; box-shadow: 0 0 30px rgba(14, 165, 233, 0.3); margin-bottom: 20px;">
                    <div class="radar" style="position: relative; width: 120px; height: 120px; border-radius: 50%; border: 2px solid rgba(14,165,233,0.5); overflow: hidden; background: radial-gradient(circle, rgba(14,165,233,0.15) 0%, rgba(15,23,42,0) 100%);">
                        <div style="position: absolute; width: 50%; height: 50%; top: 0; left: 50%; transform-origin: bottom left; background: linear-gradient(45deg, rgba(14,165,233,0.9) 0%, transparent 50%); animation: radar-spin 1.5s linear infinite;"></div>
                        <div style="position: absolute; top: 50%; left: 50%; transform: translate(-50%, -50%); color: #fff; font-weight: bold; font-size: 14px; letter-spacing: 2px; text-shadow: 0 0 10px #0ea5e9; background: #0f172a; padding: 5px; border-radius: 5px;">ENESA</div>
                        <div style="position: absolute; top: 0; bottom: 0; left: 50%; width: 1px; background: rgba(14,165,233,0.4);"></div>
                        <div style="position: absolute; left: 0; right: 0; top: 50%; height: 1px; background: rgba(14,165,233,0.4);"></div>
                        <div style="position: absolute; top: 20%; left: 20%; width: 6px; height: 6px; background: #4ade80; border-radius: 50%; box-shadow: 0 0 10px #4ade80; animation: blip 1.5s infinite;"></div>
                        <div style="position: absolute; top: 70%; left: 60%; width: 4px; height: 4px; background: #4ade80; border-radius: 50%; box-shadow: 0 0 10px #4ade80; animation: blip 1.5s infinite 0.7s;"></div>
                    </div>
                    <p style="color: #0ea5e9; margin-top: 20px; font-weight: bold; font-size: 16px; animation: pulse 1s infinite; margin-bottom: 0;">🤖 IA Processando Documentos...</p>
                    <style>
                        @keyframes radar-spin { 100% { transform: rotate(360deg); } }
                        @keyframes pulse { 0%, 100% { opacity: 1; } 50% { opacity: 0.6; } }
                        @keyframes blip { 0%, 100% { opacity: 0; } 10% { opacity: 1; } }
                    </style>
                </div>
                """
                animacao_placeholder = st.empty()
                animacao_placeholder.markdown(animacao_html, unsafe_allow_html=True)
                
                with st.status("🤖 Robô iniciando análise...", expanded=True) as status:
                    progresso = st.progress(0)
                    total_arquivos = len(arquivos_scan)
                    
                    # === PRE-PROCESSAMENTO E CHUNKING ===
                    import fitz
                    arquivos_processar = []
                    for arquivo_scan in arquivos_scan:
                        nome = arquivo_scan.name
                        if nome.lower().endswith('.pdf'):
                            pdf_bytes = arquivo_scan.getvalue()
                            
                            # Tentar abrir via stream
                            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
                            num_pages = len(doc)
                            
                            # FALLBACK 1: Se stream falhou, salvar no disco e abrir do disco
                            if num_pages == 0:
                                doc.close()
                                tmp_raw = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
                                tmp_raw.write(pdf_bytes)
                                tmp_raw.close()
                                doc = fitz.open(tmp_raw.name)
                                num_pages = len(doc)
                                os.unlink(tmp_raw.name)
                            
                            # FALLBACK 2: Se ainda 0 páginas, verificar se é uma imagem renomeada para PDF (ex: TIFF)
                            if num_pages == 0:
                                doc.close()
                                try:
                                    from PIL import Image
                                    import io
                                    
                                    # Tenta abrir com PIL
                                    img = Image.open(io.BytesIO(pdf_bytes))
                                    # Se conseguiu, é uma imagem! Vamos converter pra JPG e mandar pra IA
                                    tmp_img = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
                                    # Converte pra RGB caso seja TIFF CMYK ou RGBA
                                    if img.mode != 'RGB':
                                        img = img.convert('RGB')
                                    img.save(tmp_img.name, "JPEG")
                                    tmp_img.close()
                                    arquivos_processar.append({'name': f"{nome} (Recuperado)", 'tmp_path': tmp_img.name})
                                    st.info(f"⚙️ {nome}: Arquivo corrigido (era uma imagem salva como PDF).")
                                except Exception:
                                    # Se nem o PIL abrir, o arquivo está realmente quebrado
                                    st.error(f"❌ O arquivo {nome} está corrompido ou o scanner falhou ao gerar o PDF (0 páginas válidas).")
                                    continue
                            elif num_pages > 15:
                                chunk_size = 15
                                for start_idx in range(0, num_pages, chunk_size):
                                    chunk_doc = fitz.open()
                                    chunk_doc.insert_pdf(doc, from_page=start_idx, to_page=min(start_idx + chunk_size - 1, num_pages - 1))
                                    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
                                    tmp.close()
                                    chunk_doc.save(tmp.name, deflate=True)
                                    chunk_doc.close()
                                    arquivos_processar.append({'name': f"{nome} (Pág {start_idx+1}-{min(start_idx+chunk_size, num_pages)})", 'tmp_path': tmp.name})
                                doc.close()
                            else:
                                # Rebuild PDF limpo
                                chunk_doc = fitz.open()
                                chunk_doc.insert_pdf(doc)
                                tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
                                tmp.close()
                                chunk_doc.save(tmp.name, deflate=True)
                                chunk_doc.close()
                                doc.close()
                                
                                # Verificar se o rebuild funcionou
                                check = fitz.open(tmp.name)
                                if len(check) == 0:
                                    check.close()
                                    os.unlink(tmp.name)
                                    # Último recurso: enviar bytes originais
                                    tmp2 = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
                                    tmp2.write(pdf_bytes)
                                    tmp2.close()
                                    arquivos_processar.append({'name': nome, 'tmp_path': tmp2.name})
                                else:
                                    check.close()
                                    arquivos_processar.append({'name': nome, 'tmp_path': tmp.name})
                        else:
                            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=f".{nome.split('.')[-1]}")
                            tmp.write(arquivo_scan.getvalue())
                            tmp.close()
                            arquivos_processar.append({'name': nome, 'tmp_path': tmp.name})
                            
                    total_arquivos = len(arquivos_processar)
                    
                    def processar_chunk_ia(arquivo_dict, chaves_api, modelo_gemini, prompt_ia, schema, _old_cred):
                        import time
                        import json
                        import os
                        import ast
                        import re
                        from google import genai
                        
                        tmp_path = arquivo_dict['tmp_path']
                        max_tentativas = 3
                        idx_chave_atual_local = 0
                        client_local = genai.Client(api_key=chaves_api[idx_chave_atual_local])
                        
                        for tentativa in range(max_tentativas):
                            try:
                                arquivo_up = client_local.files.upload(file=tmp_path)
                                
                                tempo_espera = 0
                                while tempo_espera < 180:
                                    file_info = client_local.files.get(name=arquivo_up.name)
                                    estado = str(file_info.state).upper()
                                    if "ACTIVE" in estado:
                                        break
                                    elif "FAILED" in estado:
                                        raise Exception("Falha interna do Google ao processar este arquivo. Tente um arquivo menor.")
                                    time.sleep(3)
                                    tempo_espera += 3
                                    
                                if tempo_espera >= 180:
                                    raise Exception("Tempo limite esgotado aguardando o Google processar o PDF (demorou mais de 3 minutos).")
                                
                                resposta = client_local.models.generate_content(
                                    model=modelo_gemini,
                                    contents=[arquivo_up, prompt_ia],
                                    config=genai.types.GenerateContentConfig(
                                        response_mime_type="application/json",
                                        response_schema=list[schema],
                                        temperature=0.0
                                    )
                                )
                                
                                if _old_cred:
                                    os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = _old_cred

                                texto_resposta = resposta.text.strip()
                                if "```json" in texto_resposta:
                                    texto_resposta = texto_resposta.split("```json")[1].split("```")[0].strip()
                                elif "```" in texto_resposta:
                                    texto_resposta = texto_resposta.split("```")[1].split("```")[0].strip()
                                    
                                start_idx = max(0, texto_resposta.find('[')) if '[' in texto_resposta else max(0, texto_resposta.find('{'))
                                end_idx = max(texto_resposta.rfind(']'), texto_resposta.rfind('}'))
                                if end_idx > start_idx:
                                    texto_resposta = texto_resposta[start_idx:end_idx+1]

                                try:
                                    dados_extraidos_lista = json.loads(texto_resposta)
                                except json.JSONDecodeError as err_json:
                                    recuperado = False
                                    last_brace = texto_resposta.rfind('}')
                                    if last_brace != -1:
                                        texto_recuperado = texto_resposta[:last_brace+1] + ']'
                                        try:
                                            dados_extraidos_lista = json.loads(texto_recuperado)
                                            recuperado = True
                                        except:
                                            pass
                                            
                                    if not recuperado:
                                        texto_fix = texto_resposta.replace("null", "None").replace("true", "True").replace("false", "False")
                                        texto_fix = re.sub(r'\}\s*\{', '}, {', texto_fix)
                                        texto_fix = re.sub(r'\]\s*\[', '], [', texto_fix)
                                        texto_fix = re.sub(r'("|\]|\})\s+(")', r'\1, \2', texto_fix)
                                        try:
                                            dados_extraidos_lista = ast.literal_eval(texto_fix)
                                        except SyntaxError:
                                            try:
                                                dados_extraidos_lista = ast.literal_eval(texto_fix + '"}]')
                                            except:
                                                if last_brace != -1:
                                                    try:
                                                        dados_extraidos_lista = ast.literal_eval(texto_fix[:last_brace+1] + ']')
                                                    except:
                                                        raise err_json
                                                else:
                                                    raise err_json
                                        except:
                                            raise err_json

                                if isinstance(dados_extraidos_lista, dict):
                                    dados_extraidos_lista = [dados_extraidos_lista]

                                datas_encontradas = [str(d.get("DATA")).strip() for d in dados_extraidos_lista if d.get("DATA") and str(d.get("DATA")).strip() != ""]
                                if datas_encontradas:
                                    data_consenso = max(set(datas_encontradas), key=datas_encontradas.count)
                                    for d in dados_extraidos_lista:
                                        d["DATA"] = data_consenso

                                for dados in dados_extraidos_lista:
                                    if 'LOCAL' not in dados:
                                        dados['LOCAL'] = ''
                                    if 'AREA' not in dados:
                                        dados['AREA'] = ''
                                    dados['SUB'] = 1
                                    dados['SUB_ATIVIDADE'] = dados.get('ATIVIDADE', '')
                                    dados['LOCAL_ESPECIFICO'] = ''
                                    dados['EFETIVO_ATIVIDADE'] = ''
                                
                                return dados_extraidos_lista
                                
                            except Exception as inner_e:
                                erro_str = str(inner_e)
                                if '429' in erro_str or 'RESOURCE_EXHAUSTED' in erro_str:
                                    if tentativa < max_tentativas - 1:
                                        if idx_chave_atual_local < len(chaves_api) - 1:
                                            idx_chave_atual_local += 1
                                            client_local = genai.Client(api_key=chaves_api[idx_chave_atual_local])
                                            time.sleep(2)
                                            continue
                                        else:
                                            time.sleep(60)
                                            continue
                                elif '503' in erro_str or 'UNAVAILABLE' in erro_str:
                                    if tentativa < max_tentativas - 1:
                                        time.sleep(10)
                                        continue
                                raise Exception(f"Erro detalhado na IA: {inner_e}")
                        raise Exception("Falha após múltiplas tentativas.")

                    import concurrent.futures
                    modelo_usado = st.session_state.get('modelo_gemini', 'gemini-2.5-flash')
                    
                    with concurrent.futures.ThreadPoolExecutor(max_workers=5) as executor:
                        future_to_chunk = {
                            executor.submit(processar_chunk_ia, chunk, lista_chaves, modelo_usado, prompt_ia, RDC_Schema, old_cred): chunk
                            for chunk in arquivos_processar
                        }
                        
                        completed = 0
                        for future in concurrent.futures.as_completed(future_to_chunk):
                            chunk = future_to_chunk[future]
                            nome_atual = chunk['name']
                            tmp_path = chunk['tmp_path']
                            
                            try:
                                dados_extraidos_lista = future.result()
                                
                                for dados in dados_extraidos_lista:
                                    ultimo_item = st.session_state.df_ia['ITEM'].max() if not st.session_state.df_ia.empty and pd.notna(st.session_state.df_ia['ITEM'].max()) else 0
                                    dados['ITEM'] = int(ultimo_item) + 1
                                    st.session_state.df_ia = pd.concat([st.session_state.df_ia, pd.DataFrame([dados])], ignore_index=True)

                                
                                st.toast(f"✅ {nome_atual} processado com sucesso!")
                            except Exception as e:
                                st.error(f"Erro no envio de {nome_atual}: {e}")
                                st.session_state.teve_falha_ia = True
                                
                            try:
                                os.remove(tmp_path)
                            except:
                                pass
                                
                            completed += 1
                            progresso.progress(completed / total_arquivos)
                            status.update(label=f"Processando blocos simultaneamente... ({completed}/{total_arquivos} concluídos)", state="running")

                    expandir_status = st.session_state.get('teve_falha_ia', False)
                    status.update(label="🎉 Leitura concluída!" if not expandir_status else "⚠️ Leitura finalizada com erros", state="complete", expanded=expandir_status)
                    st.session_state.teve_falha_ia = False
                    
                    pass
                animacao_placeholder.empty()
                st.session_state.force_use_local = True
                
            if not st.session_state.df_ia.empty:
                st.markdown("#### Dados Extraídos")
                
                lista_com_alerta = lista_encarregados_base + ["AJUSTAR NOME"]
                df_filtrado = st.session_state.df_ia[st.session_state.df_ia['ENCARREGADO'].isin(lista_com_alerta)]
                
                # --- NOVO FILTRO DE DATA ---
                datas_disponiveis = df_filtrado['DATA'].dropna().unique().tolist()
                if datas_disponiveis:
                    datas_sel = st.multiselect("📅 Filtrar Tabela por Data (Deixe em branco para ver todos):", sorted(datas_disponiveis, reverse=True), default=None)
                    if datas_sel:
                        df_filtrado = df_filtrado[df_filtrado['DATA'].isin(datas_sel)]
                # ---------------------------
                
                col_dw1, col_dw2 = st.columns([1, 1])
                with col_dw1:
                    buffer_df = io.BytesIO()
                    
                    # Preparar Excel
                    df_excel_ia = df_filtrado.copy()
                    df_excel_ia.to_excel(buffer_df, index=False, engine='openpyxl')
                    buffer_df.seek(0)
                    st.download_button(
                        label="⬇️ Baixar Planilha RDC Lida (.xlsx)",
                        data=buffer_df,
                        file_name=f"RDCs_Extraidos_{datetime.datetime.now().strftime('%d%m%Y_%H%M')}.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        use_container_width=True,
                        type="primary"
                    )
                with col_dw2:
                    if st.button("🗑️ Limpar Dados Lidos", use_container_width=True):
                        st.session_state.df_ia = pd.DataFrame(columns=['ITEM', 'SUB', 'DATA', 'DISCIPLINA', 'ENCARREGADO', 'TURNO', 'DDS', 'TRANSCRICAO', 'ATIVIDADE', 'SUB_ATIVIDADE', 'LOCAL_ESPECIFICO', 'EFETIVO_ATIVIDADE', 'PROBLEMAS', 'LOCAL', 'AREA', 'CALDEIRA'])
                        st.rerun()
                
                st.info("✏️ **Dica:** Você pode editar os dados na tabela abaixo antes de confirmar. Dê dois cliques em qualquer célula para corrigir nomes errados, datas ou locais.")
                
                # === VISUALIZAÇÃO DOS RDCs LIDOS ===
                with st.expander("🔍 Visualizar RDCs (Texto Original vs Resumo)", expanded=False):
                    itens_unicos = df_filtrado['ITEM'].unique()
                    for item_id in sorted(itens_unicos):
                        bloco = df_filtrado[df_filtrado['ITEM'] == item_id]
                        if not bloco.empty:
                            enc = bloco.iloc[0].get('ENCARREGADO', '?')
                            disc = bloco.iloc[0].get('DISCIPLINA', '?')
                            data = bloco.iloc[0].get('DATA', '?')
                            resumo = bloco.iloc[0].get('ATIVIDADE', '')
                            transcricao = bloco.iloc[0].get('TRANSCRICAO', '')
                            
                            st.markdown(f"""
                            <div style="background: rgba(14,165,233,0.08); border: 1px solid rgba(14,165,233,0.25); border-radius: 12px; padding: 15px; margin-bottom: 12px;">
                                <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 8px;">
                                    <span style="font-weight: 700; color: #0ea5e9; font-size: 15px;">📋 RDC #{int(item_id)} — {enc}</span>
                                    <span style="color: #64748b; font-size: 12px;">{disc} · {data}</span>
                                </div>
                                <p style="color: #94a3b8; font-size: 13px; margin: 0 0 5px 0;">📝 Resumo Final: {resumo}</p>
                            """, unsafe_allow_html=True)
                            
                            if transcricao and str(transcricao).strip():
                                st.markdown(f"""
                                <div style="background: rgba(245,158,11,0.08); border: 1px solid rgba(245,158,11,0.2); border-radius: 8px; padding: 10px 14px; margin-top: 10px;">
                                    <span style="color: #f59e0b; font-weight: 600; font-size: 11px; text-transform: uppercase; letter-spacing: 0.5px;">✏️ Texto Original do RDC:</span>
                                    <p style="color: #cbd5e1; font-size: 12px; margin: 6px 0 0 0; font-style: italic; line-height: 1.5;">{transcricao}</p>
                                </div>
                                """, unsafe_allow_html=True)
                                
                            st.markdown("</div>", unsafe_allow_html=True)
                
                # Preparar dados para edição
                df_editavel = df_filtrado.copy()
                colunas_mostrar = ['ITEM', 'DATA', 'DISCIPLINA', 'ENCARREGADO', 'TURNO', 'DDS', 'CALDEIRA', 'ATIVIDADE']
                df_editavel = df_editavel[[c for c in colunas_mostrar if c in df_editavel.columns]]
                
                df_editado = st.data_editor(df_editavel, hide_index=True, use_container_width=True, key="editor_ia_df")
                
                if st.button("✅ Confirmar e Salvar no Sistema", type="primary", use_container_width=True):
                    # Salvar atualizações no df_ia
                    itens_processados = df_editavel['ITEM'].unique()
                    df_restante = st.session_state.df_ia[~st.session_state.df_ia['ITEM'].isin(itens_processados)].copy()
                    
                    # Atualiza os dados do df_filtrado com as edições feitas na tabela
                    for _, row_edit in df_editado.iterrows():
                        item_id = row_edit['ITEM']
                        idx = df_filtrado[df_filtrado['ITEM'] == item_id].index
                        for col in ['DATA', 'DISCIPLINA', 'ENCARREGADO', 'TURNO', 'DDS', 'CALDEIRA', 'ATIVIDADE']:
                            if col in row_edit:
                                df_filtrado.loc[idx, col] = row_edit[col]
                                
                    st.session_state.df_ia = pd.concat([df_restante, df_filtrado], ignore_index=True)
                    
                    # === PERSISTIR RDCs NO BANCO DE DADOS ===
                    try:
                        colunas_rdc = ['DATA', 'DISCIPLINA', 'ENCARREGADO', 'TURNO', 'DDS', 'TRANSCRICAO', 'ATIVIDADE', 'PROBLEMAS', 'LOCAL', 'AREA', 'CALDEIRA']
                        df_salvar_rdc = df_editado.copy()
                        for col in colunas_rdc:
                            if col not in df_salvar_rdc.columns:
                                df_salvar_rdc[col] = ''
                        df_salvar_rdc = df_salvar_rdc[colunas_rdc]
                        df_salvar_rdc['CRIADO_EM'] = datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')
                        
                        if os.path.exists(caminho_rdc_registros_csv):
                            df_rdc_existente = pd.read_csv(caminho_rdc_registros_csv)
                            df_rdc_final = pd.concat([df_rdc_existente, df_salvar_rdc], ignore_index=True)
                        else:
                            df_rdc_final = df_salvar_rdc
                        
                        df_rdc_final.to_csv(caminho_rdc_registros_csv, index=False)
                        st.toast(f"💾 {len(df_salvar_rdc)} RDCs salvos permanentemente no Banco de Dados!", icon="✅")
                    except Exception as e_rdc_save:
                        st.warning(f"⚠️ RDCs processados, mas erro ao salvar no banco permanente: {e_rdc_save}")
                    
                    # Salva no F1
                    novos_registros = []
                    for _, row in df_editado.iterrows():
                        enc_lido = str(row.get('ENCARREGADO', '')).strip()
                        if enc_lido and enc_lido in lista_encarregados_base:
                            data_extraida = str(row.get('DATA', '')).strip()
                            try:
                                data_registro = pd.to_datetime(data_extraida).strftime('%Y-%m-%d')
                            except:
                                data_registro = datetime.date.today().strftime('%Y-%m-%d')
                            
                            ja_existe = ((st.session_state.df_historico_f1["DATA"] == data_registro) & (st.session_state.df_historico_f1["ENCARREGADO"] == enc_lido)).any()
                            if not ja_existe:
                                novos_registros.append({"DATA": data_registro, "ENCARREGADO": enc_lido})
                                
                    if novos_registros:
                        df_novos = pd.DataFrame(novos_registros)
                        if conn and not st.session_state.get('force_use_local', False):
                            try:
                                df_fresco = conn.read(worksheet="Historico_F1", ttl=0)
                                if not df_fresco.empty:
                                    df_fresco = df_fresco.dropna(how='all')
                                    df_final = pd.concat([df_fresco, df_novos], ignore_index=True).drop_duplicates(subset=["DATA", "ENCARREGADO"])
                                else:
                                    df_final = pd.concat([st.session_state.df_historico_f1, df_novos], ignore_index=True).drop_duplicates(subset=["DATA", "ENCARREGADO"])
                                
                                ok, msg = salvar_f1_seguro(conn, df_final, caminho_historico_f1_csv)
                                if ok:
                                    st.session_state.df_historico_f1 = df_final
                                st.cache_data.clear()
                                st.toast(f"{len(novos_registros)} RDCs registrados no Resumo Diário e sincronizados com a nuvem!", icon="✅")
                            except Exception as e:
                                st.error(f"Erro ao salvar na nuvem: {e}")
                                st.session_state.df_historico_f1 = pd.concat([st.session_state.df_historico_f1, df_novos], ignore_index=True).drop_duplicates(subset=["DATA", "ENCARREGADO"])
                                st.session_state.df_historico_f1.to_csv(caminho_historico_f1_csv, index=False)
                        else:
                            st.session_state.df_historico_f1 = pd.concat([st.session_state.df_historico_f1, df_novos], ignore_index=True).drop_duplicates(subset=["DATA", "ENCARREGADO"])
                            st.session_state.df_historico_f1.to_csv(caminho_historico_f1_csv, index=False)
                            st.toast(f"{len(novos_registros)} RDCs registrados localmente no Resumo Diário!", icon="✅")
                    else:
                        st.info("ℹ️ Os dados foram processados, mas os Encarregados dessa lista já haviam sido contabilizados.")

                # === DASHBOARD ANALÍTICO DO LOTE ===
                with st.expander("📊 Analytics do Lote Lido", expanded=True):
                    if not df_filtrado.empty:
                        col_dash1, col_dash2 = st.columns(2)
                        with col_dash1:
                            st.markdown("**Frentes de Serviço por Encarregado**")
                            enc_count = df_filtrado['ENCARREGADO'].replace('', 'Não Informado').value_counts()
                            st.bar_chart(enc_count, color="#0ea5e9")
                            
                        with col_dash2:
                            st.markdown("**Frentes por Disciplina**")
                            disc_count = df_filtrado['DISCIPLINA'].replace('', 'Não Informado').value_counts()
                            st.bar_chart(disc_count, color="#4ade80")
                            
                        st.markdown("**Distribuição por Caldeira (PB/RB)**")
                        caldeira_count = df_filtrado['CALDEIRA'].copy()
                        caldeira_count = caldeira_count.replace('', 'Não Identificada').value_counts()
                        st.bar_chart(caldeira_count, color="#f59e0b")

    with tab_ia_cc:
        st.markdown("### Robô Atualizador de C.C (Google Gemini)")
        st.markdown("Faça o upload dos PDFs aqui para o robô identificar o Local (PB/RB) e a Área (Estrutura, Tubulação, etc) e atualizar automaticamente o C.C. das equipes na base global do Google Sheets.")
        
        if HAS_GENAI:
            chave_padrao = ""
            try:
                chave_padrao = st.secrets.get("GEMINI_API_KEY", "")
            except Exception:
                pass
            
            st.markdown("#### Configuração e Upload")
            if not chave_padrao:
                chave_padrao = st.text_input("🔑 Cole suas Chaves da API Gemini (separadas por vírgula):", type="password", help="Chave oculta e protegida.", key="chave_cc")
            
            arquivos_scan_cc = st.file_uploader("Upload de RDCs para atualização de C.C (PDF, JPG, PNG)", type=["png", "jpg", "jpeg", "pdf"], accept_multiple_files=True, key="uploader_cc")
                
            btn_processar_cc = st.button("🚀 Atualizar C.C das Equipes com IA", type="primary", use_container_width=True)
            
            if btn_processar_cc and arquivos_scan_cc and chave_padrao:
                old_cred = os.environ.pop("GOOGLE_APPLICATION_CREDENTIALS", None)
                lista_chaves = [c.strip() for c in chave_padrao.split(",") if c.strip()]
                idx_chave_atual = 0
                client = genai.Client(api_key=lista_chaves[idx_chave_atual])
                nomes_para_prompt = ", ".join(lista_encarregados_base)
                
                prompt_ia_cc = f"""
                Analise este documento. Para CADA formulário de obra (RDC) encontrado no arquivo, extraia os dados.
                REGRA IMPORTANTÍSSIMA: Retorne APENAS UM objeto JSON por formulário/página.
                Retorne APENAS um array (lista) em formato JSON válido.
                [
                  {{
                    "DISCIPLINA": "...",
                    "ENCARREGADO": "...",
                    "PROBLEMAS": "...",
                    "LOCAL": "...",
                    "AREA": "..."
                  }}
                ]

                Regras de negócio:
                - DISCIPLINA: Extraia a disciplina ou função do topo, mas RETORNE APENAS A PRIMEIRA PALAVRA OU A PALAVRA PRINCIPAL.
                - ENCARREGADO: Extraia o nome do Encarregado escrito no papel. Compare com: [{nomes_para_prompt}]. Retorne EXATAMENTE o nome correspondente. Se ilegível, retorne 'AJUSTAR NOME'.
                - CALDEIRA: Se mencionar 'caldeira de recuperação' = 'RB'. Se 'caldeira de potência' = 'PB'. Se 'PRECIPITADOR' ou 'ESP' = 'ESP'. Se nenhum = ''.
                - LOCAL: Analise a imagem CUIDADOSAMENTE. Procure as opções 'PB ( )' e 'RB ( )'. Verifique se há um 'X', rabisco, visto ou marcação (mesmo que mal desenhada) dentro, em cima ou do lado dos parênteses. Retorne APENAS 'PB' ou 'RB'. Se nenhum, retorne ''.
                - AREA: Analise as caixinhas de área na imagem com LUPA. Procure por qualquer marcação (X, visto, círculo, rabisco) dentro ou sobre os parênteses. Opções: DUTO, EQUIPAMENTO, TUBULAÇÃO, ESTRUTURA MET, PRECIPITADOR, PRESSAO - MEC, PRESSAO - TUBULACAO, PRESSAO - FORNALHA, PINTURA, SOPRAGEM, ANDAIME. Retorne EXATAMENTE a área marcada. Se nenhuma, retorne ''.

                Apenas o JSON puro começando com [ e terminando com ].
                """

                # === ANIMAÇÃO PREMIUM DE LOADING ===
                animacao_html = """
                <div style="display: flex; flex-direction: column; align-items: center; justify-content: center; padding: 25px; background: rgba(15,23,42,0.9); border: 1px solid #0ea5e9; border-radius: 15px; box-shadow: 0 0 30px rgba(14, 165, 233, 0.3); margin-bottom: 20px;">
                    <div class="radar" style="position: relative; width: 120px; height: 120px; border-radius: 50%; border: 2px solid rgba(14,165,233,0.5); overflow: hidden; background: radial-gradient(circle, rgba(14,165,233,0.15) 0%, rgba(15,23,42,0) 100%);">
                        <div style="position: absolute; width: 50%; height: 50%; top: 0; left: 50%; transform-origin: bottom left; background: linear-gradient(45deg, rgba(14,165,233,0.9) 0%, transparent 50%); animation: radar-spin 1.5s linear infinite;"></div>
                        <div style="position: absolute; top: 50%; left: 50%; transform: translate(-50%, -50%); color: #fff; font-weight: bold; font-size: 14px; letter-spacing: 2px; text-shadow: 0 0 10px #0ea5e9; background: #0f172a; padding: 5px; border-radius: 5px;">ENESA</div>
                        <div style="position: absolute; top: 0; bottom: 0; left: 50%; width: 1px; background: rgba(14,165,233,0.4);"></div>
                        <div style="position: absolute; left: 0; right: 0; top: 50%; height: 1px; background: rgba(14,165,233,0.4);"></div>
                        <div style="position: absolute; top: 20%; left: 20%; width: 6px; height: 6px; background: #4ade80; border-radius: 50%; box-shadow: 0 0 10px #4ade80; animation: blip 1.5s infinite;"></div>
                        <div style="position: absolute; top: 70%; left: 60%; width: 4px; height: 4px; background: #4ade80; border-radius: 50%; box-shadow: 0 0 10px #4ade80; animation: blip 1.5s infinite 0.7s;"></div>
                    </div>
                    <p style="color: #0ea5e9; margin-top: 20px; font-weight: bold; font-size: 16px; animation: pulse 1s infinite; margin-bottom: 0;">🤖 IA Atualizando Centros de Custo...</p>
                    <style>
                        @keyframes radar-spin { 100% { transform: rotate(360deg); } }
                        @keyframes pulse { 0%, 100% { opacity: 1; } 50% { opacity: 0.6; } }
                        @keyframes blip { 0%, 100% { opacity: 0; } 10% { opacity: 1; } }
                    </style>
                </div>
                """
                animacao_placeholder_cc = st.empty()
                animacao_placeholder_cc.markdown(animacao_html, unsafe_allow_html=True)

                with st.status("🤖 Atualizando C.C das equipes...", expanded=True) as status_cc:
                    progresso = st.progress(0)
                    total_arquivos = len(arquivos_scan_cc)
                    houve_atualizacao_global = False
                    
                    for i, arquivo_scan in enumerate(arquivos_scan_cc):
                        status_cc.update(label=f"Processando arquivo {i+1} de {total_arquivos}: {arquivo_scan.name}...", state="running")
                    
                    try:
                        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{arquivo_scan.name.split('.')[-1]}") as tmp:
                            tmp.write(arquivo_scan.getvalue())
                            tmp_path = tmp.name
                            
                        max_tentativas = 3
                        sucesso_arquivo = False
                        for tentativa in range(max_tentativas):
                            try:
                                arquivo_up = client.files.upload(file=tmp_path)
                                
                                resposta = client.models.generate_content(
                                    model=st.session_state.get('modelo_gemini', 'gemini-2.5-flash'),
                                    contents=[arquivo_up, prompt_ia_cc],
                                    config=genai.types.GenerateContentConfig(
                                        response_mime_type="application/json",
                                        response_schema=list[RDC_CC_Schema],
                                        temperature=0.0
                                    )
                                )
                                
                                if old_cred:
                                    os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = old_cred

                                texto_json = resposta.text.strip()
                                if "```json" in texto_json:
                                    texto_json = texto_json.split("```json")[1].split("```")[0].strip()
                                elif "```" in texto_json:
                                    texto_json = texto_json.split("```")[1].split("```")[0].strip()
                                
                                start_idx = max(0, texto_json.find('[')) if '[' in texto_json else max(0, texto_json.find('{'))
                                end_idx = max(texto_json.rfind(']'), texto_json.rfind('}'))
                                if end_idx > start_idx:
                                    texto_json = texto_json[start_idx:end_idx+1]

                                try:
                                    dados_extraidos_lista = json.loads(texto_json)
                                except json.JSONDecodeError as err_json:
                                    import ast
                                    import re
                                    texto_fix = texto_json.replace("null", "None").replace("true", "True").replace("false", "False")
                                    texto_fix = re.sub(r'\}\s*\{', '}, {', texto_fix)
                                    texto_fix = re.sub(r'\]\s*\[', '], [', texto_fix)
                                    texto_fix = re.sub(r'("|\]|\})\s+(")', r'\1, \2', texto_fix)
                                    try:
                                        dados_extraidos_lista = ast.literal_eval(texto_fix)
                                    except SyntaxError:
                                        try:
                                            dados_extraidos_lista = ast.literal_eval(texto_fix + '"}]')
                                        except:
                                            raise err_json
                                    except:
                                        raise err_json
                                if isinstance(dados_extraidos_lista, dict):
                                    dados_extraidos_lista = [dados_extraidos_lista]
                                    
                                for dados in dados_extraidos_lista:
                                    # === ATUALIZAR C.C. COMPLETO NA BASE ===
                                    local_bruto = str(dados.get('LOCAL', '')).strip().upper()
                                    area_bruta = str(dados.get('AREA', '')).strip().upper()
                                    disciplina_lida = str(dados.get('DISCIPLINA', '')).strip().upper()
                                    enc_lido = str(dados.get('ENCARREGADO', '')).strip().upper()
                                    
                                    local_lido = ''
                                    if 'PB' in local_bruto: local_lido = 'PB'
                                    elif 'RB' in local_bruto: local_lido = 'RB'
                                    if not local_lido:
                                        cald = str(dados.get('CALDEIRA', '')).strip().upper()
                                        if 'PB' in cald: local_lido = 'PB'
                                        elif 'RB' in cald: local_lido = 'RB'
                                    
                                    area_lida = ''
                                    # 1. Tenta achar na área bruta (exato ou contendo)
                                    chaves_ordenadas = sorted(mapa_area_sufixo.keys(), key=len, reverse=True)
                                    for k in chaves_ordenadas:
                                        if k in area_bruta:
                                            area_lida = k
                                            break
                                            
                                    # 2. Se não achar, procura na disciplina (cuidado com falsos positivos de 'ESP')
                                    if not area_lida:
                                        import re
                                        for k in chaves_ordenadas:
                                            if k == 'ESP':
                                                if re.search(r'\bESP\b', disciplina_lida):
                                                    area_lida = k
                                                    break
                                            elif k in disciplina_lida:
                                                area_lida = k
                                                break
                                            
                                    if enc_lido and enc_lido != 'AJUSTAR NOME' and 'C.C' in df_atual.columns:
                                        encarregados_unicos = df_atual['ENCARREGADO'].dropna().unique()
                                        enc_encontrado = None
                                        
                                        for e in encarregados_unicos:
                                            if str(e).strip().upper() == enc_lido:
                                                enc_encontrado = e
                                                break
                                        if not enc_encontrado:
                                            for e in encarregados_unicos:
                                                if enc_lido in str(e).upper():
                                                    enc_encontrado = e
                                                    break
                                        if not enc_encontrado:
                                            import difflib
                                            matches = difflib.get_close_matches(enc_lido, [str(e).upper() for e in encarregados_unicos], n=1, cutoff=0.6)
                                            if matches:
                                                for e in encarregados_unicos:
                                                    if str(e).upper() == matches[0]:
                                                        enc_encontrado = e
                                                        break
                                                        
                                        if enc_encontrado:
                                            mask_enc = df_atual['ENCARREGADO'] == enc_encontrado
                                            atualizado = False
                                            
                                            if local_lido in ['PB', 'RB']:
                                                prefixo_novo = '125.02' if local_lido == 'PB' else '125.01'
                                                sufixo = mapa_area_sufixo.get(area_lida, '')
                                                
                                                if sufixo:
                                                    cc_novo = f"{prefixo_novo}.{sufixo}"
                                                    df_atual.loc[mask_enc, 'C.C'] = cc_novo
                                                    atualizado = True
                                                    st.toast(f"✅ C.C. de TODA A EQUIPE de {enc_encontrado} → {cc_novo}")
                                                else:
                                                    if local_lido == 'PB':
                                                        df_atual.loc[mask_enc, 'C.C'] = df_atual.loc[mask_enc, 'C.C'].str.replace('125.01.', '125.02.', regex=False)
                                                    else:
                                                        df_atual.loc[mask_enc, 'C.C'] = df_atual.loc[mask_enc, 'C.C'].str.replace('125.02.', '125.01.', regex=False)
                                                    atualizado = True
                                                    st.toast(f"⚠️ C.C. de TODA A EQUIPE de {enc_encontrado} atualizado parcialmente → {local_lido} (manteve sufixo)")
                                            else:
                                                st.warning(f"❌ C.C não atualizado para a equipe de {enc_encontrado}: O robô não conseguiu identificar se o local era PB ou RB.")
                                            
                                            if atualizado:
                                                st.session_state.df = df_atual.copy()
                                                houve_atualizacao_global = True
                                        else:
                                            st.error(f"❌ Encarregado '{enc_lido}' não encontrado na base. Equipe não atualizada.")

                                sucesso_arquivo = True
                                break 

                            except Exception as inner_e:
                                erro_str = str(inner_e)
                                if '429' in erro_str or 'RESOURCE_EXHAUSTED' in erro_str:
                                    if tentativa < max_tentativas - 1:
                                        if idx_chave_atual < len(lista_chaves) - 1:
                                            idx_chave_atual += 1
                                            client = genai.Client(api_key=lista_chaves[idx_chave_atual])
                                            st.warning(f"🔄 Limite atingido na chave atual. Trocando para a chave reserva {idx_chave_atual + 1}/{len(lista_chaves)}...")
                                            time.sleep(2)
                                            continue
                                        else:
                                            st.warning(f"⏳ Cota do Google atingida em todas as chaves. Aguardando 60 segundos... (Tentativa {tentativa+1}/{max_tentativas})")
                                            time.sleep(60)
                                            continue
                                elif '503' in erro_str or 'UNAVAILABLE' in erro_str:
                                    if tentativa < max_tentativas - 1:
                                        st.warning(f"⏳ Servidores da IA sobrecarregados. Tentando novamente em 10 segundos... (Tentativa {tentativa+1}/{max_tentativas})")
                                        time.sleep(10)
                                        continue
                                        
                                msg_erro = f"Erro detalhado na IA: {inner_e}"
                                try:
                                    modelos = [m.name for m in client.models.list()]
                                    msg_erro += f" | Modelos liberados: {modelos}"
                                except:
                                    pass
                                st.error(msg_erro)
                                break
                                    
                        os.remove(tmp_path)
                        
                        if sucesso_arquivo:
                            st.toast(f"✅ {arquivo_scan.name} processado com sucesso!")
                        else:
                            st.toast(f"❌ Falha ao processar {arquivo_scan.name}.")
                            st.session_state.teve_falha_ia_cc = True
                            
                    except Exception as e:
                        st.error(f"Erro no envio do arquivo {arquivo_scan.name}: {e}")
                        st.session_state.teve_falha_ia_cc = True
                        
                    progresso.progress((i + 1) / total_arquivos)

                expandir_status = st.session_state.get('teve_falha_ia_cc', False)
                status_cc.update(label="✅ Atualização de C.Cs concluída!" if not expandir_status else "⚠️ Leitura finalizada com erros", state="complete", expanded=expandir_status)
                animacao_placeholder_cc.empty()
                st.session_state.teve_falha_ia_cc = False
                
                if houve_atualizacao_global:
                    try:
                        df_atual = preparar_dataframe(df_atual)
                        st.session_state.df = df_atual.copy()
                        
                        status_cc.update(label="Sincronizando C.Cs atualizados com a nuvem...", state="running")
                        conn_update = st.connection("gsheets", type=GSheetsConnection)
                        conn_update.update(worksheet="PDE", data=df_atual)
                        st.cache_data.clear()
                    except Exception as e:
                        st.error(f"Erro ao salvar na nuvem: {e}")

                status_cc.update(label="🎉 Atualização de C.Cs concluída!", state="complete", expanded=False)
                time.sleep(2)
                st.session_state.force_use_local = True
                st.rerun()
                
                st.dataframe(df_filtrado, use_container_width=True)

    with tab_cc:
        st.markdown("### 💰 Controle de Centro de Custo (C.C)")
        
        # === ÚLTIMA ATUALIZAÇÃO ===
        ultima_base = ""
        ultima_cc = ""
        try:
            if os.path.exists(caminho_base_salva_csv):
                ts_base = os.path.getmtime(caminho_base_salva_csv)
                ultima_base = datetime.datetime.fromtimestamp(ts_base).strftime("%d/%m/%Y às %H:%M")
            elif os.path.exists(caminho_base_salva_xlsx):
                ts_base = os.path.getmtime(caminho_base_salva_xlsx)
                ultima_base = datetime.datetime.fromtimestamp(ts_base).strftime("%d/%m/%Y às %H:%M")
        except Exception:
            pass
        try:
            if os.path.exists(caminho_hist_cc):
                ts_cc = os.path.getmtime(caminho_hist_cc)
                ultima_cc = datetime.datetime.fromtimestamp(ts_cc).strftime("%d/%m/%Y às %H:%M")
        except Exception:
            pass
        
        html_update = f"""
        <div style="display: flex; gap: 15px; margin-bottom: 20px; flex-wrap: wrap;">
            <div style="background: rgba(16, 185, 129, 0.1); border: 1px solid rgba(16, 185, 129, 0.3); border-radius: 10px; padding: 10px 18px; display: flex; align-items: center; gap: 10px;">
                <div style="width: 8px; height: 8px; border-radius: 50%; background: #10b981; box-shadow: 0 0 8px #10b981; animation: pulse_dot 2s infinite;"></div>
                <span style="font-size: 13px; color: #94a3b8;">Base PDE atualizada em: <b style="color: #10b981;">{ultima_base if ultima_base else 'N/A'}</b></span>
            </div>
            <div style="background: rgba(139, 92, 246, 0.1); border: 1px solid rgba(139, 92, 246, 0.3); border-radius: 10px; padding: 10px 18px; display: flex; align-items: center; gap: 10px;">
                <div style="width: 8px; height: 8px; border-radius: 50%; background: #8b5cf6; box-shadow: 0 0 8px #8b5cf6; animation: pulse_dot 2s infinite;"></div>
                <span style="font-size: 13px; color: #94a3b8;">Histórico C.C salvo em: <b style="color: #8b5cf6;">{ultima_cc if ultima_cc else 'N/A'}</b></span>
            </div>
        </div>
        <style>
            @keyframes pulse_dot {{
                0%, 100% {{ opacity: 1; }}
                50% {{ opacity: 0.3; }}
            }}
        </style>
        """
        st.markdown(html_update, unsafe_allow_html=True)
        
        if "C.C" not in df_atual.columns or df_atual["C.C"].str.strip().eq("").all():
            st.warning("⚠️ A coluna de Centro de Custo (C.C) não foi encontrada na base de dados atual. Verifique se a planilha possui essa coluna.")
        else:
            # === ALERTA DE C.C INVÁLIDO ===
            df_em_branco = df_atual[df_atual["C.C"].isna() | df_atual["C.C"].str.strip().eq("")]
            if not df_em_branco.empty:
                st.error(f"⚠️ **ALERTA DE SISTEMA:** Existem **{len(df_em_branco)} colaboradores** na base atual **sem Centro de Custo** (C.C em branco). Eles não aparecerão nos cálculos de custo!")
            
            valid_prefixes = ["125.01.", "125.02."]
            valid_suffixes = ['001', '002', '003', '004', '005', '006', '007', '008', '009', '010', '011', '012', '013', '014', '015', '016', '101', '102', '103', '104', '105', '106', '107', '108', '109', '110', '111', '112', '113']
            
            invalid_cc_list = []
            df_preenchido = df_atual[~df_atual["C.C"].isna() & (df_atual["C.C"].str.strip() != "")]
            for _, row in df_preenchido.iterrows():
                cc_val = str(row["C.C"]).strip()
                is_valid = False
                for prefix in valid_prefixes:
                    if cc_val.startswith(prefix):
                        suf = cc_val.replace(prefix, "")
                        if suf in valid_suffixes:
                            is_valid = True
                            break
                if not is_valid:
                    invalid_cc_list.append(cc_val)
                    
            if invalid_cc_list:
                invalid_cc_count = len(invalid_cc_list)
                unique_invalids = list(set(invalid_cc_list))
                st.warning(f"⚠️ **ATENÇÃO:** Foram encontrados **{invalid_cc_count} colaboradores** com C.C **inválido** (não existe no mapa oficial). Exemplos: {', '.join(unique_invalids[:5])}")

            # Filtro PB/RB/ESP, Turno e Status Global para a aba C.C
            col_cc_filt1, col_cc_filt2, col_cc_filt3 = st.columns(3)
            with col_cc_filt1:
                filtro_local = st.segmented_control(
                    "Filtrar Dados por Local:", 
                    ["Ambas", "PB", "RB", "ESP"], 
                    default="Ambas",
                    key="filtro_cc_local_key"
                )
                if not filtro_local:
                    filtro_local = "Ambas"
                    
            with col_cc_filt2:
                turnos_cc_disponiveis = ["Todos"]
                if "TURNO" in df_atual.columns:
                    turnos_reais = [t for t in df_atual["TURNO"].unique() if str(t).strip() and str(t) != "nan"]
                    turnos_cc_disponiveis.extend(sorted(turnos_reais))
                
                filtro_cc_turno = st.selectbox(
                    "Filtrar por Turno:", 
                    turnos_cc_disponiveis,
                    index=0,
                    key="filtro_cc_turno_key"
                )
                
            with col_cc_filt3:
                status_cc_disponiveis = ["Todos"]
                if "STATUS" in df_atual.columns:
                    status_reais = [s for s in df_atual["STATUS"].unique() if str(s).strip() and str(s) != "nan"]
                    status_cc_disponiveis.extend(sorted(status_reais))
                
                filtro_cc_status = st.selectbox(
                    "Filtrar por Status:", 
                    status_cc_disponiveis,
                    index=status_cc_disponiveis.index("ATIVO") if "ATIVO" in status_cc_disponiveis else 0,
                    key="filtro_cc_status_key"
                )
                
            df_cc_aba = df_atual[df_atual["C.C"].str.strip() != ""]
            if filtro_local == "PB":
                df_cc_aba = df_cc_aba[df_cc_aba["C.C"].apply(lambda x: "125.02" in str(x) and ".005" not in str(x))]
            elif filtro_local == "RB":
                df_cc_aba = df_cc_aba[df_cc_aba["C.C"].apply(lambda x: "125.01" in str(x) and ".005" not in str(x))]
            elif filtro_local == "ESP":
                df_cc_aba = df_cc_aba[df_cc_aba["C.C"].apply(lambda x: ".005" in str(x))]
                
            if filtro_cc_turno != "Todos" and "TURNO" in df_cc_aba.columns:
                df_cc_aba = df_cc_aba[df_cc_aba["TURNO"] == filtro_cc_turno]
                
                
            if filtro_cc_status != "Todos" and "STATUS" in df_cc_aba.columns:
                df_cc_aba = df_cc_aba[df_cc_aba["STATUS"] == filtro_cc_status]

            lista_cc = sorted([str(cc) for cc in df_cc_aba["C.C"].unique()])
            
            def format_cc(cc_code):
                if cc_code == "TODOS": return "TODOS"
                local = "PB" if "125.02" in cc_code else ("RB" if "125.01" in cc_code else "")
                sufixo = str(cc_code).split('.')[-1] if '.' in str(cc_code) else str(cc_code)
                mapa_sufixo_nome = {
                    '001': 'Equipamentos', '002': 'Dutos', '003': 'Tubulação', 
                    '004': 'Estrutura Metálica', '005': 'Precipitador', '006': 'Pressão - Mecânica', 
                    '007': 'Pressão - Tubulação', '008': 'Pressão - Fornalha', '009': 'Pintura', 
                    '010': 'Comissionamento', '011': 'Op. Assistida', '012': 'Lavagem Química', 
                    '013': 'Sopragem', '014': 'Andaime', '015': 'Operadores', '016': 'Fora de Escopo',
                    '101': 'Gerência', '102': 'Produção', '103': 'Garantia da Qualidade',
                    '104': 'Planejamento', '105': 'Administração', '106': 'Segurança e Medicina',
                    '107': 'Infraestrutura', '108': 'Almoxarifado ENESA', '109': 'Almoxarifado Materiais',
                    '110': 'Manut. Elétrica Provisória', '111': 'Topografia', '112': 'Movimentação de Cargas',
                    '113': 'Medição/Contratos'
                }
                nome = mapa_sufixo_nome.get(sufixo, '')
                
                if nome and local: return f"{cc_code} - {nome} ({local})"
                elif nome: return f"{cc_code} - {nome}"
                elif local: return f"{cc_code} ({local})"
                else: return str(cc_code)
            
            # Métricas gerais Customizadas
            def card_kpi_cc(titulo, valor, cor):
                return f"""
                <div style="background: rgba(30, 41, 59, 0.45); backdrop-filter: blur(10px); border-radius: 16px; border: 1px solid rgba(255,255,255,0.05); padding: 18px; box-shadow: 0 4px 15px rgba(0,0,0,0.2); position: relative; overflow: hidden; height: 110px; transition: transform 0.3s ease;" onmouseover="this.style.transform='translateY(-5px)'" onmouseout="this.style.transform='translateY(0px)'">
                    <p style="margin: 0; font-size: 13px; color: #94a3b8; font-weight: 600; text-transform: uppercase; letter-spacing: 0.5px;">{titulo}</p>
                    <h2 style="margin: 5px 0 0 0; font-size: 34px; font-weight: 700; color: #f8fafc; text-shadow: 0 0 15px {cor}60;">{valor}</h2>
                    <div style="position: absolute; bottom: 0; left: 0; width: 100%; height: 4px; background: linear-gradient(90deg, {cor}, transparent); box-shadow: 0 -2px 10px {cor}80;"></div>
                </div>
                """
                
            mc1, mc2, mc3, mc4, mc5 = st.columns(5)
            with mc1: st.markdown(card_kpi_cc(t("Centros de Custo"), len(lista_cc), "#8b5cf6"), unsafe_allow_html=True)
            with mc2: st.markdown(card_kpi_cc(t("Total Alocados"), len(df_cc_aba), "#3b82f6"), unsafe_allow_html=True)
            with mc3: st.markdown(card_kpi_cc(t("Funções Distintas"), df_cc_aba["FUNÇÃO"].nunique(), "#f59e0b"), unsafe_allow_html=True)
            
            qtd_encarregados = len([e for e in df_cc_aba["ENCARREGADO"].unique() if str(e).strip() != "" and str(e) in lista_completa_encarregados])
            with mc4: st.markdown(card_kpi_cc(t("Encarregados"), qtd_encarregados, "#10b981"), unsafe_allow_html=True)
            
            span_of_control = round(len(df_cc_aba) / qtd_encarregados, 1) if qtd_encarregados > 0 else 0
            with mc5: st.markdown(card_kpi_cc(t("Span of Control"), span_of_control, "#0ea5e9"), unsafe_allow_html=True)
            st.markdown("<br>", unsafe_allow_html=True)

            # Gráficos lado a lado
            col_graf1, col_graf2 = st.columns([6, 4])
            
            with col_graf1:
                # Gráfico de distribuição por C.C.
                st.markdown("**Distribuição de Efetivo por Centro de Custo**")
                
                cc_contagem = df_cc_aba["C.C"].value_counts().reset_index()
                cc_contagem.columns = ["Centro de Custo", "Quantidade"]
                cc_contagem["Nome C.C"] = cc_contagem["Centro de Custo"].apply(format_cc)
                
                if len(cc_contagem) > 0:
                    fig_cc = px.bar(cc_contagem, x="Quantidade", y="Nome C.C", orientation="h", color="Quantidade", color_continuous_scale=[(0, "#0f172a"), (1, "#8b5cf6")], text="Quantidade")
                    fig_cc.update_layout(showlegend=False, xaxis_title="", yaxis_title="", margin=dict(l=0, r=40, t=10, b=0), paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", font=dict(color="#e0e4ea"), height=max(300, len(cc_contagem) * 35))
                    fig_cc.update_yaxes(categoryorder="total ascending")
                    fig_cc.update_xaxes(visible=False)
                    fig_cc.update_coloraxes(showscale=False)
                    fig_cc.update_traces(textposition='outside', cliponaxis=False)
                    
                    st.plotly_chart(fig_cc, use_container_width=True)
                else:
                    st.info("Nenhum dado encontrado para gerar gráfico de C.C.")
                    
            with col_graf2:
                # Gráfico: MOD vs MOI
                st.markdown("**Proporção MOD vs MOI**")
                df_mod = df_cc_aba[df_cc_aba["MÃO DE OBRA"].str.strip() != ""]
                if not df_mod.empty:
                    mo_contagem = df_mod["MÃO DE OBRA"].value_counts().reset_index()
                    mo_contagem.columns = ["Tipo", "Quantidade"]
                    fig_mo = px.pie(mo_contagem, values="Quantidade", names="Tipo", hole=0.65, color_discrete_sequence=["#4a9eed", "#f39c12", "#e74c3c"])
                    fig_mo.update_layout(margin=dict(l=20, r=20, t=40, b=20), paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", font=dict(color="#e0e4ea"), height=350, showlegend=True, legend=dict(orientation="h", yanchor="bottom", y=-0.2, xanchor="center", x=0.5))
                    
                    st.plotly_chart(fig_mo, use_container_width=True)
                else:
                    st.info("Dados de Mão de Obra não disponíveis.")
            
            st.markdown("---")
            
            # --- Seção de Histórico (Máquina do Tempo) ---
            if os.path.exists(caminho_hist_cc):
                try:
                    df_hist = pd.read_csv(caminho_hist_cc)
                    if not df_hist.empty and "DATA" in df_hist.columns:
                        st.markdown("**📈 Máquina do Tempo: Evolução do Efetivo**")
                        # Filtro para escolher o CC ou Área
                        opcoes_historico = ["Geral (Todos)", "Resumo: PB (Caldeira)", "Resumo: RB (Retorta)", "Resumo: ESP (Precipitador)"] + lista_cc
                        cc_selecionado = st.selectbox("Selecione a equipe ou área para analisar o crescimento:", opcoes_historico)
                        
                        if cc_selecionado == "Geral (Todos)":
                            df_plot = df_hist.groupby("DATA")["Efetivo"].sum().reset_index()
                            titulo_graf = "Crescimento Geral da Obra"
                        elif cc_selecionado == "Resumo: PB (Caldeira)":
                            df_plot = df_hist[df_hist["C.C"].apply(lambda x: "125.02" in str(x) and ".005" not in str(x))].groupby("DATA")["Efetivo"].sum().reset_index()
                            titulo_graf = "Crescimento - Área PB (Caldeira)"
                        elif cc_selecionado == "Resumo: RB (Retorta)":
                            df_plot = df_hist[df_hist["C.C"].apply(lambda x: "125.01" in str(x) and ".005" not in str(x))].groupby("DATA")["Efetivo"].sum().reset_index()
                            titulo_graf = "Crescimento - Área RB (Retorta)"
                        elif cc_selecionado == "Resumo: ESP (Precipitador)":
                            df_plot = df_hist[df_hist["C.C"].apply(lambda x: ".005" in str(x))].groupby("DATA")["Efetivo"].sum().reset_index()
                            titulo_graf = "Crescimento - Área ESP (Precipitador)"
                        else:
                            df_plot = df_hist[df_hist["C.C"] == cc_selecionado].copy()
                            titulo_graf = f"Evolução - C.C {cc_selecionado}"
                            
                        if not df_plot.empty:
                            df_plot["DATA_DT"] = pd.to_datetime(df_plot["DATA"], errors='coerce')
                            df_plot = df_plot.sort_values("DATA_DT")
                            
                            fig_hist = px.line(df_plot, x="DATA", y="Efetivo", markers=True, title=titulo_graf, line_shape="spline")
                            fig_hist.update_layout(xaxis_title="", yaxis_title="Quantidade de Colaboradores", paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", font=dict(color="#e0e4ea"), height=300)
                            fig_hist.update_xaxes(type='category')
                            fig_hist.update_yaxes(tickformat="d")
                            fig_hist.update_traces(line=dict(width=3, color="#0ea5e9"), marker=dict(size=8, color="#10b981"))
                            st.plotly_chart(fig_hist, use_container_width=True)
                        else:
                            st.info("Aguardando acumular mais dias de dados para gerar a curva.")
                except Exception:
                    pass
            # ---------------------------------------------
            st.markdown("**Liderança: Efetivo de Encarregados por C.C**")
            # Tabela sumarizando quantos encarregados tem em cada CC (Apenas encarregados reais da lista oficial)
            df_lideres = df_cc_aba[(df_cc_aba["ENCARREGADO"].str.strip() != "") & (df_cc_aba["ENCARREGADO"].isin(lista_completa_encarregados))].copy()
            
            if len(df_lideres) > 0:
                df_agrupado = df_lideres.groupby(["C.C", "ENCARREGADO"]).size().reset_index(name="QTD. COLABORADORES")
                df_agrupado["LOCAL"] = df_agrupado["C.C"].apply(lambda x: "PB" if "125.02" in str(x) else ("RB" if "125.01" in str(x) else "OUTROS"))
                df_agrupado["NOME C.C"] = df_agrupado["C.C"].apply(format_cc)
                
                # Reorganizar colunas
                df_agrupado = df_agrupado[["LOCAL", "NOME C.C", "ENCARREGADO", "QTD. COLABORADORES"]].sort_values(by=["LOCAL", "NOME C.C", "QTD. COLABORADORES"], ascending=[True, True, False])
                
                # Exibe a tabela agrupada
                st.dataframe(df_agrupado, hide_index=True, use_container_width=True)
            else:
                st.info("Nenhum encarregado da lista oficial vinculado a um Centro de Custo para o filtro selecionado.")
                
            st.markdown("---")
            
            # Filtros por C.C. e Equipe
            st.markdown("**Consulta Detalhada**")
            
            # Filtramos a lista de encarregados para exibir APENAS quem realmente é encarregado da lista oficial
            lista_encarregados_detalhada = sorted([str(e) for e in df_cc_aba["ENCARREGADO"].unique() if str(e).strip() != "" and str(e) in lista_completa_encarregados])
            
            col_f1, col_f2 = st.columns(2)
            with col_f1:
                cc_selecionado = st.selectbox("Selecione o Centro de Custo:", ["TODOS"] + lista_cc, format_func=format_cc)
            with col_f2:
                enc_selecionado = st.selectbox("Selecione a Equipe (Encarregado):", ["TODAS AS EQUIPES"] + lista_encarregados_detalhada)
            
            df_cc_filtrado = df_cc_aba.copy()
            
            if cc_selecionado != "TODOS":
                df_cc_filtrado = df_cc_filtrado[df_cc_filtrado["C.C"] == cc_selecionado]
                
            if enc_selecionado != "TODAS AS EQUIPES":
                df_cc_filtrado = df_cc_filtrado[df_cc_filtrado["ENCARREGADO"] == enc_selecionado]
            
            if len(df_cc_filtrado) > 0:
                # Resumo de funções no C.C. selecionado
                st.markdown(f"**Funções no C.C. selecionado** ({len(df_cc_filtrado)} colaboradores)")
                func_cc = df_cc_filtrado["FUNÇÃO"].value_counts().reset_index()
                func_cc.columns = ["Função", "Quantidade"]
                
                fig_func = px.bar(func_cc, x="Quantidade", y="Função", orientation="h", color="Quantidade", color_continuous_scale="Oranges", text="Quantidade")
                fig_func.update_layout(showlegend=False, xaxis_title="", yaxis_title="", margin=dict(l=0, r=0, t=10, b=0), paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", font=dict(color="#e0e4ea"), height=max(200, len(func_cc) * 35))
                fig_func.update_yaxes(categoryorder="total ascending")
                fig_func.update_xaxes(visible=False)
                fig_func.update_coloraxes(showscale=False)
                fig_func.update_traces(textposition='outside')
                if st.toggle("📊 Visualizar Gráfico de Funções"):
                    st.plotly_chart(fig_func, use_container_width=True)
                
                # Tabela detalhada
                colunas_exibir = ["MATRICULA", "NOME", "FUNÇÃO", "C.C", "ENCARREGADO"]
                if "DISCIPLINA" in df_cc_filtrado.columns:
                    colunas_exibir.append("DISCIPLINA")
                if "MÃO DE OBRA" in df_cc_filtrado.columns:
                    colunas_exibir.append("MÃO DE OBRA")
                colunas_exibir = [c for c in colunas_exibir if c in df_cc_filtrado.columns]
                
                # Botão de download
                buf_cc = io.BytesIO()
                df_cc_filtrado[colunas_exibir].to_excel(buf_cc, index=False)
                buf_cc.seek(0)
                nome_cc_arq = cc_selecionado.replace(".", "_") if cc_selecionado != "TODOS" else "TODOS"
                st.download_button("⬇️ Baixar Relatório C.C (.xlsx)", data=buf_cc, file_name=f"CC_{nome_cc_arq}.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", type="primary", use_container_width=True)
                
                st.dataframe(df_cc_filtrado[colunas_exibir].reset_index(drop=True), hide_index=True, use_container_width=True)
            else:
                st.info("Nenhum colaborador encontrado para este Centro de Custo.")

    with tab_rdc_digital:
        st.markdown("### <span class='material-symbols-rounded' style='vertical-align: middle; color: #0ea5e9; font-size: 32px;'>edit_document</span> Lançamento de RDC Digital", unsafe_allow_html=True)
        st.caption("Preencha as informações do seu dia de trabalho seguindo as 3 etapas abaixo. Os dados serão salvos na nuvem.")
        
        with st.form("form_rdc_digital"):
            tab_id, tab_local, tab_ativ = st.tabs(["1️⃣ Identificação", "2️⃣ Localização", "3️⃣ Atividades e Envio"])
            
            with tab_id:
                st.markdown("<p style='color: #94a3b8; font-size: 14px;'>Quem é você e qual seu turno?</p>", unsafe_allow_html=True)
                rdc_encarregado = st.selectbox("Selecione seu Nome (Encarregado):", [""] + lista_completa_encarregados)
                rdc_turno = st.selectbox("Turno de Trabalho:", ["DIURNO", "NOTURNO", "MISTO"])
                
            with tab_local:
                import datetime
                st.markdown("<p style='color: #94a3b8; font-size: 14px;'>Onde você trabalhou hoje?</p>", unsafe_allow_html=True)
                
                rdc_data = st.date_input("Data do Relatório:", datetime.date.today())
                
                area_options = ["PB", "RB", "ESP", "LAYDOWN 1", "LAYDOWN 2", "OUTRO (DIGITAR)"]
                area_sel = st.selectbox("Área / Local de Trabalho:", area_options)
                rdc_area = area_sel
                if area_sel == "OUTRO (DIGITAR)":
                    rdc_area = st.text_input("Qual Área/Local?", placeholder="Ex: Escritório, Almoxarifado...")
                
                disc_options = [
                    "EQUIPAMENTOS", "DUTOS", "TUBULACAO", "ESTRUTURA METALICA", "PRECIPITADOR", 
                    "PRESSAO - MECANICA", "PRESSAO - TUBULACAO", "PRESSAO - FORNALHA", "PINTURA", 
                    "COMISSIONAMENTO", "OP. ASSISTIDA", "LAVAGEM QUIMICA", "SOPRAGEM", "ANDAIME", 
                    "OPERADORES", "FORA DE ESCOPO", "GERENCIA", "PRODUCAO", "GARANTIA DA QUALIDADE", 
                    "PLANEJAMENTO", "ADMINISTRACAO", "SEGURANCA E MEDICINA DO TRABALHO", "INFRAESTRUTURA", 
                    "ALMOXARIFADO ENESA", "ALMOXARIFADO MATERIAIS", "MANUT. ELETRICA PROVISORIA", 
                    "TOPOGRAFIA", "MOVIMENTACAO DE CARGAS", "MEDICAO/CUSTO/CONTRATOS", "CIVIL", "MECÂNICA", "ELÉTRICA", "INSTRUMENTAÇÃO", "ISOLAMENTO", "OUTRA (DIGITAR)"
                ]
                disc_sel = st.selectbox("Disciplina Principal:", disc_options)
                
                rdc_disciplina = disc_sel
                if disc_sel == "OUTRA (DIGITAR)":
                    rdc_disciplina = st.text_input("Qual Disciplina?", placeholder="Ex: Tubulação, Solda...")
                    
            with tab_ativ:
                st.markdown("<p style='color: #94a3b8; font-size: 14px;'>O que foi executado?</p>", unsafe_allow_html=True)
                rdc_dds = st.text_input("Tópico do DDS do dia:")
                rdc_atividades = st.text_area("Atividades Executadas (Detalhe os serviços feitos pela equipe):", height=150)
                rdc_problemas = st.text_area("Problemas / Interrupções / Ocorrências (Opcional):", height=68)
                
                st.markdown("<br>", unsafe_allow_html=True)
                submit_rdc = st.form_submit_button("🚀 Salvar e Enviar RDC na Nuvem", use_container_width=True, type="primary")
            
            if submit_rdc:
                if not rdc_encarregado:
                    st.error("⚠️ Por favor, selecione o nome do Encarregado.")
                elif not rdc_atividades.strip():
                    st.error("⚠️ Por favor, preencha as Atividades Executadas.")
                elif disc_sel == "OUTRA (DIGITAR)" and not rdc_disciplina.strip():
                    st.error("⚠️ Digite a disciplina na caixa 'Qual Disciplina?'.")
                else:
                    rdc_json = [{
                        "ENCARREGADO": rdc_encarregado,
                        "DATA": rdc_data.strftime("%Y/%m/%d"),
                        "TURNO": rdc_turno,
                        "AREA": rdc_area.strip().upper(),
                        "DISCIPLINA": rdc_disciplina.strip().upper(),
                        "DDS": rdc_dds.strip(),
                        "ATIVIDADE": rdc_atividades.strip(),
                        "CALDEIRA": rdc_problemas.strip(),
                        "PROBLEMAS": rdc_problemas.strip()
                    }]
                    
                    import json
                    import requests
                    
                    WEBHOOK_URL = "https://script.google.com/macros/s/AKfycbxfE96gE7ckdmapBLBHJuoX2bvAt-2d76OUJNiSRsLgFCOiySeQhFOopp3DoC5Fn95D/exec"
                    
                    try:
                        with st.spinner("Enviando dados para a nuvem..."):
                            res = requests.post(WEBHOOK_URL, json=rdc_json, allow_redirects=True)
                        if res.status_code == 200:
                            st.toast(f"RDC Digital de {rdc_encarregado} salvo com sucesso na Nuvem!", icon="✅")
                            st.info("Para visualizar na tabela da IA, clique em 'Puxar Dados Automáticos' abaixo.")
                        else:
                            st.error(f"❌ Erro ao enviar. Servidor retornou: {res.text}")
                    except Exception as e:
                        st.error(f"❌ Falha de conexão: {e}")
        

        st.markdown("---")
        st.markdown("### 📥 Sincronização de RDCs (Nuvem)")
        st.caption("Clique no botão abaixo para puxar todos os RDCs lançados pelos encarregados no sistema.")
        
        WEBHOOK_URL = "https://script.google.com/macros/s/AKfycbxfE96gE7ckdmapBLBHJuoX2bvAt-2d76OUJNiSRsLgFCOiySeQhFOopp3DoC5Fn95D/exec"
        
        if st.button("🔄 Puxar Dados Automáticos (Google Sheets)", type="primary", use_container_width=True):
            with st.spinner("Conectando ao Banco de Dados na Nuvem..."):
                try:
                    import requests
                    response = requests.get(WEBHOOK_URL, timeout=15)
                    
                    if response.status_code == 200:
                        dados_offline = response.json()
                        
                        if isinstance(dados_offline, list) and len(dados_offline) > 0:
                            if 'df_ia' not in st.session_state:
                                st.session_state.df_ia = pd.DataFrame(columns=['ITEM', 'SUB', 'DATA', 'DISCIPLINA', 'ENCARREGADO', 'TURNO', 'DDS', 'TRANSCRICAO', 'ATIVIDADE', 'SUB_ATIVIDADE', 'LOCAL_ESPECIFICO', 'EFETIVO_ATIVIDADE', 'PROBLEMAS', 'LOCAL', 'AREA', 'CALDEIRA'])
                                
                            ultimo_item = st.session_state.df_ia['ITEM'].max() if not st.session_state.df_ia.empty and pd.notna(st.session_state.df_ia['ITEM'].max()) else 0
                            
                            novos_registros = []
                            for r in dados_offline:
                                ultimo_item += 1
                                novo_reg = {
                                    'ITEM': ultimo_item,
                                    'DATA': r.get('DATA', ''),
                                    'DISCIPLINA': str(r.get('DISCIPLINA', '')).strip().upper(),
                                    'ENCARREGADO': r.get('ENCARREGADO', ''),
                                    'TURNO': r.get('TURNO', ''),
                                    'DDS': r.get('TOPICO_DDS', r.get('DDS', '')),
                                    'ATIVIDADE': r.get('ATIVIDADES', r.get('ATIVIDADE', '')),
                                    'PROBLEMAS': r.get('PROBLEMAS', r.get('CALDEIRA', '')),
                                    'LOCAL': str(r.get('AREA', '')).strip().upper(),
                                    'AREA': str(r.get('AREA', '')).strip().upper()
                                }
                                novos_registros.append(novo_reg)
                                
                            st.session_state.df_ia = pd.concat([st.session_state.df_ia, pd.DataFrame(novos_registros)], ignore_index=True)
                            st.success(f"📦 Sincronização Automática concluída! {len(novos_registros)} RDCs puxados do Google Sheets com sucesso.")
                            st.balloons()
                        else:
                            st.info("👍 Nenhum RDC novo pendente no Google Sheets no momento.")
                    else:
                        st.error(f"❌ Erro de conexão. Código HTTP: {response.status_code}")
                except Exception as e:
                    st.error(f"❌ Falha de rede ao tentar conectar com a nuvem: {e}")

    # ==============================================================
    # ABA 10: GERENCIAR PDE
    # ==============================================================
    with tab_pde:
        st.markdown("### 👷 Gerenciar PDE — Base de Funcionários")
        st.markdown("Visualize, edite e exporte a base completa de efetivo.")

        # Filters
        col1, col2, col3 = st.columns(3)
        with col1:
            encarregado_filtro = st.selectbox("Encarregado", options=["Todos"] + sorted(list(df_atual["ENCARREGADO"].dropna().unique())), key="pde_enc_filtro")
        with col2:
            disciplina_filtro = st.selectbox("Disciplina", options=["Todos"] + sorted(list(df_atual["DISCIPLINA"].dropna().unique())), key="pde_disc_filtro")
        with col3:
            status_filtro = st.selectbox("Status", options=["Todos"] + sorted(list(df_atual["STATUS"].dropna().unique())), key="pde_status_filtro")

        # Apply filters
        df_pde_filtrado = df_atual.copy()
        if encarregado_filtro != "Todos":
            df_pde_filtrado = df_pde_filtrado[df_pde_filtrado["ENCARREGADO"] == encarregado_filtro]
        if disciplina_filtro != "Todos":
            df_pde_filtrado = df_pde_filtrado[df_pde_filtrado["DISCIPLINA"] == disciplina_filtro]
        if status_filtro != "Todos":
            df_pde_filtrado = df_pde_filtrado[df_pde_filtrado["STATUS"] == status_filtro]

        # Summary Metrics
        m1, m2, m3, m4 = st.columns(4)
        m1.metric("Total de Funcionários", len(df_pde_filtrado))
        m2.metric("Total de Encarregados", df_pde_filtrado["ENCARREGADO"].nunique())
        m3.metric("Total de C.C", df_pde_filtrado["C.C"].nunique())
        
        # Detectar colaboradores sem encarregado
        mask_sem_enc = (
            df_atual["ENCARREGADO"].isna() | 
            (df_atual["ENCARREGADO"].astype(str).str.strip() == "") | 
            (df_atual["ENCARREGADO"].astype(str).str.strip().str.upper() == "NAN") |
            (df_atual["ENCARREGADO"].astype(str).str.strip() == "-") |
            (df_atual["ENCARREGADO"].astype(str).str.strip() == "0")
        )
        df_sem_encarregado = df_atual[mask_sem_enc]
        qtd_sem = len(df_sem_encarregado)
        
        # Métrica com destaque vermelho se houver
        if qtd_sem > 0:
            m4.metric("⚠️ Sem Encarregado", qtd_sem)
        else:
            m4.metric("✅ Sem Encarregado", 0)
        
        # Alerta e tabela expansível
        if qtd_sem > 0:
            st.warning(f"⚠️ **{qtd_sem} colaborador(es) estão SEM ENCARREGADO definido!** Clique abaixo para ver a lista.")
            with st.expander(f"👁️ Ver {qtd_sem} Colaboradores sem Encarregado", expanded=False):
                colunas_mostrar = ["MATRICULA", "NOME", "FUNÇÃO", "C.C", "DISCIPLINA", "STATUS"]
                colunas_existentes = [c for c in colunas_mostrar if c in df_sem_encarregado.columns]
                st.dataframe(
                    df_sem_encarregado[colunas_existentes].reset_index(drop=True),
                    use_container_width=True,
                    hide_index=True
                )
                
                # Botão para exportar a lista
                csv_sem_enc = df_sem_encarregado[colunas_existentes].to_csv(index=False).encode("utf-8")
                st.download_button(
                    "📥 Baixar Lista (CSV)",
                    data=csv_sem_enc,
                    file_name="colaboradores_sem_encarregado.csv",
                    mime="text/csv",
                    use_container_width=True
                )

        # ==============================================================
        # QUADRO DE FUNÇÕES POR DISCIPLINA
        # ==============================================================
        st.markdown("---")
        st.markdown("#### 📊 Quantidade de Funções por Disciplina")
        
        if "DISCIPLINA" in df_atual.columns and "FUNÇÃO" in df_atual.columns:
            # Tabela pivot: Disciplina x Função com contagem
            df_pivot = df_atual.groupby(["DISCIPLINA", "FUNÇÃO"]).size().reset_index(name="QTD")
            df_pivot = df_pivot.sort_values(["DISCIPLINA", "QTD"], ascending=[True, False])
            
            # Resumo por disciplina (total de pessoas e total de funções distintas)
            df_resumo_disc = df_atual.groupby("DISCIPLINA").agg(
                Total_Pessoas=("NOME", "count"),
                Total_Funcoes=("FUNÇÃO", "nunique")
            ).reset_index().sort_values("Total_Pessoas", ascending=False)
            
            # Mostrar tabela resumo
            st.dataframe(
                df_resumo_disc.rename(columns={
                    "DISCIPLINA": "Disciplina",
                    "Total_Pessoas": "👷 Pessoas",
                    "Total_Funcoes": "🔧 Funções Distintas"
                }),
                use_container_width=True,
                hide_index=True
            )
            
            # Detalhamento por disciplina (expansível)
            disciplinas_unicas = df_resumo_disc["DISCIPLINA"].tolist()
            
            with st.expander(f"🔍 Ver Detalhamento ({len(disciplinas_unicas)} disciplinas)", expanded=False):
                disc_selecionada = st.selectbox("Selecione a Disciplina:", disciplinas_unicas, key="sel_disc_detalhe")
                
                df_detalhe = df_pivot[df_pivot["DISCIPLINA"] == disc_selecionada][["FUNÇÃO", "QTD"]].reset_index(drop=True)
                total_disc = df_detalhe["QTD"].sum()
                
                st.markdown(f"**{disc_selecionada}** — {total_disc} pessoas em {len(df_detalhe)} função(ões)")
                st.dataframe(
                    df_detalhe.rename(columns={"FUNÇÃO": "Função", "QTD": "Quantidade"}),
                    use_container_width=True,
                    hide_index=True
                )
            
            # ========== BOTÃO BAIXAR EXCEL COM TABELA DINÂMICA ==========
            st.markdown("")
            if st.button("📥 Baixar Tabela Dinâmica em Excel", use_container_width=True, key="btn_baixar_pivot_excel"):
                try:
                    from openpyxl import Workbook
                    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
                    from openpyxl.utils import get_column_letter
                    
                    wb_pv = Workbook()
                    
                    # ---- ABA 1: TABELA DINÂMICA (Matriz Disciplina x Função) ----
                    ws_pivot = wb_pv.active
                    ws_pivot.title = "Tabela Dinamica"
                    
                    todas_funcoes = sorted(df_atual["FUNÇÃO"].dropna().unique())
                    todas_disciplinas = sorted(df_atual["DISCIPLINA"].dropna().unique())
                    
                    header_fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
                    header_font = Font(bold=True, color="FFFFFF", size=11)
                    total_fill = PatternFill(start_color="D6E4F0", end_color="D6E4F0", fill_type="solid")
                    total_font = Font(bold=True, size=11)
                    borda = Border(
                        left=Side(style="thin"), right=Side(style="thin"),
                        top=Side(style="thin"), bottom=Side(style="thin")
                    )
                    
                    # Título
                    ws_pivot.merge_cells(start_row=1, start_column=1, end_row=1, end_column=len(todas_disciplinas) + 2)
                    ws_pivot.cell(row=1, column=1, value="TABELA DINAMICA - FUNCOES POR DISCIPLINA")
                    ws_pivot.cell(row=1, column=1).font = Font(bold=True, size=14, color="1F4E79")
                    ws_pivot.cell(row=1, column=1).alignment = Alignment(horizontal="center")
                    
                    # Cabeçalhos
                    ws_pivot.cell(row=3, column=1, value="FUNCAO")
                    ws_pivot.cell(row=3, column=1).font = header_font
                    ws_pivot.cell(row=3, column=1).fill = header_fill
                    ws_pivot.cell(row=3, column=1).border = borda
                    ws_pivot.column_dimensions["A"].width = 40
                    
                    for j, disc in enumerate(todas_disciplinas):
                        col = j + 2
                        ws_pivot.cell(row=3, column=col, value=disc)
                        ws_pivot.cell(row=3, column=col).font = header_font
                        ws_pivot.cell(row=3, column=col).fill = header_fill
                        ws_pivot.cell(row=3, column=col).alignment = Alignment(horizontal="center")
                        ws_pivot.cell(row=3, column=col).border = borda
                        ws_pivot.column_dimensions[get_column_letter(col)].width = 18
                    
                    col_total = len(todas_disciplinas) + 2
                    ws_pivot.cell(row=3, column=col_total, value="TOTAL")
                    ws_pivot.cell(row=3, column=col_total).font = header_font
                    ws_pivot.cell(row=3, column=col_total).fill = PatternFill(start_color="C00000", end_color="C00000", fill_type="solid")
                    ws_pivot.cell(row=3, column=col_total).alignment = Alignment(horizontal="center")
                    ws_pivot.cell(row=3, column=col_total).border = borda
                    ws_pivot.column_dimensions[get_column_letter(col_total)].width = 12
                    
                    # Dados
                    lookup = df_pivot.set_index(["FUNÇÃO", "DISCIPLINA"])["QTD"]
                    
                    for i, func in enumerate(todas_funcoes):
                        row = i + 4
                        ws_pivot.cell(row=row, column=1, value=func)
                        ws_pivot.cell(row=row, column=1).border = borda
                        total_linha = 0
                        
                        for j, disc in enumerate(todas_disciplinas):
                            col = j + 2
                            try:
                                qtd = int(lookup.get((func, disc), 0))
                            except Exception:
                                qtd = 0
                            if qtd > 0:
                                ws_pivot.cell(row=row, column=col, value=qtd)
                            ws_pivot.cell(row=row, column=col).alignment = Alignment(horizontal="center")
                            ws_pivot.cell(row=row, column=col).border = borda
                            total_linha += qtd
                        
                        ws_pivot.cell(row=row, column=col_total, value=total_linha)
                        ws_pivot.cell(row=row, column=col_total).font = Font(bold=True)
                        ws_pivot.cell(row=row, column=col_total).alignment = Alignment(horizontal="center")
                        ws_pivot.cell(row=row, column=col_total).border = borda
                    
                    # Linha TOTAL
                    row_total = len(todas_funcoes) + 4
                    ws_pivot.cell(row=row_total, column=1, value="TOTAL GERAL")
                    ws_pivot.cell(row=row_total, column=1).font = total_font
                    ws_pivot.cell(row=row_total, column=1).fill = total_fill
                    ws_pivot.cell(row=row_total, column=1).border = borda
                    grande_total = 0
                    
                    for j, disc in enumerate(todas_disciplinas):
                        col = j + 2
                        soma_col = sum(int(lookup.get((f, disc), 0)) for f in todas_funcoes)
                        ws_pivot.cell(row=row_total, column=col, value=soma_col)
                        ws_pivot.cell(row=row_total, column=col).font = total_font
                        ws_pivot.cell(row=row_total, column=col).fill = total_fill
                        ws_pivot.cell(row=row_total, column=col).alignment = Alignment(horizontal="center")
                        ws_pivot.cell(row=row_total, column=col).border = borda
                        grande_total += soma_col
                    
                    ws_pivot.cell(row=row_total, column=col_total, value=grande_total)
                    ws_pivot.cell(row=row_total, column=col_total).font = Font(bold=True, color="FFFFFF", size=12)
                    ws_pivot.cell(row=row_total, column=col_total).fill = PatternFill(start_color="C00000", end_color="C00000", fill_type="solid")
                    ws_pivot.cell(row=row_total, column=col_total).alignment = Alignment(horizontal="center")
                    ws_pivot.cell(row=row_total, column=col_total).border = borda
                    
                    # ---- ABA 2: DETALHADO ----
                    ws_det = wb_pv.create_sheet("Detalhado")
                    ws_det.append(["DISCIPLINA", "FUNCAO", "QUANTIDADE"])
                    for cell in ws_det[1]:
                        cell.font = header_font
                        cell.fill = header_fill
                        cell.border = borda
                    for _, r in df_pivot.iterrows():
                        ws_det.append([r["DISCIPLINA"], r["FUNÇÃO"], r["QTD"]])
                    ws_det.column_dimensions["A"].width = 30
                    ws_det.column_dimensions["B"].width = 40
                    ws_det.column_dimensions["C"].width = 15
                    
                    # ---- ABA 3: RESUMO ----
                    ws_res = wb_pv.create_sheet("Resumo")
                    ws_res.append(["DISCIPLINA", "TOTAL PESSOAS", "FUNCOES DISTINTAS"])
                    for cell in ws_res[1]:
                        cell.font = header_font
                        cell.fill = header_fill
                        cell.border = borda
                    for _, r in df_resumo_disc.iterrows():
                        ws_res.append([r["DISCIPLINA"], r["Total_Pessoas"], r["Total_Funcoes"]])
                    ws_res.column_dimensions["A"].width = 30
                    ws_res.column_dimensions["B"].width = 18
                    ws_res.column_dimensions["C"].width = 22
                    
                    buf_pivot = io.BytesIO()
                    wb_pv.save(buf_pivot)
                    wb_pv.close()
                    buf_pivot.seek(0)
                    
                    st.download_button(
                        label="⬇️ Clique aqui para Baixar",
                        data=buf_pivot,
                        file_name=f"Tabela_Dinamica_Funcoes_{datetime.datetime.now().strftime('%d_%m_%Y')}.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        use_container_width=True,
                        key="btn_download_pivot_file"
                    )
                    st.success("✅ Planilha gerada com 3 abas: Tabela Dinâmica, Detalhado e Resumo!")
                except Exception as e:
                    st.error(f"Erro ao gerar Excel: {e}")
        else:
            st.info("ℹ️ As colunas DISCIPLINA e FUNÇÃO não foram encontradas no PDE.")
        
        st.markdown("---")

        # Editable data
        st.markdown("#### Base Atual")
        colunas_pde = ["MATRICULA", "NOME", "FUNÇÃO", "C.C", "ENCARREGADO", "TURNO", "STATUS", "DISCIPLINA", "MÃO DE OBRA"]
        colunas_pde_existentes = [c for c in colunas_pde if c in df_pde_filtrado.columns]
        df_editado_pde = st.data_editor(
            df_pde_filtrado[colunas_pde_existentes],
            key="editor_pde",
            use_container_width=True
        )

        # Add employee
        with st.expander("➕ Adicionar Funcionário"):
            with st.form("form_add_func"):
                c1, c2, c3 = st.columns(3)
                with c1:
                    new_mat = st.text_input("MATRICULA")
                    new_nome = st.text_input("NOME")
                    new_funcao = st.text_input("FUNÇÃO")
                with c2:
                    new_cc = st.text_input("C.C")
                    new_enc = st.text_input("ENCARREGADO")
                    new_turno = st.text_input("TURNO")
                with c3:
                    new_status = st.text_input("STATUS", value="ATIVO")
                    new_disc = st.text_input("DISCIPLINA")
                    new_mo = st.text_input("MÃO DE OBRA")

                submitted_add = st.form_submit_button("➕ Adicionar ao Sistema")
                if submitted_add:
                    try:
                        novo_dado = pd.DataFrame([{
                            "MATRICULA": new_mat, "NOME": new_nome, "FUNÇÃO": new_funcao,
                            "C.C": new_cc, "ENCARREGADO": new_enc, "TURNO": new_turno,
                            "STATUS": new_status, "DISCIPLINA": new_disc, "MÃO DE OBRA": new_mo
                        }])
                        st.session_state.df = pd.concat([df_atual, novo_dado], ignore_index=True)
                        st.session_state.df.to_csv(caminho_base_salva_csv, index=False)
                        st.success("✅ Funcionário adicionado com sucesso!")
                        time.sleep(1)
                        st.rerun()
                    except Exception as e:
                        st.error(f"Erro ao adicionar: {e}")

        # Action Buttons
        b1, b2, b3 = st.columns(3)
        with b1:
            if st.button("💾 Salvar Alterações na PDE", key="btn_salvar_pde", type="primary", use_container_width=True):
                try:
                    # Atualizar o DataFrame principal com as edições
                    for col in colunas_pde_existentes:
                        if col in df_editado_pde.columns:
                            df_pde_filtrado[col] = df_editado_pde[col].values
                    # Aplicar de volta no df_atual
                    df_atual.update(df_pde_filtrado)
                    df_atual.to_csv(caminho_base_salva_csv, index=False)
                    st.session_state.df = df_atual
                    try:
                        if conn:
                            conn.update(worksheet="PDE", data=df_atual)
                            st.toast("☁️ Sincronizado com Google Sheets!", icon="✅")
                    except Exception as e_gs:
                        st.error(f"⚠️ Erro ao salvar no Google Sheets: {e_gs}")
                    st.success("✅ Alterações salvas com sucesso!")
                except Exception as e:
                    st.error(f"Erro ao salvar: {e}")

        with b2:
            if st.button("🔄 Puxar do Google Sheets", key="btn_sync_pde", use_container_width=True):
                try:
                    if conn:
                        df_gs = conn.read(worksheet="PDE", ttl=0)
                        df_gs = df_gs.dropna(how='all')
                        df_gs.to_csv(caminho_base_salva_csv, index=False)
                        st.success("✅ Dados sincronizados do Google Sheets!")
                        time.sleep(1)
                        st.rerun()
                    else:
                        st.warning("Conexão com Google Sheets não disponível.")
                except Exception as e:
                    st.error(f"Erro ao sincronizar: {e}")

        with b3:
            try:
                buffer_pde = io.BytesIO()
                with pd.ExcelWriter(buffer_pde, engine='openpyxl') as writer:
                    df_pde_filtrado.to_excel(writer, index=False, sheet_name='PDE')
                st.download_button(
                    label="⬇️ Exportar Excel",
                    data=buffer_pde.getvalue(),
                    file_name="Base_PDE_Filtrada.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    key="btn_export_pde",
                    use_container_width=True
                )
            except Exception as e:
                st.error(f"Erro ao gerar Excel: {e}")

    # ==============================================================
    # ABA 11: BANCO DE RDCs
    # ==============================================================
    with tab_banco_rdc:
        st.markdown("### 📑 Banco de RDCs — Histórico Completo")
        st.markdown("Todos os RDCs lidos pela IA, salvos permanentemente no sistema.")

        try:
            if os.path.exists(caminho_rdc_registros_csv):
                df_rdc_banco = pd.read_csv(caminho_rdc_registros_csv)

                if not df_rdc_banco.empty:
                    # Filters
                    c1, c2, c3 = st.columns(3)
                    with c1:
                        data_inicio_rdc = st.date_input("De:", value=None, key="rdc_data_de")
                    with c2:
                        data_fim_rdc = st.date_input("Até:", value=None, key="rdc_data_ate")
                    with c3:
                        encarregados_rdc = ["Todos"] + sorted(list(df_rdc_banco["ENCARREGADO"].dropna().unique()))
                        enc_filtro_rdc = st.selectbox("Encarregado", options=encarregados_rdc, key="rdc_enc_filtro")

                    c4, c5, c6 = st.columns(3)
                    with c4:
                        disc_rdc = ["Todos"] + sorted(list(df_rdc_banco["DISCIPLINA"].dropna().unique()))
                        disc_filtro_rdc = st.selectbox("Disciplina", options=disc_rdc, key="rdc_disc_filtro")
                    with c5:
                        if "CALDEIRA" in df_rdc_banco.columns:
                            caldeira_rdc = ["Todos"] + sorted([x for x in df_rdc_banco["CALDEIRA"].dropna().unique() if str(x).strip()])
                        else:
                            caldeira_rdc = ["Todos"]
                        cald_filtro_rdc = st.selectbox("Caldeira", options=caldeira_rdc, key="rdc_cald_filtro")
                    with c6:
                        busca_atividade = st.text_input("🔎 Buscar na atividade", key="rdc_busca")

                    # Apply filters
                    df_rdc_filtrado = df_rdc_banco.copy()
                    if "DATA" in df_rdc_filtrado.columns and (data_inicio_rdc or data_fim_rdc):
                        df_rdc_filtrado["_DATA_DT"] = pd.to_datetime(df_rdc_filtrado["DATA"], errors='coerce')
                        if data_inicio_rdc:
                            df_rdc_filtrado = df_rdc_filtrado[df_rdc_filtrado["_DATA_DT"].dt.date >= data_inicio_rdc]
                        if data_fim_rdc:
                            df_rdc_filtrado = df_rdc_filtrado[df_rdc_filtrado["_DATA_DT"].dt.date <= data_fim_rdc]
                        df_rdc_filtrado = df_rdc_filtrado.drop(columns=["_DATA_DT"])

                    if enc_filtro_rdc != "Todos":
                        df_rdc_filtrado = df_rdc_filtrado[df_rdc_filtrado["ENCARREGADO"] == enc_filtro_rdc]
                    if disc_filtro_rdc != "Todos":
                        df_rdc_filtrado = df_rdc_filtrado[df_rdc_filtrado["DISCIPLINA"] == disc_filtro_rdc]
                    if cald_filtro_rdc != "Todos" and "CALDEIRA" in df_rdc_filtrado.columns:
                        df_rdc_filtrado = df_rdc_filtrado[df_rdc_filtrado["CALDEIRA"] == cald_filtro_rdc]
                    if busca_atividade and "ATIVIDADE" in df_rdc_filtrado.columns:
                        df_rdc_filtrado = df_rdc_filtrado[df_rdc_filtrado["ATIVIDADE"].astype(str).str.contains(busca_atividade, case=False, na=False)]

                    # Summary
                    st.caption(f"📊 Mostrando **{len(df_rdc_filtrado)}** de **{len(df_rdc_banco)}** RDCs")

                    # Data Editor
                    cols_to_show = [c for c in ["DATA", "DISCIPLINA", "ENCARREGADO", "TURNO", "CALDEIRA", "ATIVIDADE", "DDS", "PROBLEMAS"] if c in df_rdc_filtrado.columns]
                    df_editado_rdc = st.data_editor(
                        df_rdc_filtrado[cols_to_show] if cols_to_show else df_rdc_filtrado,
                        key="editor_banco_rdc",
                        use_container_width=True
                    )

                    with st.expander("📝 Ver Transcrições Completas"):
                        if "TRANSCRICAO" in df_rdc_filtrado.columns:
                            for _, row in df_rdc_filtrado.iterrows():
                                st.markdown(f"**{row.get('ENCARREGADO', '?')}** — {row.get('DATA', '?')}")
                                st.text(str(row.get('TRANSCRICAO', '')))
                                st.markdown("---")
                        else:
                            st.info("Coluna de transcrição não disponível.")

                    # Action buttons
                    b1, b2, b3 = st.columns(3)
                    with b1:
                        if st.button("💾 Salvar Edições", key="btn_salvar_rdc", type="primary", use_container_width=True):
                            try:
                                # Aplicar edições de volta no dataframe completo
                                for col in cols_to_show:
                                    if col in df_editado_rdc.columns:
                                        df_rdc_filtrado[col] = df_editado_rdc[col].values
                                df_rdc_banco.update(df_rdc_filtrado)
                                df_rdc_banco.to_csv(caminho_rdc_registros_csv, index=False)
                                st.success("✅ Registros atualizados!")
                            except Exception as e:
                                st.error(f"Erro ao salvar: {e}")
                    with b2:
                        try:
                            buffer_rdc_xls = io.BytesIO()
                            with pd.ExcelWriter(buffer_rdc_xls, engine='openpyxl') as writer:
                                df_rdc_filtrado.to_excel(writer, index=False, sheet_name='RDCs')
                            st.download_button(
                                label="⬇️ Excel",
                                data=buffer_rdc_xls.getvalue(),
                                file_name="Banco_RDC_Filtrado.xlsx",
                                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                key="btn_export_rdc_xls",
                                use_container_width=True
                            )
                        except Exception as e:
                            st.error(f"Erro ao gerar Excel: {e}")
                    with b3:
                        csv_rdc = df_rdc_filtrado.to_csv(index=False).encode('utf-8')
                        st.download_button(
                            label="⬇️ CSV",
                            data=csv_rdc,
                            file_name="Banco_RDC_Filtrado.csv",
                            mime="text/csv",
                            key="btn_export_rdc_csv",
                            use_container_width=True
                        )

                    with st.expander("🗑️ Excluir Registros"):
                        st.warning("⚠️ Atenção: Esta ação é irreversível.")
                        id_delete = st.number_input("Índice do registro para excluir:", min_value=0, max_value=max(0, len(df_rdc_banco)-1), step=1, key="rdc_id_delete")
                        if st.button("🗑️ Confirmar Exclusão", type="primary", key="btn_delete_rdc"):
                            try:
                                df_rdc_banco = df_rdc_banco.drop(index=id_delete).reset_index(drop=True)
                                df_rdc_banco.to_csv(caminho_rdc_registros_csv, index=False)
                                st.success(f"✅ Registro {id_delete} excluído com sucesso!")
                                time.sleep(1)
                                st.rerun()
                            except Exception as e:
                                st.error(f"Erro ao excluir: {e}")
                else:
                    st.info("📭 Banco de RDCs vazio. Processe alguns RDCs na aba 'Leitor de RDC (IA)' e clique em 'Confirmar e Salvar' para popular esta tabela.")
            else:
                st.info("📭 Nenhum RDC salvo ainda. Processe RDCs na aba 'Leitor de RDC (IA)' e clique em 'Confirmar e Salvar' para começar a guardar.")
        except Exception as e:
            st.error(f"Erro ao carregar banco de RDCs: {e}")

    # ==============================================================
    # ABA 12: ANÁLISE DE GARGALOS
    # ==============================================================
    with tab_gargalos:
        st.markdown("### 🔍 Análise de Gargalos — Inteligência dos RDCs")
        st.markdown("Dashboard analítico gerado automaticamente a partir dos RDCs processados pela IA. Identifique os maiores gargalos da obra em segundos.")

        try:
            if os.path.exists(caminho_rdc_registros_csv):
                df_gargalos = pd.read_csv(caminho_rdc_registros_csv)

                if not df_gargalos.empty and len(df_gargalos) > 0:
                    # === PREPARAR COLUNA DE DATA ===
                    df_gargalos["_DATA_DT"] = pd.to_datetime(df_gargalos["DATA"], errors='coerce')

                    # === FILTROS ===
                    st.markdown("#### 🎛️ Filtros")
                    datas_disponiveis_g = sorted(df_gargalos["_DATA_DT"].dropna().dt.strftime("%d/%m/%Y").unique().tolist(), key=lambda x: pd.to_datetime(x, format="%d/%m/%Y"), reverse=True)
                    col_f1, col_f2 = st.columns(2)
                    with col_f1:
                        datas_sel_g = st.multiselect("📅 Selecione a(s) Data(s):", datas_disponiveis_g, default=[], key="garg_datas", help="Deixe vazio para ver todos os RDCs")
                    with col_f2:
                        disc_opcoes_g = sorted(df_gargalos["DISCIPLINA"].dropna().unique().tolist())
                        disc_sel_g = st.multiselect("Disciplina:", disc_opcoes_g, default=[], key="garg_disc")

                    df_g = df_gargalos.copy()
                    if datas_sel_g:
                        datas_filtro = [pd.to_datetime(d, format="%d/%m/%Y").date() for d in datas_sel_g]
                        df_g = df_g[df_g["_DATA_DT"].dt.date.isin(datas_filtro)]
                    if disc_sel_g:
                        df_g = df_g[df_g["DISCIPLINA"].isin(disc_sel_g)]

                    # === CATEGORIZAR PROBLEMAS ===
                    def categorizar_problema(texto):
                        if pd.isna(texto) or str(texto).strip() == "" or str(texto).strip().lower() in ["nan", "não informado", "nenhum", "n/a", "nao informado", "sem problemas", "-", "nenhum problema"]:
                            return None
                        texto = str(texto).upper()
                        if any(p in texto for p in ["CHUVA", "TEMPORAL", "INTEMPÉRIE", "INTEMPERIE", "TEMPO", "CLIMÁT", "CLIMAT"]):
                            return "☁️ Chuva / Intempérie"
                        if any(p in texto for p in ["MATERIAL", "CONSUMÍVEL", "CONSUMIVEL", "ELETRODO", "FALTA DE PEÇA", "FALTA DE PECA", "INSUMO", "SUPRIMENTO"]):
                            return "🔧 Falta de Material"
                        if any(p in texto for p in ["EQUIPAMENTO", "MÁQUINA", "MAQUINA", "QUEBRA", "DEFEITO", "PANE", "MANUTENÇÃO", "MANUTENCAO", "GUINDASTE", "GUINCHO", "MUNCK"]):
                            return "⚡ Quebra / Falta de Equipamento"
                        if any(p in texto for p in ["EFETIVO", "MÃO DE OBRA", "MAO DE OBRA", "FALTA DE PESSOAL", "FALTA PESSOAL", "ABSENTEÍSMO", "ABSENTEISMO", "ATESTADO", "AFASTADO"]):
                            return "👷 Falta de Efetivo"
                        if any(p in texto for p in ["ACESSO", "LIBERAÇÃO", "LIBERACAO", "PERMISSÃO", "PERMISSAO", "ANDAIME", "BLOQUEIO", "ISOLAMENTO", "SEGURANÇA", "SEGURANCA", "PT"]):
                            return "🚫 Falta de Acesso / Liberação"
                        if any(p in texto for p in ["ENERGIA", "ELÉTRIC", "ELETRIC", "SOLDA", "REDE", "TOMADA", "EXTENSÃO", "EXTENSAO"]):
                            return "⚡ Energia / Elétrica"
                        if any(p in texto for p in ["PROJETO", "DESENHO", "INFORMAÇÃO", "INFORMACAO", "ENGENHARIA"]):
                            return "📐 Falta de Projeto / Info"
                        return "📋 Outros"

                    df_g["_CATEGORIA_PROB"] = df_g["PROBLEMAS"].apply(categorizar_problema)
                    df_com_problema = df_g[df_g["_CATEGORIA_PROB"].notna()].copy()
                    total_rdcs = len(df_g)
                    total_com_prob = len(df_com_problema)
                    pct_prob = (total_com_prob / total_rdcs * 100) if total_rdcs > 0 else 0

                    # === MÉTRICAS ===
                    st.markdown("---")
                    col_m1, col_m2, col_m3, col_m4 = st.columns(4)
                    with col_m1:
                        st.metric("📑 RDCs Analisados", f"{total_rdcs}")
                    with col_m2:
                        st.metric("⚠️ RDCs com Problemas", f"{total_com_prob}")
                    with col_m3:
                        st.metric("📊 % com Problemas", f"{pct_prob:.1f}%")
                    with col_m4:
                        if not df_com_problema.empty and "DISCIPLINA" in df_com_problema.columns:
                            disc_top = df_com_problema["DISCIPLINA"].value_counts().idxmax()
                            st.metric("🏗️ Disciplina Crítica", disc_top)
                        else:
                            st.metric("🏗️ Disciplina Crítica", "N/A")

                    st.markdown("---")

                    # === GRÁFICOS ===
                    if not df_com_problema.empty:
                        col_g1, col_g2 = st.columns(2)

                        # PIZZA: Categorias de Problemas
                        with col_g1:
                            st.markdown("#### 🥧 Categorias de Problemas")
                            cat_counts = df_com_problema["_CATEGORIA_PROB"].value_counts().reset_index()
                            cat_counts.columns = ["Categoria", "Quantidade"]
                            fig_pie = px.pie(
                                cat_counts,
                                names="Categoria",
                                values="Quantidade",
                                hole=0.4,
                                color_discrete_sequence=px.colors.qualitative.Set2
                            )
                            fig_pie.update_layout(
                                paper_bgcolor="rgba(0,0,0,0)",
                                plot_bgcolor="rgba(0,0,0,0)",
                                font_color="#e0e4ea",
                                legend=dict(font=dict(size=11)),
                                margin=dict(t=30, b=30, l=10, r=10),
                                height=400
                            )
                            st.plotly_chart(fig_pie, use_container_width=True)

                        # BARRAS: Problemas por Disciplina
                        with col_g2:
                            st.markdown("#### 📊 Problemas por Disciplina")
                            disc_prob = df_com_problema.groupby("DISCIPLINA")["_CATEGORIA_PROB"].count().reset_index()
                            disc_prob.columns = ["Disciplina", "Ocorrências"]
                            disc_prob = disc_prob.sort_values("Ocorrências", ascending=True)
                            fig_bar = px.bar(
                                disc_prob,
                                x="Ocorrências",
                                y="Disciplina",
                                orientation="h",
                                color="Ocorrências",
                                color_continuous_scale=["#0ea5e9", "#f59e0b", "#ef4444"]
                            )
                            fig_bar.update_layout(
                                paper_bgcolor="rgba(0,0,0,0)",
                                plot_bgcolor="rgba(0,0,0,0)",
                                font_color="#e0e4ea",
                                showlegend=False,
                                margin=dict(t=30, b=30, l=10, r=10),
                                height=400,
                                yaxis=dict(tickfont=dict(size=11))
                            )
                            fig_bar.update_coloraxes(showscale=False)
                            st.plotly_chart(fig_bar, use_container_width=True)

                        # LINHA: Evolução Temporal de Problemas
                        st.markdown("#### 📈 Evolução de Problemas ao Longo do Tempo")
                        df_tempo = df_com_problema.copy()
                        df_tempo["_SEMANA"] = df_tempo["_DATA_DT"].dt.to_period("W").apply(lambda r: r.start_time)
                        evolucao = df_tempo.groupby("_SEMANA").size().reset_index(name="Problemas")
                        evolucao.columns = ["Semana", "Problemas"]
                        if len(evolucao) > 1:
                            fig_line = px.area(
                                evolucao,
                                x="Semana",
                                y="Problemas",
                                markers=True,
                                color_discrete_sequence=["#0ea5e9"]
                            )
                            fig_line.update_layout(
                                paper_bgcolor="rgba(0,0,0,0)",
                                plot_bgcolor="rgba(0,0,0,0)",
                                font_color="#e0e4ea",
                                margin=dict(t=30, b=30, l=10, r=10),
                                height=350,
                                xaxis_title="Semana",
                                yaxis_title="Nº de Problemas"
                            )
                            st.plotly_chart(fig_line, use_container_width=True)
                        else:
                            st.info("📅 Dados insuficientes para gerar a evolução temporal. Continue processando RDCs para ver a tendência.")

                        # BARRAS: Categorias cruzadas por Disciplina (Stacked)
                        st.markdown("#### 🧩 Mapa de Calor: Problema × Disciplina")
                        cross = df_com_problema.groupby(["DISCIPLINA", "_CATEGORIA_PROB"]).size().reset_index(name="Qtd")
                        cross.columns = ["Disciplina", "Categoria", "Qtd"]
                        fig_stack = px.bar(
                            cross,
                            x="Disciplina",
                            y="Qtd",
                            color="Categoria",
                            barmode="stack",
                            color_discrete_sequence=px.colors.qualitative.Set2
                        )
                        fig_stack.update_layout(
                            paper_bgcolor="rgba(0,0,0,0)",
                            plot_bgcolor="rgba(0,0,0,0)",
                            font_color="#e0e4ea",
                            margin=dict(t=30, b=30, l=10, r=10),
                            height=400,
                            legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1, font=dict(size=10))
                        )
                        st.plotly_chart(fig_stack, use_container_width=True)

                        # TABELA: Top 10 Problemas Recentes
                        st.markdown("#### 📋 Últimos Problemas Reportados")
                        df_recentes = df_com_problema.sort_values("_DATA_DT", ascending=False).head(10)
                        cols_mostrar = ["DATA", "ENCARREGADO", "DISCIPLINA", "_CATEGORIA_PROB", "PROBLEMAS"]
                        cols_existe = [c for c in cols_mostrar if c in df_recentes.columns]
                        df_mostrar = df_recentes[cols_existe].copy()
                        df_mostrar = df_mostrar.rename(columns={"_CATEGORIA_PROB": "CATEGORIA"})
                        st.dataframe(df_mostrar, hide_index=True, use_container_width=True)

                    else:
                        st.success("🎉 Nenhum problema identificado nos RDCs filtrados! Excelente resultado operacional.")

                    # === NUVEM DE ATIVIDADES (Barras Horizontais) ===
                    st.markdown("---")
                    st.markdown("#### 🏗️ Atividades Mais Frequentes nos RDCs")
                    if "ATIVIDADE" in df_g.columns:
                        import re as re_mod
                        stopwords = {"DE", "DO", "DA", "DOS", "DAS", "E", "EM", "NO", "NA", "NOS", "NAS", "COM", "PARA", "POR", "UM", "UMA", "O", "A", "OS", "AS", "AO", "À", "SE", "QUE", "SÃO", "FOI", "SER", "TER", "ESTÁ", "COMO"}
                        all_text = " ".join(df_g["ATIVIDADE"].dropna().astype(str).tolist()).upper()
                        words = re_mod.findall(r"[A-ZÁÉÍÓÚÂÊÔÃÕÇ]{4,}", all_text)
                        words = [w for w in words if w not in stopwords and len(w) > 3]
                        if words:
                            from collections import Counter
                            word_counts = Counter(words).most_common(15)
                            df_words = pd.DataFrame(word_counts, columns=["Atividade", "Frequência"])
                            df_words = df_words.sort_values("Frequência", ascending=True)
                            fig_words = px.bar(
                                df_words,
                                x="Frequência",
                                y="Atividade",
                                orientation="h",
                                color="Frequência",
                                color_continuous_scale=["#22c55e", "#0ea5e9", "#8b5cf6"]
                            )
                            fig_words.update_layout(
                                paper_bgcolor="rgba(0,0,0,0)",
                                plot_bgcolor="rgba(0,0,0,0)",
                                font_color="#e0e4ea",
                                showlegend=False,
                                margin=dict(t=10, b=30, l=10, r=10),
                                height=450,
                                yaxis=dict(tickfont=dict(size=12))
                            )
                            fig_words.update_coloraxes(showscale=False)
                            st.plotly_chart(fig_words, use_container_width=True)
                        else:
                            st.info("Sem dados de atividades suficientes para gerar o gráfico.")
                    else:
                        st.info("Coluna ATIVIDADE não encontrada nos dados.")

                    # === DDS: Temas mais abordados ===
                    st.markdown("#### 🦺 Temas de DDS Mais Abordados")
                    if "DDS" in df_g.columns:
                        dds_validos = df_g["DDS"].dropna().astype(str)
                        dds_validos = dds_validos[~dds_validos.str.strip().str.lower().isin(["nan", "", "não informado", "nao informado", "-", "n/a"])]
                        if len(dds_validos) > 0:
                            dds_counts = dds_validos.str.upper().value_counts().head(10).reset_index()
                            dds_counts.columns = ["Tema DDS", "Frequência"]
                            fig_dds = px.bar(
                                dds_counts,
                                x="Frequência",
                                y="Tema DDS",
                                orientation="h",
                                color_discrete_sequence=["#f59e0b"]
                            )
                            fig_dds.update_layout(
                                paper_bgcolor="rgba(0,0,0,0)",
                                plot_bgcolor="rgba(0,0,0,0)",
                                font_color="#e0e4ea",
                                showlegend=False,
                                margin=dict(t=10, b=30, l=10, r=10),
                                height=350,
                                yaxis=dict(tickfont=dict(size=11))
                            )
                            st.plotly_chart(fig_dds, use_container_width=True)
                        else:
                            st.info("Nenhum tema de DDS registrado nos RDCs filtrados.")
                    else:
                        st.info("Coluna DDS não encontrada nos dados.")

                    # === BOTÃO: GERAR RELATÓRIO EXECUTIVO EM POWERPOINT ===
                    st.markdown("---")
                    st.markdown("#### 📑 Relatório Executivo")
                    if st.button("📑 GERAR APRESENTAÇÃO POWERPOINT (.PPTX)", type="primary", use_container_width=True, key="btn_gerar_pptx"):
                        with st.spinner("Gerando apresentação executiva..."):
                            try:
                                from pptx import Presentation as PptxPresentation
                                from pptx.util import Inches, Pt, Emu
                                from pptx.dml.color import RGBColor
                                from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
                                from pptx.enum.chart import XL_CHART_TYPE
                                
                                prs = PptxPresentation()
                                prs.slide_width = Inches(13.333)
                                prs.slide_height = Inches(7.5)
                                
                                # Cores padrão
                                COR_FUNDO = RGBColor(15, 23, 42)
                                COR_AZUL = RGBColor(14, 165, 233)
                                COR_BRANCO = RGBColor(224, 228, 234)
                                COR_CINZA = RGBColor(148, 163, 184)
                                COR_VERDE = RGBColor(34, 197, 94)
                                COR_AMARELO = RGBColor(245, 158, 11)
                                COR_VERMELHO = RGBColor(239, 68, 68)
                                COR_CARD_BG = RGBColor(30, 41, 59)
                                
                                def set_slide_bg(slide, cor):
                                    background = slide.background
                                    fill = background.fill
                                    fill.solid()
                                    fill.fore_color.rgb = cor
                                
                                def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=COR_BRANCO, alignment=PP_ALIGN.LEFT):
                                    txBox = slide.shapes.add_textbox(left, top, width, height)
                                    tf = txBox.text_frame
                                    tf.word_wrap = True
                                    p = tf.paragraphs[0]
                                    p.text = text
                                    p.font.size = Pt(font_size)
                                    p.font.bold = bold
                                    p.font.color.rgb = color
                                    p.alignment = alignment
                                    return txBox
                                
                                def add_card(slide, left, top, width, height, titulo, valor, cor_valor=COR_AZUL):
                                    shape = slide.shapes.add_shape(1, left, top, width, height)
                                    shape.fill.solid()
                                    shape.fill.fore_color.rgb = COR_CARD_BG
                                    shape.line.color.rgb = RGBColor(51, 65, 85)
                                    shape.line.width = Pt(1)
                                    
                                    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4))
                                    tf = txBox.text_frame
                                    p = tf.paragraphs[0]
                                    p.text = titulo
                                    p.font.size = Pt(11)
                                    p.font.color.rgb = COR_CINZA
                                    
                                    txBox2 = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), Inches(0.6))
                                    tf2 = txBox2.text_frame
                                    p2 = tf2.paragraphs[0]
                                    p2.text = str(valor)
                                    p2.font.size = Pt(28)
                                    p2.font.bold = True
                                    p2.font.color.rgb = cor_valor
                                
                                # ============================================
                                # SLIDE 1: CAPA
                                # ============================================
                                slide1 = prs.slides.add_slide(prs.slide_layouts[6])
                                set_slide_bg(slide1, COR_FUNDO)
                                
                                # Logo
                                if os.path.exists(caminho_logo):
                                    try:
                                        slide1.shapes.add_picture(caminho_logo, Inches(5.4), Inches(0.8), height=Inches(1.5))
                                    except:
                                        pass
                                
                                # Linha decorativa
                                line = slide1.shapes.add_shape(1, Inches(3), Inches(2.8), Inches(7.333), Inches(0.04))
                                line.fill.solid()
                                line.fill.fore_color.rgb = COR_AZUL
                                line.line.fill.background()
                                
                                add_text_box(slide1, Inches(1.5), Inches(3.0), Inches(10.333), Inches(1.2),
                                    f"RELATÓRIO EXECUTIVO DE OBRA", font_size=36, bold=True, color=COR_BRANCO, alignment=PP_ALIGN.CENTER)
                                add_text_box(slide1, Inches(1.5), Inches(4.0), Inches(10.333), Inches(0.8),
                                    nome_site, font_size=24, bold=False, color=COR_AZUL, alignment=PP_ALIGN.CENTER)
                                
                                datas_texto = ""
                                if datas_sel_g:
                                    datas_texto = ", ".join(datas_sel_g)
                                else:
                                    datas_texto = f"Todos os dados até {datetime.datetime.now().strftime('%d/%m/%Y')}"
                                add_text_box(slide1, Inches(1.5), Inches(4.8), Inches(10.333), Inches(0.6),
                                    f"Período: {datas_texto}", font_size=16, color=COR_CINZA, alignment=PP_ALIGN.CENTER)
                                add_text_box(slide1, Inches(1.5), Inches(5.5), Inches(10.333), Inches(0.5),
                                    f"Gerado em: {datetime.datetime.now().strftime('%d/%m/%Y às %H:%M')}", font_size=12, color=COR_CINZA, alignment=PP_ALIGN.CENTER)
                                
                                # ============================================
                                # SLIDE 2: MÉTRICAS GERAIS
                                # ============================================
                                slide2 = prs.slides.add_slide(prs.slide_layouts[6])
                                set_slide_bg(slide2, COR_FUNDO)
                                
                                add_text_box(slide2, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                                    "INDICADORES GERAIS", font_size=28, bold=True, color=COR_BRANCO)
                                
                                line2 = slide2.shapes.add_shape(1, Inches(0.5), Inches(1.0), Inches(12.333), Inches(0.03))
                                line2.fill.solid()
                                line2.fill.fore_color.rgb = COR_AZUL
                                line2.line.fill.background()
                                
                                card_w = Inches(2.8)
                                card_h = Inches(1.2)
                                card_y = Inches(1.4)
                                gap = Inches(0.3)
                                start_x = Inches(0.7)
                                
                                disc_top_val = "N/A"
                                if not df_com_problema.empty and "DISCIPLINA" in df_com_problema.columns:
                                    try:
                                        disc_top_val = df_com_problema["DISCIPLINA"].value_counts().idxmax()
                                    except:
                                        pass
                                
                                add_card(slide2, start_x, card_y, card_w, card_h, "RDCs Analisados", total_rdcs, COR_AZUL)
                                add_card(slide2, start_x + card_w + gap, card_y, card_w, card_h, "RDCs com Problemas", total_com_prob, COR_AMARELO)
                                add_card(slide2, start_x + 2*(card_w + gap), card_y, card_w, card_h, "% com Problemas", f"{pct_prob:.1f}%", COR_VERMELHO if pct_prob > 30 else COR_VERDE)
                                add_card(slide2, start_x + 3*(card_w + gap), card_y, card_w, card_h, "Disciplina Crítica", disc_top_val, COR_AMARELO)
                                
                                # Disciplinas resumo
                                if "DISCIPLINA" in df_g.columns:
                                    disc_resumo = df_g["DISCIPLINA"].value_counts().head(8)
                                    add_text_box(slide2, Inches(0.5), Inches(3.0), Inches(6), Inches(0.6),
                                        "RDCS POR DISCIPLINA", font_size=18, bold=True, color=COR_AZUL)
                                    
                                    tbl_rows = len(disc_resumo) + 1
                                    tbl = slide2.shapes.add_table(tbl_rows, 2, Inches(0.5), Inches(3.6), Inches(5.5), Inches(0.4 * tbl_rows)).table
                                    tbl.columns[0].width = Inches(3.5)
                                    tbl.columns[1].width = Inches(2.0)
                                    
                                    # Header
                                    for j, hdr in enumerate(["Disciplina", "Quantidade"]):
                                        cell = tbl.cell(0, j)
                                        cell.text = hdr
                                        cell.fill.solid()
                                        cell.fill.fore_color.rgb = COR_AZUL
                                        for paragraph in cell.text_frame.paragraphs:
                                            paragraph.font.size = Pt(11)
                                            paragraph.font.bold = True
                                            paragraph.font.color.rgb = COR_FUNDO
                                    
                                    for i, (disc_name, qtd) in enumerate(disc_resumo.items()):
                                        row_idx = i + 1
                                        for j, val in enumerate([str(disc_name), str(qtd)]):
                                            cell = tbl.cell(row_idx, j)
                                            cell.text = val
                                            cell.fill.solid()
                                            cell.fill.fore_color.rgb = COR_CARD_BG if row_idx % 2 == 0 else COR_FUNDO
                                            for paragraph in cell.text_frame.paragraphs:
                                                paragraph.font.size = Pt(10)
                                                paragraph.font.color.rgb = COR_BRANCO
                                
                                # Caldeira resumo
                                if "CALDEIRA" in df_g.columns:
                                    cald_resumo = df_g["CALDEIRA"].replace("", "N/I").value_counts().head(5)
                                    add_text_box(slide2, Inches(7), Inches(3.0), Inches(6), Inches(0.6),
                                        "DISTRIBUIÇÃO POR CALDEIRA", font_size=18, bold=True, color=COR_AZUL)
                                    
                                    tbl_rows2 = len(cald_resumo) + 1
                                    tbl2 = slide2.shapes.add_table(tbl_rows2, 2, Inches(7), Inches(3.6), Inches(5.5), Inches(0.4 * tbl_rows2)).table
                                    tbl2.columns[0].width = Inches(3.5)
                                    tbl2.columns[1].width = Inches(2.0)
                                    
                                    for j, hdr in enumerate(["Caldeira", "RDCs"]):
                                        cell = tbl2.cell(0, j)
                                        cell.text = hdr
                                        cell.fill.solid()
                                        cell.fill.fore_color.rgb = COR_AMARELO
                                        for paragraph in cell.text_frame.paragraphs:
                                            paragraph.font.size = Pt(11)
                                            paragraph.font.bold = True
                                            paragraph.font.color.rgb = COR_FUNDO
                                    
                                    for i, (cald_name, qtd) in enumerate(cald_resumo.items()):
                                        row_idx = i + 1
                                        for j, val in enumerate([str(cald_name), str(qtd)]):
                                            cell = tbl2.cell(row_idx, j)
                                            cell.text = val
                                            cell.fill.solid()
                                            cell.fill.fore_color.rgb = COR_CARD_BG if row_idx % 2 == 0 else COR_FUNDO
                                            for paragraph in cell.text_frame.paragraphs:
                                                paragraph.font.size = Pt(10)
                                                paragraph.font.color.rgb = COR_BRANCO
                                
                                # ============================================
                                # SLIDE 3: ANÁLISE DE GARGALOS
                                # ============================================
                                slide3 = prs.slides.add_slide(prs.slide_layouts[6])
                                set_slide_bg(slide3, COR_FUNDO)
                                
                                add_text_box(slide3, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                                    "ANÁLISE DE GARGALOS", font_size=28, bold=True, color=COR_BRANCO)
                                
                                line3 = slide3.shapes.add_shape(1, Inches(0.5), Inches(1.0), Inches(12.333), Inches(0.03))
                                line3.fill.solid()
                                line3.fill.fore_color.rgb = COR_AZUL
                                line3.line.fill.background()
                                
                                if not df_com_problema.empty:
                                    # Tabela de categorias de problemas
                                    cat_counts_pptx = df_com_problema["_CATEGORIA_PROB"].value_counts()
                                    add_text_box(slide3, Inches(0.5), Inches(1.3), Inches(6), Inches(0.5),
                                        "CATEGORIAS DE PROBLEMAS IDENTIFICADOS", font_size=16, bold=True, color=COR_AZUL)
                                    
                                    tbl_rows3 = len(cat_counts_pptx) + 1
                                    tbl3 = slide3.shapes.add_table(tbl_rows3, 3, Inches(0.5), Inches(1.9), Inches(6), Inches(0.4 * tbl_rows3)).table
                                    tbl3.columns[0].width = Inches(3.2)
                                    tbl3.columns[1].width = Inches(1.4)
                                    tbl3.columns[2].width = Inches(1.4)
                                    
                                    for j, hdr in enumerate(["Categoria", "Ocorrências", "% do Total"]):
                                        cell = tbl3.cell(0, j)
                                        cell.text = hdr
                                        cell.fill.solid()
                                        cell.fill.fore_color.rgb = COR_VERMELHO
                                        for paragraph in cell.text_frame.paragraphs:
                                            paragraph.font.size = Pt(11)
                                            paragraph.font.bold = True
                                            paragraph.font.color.rgb = RGBColor(255, 255, 255)
                                    
                                    for i, (cat_name, qtd) in enumerate(cat_counts_pptx.items()):
                                        row_idx = i + 1
                                        pct = (qtd / total_com_prob * 100) if total_com_prob > 0 else 0
                                        for j, val in enumerate([str(cat_name), str(qtd), f"{pct:.1f}%"]):
                                            cell = tbl3.cell(row_idx, j)
                                            cell.text = val
                                            cell.fill.solid()
                                            cell.fill.fore_color.rgb = COR_CARD_BG if row_idx % 2 == 0 else COR_FUNDO
                                            for paragraph in cell.text_frame.paragraphs:
                                                paragraph.font.size = Pt(10)
                                                paragraph.font.color.rgb = COR_BRANCO
                                    
                                    # Problemas por disciplina
                                    disc_prob_pptx = df_com_problema["DISCIPLINA"].value_counts().head(6)
                                    if not disc_prob_pptx.empty:
                                        add_text_box(slide3, Inches(7), Inches(1.3), Inches(6), Inches(0.5),
                                            "PROBLEMAS POR DISCIPLINA", font_size=16, bold=True, color=COR_AZUL)
                                        
                                        tbl_rows4 = len(disc_prob_pptx) + 1
                                        tbl4 = slide3.shapes.add_table(tbl_rows4, 2, Inches(7), Inches(1.9), Inches(5.5), Inches(0.4 * tbl_rows4)).table
                                        tbl4.columns[0].width = Inches(3.5)
                                        tbl4.columns[1].width = Inches(2.0)
                                        
                                        for j, hdr in enumerate(["Disciplina", "Problemas"]):
                                            cell = tbl4.cell(0, j)
                                            cell.text = hdr
                                            cell.fill.solid()
                                            cell.fill.fore_color.rgb = COR_AMARELO
                                            for paragraph in cell.text_frame.paragraphs:
                                                paragraph.font.size = Pt(11)
                                                paragraph.font.bold = True
                                                paragraph.font.color.rgb = COR_FUNDO
                                        
                                        for i, (d_name, qtd) in enumerate(disc_prob_pptx.items()):
                                            row_idx = i + 1
                                            for j, val in enumerate([str(d_name), str(qtd)]):
                                                cell = tbl4.cell(row_idx, j)
                                                cell.text = val
                                                cell.fill.solid()
                                                cell.fill.fore_color.rgb = COR_CARD_BG if row_idx % 2 == 0 else COR_FUNDO
                                                for paragraph in cell.text_frame.paragraphs:
                                                    paragraph.font.size = Pt(10)
                                                    paragraph.font.color.rgb = COR_BRANCO
                                    
                                    # Últimos problemas
                                    df_ult_prob = df_com_problema.sort_values("_DATA_DT", ascending=False).head(6)
                                    if not df_ult_prob.empty:
                                        add_text_box(slide3, Inches(0.5), Inches(1.9 + 0.4 * tbl_rows3 + 0.3), Inches(12), Inches(0.5),
                                            "ÚLTIMOS PROBLEMAS REPORTADOS", font_size=16, bold=True, color=COR_AMARELO)
                                        
                                        y_tbl5 = Inches(1.9 + 0.4 * tbl_rows3 + 0.8)
                                        rows5 = min(len(df_ult_prob), 6) + 1
                                        tbl5 = slide3.shapes.add_table(rows5, 4, Inches(0.5), y_tbl5, Inches(12.333), Inches(0.38 * rows5)).table
                                        tbl5.columns[0].width = Inches(1.5)
                                        tbl5.columns[1].width = Inches(3.0)
                                        tbl5.columns[2].width = Inches(2.5)
                                        tbl5.columns[3].width = Inches(5.333)
                                        
                                        for j, hdr in enumerate(["Data", "Encarregado", "Disciplina", "Problema"]):
                                            cell = tbl5.cell(0, j)
                                            cell.text = hdr
                                            cell.fill.solid()
                                            cell.fill.fore_color.rgb = COR_AZUL
                                            for paragraph in cell.text_frame.paragraphs:
                                                paragraph.font.size = Pt(10)
                                                paragraph.font.bold = True
                                                paragraph.font.color.rgb = COR_FUNDO
                                        
                                        for i, (_, row_p) in enumerate(df_ult_prob.head(6).iterrows()):
                                            row_idx = i + 1
                                            vals = [
                                                str(row_p.get("DATA", ""))[:10],
                                                str(row_p.get("ENCARREGADO", "")),
                                                str(row_p.get("DISCIPLINA", "")),
                                                str(row_p.get("PROBLEMAS", ""))[:80]
                                            ]
                                            for j, val in enumerate(vals):
                                                cell = tbl5.cell(row_idx, j)
                                                cell.text = val
                                                cell.fill.solid()
                                                cell.fill.fore_color.rgb = COR_CARD_BG if row_idx % 2 == 0 else COR_FUNDO
                                                for paragraph in cell.text_frame.paragraphs:
                                                    paragraph.font.size = Pt(9)
                                                    paragraph.font.color.rgb = COR_BRANCO
                                else:
                                    add_text_box(slide3, Inches(1.5), Inches(2.5), Inches(10), Inches(1),
                                        "Nenhum problema identificado nos RDCs do período selecionado.\nExcelente resultado operacional!",
                                        font_size=20, color=COR_VERDE, alignment=PP_ALIGN.CENTER)
                                
                                # ============================================
                                # SLIDE 4: ATIVIDADES REALIZADAS
                                # ============================================
                                slide4 = prs.slides.add_slide(prs.slide_layouts[6])
                                set_slide_bg(slide4, COR_FUNDO)
                                
                                add_text_box(slide4, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                                    "ATIVIDADES REALIZADAS", font_size=28, bold=True, color=COR_BRANCO)
                                
                                line4 = slide4.shapes.add_shape(1, Inches(0.5), Inches(1.0), Inches(12.333), Inches(0.03))
                                line4.fill.solid()
                                line4.fill.fore_color.rgb = COR_AZUL
                                line4.line.fill.background()
                                
                                # Top atividades
                                if "ATIVIDADE" in df_g.columns:
                                    import re as re_mod2
                                    stopwords2 = {"DE", "DO", "DA", "DOS", "DAS", "E", "EM", "NO", "NA", "NOS", "NAS", "COM", "PARA", "POR", "UM", "UMA", "O", "A", "OS", "AS", "AO", "SE", "QUE", "FOI", "SER", "TER", "COMO", "ESTÁ", "SÃO"}
                                    all_text2 = " ".join(df_g["ATIVIDADE"].dropna().astype(str).tolist()).upper()
                                    words2 = re_mod2.findall(r"[A-ZÁÉÍÓÚÂÊÔÃÕÇ]{4,}", all_text2)
                                    words2 = [w for w in words2 if w not in stopwords2 and len(w) > 3]
                                    if words2:
                                        from collections import Counter as Counter2
                                        word_counts2 = Counter2(words2).most_common(12)
                                        
                                        add_text_box(slide4, Inches(0.5), Inches(1.3), Inches(6), Inches(0.5),
                                            "PALAVRAS-CHAVE MAIS FREQUENTES", font_size=16, bold=True, color=COR_AZUL)
                                        
                                        tbl_rows6 = len(word_counts2) + 1
                                        tbl6 = slide4.shapes.add_table(tbl_rows6, 2, Inches(0.5), Inches(1.9), Inches(5.5), Inches(0.35 * tbl_rows6)).table
                                        tbl6.columns[0].width = Inches(3.5)
                                        tbl6.columns[1].width = Inches(2.0)
                                        
                                        for j, hdr in enumerate(["Atividade", "Frequência"]):
                                            cell = tbl6.cell(0, j)
                                            cell.text = hdr
                                            cell.fill.solid()
                                            cell.fill.fore_color.rgb = COR_VERDE
                                            for paragraph in cell.text_frame.paragraphs:
                                                paragraph.font.size = Pt(11)
                                                paragraph.font.bold = True
                                                paragraph.font.color.rgb = COR_FUNDO
                                        
                                        for i, (word, cnt) in enumerate(word_counts2):
                                            row_idx = i + 1
                                            for j, val in enumerate([word, str(cnt)]):
                                                cell = tbl6.cell(row_idx, j)
                                                cell.text = val
                                                cell.fill.solid()
                                                cell.fill.fore_color.rgb = COR_CARD_BG if row_idx % 2 == 0 else COR_FUNDO
                                                for paragraph in cell.text_frame.paragraphs:
                                                    paragraph.font.size = Pt(10)
                                                    paragraph.font.color.rgb = COR_BRANCO
                                
                                # DDS resumo
                                if "DDS" in df_g.columns:
                                    dds_v2 = df_g["DDS"].dropna().astype(str)
                                    dds_v2 = dds_v2[~dds_v2.str.strip().str.lower().isin(["nan", "", "não informado", "nao informado", "-", "n/a"])]
                                    if len(dds_v2) > 0:
                                        dds_top = dds_v2.str.upper().value_counts().head(8)
                                        add_text_box(slide4, Inches(7), Inches(1.3), Inches(6), Inches(0.5),
                                            "TEMAS DE DDS ABORDADOS", font_size=16, bold=True, color=COR_AMARELO)
                                        
                                        tbl_rows7 = len(dds_top) + 1
                                        tbl7 = slide4.shapes.add_table(tbl_rows7, 2, Inches(7), Inches(1.9), Inches(5.5), Inches(0.35 * tbl_rows7)).table
                                        tbl7.columns[0].width = Inches(3.5)
                                        tbl7.columns[1].width = Inches(2.0)
                                        
                                        for j, hdr in enumerate(["Tema DDS", "Vezes"]):
                                            cell = tbl7.cell(0, j)
                                            cell.text = hdr
                                            cell.fill.solid()
                                            cell.fill.fore_color.rgb = COR_AMARELO
                                            for paragraph in cell.text_frame.paragraphs:
                                                paragraph.font.size = Pt(11)
                                                paragraph.font.bold = True
                                                paragraph.font.color.rgb = COR_FUNDO
                                        
                                        for i, (dds_name, qtd) in enumerate(dds_top.items()):
                                            row_idx = i + 1
                                            for j, val in enumerate([str(dds_name), str(qtd)]):
                                                cell = tbl7.cell(row_idx, j)
                                                cell.text = val
                                                cell.fill.solid()
                                                cell.fill.fore_color.rgb = COR_CARD_BG if row_idx % 2 == 0 else COR_FUNDO
                                                for paragraph in cell.text_frame.paragraphs:
                                                    paragraph.font.size = Pt(10)
                                                    paragraph.font.color.rgb = COR_BRANCO
                                
                                # Rodapé no último slide
                                add_text_box(slide4, Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4),
                                    f"{nome_site} — Relatório gerado automaticamente pelo Sistema RDC & PDE — {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}",
                                    font_size=9, color=COR_CINZA, alignment=PP_ALIGN.CENTER)
                                
                                # === SALVAR E DISPONIBILIZAR ===
                                buffer_pptx = io.BytesIO()
                                prs.save(buffer_pptx)
                                buffer_pptx.seek(0)
                                nome_pptx = f"Relatorio_Executivo_{datetime.datetime.now().strftime('%d_%m_%Y_%H%M')}.pptx"
                                
                                st.download_button(
                                    label="⬇️ Baixar Apresentação (.pptx)",
                                    data=buffer_pptx,
                                    file_name=nome_pptx,
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                    use_container_width=True
                                )
                                st.success(f"✅ Apresentação gerada com {len(prs.slides)} slides!")
                                
                            except ImportError:
                                st.error("A biblioteca `python-pptx` não está instalada. Instale com `pip install python-pptx`.")
                            except Exception as e:
                                st.error(f"Erro ao gerar apresentação: {e}")

                else:
                    st.info("📭 Banco de RDCs vazio. Processe alguns RDCs na aba 'Leitor de RDC (IA)' e clique em 'Confirmar e Salvar' para popular os gráficos de análise.")
            else:
                st.info("📭 Nenhum RDC salvo ainda. Processe RDCs na aba 'Leitor de RDC (IA)' e clique em 'Confirmar e Salvar' para começar a análise.")
        except Exception as e:
            st.error(f"Erro ao carregar Análise de Gargalos: {e}")

    # ==============================================================
    # ABA 13: BANCO DE DADOS (GOOGLE SHEETS EMBUTIDO)
    # ==============================================================
    with tab_banco_dados:
        st.markdown("### 📊 Banco de Dados (Planilha ao Vivo)")
        
        sheets_url_edit = "https://docs.google.com/spreadsheets/d/1ajWLKG4I56_QAwc1VoZmi8w4YSGbmHf6oEho_yWmsYY/edit?usp=sharing"
        
        # Painel principal com instruções
        st.markdown("""
        <div style="background: linear-gradient(135deg, #1e3a5f, #0f2027); border-radius: 16px; padding: 25px; border: 1px solid rgba(14,165,233,0.3); margin-bottom: 20px;">
            <h3 style="color: #0ea5e9; margin-top: 0;">📝 Como editar o Banco de Dados</h3>
            <ol style="color: #94a3b8; font-size: 15px; line-height: 1.8;">
                <li>Clique no botão <b style="color: #22c55e;">verde</b> abaixo para abrir a planilha no Google Sheets</li>
                <li>Faça suas edições normalmente (as fórmulas PROCV funcionam!)</li>
                <li>Volte aqui e clique em <b style="color: #0ea5e9;">"🔄 Atualizar Dados no Site"</b></li>
            </ol>
        </div>
        """, unsafe_allow_html=True)
        
        col_open, col_refresh = st.columns(2)
        with col_open:
            st.link_button("📝 ABRIR PLANILHA PARA EDITAR", sheets_url_edit, use_container_width=True, type="primary")
        with col_refresh:
            if st.button("🔄 Atualizar Dados no Site", use_container_width=True, key="btn_refresh_sheets", type="secondary"):
                st.session_state.df = None
                st.cache_data.clear()
                st.rerun()
        
        st.markdown("---")
        
        # Mostrar preview dos dados atuais do PDE
        st.markdown("#### 👁️ Visualização Atual do PDE no Sistema")
        if st.session_state.df is not None and not st.session_state.df.empty:
            df_preview = st.session_state.df.copy()
            
            col_info1, col_info2, col_info3 = st.columns(3)
            with col_info1:
                st.metric("👷 Total de Funcionários", len(df_preview))
            with col_info2:
                if 'DISCIPLINA' in df_preview.columns:
                    st.metric("📂 Disciplinas", df_preview['DISCIPLINA'].nunique())
                elif 'C.C' in df_preview.columns:
                    st.metric("📂 C.C", df_preview['C.C'].nunique())
            with col_info3:
                if 'FUNCAO' in df_preview.columns:
                    st.metric("🔧 Funções", df_preview['FUNCAO'].nunique())
            
            st.markdown("")
            
            # Filtro de busca
            busca = st.text_input("🔍 Buscar por nome, matrícula ou função:", key="busca_banco_dados")
            if busca:
                mask = df_preview.apply(lambda row: busca.upper() in str(row.values).upper(), axis=1)
                df_preview = df_preview[mask]
                st.caption(f"Mostrando {len(df_preview)} resultado(s) para '{busca}'")
            
            st.dataframe(df_preview, use_container_width=True, height=500)
        else:
            st.info("ℹ️ Nenhum dado carregado. Clique em '🔄 Atualizar Dados no Site' para carregar.")

    # ==============================================================
    # ABA 13: ADMIN
    # ==============================================================
    with tab_admin:
        st.markdown("### ⚙️ Painel Administrativo")
        st.markdown("Controle central do banco de dados e configurações do sistema.")

        try:
            # Database Status Table
            st.markdown("#### 📦 Status do Banco de Dados")
            tabelas_info = []
            arquivos_banco = {
                "Colaboradores (PDE)": caminho_base_salva_csv,
                "Escala Diária": caminho_escala_csv,
                "Registros RDC": caminho_rdc_registros_csv,
                "Histórico F1": caminho_historico_f1_csv,
                "Histórico C.C": caminho_hist_cc,
                "Exceções F1": os.path.join(pasta_base, "f1_excecoes.csv")
            }

            for nome_tabela, caminho_tabela in arquivos_banco.items():
                if os.path.exists(caminho_tabela):
                    tamanho = os.path.getsize(caminho_tabela) / 1024
                    modificado = datetime.datetime.fromtimestamp(os.path.getmtime(caminho_tabela)).strftime('%d/%m/%Y %H:%M')
                    try:
                        linhas = len(pd.read_csv(caminho_tabela))
                    except:
                        linhas = 0
                    status_tb = "✅ Ativo"
                else:
                    tamanho = 0
                    modificado = "-"
                    linhas = 0
                    status_tb = "❌ Não encontrado"

                tabelas_info.append({
                    "Tabela": nome_tabela,
                    "Status": status_tb,
                    "Registros": linhas,
                    "Tamanho (KB)": round(tamanho, 1),
                    "Última Atualização": modificado
                })

            df_status = pd.DataFrame(tabelas_info)
            st.dataframe(df_status, use_container_width=True, hide_index=True)

            # Action Buttons
            st.markdown("#### 🔧 Ações")
            c1, c2, c3 = st.columns(3)
            with c1:
                try:
                    buffer_backup = io.BytesIO()
                    with pd.ExcelWriter(buffer_backup, engine='openpyxl') as writer:
                        for item in tabelas_info:
                            caminho_item = arquivos_banco[item["Tabela"]]
                            if os.path.exists(caminho_item):
                                try:
                                    df_temp = pd.read_csv(caminho_item)
                                    sheet_name = item["Tabela"][:31]
                                    df_temp.to_excel(writer, index=False, sheet_name=sheet_name)
                                except:
                                    pass

                    st.download_button(
                        label="📥 Backup Geral (Excel)",
                        data=buffer_backup.getvalue(),
                        file_name=f"Backup_Geral_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        key="btn_backup_geral",
                        use_container_width=True
                    )
                except Exception as e:
                    st.error(f"Erro ao gerar backup: {e}")

            with c2:
                if st.button("🔄 Sincronizar Google Sheets", key="btn_sync_admin", use_container_width=True):
                    try:
                        if conn:
                            conn.update(worksheet="PDE", data=df_atual)
                            if st.session_state.get('df_historico_f1') is not None:
                                salvar_f1_seguro(conn, st.session_state.df_historico_f1, caminho_historico_f1_csv)
                            st.success("✅ Sincronização com Google Sheets concluída!")
                        else:
                            st.warning("Conexão com Google Sheets não disponível.")
                    except Exception as e:
                        st.error(f"Erro ao sincronizar: {e}")

            with c3:
                with st.popover("🗑️ Limpar Tabela"):
                    st.warning("⚠️ Cuidado! Isso apagará todos os dados da tabela selecionada.")
                    tabela_limpar = st.selectbox("Selecione a tabela:", options=list(arquivos_banco.keys()), key="admin_tabela_limpar")
                    confirmacao = st.text_input("Digite 'CONFIRMAR' para prosseguir:", key="admin_confirmacao")
                    if st.button("🗑️ Apagar Tabela", type="primary", key="btn_apagar_tabela"):
                        if confirmacao == "CONFIRMAR":
                            caminho_apagar = arquivos_banco[tabela_limpar]
                            if os.path.exists(caminho_apagar):
                                try:
                                    df_empty = pd.read_csv(caminho_apagar).head(0)
                                    df_empty.to_csv(caminho_apagar, index=False)
                                    st.success(f"✅ Tabela '{tabela_limpar}' limpa com sucesso!")
                                    time.sleep(1)
                                    st.rerun()
                                except Exception as e:
                                    st.error(f"Erro ao limpar: {e}")
                            else:
                                st.error("Arquivo não encontrado.")
                        else:
                            st.error("Confirmação incorreta. Digite exatamente 'CONFIRMAR'.")

            # Encarregados Manager
            with st.expander("👥 Gerenciar Encarregados (Lista Oficial)"):
                if "ENCARREGADO" in df_atual.columns:
                    encarregados_ativos = sorted(df_atual["ENCARREGADO"].dropna().unique().tolist())
                    df_enc = pd.DataFrame({"Nome do Encarregado": encarregados_ativos, "Efetivo": [len(df_atual[df_atual["ENCARREGADO"] == e]) for e in encarregados_ativos]})
                    st.dataframe(df_enc, use_container_width=True, hide_index=True)
                    st.caption(f"Total: {len(encarregados_ativos)} encarregados ativos")
                else:
                    st.info("Coluna 'ENCARREGADO' não encontrada na base atual.")

            # Activity Logs
            st.markdown("#### 🕒 Atividade Recente")
            logs_ordenados = sorted(tabelas_info, key=lambda x: x["Última Atualização"], reverse=True)
            for item in logs_ordenados:
                if item["Última Atualização"] != "-":
                    st.markdown(f"• **{item['Tabela']}** — {item['Registros']} registros — atualizada em {item['Última Atualização']}")

        except Exception as e:
            st.error(f"Erro no painel administrativo: {e}")


else:
    st.markdown(f"""
    <style>
    @keyframes slowPulse {{
        0% {{ box-shadow: 0 0 20px rgba(14, 165, 233, 0.1); transform: scale(1); }}
        50% {{ box-shadow: 0 0 50px rgba(14, 165, 233, 0.3); transform: scale(1.02); }}
        100% {{ box-shadow: 0 0 20px rgba(14, 165, 233, 0.1); transform: scale(1); }}
    }}
    .empty-state-card {{
        background: rgba(15, 23, 42, 0.6);
        backdrop-filter: blur(20px);
        padding: 60px 40px;
        border-radius: 24px;
        text-align: center;
        border: 2px dashed rgba(14, 165, 233, 0.4);
        margin-top: 60px;
        animation: slowPulse 3s infinite ease-in-out;
        max-width: 600px;
        margin-left: auto;
        margin-right: auto;
    }}
    .empty-state-icon {{
        font-size: 64px;
        margin-bottom: 20px;
        background: linear-gradient(135deg, #0ea5e9, #8b5cf6);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        display: inline-block;
    }}
    </style>
    <div class="empty-state-card">
        <div class="empty-state-icon">📁</div>
        <h2 style="color: {cor_texto} !important; font-family: 'Outfit', sans-serif; font-size: 2.2rem; font-weight: 700; margin-bottom: 15px;">{t('Aguardando Base de Dados')}</h2>
        <p style="font-size: 1.1rem; color: {cor_texto_sub}; letter-spacing: 0.5px; line-height: 1.6;">{t('O sistema está pronto.<br>Para iniciar a gestão, <b>arraste o arquivo de Efetivo (.csv ou .xlsx)</b><br>para a área de upload na barra lateral.')}</p>
    </div>
    """, unsafe_allow_html=True)
